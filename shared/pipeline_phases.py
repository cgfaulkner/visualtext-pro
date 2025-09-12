#!/usr/bin/env python3
"""
Three-Phase Pipeline Implementation
==================================

Implements the clean three-phase pipeline:

Phase 1: Scan - Extract visual_index and current_alt_by_key from PPTX
Phase 2: Generate - Create generated_alt_by_key for missing ALT text  
Phase 3: Resolve - Merge current + generated into final_alt_map

This ensures single source of truth and eliminates double LLaVA calls.
"""

from __future__ import annotations
import hashlib
import logging
import time
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple

from pipeline_artifacts import RunArtifacts


logger = logging.getLogger(__name__)


def phase1_scan(pptx_path: Path, artifacts: RunArtifacts, 
               generate_thumbnails: bool = True) -> Dict[str, Any]:
    """
    Phase 1: Scan PPTX and build visual_index + current_alt_by_key.
    
    This phase extracts:
    - visual_index: Complete catalog of all images with stable keys, metadata
    - current_alt_by_key: Existing ALT text already in the deck
    - thumbnails: Generated for DOCX review (optional)
    
    Args:
        pptx_path: Path to PPTX file
        artifacts: RunArtifacts for managing file paths
        generate_thumbnails: Whether to generate thumbnail images
        
    Returns:
        Dictionary with scan results
    """
    logger.info(f"Phase 1: Scanning {pptx_path.name}")
    
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import base64
        from PIL import Image
        import io
        
        # Load presentation
        prs = Presentation(str(pptx_path))
        
        visual_index = {}
        current_alt_by_key = {}
        thumbnails_generated = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # Create stable key
                        image_data = shape.image.blob if hasattr(shape, 'image') else b''
                        hash_suffix = hashlib.sha256(image_data).hexdigest()[:8]
                        
                        # Try to get shape ID from cNvPr
                        shape_id = getattr(shape, 'shape_id', shape_idx)
                        stable_key = f"slide_{slide_idx}_shapeid_{shape_id}_hash_{hash_suffix}"
                        
                        # Extract current ALT text
                        current_alt = ""
                        try:
                            # Try descr first (primary ALT text)
                            if hasattr(shape, '_element'):
                                pic_element = shape._element
                                nvpicpr = pic_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
                                if nvpicpr is not None:
                                    current_alt = nvpicpr.get('descr', '') or nvpicpr.get('title', '')
                        except Exception as e:
                            logger.debug(f"Could not extract ALT text from {stable_key}: {e}")
                        
                        # Clean current ALT text
                        if current_alt:
                            current_alt = current_alt.strip()
                            if current_alt:
                                current_alt_by_key[stable_key] = current_alt
                        
                        # Build visual index entry
                        visual_entry = {
                            'stable_key': stable_key,
                            'slide_idx': slide_idx,
                            'shape_idx': shape_idx,
                            'shape_id': shape_id,
                            'slide_number': slide_idx + 1,
                            'image_number': len(visual_index) + 1,
                            'has_current_alt': bool(current_alt),
                            'current_alt_text': current_alt
                        }
                        
                        # Add image metadata if available
                        if hasattr(shape, 'image'):
                            try:
                                visual_entry.update({
                                    'width_px': shape.image.size[0] if shape.image.size else 0,
                                    'height_px': shape.image.size[1] if shape.image.size else 0,
                                    'file_size': len(image_data)
                                })
                            except Exception as e:
                                logger.debug(f"Could not extract image metadata: {e}")
                        
                        # Generate thumbnail if requested
                        if generate_thumbnails and image_data:
                            try:
                                thumbnail_path = artifacts.thumbs_dir / f"{stable_key}.jpg"
                                
                                # Create thumbnail
                                img = Image.open(io.BytesIO(image_data))
                                img.thumbnail((200, 200), Image.Resampling.LANCZOS)
                                img.save(thumbnail_path, 'JPEG', quality=85)
                                
                                visual_entry['thumbnail_path'] = str(thumbnail_path)
                                thumbnails_generated += 1
                                
                            except Exception as e:
                                logger.debug(f"Could not generate thumbnail for {stable_key}: {e}")
                                visual_entry['thumbnail_error'] = str(e)
                        
                        visual_index[stable_key] = visual_entry
                        
                    except Exception as e:
                        logger.warning(f"Error processing image at slide {slide_idx}, shape {shape_idx}: {e}")
        
        # Save artifacts
        artifacts.save_visual_index(visual_index)
        artifacts.save_current_alt_by_key(current_alt_by_key)
        
        result = {
            'success': True,
            'total_images': len(visual_index),
            'images_with_current_alt': len(current_alt_by_key),
            'thumbnails_generated': thumbnails_generated,
            'visual_index_path': str(artifacts.visual_index_path),
            'current_alt_by_key_path': str(artifacts.current_alt_by_key_path)
        }
        
        logger.info(f"Phase 1 complete: {len(visual_index)} images, {len(current_alt_by_key)} with existing ALT")
        return result
        
    except Exception as e:
        logger.error(f"Phase 1 failed: {e}", exc_info=True)
        return {
            'success': False,
            'error': str(e),
            'total_images': 0
        }


def phase2_generate(artifacts: RunArtifacts, alt_generator, 
                   force_regenerate: bool = False) -> Dict[str, Any]:
    """
    Phase 2: Generate ALT text for images that need it.
    
    Only generates for keys where current_alt_by_key[key] is empty.
    Uses caching to avoid duplicate LLaVA calls.
    
    Args:
        artifacts: RunArtifacts for managing file paths
        alt_generator: ALT text generation provider
        force_regenerate: If True, regenerate even if cache exists
        
    Returns:
        Dictionary with generation results
    """
    logger.info("Phase 2: Generating missing ALT text")
    
    try:
        # Load Phase 1 results
        visual_index = artifacts.load_visual_index()
        current_alt_by_key = artifacts.load_current_alt_by_key()
        
        if not visual_index:
            logger.warning("No visual_index found - Phase 1 may not have run")
            return {'success': False, 'error': 'No visual index available'}
        
        # Load existing generated ALT (for caching)
        generated_alt_by_key = {}
        if not force_regenerate:
            generated_alt_by_key = artifacts.load_generated_alt_by_key()
        
        # Find keys that need generation
        keys_needing_generation = []
        for key, visual_info in visual_index.items():
            has_current = key in current_alt_by_key and current_alt_by_key[key].strip()
            has_generated = key in generated_alt_by_key and generated_alt_by_key[key].strip()
            
            if not has_current and not has_generated:
                keys_needing_generation.append(key)
        
        logger.info(f"Found {len(keys_needing_generation)} images needing ALT text generation")
        
        if not keys_needing_generation:
            # Nothing to generate, save empty result
            artifacts.save_generated_alt_by_key(generated_alt_by_key)
            return {
                'success': True,
                'generated_count': 0,
                'cached_count': len(generated_alt_by_key),
                'skipped_count': len(current_alt_by_key)
            }
        
        # Generate ALT text for needed keys
        generation_errors = []
        newly_generated = 0
        
        for key in keys_needing_generation:
            try:
                visual_info = visual_index[key]
                
                # Check if we have a thumbnail to use
                thumbnail_path = visual_info.get('thumbnail_path')
                if thumbnail_path and Path(thumbnail_path).exists():
                    # Generate using thumbnail
                    alt_text = alt_generator.generate_alt_text(thumbnail_path)
                else:
                    # Would need original image data - this is a fallback
                    logger.warning(f"No thumbnail available for {key}, skipping generation")
                    continue
                
                if alt_text and alt_text.strip():
                    generated_alt_by_key[key] = alt_text.strip()
                    newly_generated += 1
                    logger.debug(f"Generated ALT for {key}: {alt_text[:50]}...")
                else:
                    generation_errors.append(f"Empty result for {key}")
                    
            except Exception as e:
                error_msg = f"Generation failed for {key}: {e}"
                generation_errors.append(error_msg)
                logger.warning(error_msg)
        
        # Save generated ALT text
        artifacts.save_generated_alt_by_key(generated_alt_by_key)
        
        result = {
            'success': True,
            'generated_count': newly_generated,
            'total_generated': len(generated_alt_by_key),
            'cached_count': len(generated_alt_by_key) - newly_generated,
            'skipped_count': len(current_alt_by_key),
            'error_count': len(generation_errors),
            'errors': generation_errors[:10]  # Limit error list
        }
        
        logger.info(f"Phase 2 complete: {newly_generated} newly generated, "
                   f"{len(generated_alt_by_key)} total available")
        return result
        
    except Exception as e:
        logger.error(f"Phase 2 failed: {e}", exc_info=True)
        return {
            'success': False,
            'error': str(e),
            'generated_count': 0
        }


def phase3_resolve(artifacts: RunArtifacts) -> Dict[str, Any]:
    """
    Phase 3: Resolve final_alt_map by merging current + generated.
    
    Deterministic merge strategy:
    - If current_alt_by_key[key] exists → use that (preserve existing)
    - Else → use generated_alt_by_key[key] if available
    
    Args:
        artifacts: RunArtifacts for managing file paths
        
    Returns:
        Dictionary with resolution results
    """
    logger.info("Phase 3: Resolving final ALT text mappings")
    
    try:
        # Load Phase 1 & 2 results
        current_alt_by_key = artifacts.load_current_alt_by_key()
        generated_alt_by_key = artifacts.load_generated_alt_by_key()
        visual_index = artifacts.load_visual_index()
        
        if not visual_index:
            logger.warning("No visual_index found - Phase 1 may not have run")
            return {'success': False, 'error': 'No visual index available'}
        
        # Build final mappings using preserve-first strategy
        final_alt_map = {}
        
        # Stats tracking
        preserved_count = 0
        generated_count = 0
        missing_count = 0
        
        for key in visual_index.keys():
            if key in current_alt_by_key and current_alt_by_key[key].strip():
                # Preserve existing ALT text
                final_alt_map[key] = current_alt_by_key[key].strip()
                preserved_count += 1
            elif key in generated_alt_by_key and generated_alt_by_key[key].strip():
                # Use generated ALT text
                final_alt_map[key] = generated_alt_by_key[key].strip()
                generated_count += 1
            else:
                # No ALT text available
                missing_count += 1
        
        # Save final mappings
        artifacts.save_final_alt_map(final_alt_map)
        
        result = {
            'success': True,
            'total_images': len(visual_index),
            'final_mappings': len(final_alt_map),
            'preserved_count': preserved_count,
            'generated_count': generated_count,
            'missing_count': missing_count,
            'coverage_percent': (len(final_alt_map) / len(visual_index) * 100) if visual_index else 0
        }
        
        logger.info(f"Phase 3 complete: {len(final_alt_map)}/{len(visual_index)} images have ALT text "
                   f"({result['coverage_percent']:.1f}% coverage)")
        return result
        
    except Exception as e:
        logger.error(f"Phase 3 failed: {e}", exc_info=True)
        return {
            'success': False,
            'error': str(e),
            'final_mappings': 0
        }


def run_pipeline(pptx_path: Path, config, alt_generator, 
                force_regenerate: bool = False) -> RunArtifacts:
    """
    Run the complete three-phase pipeline.
    
    Args:
        pptx_path: Path to PPTX file  
        config: Configuration object
        alt_generator: ALT text generation provider
        force_regenerate: Force regeneration even if cache exists
        
    Returns:
        RunArtifacts with all pipeline outputs
    """
    logger.info(f"Starting pipeline for {pptx_path.name}")
    start_time = time.time()
    
    # Create artifacts structure
    artifacts = RunArtifacts.create_for_run(pptx_path)
    
    try:
        # Phase 1: Scan
        scan_result = phase1_scan(pptx_path, artifacts)
        if not scan_result['success']:
            raise RuntimeError(f"Phase 1 failed: {scan_result.get('error', 'Unknown error')}")
        
        # Phase 2: Generate (only if images need ALT text)
        if scan_result['total_images'] > scan_result['images_with_current_alt']:
            generate_result = phase2_generate(artifacts, alt_generator, force_regenerate)
            if not generate_result['success']:
                logger.warning(f"Phase 2 had issues: {generate_result.get('error', 'Unknown error')}")
        else:
            logger.info("Phase 2 skipped: All images already have ALT text")
            artifacts.save_generated_alt_by_key({})
        
        # Phase 3: Resolve
        resolve_result = phase3_resolve(artifacts)
        if not resolve_result['success']:
            raise RuntimeError(f"Phase 3 failed: {resolve_result.get('error', 'Unknown error')}")
        
        elapsed = time.time() - start_time
        logger.info(f"Pipeline completed in {elapsed:.2f}s")
        
        return artifacts
        
    except Exception as e:
        logger.error(f"Pipeline failed: {e}", exc_info=True)
        # Clean up on failure
        artifacts.cleanup(keep_finals=False)
        raise