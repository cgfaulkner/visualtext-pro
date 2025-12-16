#!/usr/bin/env python3
"""
Shape Renderer for PPTX Processing
==================================

Provides shape-to-image rendering functionality for PowerPoint shapes.
Used for generating thumbnails when slide rendering is unavailable (e.g., macOS).
"""

from typing import Optional, List, Tuple
import logging
import math

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    Image = None
    ImageDraw = None
    ImageFont = None

logger = logging.getLogger(__name__)


class ShapeRenderer:
    """
    Renders PowerPoint shapes to PIL Images.
    
    Extracted from core/pptx_processor.py to enable reuse across modules
    without external dependencies.
    """
    
    def __init__(self):
        """Initialize the shape renderer."""
        if not PIL_AVAILABLE:
            logger.warning("PIL not available - shape rendering will be disabled")
    
    def render_shape_to_image(self, shape, width_px: int, height_px: int) -> Optional[Image.Image]:
        """
        Render a PowerPoint shape to a PIL Image.
        
        Args:
            shape: Shape to render (python-pptx shape object)
            width_px: Width in pixels
            height_px: Height in pixels
            
        Returns:
            PIL Image or None if rendering failed
        """
        if not PIL_AVAILABLE:
            logger.warning("PIL not available - cannot render shapes to images")
            return None
        
        if width_px <= 0 or height_px <= 0:
            return None
        
        try:
            # Create image canvas with white background
            img = Image.new('RGB', (width_px, height_px), 'white')
            draw = ImageDraw.Draw(img)
            
            # Render shape based on type
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            shape_type = getattr(shape, 'shape_type', None)
            
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                self._render_auto_shape(draw, shape, width_px, height_px)
            elif shape_type == MSO_SHAPE_TYPE.LINE:
                self._render_line_shape(draw, shape, width_px, height_px)
            elif shape_type == MSO_SHAPE_TYPE.FREEFORM:
                self._render_freeform_shape(draw, shape, width_px, height_px)
            elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                self._render_text_box(draw, shape, width_px, height_px)
            else:
                # Generic shape rendering
                self._render_generic_shape(draw, shape, width_px, height_px)
            
            return img
            
        except Exception as e:
            logger.warning(f"Failed to render shape: {e}")
            return None
    
    def render_group_shape(self, group_shape, bbox: dict) -> Optional[Image.Image]:
        """
        Render a group shape by compositing its child shapes.
        
        Args:
            group_shape: Group shape object with .shapes attribute
            bbox: Bounding box dict with 'width' and 'height' in points
            
        Returns:
            PIL Image with composited child shapes or None if rendering failed
        """
        if not PIL_AVAILABLE:
            return None
        
        try:
            if not hasattr(group_shape, 'shapes') or not group_shape.shapes:
                return None
            
            # Convert bbox from points to pixels (96 DPI)
            width_px = max(int(bbox.get('width', 0) * 96 / 72), 50)
            height_px = max(int(bbox.get('height', 0) * 96 / 72), 50)
            
            # Get group position in EMU
            group_left_emu = getattr(group_shape, 'left', None)
            group_top_emu = getattr(group_shape, 'top', None)
            
            if group_left_emu is None or group_top_emu is None:
                group_left_emu = 0
                group_top_emu = 0
            else:
                group_left_emu = group_left_emu.emu if hasattr(group_left_emu, 'emu') else group_left_emu
                group_top_emu = group_top_emu.emu if hasattr(group_top_emu, 'emu') else group_top_emu
            
            # Create canvas
            canvas = Image.new('RGB', (width_px, height_px), 'white')
            
            # Render each child shape onto canvas
            for child in group_shape.shapes:
                try:
                    if not hasattr(child, 'left') or not hasattr(child, 'top'):
                        continue
                    
                    # Get child position and size in EMU
                    child_left_emu = child.left.emu if hasattr(child.left, 'emu') else child.left
                    child_top_emu = child.top.emu if hasattr(child.top, 'emu') else child.top
                    child_width_emu = child.width.emu if hasattr(child.width, 'emu') else child.width
                    child_height_emu = child.height.emu if hasattr(child.height, 'emu') else child.height
                    
                    # Calculate position relative to group (in pixels)
                    child_left_px = int((child_left_emu - group_left_emu) / 914400 * 96)
                    child_top_px = int((child_top_emu - group_top_emu) / 914400 * 96)
                    child_width_px = max(int(child_width_emu / 914400 * 96), 1)
                    child_height_px = max(int(child_height_emu / 914400 * 96), 1)
                    
                    # Skip if child is outside canvas bounds
                    if (child_left_px >= width_px or child_top_px >= height_px or
                        child_left_px + child_width_px <= 0 or child_top_px + child_height_px <= 0):
                        continue
                    
                    # Render child shape
                    child_img = self.render_shape_to_image(child, child_width_px, child_height_px)
                    if child_img:
                        # Composite onto canvas (handle transparency if present)
                        if child_img.mode == 'RGBA':
                            canvas.paste(child_img, (child_left_px, child_top_px), child_img)
                        else:
                            canvas.paste(child_img, (child_left_px, child_top_px))
                except Exception as e:
                    logger.debug(f"Failed to render child shape in group: {e}")
                    continue
            
            return canvas
            
        except Exception as e:
            logger.warning(f"Failed to render group shape: {e}")
            return None
    
    def _render_auto_shape(self, draw: ImageDraw.Draw, shape, width: int, height: int):
        """Render AutoShape (circles, rectangles, etc.)"""
        try:
            # Get shape fill color
            fill_color = self._get_shape_fill_color(shape)
            line_color = self._get_shape_line_color(shape)
            line_width = self._get_shape_line_width(shape)
            
            # Get shape type for specific rendering
            auto_shape_type = getattr(shape, 'auto_shape_type', None)
            
            if auto_shape_type:
                from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                
                if auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                    # Draw circle/ellipse
                    draw.ellipse([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
                elif auto_shape_type in [MSO_AUTO_SHAPE_TYPE.RECTANGLE, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE]:
                    # Draw rectangle
                    draw.rectangle([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
                elif auto_shape_type == MSO_AUTO_SHAPE_TYPE.HEXAGON:
                    # Draw hexagon
                    points = self._get_hexagon_points(width, height)
                    draw.polygon(points, fill=fill_color, outline=line_color)
                else:
                    # Generic rectangle for unknown shapes
                    draw.rectangle([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
            else:
                # Default rectangle
                draw.rectangle([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
                
        except Exception as e:
            logger.warning(f"Failed to render auto shape: {e}")
            # Fallback: draw a simple rectangle
            draw.rectangle([0, 0, width-1, height-1], fill='lightgray', outline='black')
    
    def _render_line_shape(self, draw: ImageDraw.Draw, shape, width: int, height: int):
        """Render line shape"""
        try:
            line_color = self._get_shape_line_color(shape)
            line_width = self._get_shape_line_width(shape)
            
            # Draw line from top-left to bottom-right (simplified)
            draw.line([0, 0, width-1, height-1], fill=line_color, width=line_width)
            
        except Exception as e:
            logger.warning(f"Failed to render line shape: {e}")
            draw.line([0, 0, width-1, height-1], fill='black', width=2)
    
    def _render_freeform_shape(self, draw: ImageDraw.Draw, shape, width: int, height: int):
        """Render freeform shape"""
        try:
            fill_color = self._get_shape_fill_color(shape)
            line_color = self._get_shape_line_color(shape)
            
            # For freeform, draw a polygon (simplified)
            # This is a basic implementation - complex freeforms would need path parsing
            points = [(0, height//2), (width//4, 0), (3*width//4, 0), (width-1, height//2), (width-1, height-1), (0, height-1)]
            draw.polygon(points, fill=fill_color, outline=line_color)
            
        except Exception as e:
            logger.warning(f"Failed to render freeform shape: {e}")
            draw.rectangle([0, 0, width-1, height-1], fill='lightgray', outline='black')
    
    def _render_text_box(self, draw: ImageDraw.Draw, shape, width: int, height: int):
        """Render text box with background"""
        try:
            fill_color = self._get_shape_fill_color(shape)
            line_color = self._get_shape_line_color(shape)
            line_width = self._get_shape_line_width(shape)
            
            # Draw background
            draw.rectangle([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
            
            # Add text if available
            if hasattr(shape, 'text') and shape.text:
                try:
                    font = ImageFont.load_default()
                    text_color = 'black'
                    
                    # Simple text positioning (centered)
                    text_width = len(shape.text) * 8  # Rough estimate
                    text_height = 12
                    x = max(0, (width - text_width) // 2)
                    y = max(0, (height - text_height) // 2)
                    
                    draw.text((x, y), shape.text[:50], fill=text_color, font=font)  # Limit text length
                except:
                    pass  # Text rendering is optional
                    
        except Exception as e:
            logger.warning(f"Failed to render text box: {e}")
            draw.rectangle([0, 0, width-1, height-1], fill='white', outline='black')
    
    def _render_generic_shape(self, draw: ImageDraw.Draw, shape, width: int, height: int):
        """Generic shape rendering fallback"""
        try:
            fill_color = self._get_shape_fill_color(shape)
            line_color = self._get_shape_line_color(shape)
            line_width = self._get_shape_line_width(shape)
            
            # Draw as rectangle with visual indication
            draw.rectangle([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
            
            # Add X marks to indicate it's a generic shape
            draw.line([0, 0, width-1, height-1], fill=line_color, width=1)
            draw.line([0, height-1, width-1, 0], fill=line_color, width=1)
            
        except Exception as e:
            logger.warning(f"Failed to render generic shape: {e}")
            draw.rectangle([0, 0, width-1, height-1], fill='lightgray', outline='black')
    
    def _get_shape_fill_color(self, shape) -> str:
        """Get shape fill color"""
        try:
            if hasattr(shape, 'fill') and shape.fill:
                fill = shape.fill
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    # Try to get RGB color
                    try:
                        rgb = fill.fore_color.rgb
                        return f"#{rgb:06x}"
                    except:
                        pass
                
                # Check for solid fill
                from pptx.enum.dml import MSO_FILL_TYPE
                if hasattr(fill, 'type') and fill.type == MSO_FILL_TYPE.SOLID:
                    return 'lightblue'  # Default solid color
                    
            # Shape-specific defaults
            shape_type = getattr(shape, 'shape_type', None)
            auto_shape_type = getattr(shape, 'auto_shape_type', None)
            
            if auto_shape_type:
                from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                if auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                    return 'lightblue'  # Blue circle default
                elif auto_shape_type == MSO_AUTO_SHAPE_TYPE.HEXAGON:
                    return 'purple'  # Purple hexagon default
                    
            return 'lightgray'  # Generic default
            
        except Exception:
            return 'lightgray'
    
    def _get_shape_line_color(self, shape) -> str:
        """Get shape line color"""
        try:
            if hasattr(shape, 'line') and shape.line:
                line = shape.line
                if hasattr(line, 'color') and line.color:
                    try:
                        rgb = line.color.rgb
                        return f"#{rgb:06x}"
                    except:
                        pass
            return 'black'  # Default line color
        except Exception:
            return 'black'
    
    def _get_shape_line_width(self, shape) -> int:
        """Get shape line width"""
        try:
            if hasattr(shape, 'line') and shape.line:
                line = shape.line
                if hasattr(line, 'width') and line.width:
                    # Convert EMU to pixels (rough approximation)
                    width_emu = line.width
                    width_px = max(1, int(width_emu / 914400 * 96 / 72))  # Convert to reasonable pixel width
                    return min(width_px, 10)  # Cap at 10px
            return 2  # Default line width
        except Exception:
            return 2
    
    def _get_hexagon_points(self, width: int, height: int) -> List[Tuple[int, int]]:
        """Generate hexagon points"""
        cx, cy = width // 2, height // 2
        radius_x, radius_y = width // 2 - 2, height // 2 - 2
        
        points = []
        for i in range(6):
            angle = i * math.pi / 3
            x = cx + radius_x * math.cos(angle)
            y = cy + radius_y * math.sin(angle)
            points.append((int(x), int(y)))
        
        return points

