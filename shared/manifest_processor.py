#!/usr/bin/env python3
"""
Manifest-Based PPTX Processor
=============================

Processes PPTX files using the ALT manifest as single source of truth.
Handles extraction, generation with caching/idempotency, and provenance tracking.
"""

from __future__ import annotations
import logging
import time
from pathlib import Path
from typing import Dict, Any, Optional

from alt_manifest import (
    AltManifest, AltManifestEntry, compute_image_hash, create_stable_key,
    create_instance_key, create_content_key, parse_min_shape_area, MANIFEST_SCHEMA_VERSION
)
from shape_utils import is_image_like, is_decorative_shape, is_empty_placeholder_textbox
from alt_text_reader import read_existing_alt

logger = logging.getLogger(__name__)


class ManifestProcessor:
    """
    Processes PPTX files with manifest-based caching and single source of truth.
    
    Implements the three-phase pipeline:
    Step 1: Discovery & Classification (NO LLaVA YET)
    Step 2: Rendering & Thumbnails (MODEL-AGNOSTIC) 
    Step 3: ALT Generation (with caching)
    """
    
    def __init__(self, config_manager, alt_generator=None, 
                 llava_include_shapes: str = "smart",
                 max_shapes_per_slide: int = 5,
                 min_shape_area: str = "1%"):
        self.config = config_manager
        self.alt_generator = alt_generator
        self.llava_include_shapes = llava_include_shapes
        self.max_shapes_per_slide = max_shapes_per_slide
        self.min_shape_area = min_shape_area
    
    def phase1_discover_and_classify(self, pptx_path: Path, manifest: AltManifest) -> Dict[str, Any]:
        """
        Step 1: Discovery & Classification (NO LLaVA YET)
        
        Walk slides/shapes; for each reportable element:
        - Compute instance_key (slide_index, shape_id)
        - Determine shape_type, bbox, is_group_child
        - Read existing_alt from PPTX
        - Create placeholder content_key (will be replaced in Step 2 after rendering)
        """
        logger.info("Phase 1: Discovering and classifying visual elements")
        
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            prs = Presentation(str(pptx_path))
            
            # Parse minimum area threshold
            slide = prs.slides[0] if prs.slides else None
            if slide:
                slide_width = int(prs.slide_width.emu // 914400)  # Convert to points (72 DPI)
                slide_height = int(prs.slide_height.emu // 914400)
            else:
                slide_width, slide_height = 720, 540  # Default PowerPoint slide size
            
            min_area_threshold = parse_min_shape_area(self.min_shape_area, slide_width, slide_height)
            
            discovered_elements = 0
            classified_elements = 0
            
            for slide_idx, slide in enumerate(prs.slides):
                # Get slide text for context
                slide_text = self._extract_slide_text(slide)
                slide_notes = self._extract_slide_notes(slide)
                
                shapes_on_slide = []
                
                for shape_idx, shape in enumerate(slide.shapes):
                    try:
                        # Skip empty placeholder text boxes
                        if is_empty_placeholder_textbox(shape):
                            logger.debug(f"Skipping empty placeholder textbox on slide {slide_idx + 1}")
                            continue

                        # Get shape ID (unique within slide)
                        shape_id = getattr(shape, 'shape_id', shape_idx)
                        
                        # Create instance key
                        instance_key = create_instance_key(slide_idx, shape_id)
                        
                        # Classify shape type and group membership
                        shape_type, is_group_child = manifest.classify_shape_type(shape, MSO_SHAPE_TYPE)
                        
                        # Get bounding box
                        bbox = self._get_shape_bbox(shape)
                        
                        # Calculate shape area for filtering
                        shape_area = bbox["width"] * bbox["height"]
                        
                        # Determine if this shape should be included based on strategy
                        include_shape, include_reason, exclude_reason = self._should_include_shape(
                            shape, shape_type, shape_area, min_area_threshold, len(shapes_on_slide)
                        )
                        
                        if not include_shape:
                            logger.debug(f"Excluding {instance_key}: {exclude_reason}")
                            continue
                        
                        discovered_elements += 1
                        
                        # Extract existing ALT text
                        existing_alt = self._collect_existing_alt(shape)
                        
                        # Create placeholder content_key (will be replaced in Step 2)
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, 'image'):
                            # For pictures, use image blob
                            image_data = shape.image.blob if shape.image else b''
                            content_key = create_content_key(image_data) if image_data else ""
                        else:
                            # For shapes, create placeholder based on properties
                            placeholder_content = f"{shape_type}_{instance_key}_{shape_area}".encode('utf-8')
                            content_key = create_content_key(placeholder_content)
                        
                        # Determine file format
                        format_type = self._determine_format_type(shape, shape_type)
                        
                        # Create manifest entry
                        entry = AltManifestEntry(
                            instance_key=instance_key,
                            content_key=content_key,
                            shape_type=shape_type,
                            is_group_child=is_group_child,
                            bbox=bbox,
                            format=format_type,
                            existing_alt=existing_alt,
                            had_existing_alt=bool(existing_alt.strip()),
                            include_reason=include_reason,
                            slide_idx=slide_idx,
                            slide_number=slide_idx + 1,
                            image_number=discovered_elements,
                            slide_text=slide_text,
                            slide_notes=slide_notes,
                            # Legacy compatibility
                            key=create_stable_key(slide_idx, shape_id, content_key[:16]),
                            image_hash=content_key,
                            current_alt=existing_alt,
                            width_px=int(bbox["width"]),
                            height_px=int(bbox["height"])
                        )
                        
                        manifest.add_entry(entry)
                        classified_elements += 1
                        shapes_on_slide.append(entry)
                        
                        logger.debug(f"Classified {instance_key} as {shape_type} (area: {shape_area:.0f})")
                        
                    except Exception as e:
                        logger.warning(f"Error classifying shape at slide {slide_idx}, shape {shape_idx}: {e}")
            
            result = {
                'success': True,
                'discovered_elements': discovered_elements,
                'classified_elements': classified_elements,
                'min_area_threshold': min_area_threshold,
                'include_strategy': self.llava_include_shapes
            }
            
            logger.info(f"Phase 1 complete: {classified_elements} elements classified "
                       f"(strategy: {self.llava_include_shapes}, threshold: {min_area_threshold:.0f} sq pts)")
            
            return result
            
        except Exception as e:
            logger.error(f"Phase 1 failed: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e),
                'discovered_elements': 0
            }
    
    def _should_include_shape(self, shape, shape_type: str, shape_area: float, 
                            min_area_threshold: int, shapes_on_slide_count: int) -> tuple[bool, str, str]:
        """
        Determine if a shape should be included for processing based on strategy.
        
        Returns:
            (should_include, include_reason, exclude_reason)
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        # Check per-slide limit first
        if shapes_on_slide_count >= self.max_shapes_per_slide:
            return False, "", f"slide limit reached ({self.max_shapes_per_slide} shapes)"
        
        if self.llava_include_shapes == "off":
            # Only include pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True, "pictures only mode", ""
            else:
                return False, "", "non-picture in pictures-only mode"
                
        elif self.llava_include_shapes == "all":
            # Include all shapes
            return True, "include all shapes", ""
            
        elif self.llava_include_shapes == "smart":
            # Check for decorative shapes first (exclude them)
            if is_decorative_shape(shape):
                return False, "", "decorative shape (line/border/accent)"
            
            # Include image-like shapes (pictures, vectors, groups with images)
            if is_image_like(shape):
                return True, "image-like shape (picture/vector/group)", ""
            elif shape_area >= min_area_threshold:
                return True, f"shape above threshold ({shape_area:.0f} >= {min_area_threshold})", ""
            else:
                return False, "", f"shape below threshold ({shape_area:.0f} < {min_area_threshold})"
        
        return False, "", "unknown strategy"
    
    def _get_shape_bbox(self, shape) -> Dict[str, float]:
        """Extract bounding box from shape."""
        try:
            # Convert from EMUs to points (914400 EMUs = 1 inch = 72 points)
            left = shape.left.emu / 914400 * 72 if hasattr(shape, 'left') and shape.left else 0
            top = shape.top.emu / 914400 * 72 if hasattr(shape, 'top') and shape.top else 0
            width = shape.width.emu / 914400 * 72 if hasattr(shape, 'width') and shape.width else 0
            height = shape.height.emu / 914400 * 72 if hasattr(shape, 'height') and shape.height else 0
            
            return {"left": left, "top": top, "width": width, "height": height}
        except Exception as e:
            logger.debug(f"Could not extract bbox: {e}")
            return {"left": 0, "top": 0, "width": 0, "height": 0}
    
    def _determine_format_type(self, shape, shape_type: str) -> str:
        """Determine the format type of the shape."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # For pictures, try to determine from image format
            try:
                if hasattr(shape, 'image') and shape.image:
                    # Could inspect image.ext or image.content_type if available
                    return "png"  # Default assumption for pictures
            except:
                pass
            return "jpg"  # Fallback for pictures
        else:
            # Non-picture shapes will be rendered as crops
            return "shape"
    
    def phase2_render_and_generate_crops(self, pptx_path: Path, manifest: AltManifest, 
                                       artifacts) -> Dict[str, Any]:
        """
        Step 2: Rendering & Thumbnails (MODEL-AGNOSTIC)
        
        Create display thumbnails for all visuals; create model crops for shapes:
        - Render each slide to PNG at fixed width (e.g., 1920px)
        - For each element: compute crop from slide image using bbox
        - Save thumbnail (for DOCX) and crop (for model)
        - Update content_key as hash of the model crop bytes
        """
        logger.info("Phase 2: Rendering slides and generating crops")
        
        # Detect win32com availability (Windows only)
        win32_available = False
        try:
            import win32com.client
            win32_available = True
        except ImportError:
            logger.debug(
                "win32com not available — skipping PowerPoint slide rendering "
                "(non-Windows environment)"
            )
        
        try:
            from pptx import Presentation
            from PIL import Image, ImageDraw
            import io
            import time
            import tempfile
            
            entries_processed = 0
            thumbnails_created = 0
            crops_created = 0
            rendering_errors = []
            
            # Load presentation
            prs = Presentation(str(pptx_path))
            
            # Try to use PowerPoint COM for high-quality rendering (Windows)
            slide_images = {}
            if win32_available:
                try:
                    slide_images = self._render_slides_with_powerpoint(pptx_path, target_width=1920)
                    logger.info(f"Rendered {len(slide_images)} slides using PowerPoint COM")
                except Exception as e:
                    logger.warning(f"PowerPoint COM rendering failed: {e}")
                    slide_images = {}
                    logger.info("Falling back to image-only thumbnail extraction")
            else:
                slide_images = {}
            
            # Process each manifest entry
            for entry in manifest.get_all_entries():
                try:
                    entries_processed += 1
                    
                    # Get slide image
                    slide_img = slide_images.get(entry.slide_idx)
                    
                    if slide_img is None:
                        # Fallback: try to process individual shape
                        logger.debug(f"No slide image for {entry.instance_key}, trying individual shape processing")
                        
                        # For pictures, we can extract the image directly
                        if entry.shape_type == "PICTURE":
                            shape = self._find_shape_by_instance_key(prs, entry)
                            if shape and hasattr(shape, 'image') and shape.image:
                                image_data = shape.image.blob
                                
                                # Create thumbnail
                                thumbnail_path = artifacts.thumbs_dir / f"{entry.instance_key}.jpg"
                                self._create_thumbnail_from_bytes(image_data, thumbnail_path)
                                entry.thumb_path = str(thumbnail_path)
                                thumbnails_created += 1
                                
                                # Use original image as crop
                                crop_path = artifacts.crops_dir / f"{entry.instance_key}.png"
                                with open(crop_path, 'wb') as f:
                                    f.write(image_data)
                                entry.crop_path = str(crop_path)
                                crops_created += 1
                                
                                # Update content_key with actual image data
                                entry.content_key = create_content_key(image_data)
                                entry.rasterizer_info = {
                                    "engine": "direct_extract",
                                    "dpi": 0,
                                    "status": "success"
                                }
                                
                                manifest.add_entry(entry)
                                continue
                        
                        # For non-pictures without slide rendering, skip
                        logger.warning(f"Cannot process {entry.instance_key} without slide rendering")
                        entry.rasterizer_info = {
                            "engine": "none",
                            "dpi": 0,
                            "status": "no_slide_render"
                        }
                        rendering_errors.append(f"No slide rendering available for {entry.instance_key}")
                        continue
                    
                    # Create crop from slide image using bbox
                    bbox = entry.bbox
                    if not bbox or not all(bbox.get(k, 0) > 0 for k in ['width', 'height']):
                        logger.warning(f"Invalid bbox for {entry.instance_key}: {bbox}")
                        continue
                    
                    # Scale bbox coordinates to match rendered slide dimensions
                    slide_width_pts = float(prs.slide_width.emu) / 914400 * 72  # Convert to points
                    slide_height_pts = float(prs.slide_height.emu) / 914400 * 72
                    
                    scale_x = slide_img.width / slide_width_pts
                    scale_y = slide_img.height / slide_height_pts
                    
                    # Calculate crop rectangle (with small padding)
                    padding = 5  # pixels
                    crop_left = max(0, int(bbox["left"] * scale_x) - padding)
                    crop_top = max(0, int(bbox["top"] * scale_y) - padding)
                    crop_right = min(slide_img.width, int((bbox["left"] + bbox["width"]) * scale_x) + padding)
                    crop_bottom = min(slide_img.height, int((bbox["top"] + bbox["height"]) * scale_y) + padding)
                    
                    # Extract crop
                    crop_img = slide_img.crop((crop_left, crop_top, crop_right, crop_bottom))
                    
                    # Save crop (model input)
                    crop_path = artifacts.crops_dir / f"{entry.instance_key}.png"
                    crop_img.save(crop_path, 'PNG')
                    entry.crop_path = str(crop_path)
                    crops_created += 1
                    
                    # Save thumbnail (display/DOCX)
                    thumbnail_path = artifacts.thumbs_dir / f"{entry.instance_key}.jpg" 
                    thumbnail_img = crop_img.copy()
                    thumbnail_img.thumbnail((200, 200), Image.Resampling.LANCZOS)
                    thumbnail_img.save(thumbnail_path, 'JPEG', quality=85)
                    entry.thumb_path = str(thumbnail_path)
                    thumbnails_created += 1
                    
                    # Update content_key with crop image bytes
                    with open(crop_path, 'rb') as f:
                        crop_bytes = f.read()
                    entry.content_key = create_content_key(crop_bytes)
                    
                    # Record rasterizer info
                    entry.rasterizer_info = {
                        "engine": "slide_render" if slide_img else "direct_extract",
                        "dpi": int(1920 / slide_width_pts * 72),  # Approximate DPI
                        "status": "success"
                    }
                    
                    manifest.add_entry(entry)
                    
                    logger.debug(f"Processed {entry.instance_key}: crop {crop_right-crop_left}x{crop_bottom-crop_top}")
                    
                except Exception as e:
                    error_msg = f"Error processing {entry.instance_key}: {e}"
                    logger.warning(error_msg)
                    rendering_errors.append(error_msg)
                    
                    # Mark entry with error
                    entry.rasterizer_info = {
                        "engine": "error",
                        "dpi": 0,
                        "status": f"error: {str(e)[:100]}"
                    }
                    manifest.add_entry(entry)
            
            result = {
                'success': True,
                'entries_processed': entries_processed,
                'thumbnails_created': thumbnails_created,
                'crops_created': crops_created,
                'slides_rendered': len(slide_images),
                'rendering_errors': rendering_errors[:10]  # Limit error list
            }
            
            logger.info(f"Phase 2 complete: {crops_created} crops, {thumbnails_created} thumbnails created")
            return result
            
        except Exception as e:
            logger.error(f"Phase 2 failed: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e),
                'entries_processed': 0
            }
    
    def _render_slides_with_powerpoint(self, pptx_path: Path, target_width: int = 1920) -> Dict[int, Image.Image]:
        """
        Render slides to PNG using PowerPoint COM automation (Windows only).
        
        Returns:
            Dictionary mapping slide_idx -> PIL Image
        """
        import tempfile
        import win32com.client
        from PIL import Image
        
        slide_images = {}
        
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_dir_path = Path(temp_dir)
            
            try:
                # Initialize PowerPoint
                ppt = win32com.client.Dispatch("PowerPoint.Application")
                ppt.Visible = False
                
                # Open presentation
                presentation = ppt.Presentations.Open(str(pptx_path.absolute()))
                
                # Export each slide as PNG
                for slide_idx in range(presentation.Slides.Count):
                    slide_num = slide_idx + 1  # PowerPoint uses 1-based indexing
                    png_path = temp_dir_path / f"slide_{slide_idx}.png"
                    
                    # Export slide to PNG
                    presentation.Slides(slide_num).Export(str(png_path), "PNG", target_width, 0)
                    
                    # Load as PIL Image
                    slide_img = Image.open(png_path)
                    slide_images[slide_idx] = slide_img.copy()
                    
                # Close presentation
                presentation.Close()
                ppt.Quit()
                
            except Exception as e:
                logger.error(f"PowerPoint COM rendering failed: {e}")
                raise
        
        return slide_images
    
    def _find_shape_by_instance_key(self, prs, entry):
        """Find a shape in the presentation by instance key."""
        try:
            slide = prs.slides[entry.slide_idx]
            # Extract shape_id from instance_key
            parts = entry.instance_key.split('_')
            if len(parts) >= 3:
                shape_id = int(parts[2])  # slide_X_shape_Y -> Y
                for shape in slide.shapes:
                    if getattr(shape, 'shape_id', 0) == shape_id:
                        return shape
        except Exception as e:
            logger.debug(f"Could not find shape for {entry.instance_key}: {e}")
        return None
    
    def _create_thumbnail_from_bytes(self, image_data: bytes, thumbnail_path: Path):
        """Create thumbnail from image bytes."""
        try:
            from PIL import Image
            import io
            
            img = Image.open(io.BytesIO(image_data))
            img.thumbnail((200, 200), Image.Resampling.LANCZOS)
            img.save(thumbnail_path, 'JPEG', quality=85)
            
        except Exception as e:
            logger.debug(f"Could not create thumbnail at {thumbnail_path}: {e}")
    
    def phase3_inclusion_policy_and_caching(self, manifest: AltManifest, mode: str = "preserve") -> Dict[str, Any]:
        """
        Step 3: Inclusion Policy & Caching (DECISION ENGINE)
        
        Decide who goes to LLaVA; ensure idempotency:
        - Respect --llava-include-shapes (off/smart/all) 
        - Respect preserve mode with existing_alt
        - Use content_key caching for idempotency
        """
        logger.info(f"Phase 3: Applying inclusion policy and caching (mode: {mode})")
        
        try:
            preserved_count = 0
            cached_count = 0
            needs_generation_count = 0
            excluded_count = 0
            policy_decisions = []
            
            for entry in manifest.get_all_entries():
                decision = self._make_generation_decision(entry, manifest, mode)
                policy_decisions.append(decision)
                
                # Update manifest entry with decision
                entry.decision_reason = decision['decision_reason']
                entry.include_reason = decision.get('include_reason', '')
                entry.exclude_reason = decision.get('exclude_reason', '')
                
                if decision['action'] == 'preserve':
                    entry.final_alt = entry.existing_alt
                    entry.llava_called = False
                    preserved_count += 1
                    
                elif decision['action'] == 'use_cache':
                    entry.final_alt = decision['cached_alt']
                    entry.llava_called = False
                    cached_count += 1
                    
                elif decision['action'] == 'generate':
                    # Mark for generation - will be processed in Step 4
                    needs_generation_count += 1
                    
                elif decision['action'] == 'exclude':
                    # Not eligible for ALT text
                    entry.final_alt = ""
                    entry.llava_called = False
                    excluded_count += 1
                
                # Update legacy compatibility fields
                entry.suggested_alt = entry.final_alt
                entry.source = self._map_decision_to_source(decision['decision_reason'])
                
                manifest.add_entry(entry)
                
                logger.debug(f"Decision for {entry.instance_key}: {decision['action']} - {decision['decision_reason']}")
            
            result = {
                'success': True,
                'total_entries': len(manifest.get_all_entries()),
                'preserved_count': preserved_count,
                'cached_count': cached_count,
                'needs_generation_count': needs_generation_count,
                'excluded_count': excluded_count,
                'policy_decisions': policy_decisions
            }
            
            logger.info(f"Phase 3 complete: {preserved_count} preserved, {cached_count} cached, "
                       f"{needs_generation_count} need generation, {excluded_count} excluded")
            
            return result
            
        except Exception as e:
            logger.error(f"Phase 3 failed: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e),
                'total_entries': 0
            }
    
    def _make_generation_decision(self, entry, manifest: AltManifest, mode: str) -> Dict[str, Any]:
        """
        Make generation decision for a single manifest entry.
        
        Returns:
            Dictionary with action, decision_reason, and other metadata
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        # Step 1: Check preserve mode with existing ALT
        if mode == "preserve" and entry.existing_alt.strip():
            return {
                'action': 'preserve',
                'decision_reason': 'preserve_existing',
                'include_reason': f'existing ALT preserved in {mode} mode',
                'exclude_reason': ''
            }
        
        # Step 2: Check inclusion policy based on --llava-include-shapes
        inclusion_check = self._check_inclusion_policy(entry)
        if not inclusion_check['include']:
            return {
                'action': 'exclude', 
                'decision_reason': 'policy_excluded',
                'include_reason': '',
                'exclude_reason': inclusion_check['reason']
            }
        
        # Step 3: Check cache by content_key (idempotency)
        cached_entry = manifest.get_by_content_key(entry.content_key)
        if cached_entry and cached_entry != entry and cached_entry.final_alt.strip():
            return {
                'action': 'use_cache',
                'decision_reason': 'cached',
                'cached_alt': cached_entry.final_alt,
                'include_reason': f'reused from cache (content_key match)',
                'exclude_reason': ''
            }
        
        # Step 4: Check if this entry already has generated ALT (idempotency within same run)
        if entry.final_alt.strip() and entry.llava_called:
            return {
                'action': 'use_cache',
                'decision_reason': 'already_generated',
                'cached_alt': entry.final_alt,
                'include_reason': 'already processed in this run',
                'exclude_reason': ''
            }
        
        # Step 5: Needs generation
        return {
            'action': 'generate',
            'decision_reason': 'needs_generation',
            'include_reason': inclusion_check['reason'],
            'exclude_reason': ''
        }
    
    def _check_inclusion_policy(self, entry) -> Dict[str, Any]:
        """
        Check if entry should be included based on --llava-include-shapes policy.
        
        Returns:
            Dictionary with include (bool) and reason (str)
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        if self.llava_include_shapes == "off":
            # Only pictures
            if entry.shape_type == "PICTURE":
                return {'include': True, 'reason': 'picture in off mode (pictures only)'}
            else:
                return {'include': False, 'reason': f'{entry.shape_type} excluded in off mode (pictures only)'}
                
        elif self.llava_include_shapes == "all":
            # Every non-decorative shape
            if self._is_decorative_shape(entry):
                return {'include': False, 'reason': f'decorative {entry.shape_type} excluded'}
            else:
                return {'include': True, 'reason': f'{entry.shape_type} included in all mode'}
                
        elif self.llava_include_shapes == "smart":
            # Pictures + non-decorative shapes above threshold
            if entry.shape_type == "PICTURE":
                return {'include': True, 'reason': 'picture in smart mode'}
            
            if self._is_decorative_shape(entry):
                return {'include': False, 'reason': f'decorative {entry.shape_type} excluded'}
            
            # Check if shape meets size threshold
            shape_area = entry.bbox.get('width', 0) * entry.bbox.get('height', 0)
            min_area_threshold = parse_min_shape_area(self.min_shape_area)
            
            if shape_area >= min_area_threshold:
                return {'include': True, 'reason': f'{entry.shape_type} above threshold ({shape_area:.0f} >= {min_area_threshold})'}
            else:
                return {'include': False, 'reason': f'{entry.shape_type} below threshold ({shape_area:.0f} < {min_area_threshold})'}
        
        # Unknown strategy
        return {'include': False, 'reason': f'unknown strategy: {self.llava_include_shapes}'}
    
    def _is_decorative_shape(self, entry) -> bool:
        """
        Determine if a shape is obviously decorative and should be excluded.
        
        Decorative shapes typically include:
        - Simple lines and connectors
        - Very small shapes (likely bullets or decorations)
        - Placeholder shapes
        """
        # Lines and connectors are typically decorative
        if entry.shape_type in ["LINE", "CONNECTOR"]:
            return True
        
        # Very small shapes are likely decorative
        shape_area = entry.bbox.get('width', 0) * entry.bbox.get('height', 0)
        if shape_area < 100:  # Less than 100 square points
            return True
        
        # Check for common decorative patterns in shape names/text
        slide_text = entry.slide_text.lower()
        decorative_patterns = ['bullet', 'decoration', 'border', 'frame', 'divider']
        if any(pattern in slide_text for pattern in decorative_patterns):
            return True
        
        return False
    
    def _map_decision_to_source(self, decision_reason: str) -> str:
        """Map decision_reason to legacy source field."""
        mapping = {
            'preserve_existing': 'existing',
            'cached': 'cached', 
            'already_generated': 'cached',
            'needs_generation': 'generated',  # Will be updated after generation
            'policy_excluded': 'excluded',
            'generated_new': 'generated'
        }
        return mapping.get(decision_reason, 'unknown')
    
    def phase4_single_pass_llava_generation(self, manifest: AltManifest) -> Dict[str, Any]:
        """
        Step 4: Single-pass LLaVA + Normalization (ONLY WHEN NEEDED)
        
        Call LLaVA once per eligible entry; normalize; store results:
        - For entries marked needs_generation
        - Use crop_path as image input  
        - Store llm_raw, normalize, save final_alt
        - Set decision_reason=generated_new, llava_called=true, truncated_flag=false
        """
        logger.info("Phase 4: Single-pass LLaVA generation with normalization")
        
        if not self.alt_generator:
            logger.warning("No ALT generator available for Phase 4")
            return {
                'success': False,
                'error': 'No ALT generator available',
                'generated_count': 0
            }
        
        try:
            generated_count = 0
            skipped_count = 0
            error_count = 0
            generation_errors = []
            
            # Find entries that need generation
            entries_needing_generation = []
            for entry in manifest.get_all_entries():
                if entry.decision_reason == "needs_generation":
                    entries_needing_generation.append(entry)
            
            logger.info(f"Found {len(entries_needing_generation)} entries needing LLaVA generation")
            
            for entry in entries_needing_generation:
                try:
                    # Verify crop_path exists
                    if not entry.crop_path or not Path(entry.crop_path).exists():
                        error_msg = f"Crop image not found for {entry.instance_key}: {entry.crop_path}"
                        logger.warning(error_msg)
                        generation_errors.append(error_msg)
                        error_count += 1
                        
                        # Mark entry as error
                        entry.decision_reason = "generation_error"
                        entry.exclude_reason = "crop image not available"
                        entry.llava_called = False
                        manifest.add_entry(entry)
                        continue
                    
                    # Check content_key cache (idempotency)
                    cached_entry = manifest.get_by_content_key(entry.content_key)
                    if cached_entry and cached_entry != entry and cached_entry.final_alt.strip():
                        # Reuse cached result
                        entry.final_alt = cached_entry.final_alt
                        entry.llm_raw = cached_entry.llm_raw
                        entry.llava_called = False
                        entry.decision_reason = "cached"
                        entry.truncated_flag = cached_entry.truncated_flag
                        manifest.add_entry(entry)
                        skipped_count += 1
                        logger.debug(f"Reused cached result for {entry.instance_key}")
                        continue
                    
                    # Build context for LLaVA prompt
                    context = self._build_generation_context(entry)
                    
                    # Generate ALT text using crop_path
                    logger.debug(f"Generating ALT for {entry.instance_key} using {entry.crop_path}")
                    
                    # Use the flexible generator with manifest integration
                    result = self.alt_generator.generate_alt_text(
                        image_path=entry.crop_path,
                        context=context,
                        prompt_type="default",
                        return_metadata=True,
                        manifest=manifest,
                        entry_key=entry.instance_key
                    )
                    
                    if isinstance(result, tuple):
                        alt_text, metadata = result
                    else:
                        alt_text = result
                        metadata = {}
                    
                    if alt_text and alt_text.strip():
                        # Store results with proper decision tracking
                        entry.llm_raw = alt_text  # Store original before any post-processing
                        entry.final_alt = alt_text  # Normalized by the generator
                        entry.llava_called = True
                        entry.decision_reason = "generated_new"
                        entry.truncated_flag = False  # No truncation since we use complete sentences
                        
                        # Store generation metadata
                        if metadata:
                            entry.duration_ms = metadata.get("generation_time", 0) * 1000
                            entry.provider = metadata.get("successful_provider", "unknown")
                            entry.model_used = metadata.get("successful_model", "unknown")
                        
                        # Update legacy compatibility
                        entry.suggested_alt = entry.final_alt
                        entry.source = "generated"
                        
                        manifest.add_entry(entry)
                        generated_count += 1
                        
                        logger.info(f"✅ Generated ALT for {entry.instance_key}: {alt_text[:50]}...")
                        
                    else:
                        error_msg = f"Empty result from LLaVA for {entry.instance_key}"
                        logger.warning(error_msg)
                        generation_errors.append(error_msg)
                        error_count += 1
                        
                        # Mark as error
                        entry.decision_reason = "generation_failed"
                        entry.exclude_reason = "LLaVA returned empty result"
                        entry.llava_called = True
                        entry.final_alt = ""
                        manifest.add_entry(entry)
                    
                except Exception as e:
                    error_msg = f"Generation failed for {entry.instance_key}: {e}"
                    logger.error(error_msg)
                    generation_errors.append(error_msg)
                    error_count += 1
                    
                    # Mark as error
                    entry.decision_reason = "generation_error"
                    entry.exclude_reason = f"Exception: {str(e)[:100]}"
                    entry.llava_called = False
                    entry.final_alt = ""
                    manifest.add_entry(entry)
            
            result = {
                'success': True,
                'total_candidates': len(entries_needing_generation),
                'generated_count': generated_count,
                'skipped_count': skipped_count,
                'error_count': error_count,
                'generation_errors': generation_errors[:5]  # Limit error list
            }
            
            logger.info(f"Phase 4 complete: {generated_count} generated, {skipped_count} cached, {error_count} errors")
            return result
            
        except Exception as e:
            logger.error(f"Phase 4 failed: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e),
                'generated_count': 0
            }
    
    def _build_generation_context(self, entry) -> str:
        """
        Build context string for ALT text generation.
        
        Includes slide context, shape information, and any existing text.
        """
        context_parts = []
        
        # Slide information
        if entry.slide_number:
            context_parts.append(f"Slide: {entry.slide_number}")
        
        # Shape information
        if entry.shape_type:
            shape_desc = f"Shape: A {entry.shape_type.lower()}"
            if entry.bbox:
                width = int(entry.bbox.get('width', 0))
                height = int(entry.bbox.get('height', 0))
                if width > 0 and height > 0:
                    shape_desc += f" sized {width}x{height} pixels"
            context_parts.append(shape_desc)
        
        # Slide text context
        if entry.slide_text and entry.slide_text.strip():
            # Limit slide text to avoid overwhelming the prompt
            slide_text = entry.slide_text.strip()[:200]
            context_parts.append(f"Slide text: {slide_text}")
        
        # Notes context  
        if hasattr(entry, 'slide_notes') and entry.slide_notes and entry.slide_notes.strip():
            notes = entry.slide_notes.strip()[:100]
            context_parts.append(f"Notes: {notes}")
        
        return "\n".join(context_parts)
        
    def extract_and_generate(self, pptx_path: Path, manifest_path: Path,
                            mode: str = "preserve",
                            generate_thumbnails: bool = True) -> Dict[str, Any]:
        """
        Extract images from PPTX and generate/cache ALT text using manifest.
        
        Args:
            pptx_path: Path to PPTX file
            manifest_path: Path to manifest JSONL file
            mode: "preserve" or "replace" - how to handle existing ALT
            generate_thumbnails: Whether to generate thumbnail images
            
        Returns:
            Processing results with statistics
        """
        logger.info(f"Processing {pptx_path.name} with manifest (mode: {mode})")
        
        # Initialize manifest
        manifest = AltManifest(manifest_path)
        
        # Extract images and existing ALT text
        extraction_result = self._extract_images_to_manifest(
            pptx_path, manifest, generate_thumbnails
        )
        
        if not extraction_result['success']:
            return extraction_result

        # Generate missing ALT text with caching
        generation_result = self._generate_missing_alt_text(
            manifest, mode
        )

        # Finalize decisions for review/injection without requiring PPTX injection
        finalize_result = self._finalize_alt_decisions(manifest, mode)

        # Save manifest
        manifest.save()

        if logger.isEnabledFor(logging.DEBUG):
            entries = manifest.get_all_entries()
            key_sample = [e.key for e in entries[:5]]
            logger.debug("First 5 manifest entry keys: %s", key_sample)

        # Combine results
        stats = manifest.get_statistics()
        
        result = {
            'success': True,
            'manifest_path': str(manifest_path),
            'extraction': extraction_result,
            'generation': generation_result,
            'finalization': finalize_result,
            'statistics': stats,
            'total_entries': stats['total_entries'],
            'llava_calls_made': stats['llava_calls_made'],
            'with_suggested_alt': stats['with_suggested_alt']
        }
        
        logger.info(f"Manifest processing complete: {stats['total_entries']} entries, "
                   f"{stats['llava_calls_made']} LLaVA calls, "
                   f"{stats['with_suggested_alt']} with suggested ALT")
        
        return result

    def _finalize_alt_decisions(self, manifest: AltManifest, mode: str) -> Dict[str, Any]:
        """Ensure manifest entries have finalized ALT decisions for review/doc builds."""

        preserved_count = 0
        generated_count = 0
        cached_count = 0
        missing_count = 0

        for entry in manifest.get_all_entries():
            current_alt = entry.existing_alt.strip()
            suggested_alt = (entry.final_alt or entry.suggested_alt or "").strip()

            if suggested_alt:
                entry.final_alt = suggested_alt
                entry.suggested_alt = suggested_alt

                if entry.decision_reason in ["cached", "generated_new"]:
                    pass
                elif entry.llava_called:
                    entry.decision_reason = "generated_new"
                    generated_count += 1
                elif entry.source == "cached":
                    entry.decision_reason = "cached"
                    cached_count += 1
                else:
                    entry.decision_reason = entry.decision_reason or "generated_new"
                    generated_count += 1

            elif current_alt:
                entry.final_alt = current_alt
                entry.decision_reason = entry.decision_reason or "preserve_existing"
                preserved_count += 1
            else:
                entry.decision_reason = entry.decision_reason or "needs_generation"
                missing_count += 1

            manifest.add_entry(entry)

        return {
            'success': True,
            'preserved_count': preserved_count,
            'generated_count': generated_count,
            'cached_count': cached_count,
            'missing_count': missing_count,
            'mode': mode
        }
    
    def _extract_images_to_manifest(self, pptx_path: Path, manifest: AltManifest,
                                   generate_thumbnails: bool) -> Dict[str, Any]:
        """Extract images from PPTX and populate manifest entries."""
        logger.info("Extracting images and existing ALT text from PPTX")
        
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from PIL import Image
            import io
            
            prs = Presentation(str(pptx_path))
            
            entries_created = 0
            thumbnails_generated = 0
            extraction_errors = []
            
            for slide_idx, slide in enumerate(prs.slides):
                # Get slide text for context
                slide_text = self._extract_slide_text(slide)
                slide_notes = self._extract_slide_notes(slide)
                
                image_count_on_slide = 0
                
                for shape_idx, shape in enumerate(slide.shapes):
                    try:
                        # Classify shape type and group membership
                        shape_type, is_group_child = manifest.classify_shape_type(shape, MSO_SHAPE_TYPE)
                        
                        # Get shape ID
                        shape_id = getattr(shape, 'shape_id', shape_idx)
                        
                        # Extract existing ALT text from PPTX (all shapes can have ALT text)
                        current_alt = self._collect_existing_alt(shape)
                        
                        # Generate hash - for pictures use image data, for others use shape properties
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            # Handle pictures with actual image data
                            if not hasattr(shape, 'image') or not shape.image:
                                continue
                            image_data = shape.image.blob
                            if not image_data:
                                continue
                            image_hash = compute_image_hash(image_data)
                            width_px = shape.image.size[0] if shape.image.size else 0
                            height_px = shape.image.size[1] if shape.image.size else 0
                        else:
                            # Handle non-picture shapes (lines, text boxes, etc.)
                            shape_properties = f"{shape_type}_{shape_id}_{slide_idx}"
                            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                                width_px = int(shape.width.emu // 914400) if shape.width else 0  # Convert to pixels
                                height_px = int(shape.height.emu // 914400) if shape.height else 0
                                shape_properties += f"_{width_px}x{height_px}"
                            else:
                                width_px = height_px = 0
                            image_hash = compute_image_hash(shape_properties.encode('utf-8'))
                        
                        # Create stable key
                        key = create_stable_key(slide_idx, shape_id, image_hash)
                        
                        # Check if we already have this entry
                        existing_entry = manifest.get_entry(key)
                        if existing_entry:
                            logger.debug(f"Entry already exists for {key}, updating context")
                            # Update contextual information but preserve ALT decisions
                            existing_entry.slide_text = slide_text
                            existing_entry.slide_notes = slide_notes
                            if not existing_entry.existing_alt and current_alt:
                                existing_entry.existing_alt = current_alt
                                existing_entry.current_alt = current_alt  # Legacy sync
                            manifest.add_entry(existing_entry)
                            continue
                        
                        image_count_on_slide += 1
                        
                        # Create new manifest entry
                        entry = manifest.create_entry_from_shape(
                            key=key,
                            image_hash=image_hash,
                            slide_idx=slide_idx,
                            image_number=entries_created + 1,
                            current_alt=current_alt,
                            shape_type=shape_type,
                            is_group_child=is_group_child,
                            slide_text=slide_text,
                            slide_notes=slide_notes,
                            width_px=width_px,
                            height_px=height_px
                        )
                        
                        # Generate thumbnail only for picture shapes
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and generate_thumbnails:
                            thumbnail_path = self._generate_thumbnail(
                                image_data, pptx_path, key
                            )
                            if thumbnail_path:
                                entry.thumbnail_path = str(thumbnail_path)
                                thumbnails_generated += 1
                        
                        manifest.add_entry(entry)
                        entries_created += 1
                        
                    except Exception as e:
                        error_msg = f"Error extracting shape at slide {slide_idx}, shape {shape_idx}: {e}"
                        logger.warning(error_msg)
                        extraction_errors.append(error_msg)
            
            logger.info(f"Extracted {entries_created} images, generated {thumbnails_generated} thumbnails")
            
            return {
                'success': True,
                'entries_created': entries_created,
                'thumbnails_generated': thumbnails_generated,
                'extraction_errors': extraction_errors
            }
            
        except Exception as e:
            logger.error(f"Image extraction failed: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e),
                'entries_created': 0
            }
    
    def _collect_existing_alt(self, shape) -> str:
        """Collect existing ALT text for a shape using the shared reader."""
        try:
            alt_text = read_existing_alt(shape)
        except Exception as exc:  # pragma: no cover - defensive logging
            logger.debug(f"Could not extract ALT text: {exc}")
            return ""

        return (alt_text or "").strip()

    def _extract_current_alt_text(self, shape) -> str:
        """Backward-compatible alias for _collect_existing_alt."""
        return self._collect_existing_alt(shape)
    
    def _extract_slide_text(self, slide) -> str:
        """Extract text content from slide for context."""
        text_parts = []
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text.strip())
        except Exception as e:
            logger.debug(f"Could not extract slide text: {e}")
        
        return ' '.join(text_parts)[:500]  # Limit length
    
    def _extract_slide_notes(self, slide) -> str:
        """Extract slide notes for context."""
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_text = ""
                for shape in slide.notes_slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        notes_text += shape.text.strip() + " "
                return notes_text.strip()[:300]  # Limit length
        except Exception as e:
            logger.debug(f"Could not extract slide notes: {e}")
        
        return ""
    
    def _generate_thumbnail(self, image_data: bytes, pptx_path: Path, key: str) -> Optional[Path]:
        """Generate thumbnail image for DOCX review."""
        try:
            from PIL import Image
            import io
            
            # Create thumbnails directory
            thumbs_dir = pptx_path.parent / f".{pptx_path.stem}_thumbs"
            thumbs_dir.mkdir(exist_ok=True)
            
            thumbnail_path = thumbs_dir / f"{key}.jpg"
            
            # Generate thumbnail
            img = Image.open(io.BytesIO(image_data))
            img.thumbnail((200, 200), Image.Resampling.LANCZOS)
            img.save(thumbnail_path, 'JPEG', quality=85)
            
            return thumbnail_path
            
        except Exception as e:
            logger.debug(f"Could not generate thumbnail for {key}: {e}")
            return None
    
    def _generate_missing_alt_text(self, manifest: AltManifest, mode: str) -> Dict[str, Any]:
        """Generate ALT text for entries that need it, using caching and idempotency."""
        logger.info(f"Generating missing ALT text (mode: {mode})")
        
        entries_needing_generation = []
        preserved_count = 0
        cached_count = 0
        
        # Check each entry to see if generation is needed
        for entry in manifest.get_all_entries():
            should_generate, cached_alt = manifest.should_generate_alt(
                entry.key, entry.image_hash, entry.current_alt, mode
            )
            
            if not should_generate:
                # Use existing or cached ALT
                if cached_alt:
                    if entry.current_alt.strip() and mode == "preserve":
                        # Preserve existing
                        manifest.record_generation(
                            entry, entry.current_alt, "existing", False
                        )
                        preserved_count += 1
                        manifest.log_decision(entry.key, mode, "current", "preserve mode with existing ALT")
                    else:
                        # Use cached
                        manifest.record_generation(
                            entry, cached_alt, "cached", False
                        )
                        cached_count += 1
                        manifest.log_decision(entry.key, mode, "cached", "reused from cache")
            else:
                entries_needing_generation.append(entry)
        
        logger.info(f"Found {len(entries_needing_generation)} entries needing generation, "
                   f"{preserved_count} preserved, {cached_count} cached")
        
        # Generate ALT text for entries that need it
        generation_errors = []
        generated_count = 0
        
        for entry in entries_needing_generation:
            try:
                start_time = time.time()
                
                # Check if this shape type should have ALT text generated via LLaVA
                if manifest.should_generate_for_shape_type(entry.shape_type):
                    # Generate ALT text using LLaVA with thumbnail
                    if entry.thumbnail_path and Path(entry.thumbnail_path).exists():
                        alt_text = self.alt_generator.generate_alt_text(
                            entry.thumbnail_path,
                            manifest=manifest,
                            entry_key=entry.key
                        )
                        
                        if alt_text and alt_text.strip():
                            generated_count += 1
                            logger.debug(f"Generated ALT for {entry.key}: {alt_text[:50]}...")
                        else:
                            error_msg = f"Empty ALT text generated for {entry.key}"
                            generation_errors.append(error_msg)
                            logger.warning(error_msg)
                    else:
                        logger.warning(f"No thumbnail available for {entry.key}, skipping LLaVA generation")
                        continue
                else:
                    # Use shape fallback for non-picture elements
                    fallback_alt = manifest.get_shape_fallback_alt(
                        entry.shape_type,
                        entry.is_group_child,
                        entry.width_px,
                        entry.height_px
                    )
                    
                    # Normalize the fallback text
                    normalized_alt, was_truncated = manifest.normalize_alt_text(fallback_alt)
                    
                    # Update entry with fallback ALT
                    entry.llm_raw = fallback_alt
                    entry.final_alt = normalized_alt
                    entry.truncated_flag = was_truncated
                    entry.llava_called = False
                    entry.decision_reason = "shape_fallback"
                    entry.source = "shape_fallback"
                    
                    # Update legacy fields for compatibility
                    entry.suggested_alt = normalized_alt
                    
                    manifest.add_entry(entry)
                    generated_count += 1
                    
                    logger.debug(f"Applied shape fallback for {entry.key}: {normalized_alt}")
                    
                duration_ms = int((time.time() - start_time) * 1000)
                manifest.log_decision(
                    entry.key, 
                    mode, 
                    entry.final_alt,
                    f"{'LLaVA' if entry.llava_called else 'shape_fallback'} in {duration_ms}ms"
                )
                    
            except Exception as e:
                error_msg = f"Generation failed for {entry.key}: {e}"
                generation_errors.append(error_msg)
                logger.warning(error_msg)
        
        return {
            'success': True,
            'generated_count': generated_count,
            'preserved_count': preserved_count,
            'cached_count': cached_count,
            'error_count': len(generation_errors),
            'errors': generation_errors[:10]  # Limit error list
        }