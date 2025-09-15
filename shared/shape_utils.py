"""
Shape utilities for PPTX processing with enhanced vector/group detection.
Provides robust detection of image-like shapes including vectors and groups.
"""

from typing import Optional, Union, Any
import logging

logger = logging.getLogger(__name__)


def is_image_like(shape) -> bool:
    """
    Robust detector for image-like shapes in PPTX including vectors/groups.
    
    Treats the following as image-like:
    - Regular pictures (MSO_SHAPE_TYPE.PICTURE)
    - Vector shapes with picture fill (AUTO_SHAPE with picture fill)
    - Group shapes containing any picture children
    - Shapes with embedded blips in XML (WMF/EMF vectors)
    
    Args:
        shape: python-pptx shape object
        
    Returns:
        True if shape should be treated as image-like for ALT text purposes
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.warning("python-pptx not available, assuming not image-like")
        return False
    
    if not hasattr(shape, 'shape_type'):
        return False
    
    shape_type = getattr(shape, 'shape_type', None)
    
    # Regular pictures are always image-like
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    
    # Vector shapes with picture fill
    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            if hasattr(shape, 'fill') and hasattr(shape.fill, 'type'):
                # Picture fill type is typically 5 in python-pptx
                if getattr(shape.fill, 'type', None) == 5:
                    return True
        except Exception as e:
            logger.debug(f"Error checking AUTO_SHAPE fill: {e}")
    
    # Group shapes containing any picture children
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            if hasattr(shape, 'shapes'):
                for child_shape in shape.shapes:
                    child_type = getattr(child_shape, 'shape_type', None)
                    if child_type == MSO_SHAPE_TYPE.PICTURE:
                        return True
                    # Recursively check nested groups
                    if child_type == MSO_SHAPE_TYPE.GROUP and is_image_like(child_shape):
                        return True
        except Exception as e:
            logger.debug(f"Error checking GROUP children: {e}")
            return False
    
    # Fallback XML inspection for embedded blips (WMF/EMF vectors)
    try:
        element = getattr(shape, '_element', None)
        if element is not None:
            # Look for blip elements indicating embedded images/vectors
            blips = _safe_xpath(element, ".//a:blip")
            if blips:
                logger.debug(f"Found {len(blips)} blip(s) in shape, treating as image-like")
                return True
    except Exception as e:
        logger.debug(f"Error in XML blip inspection: {e}")
    
    return False


def is_low_value_alt(alt_text: str) -> bool:
    """
    Determine if ALT text is low-value boilerplate that should be overwritten.
    
    Args:
        alt_text: The ALT text to evaluate
        
    Returns:
        True if the ALT text is low-value and should be replaced
    """
    if not alt_text or not alt_text.strip():
        return True
    
    text = alt_text.strip()
    
    # Check for common boilerplate patterns
    import re
    low_value_patterns = [
        r"^\s*This is a PowerPoint shape\b",
        r"^\s*Image of\b", 
        r"^\s*Picture\b",
        r"^\s*Graphic\b",
        r"^\s*unknown\b",
        r"^\s*Shape\b",
        r"^\s*Object\b",
    ]
    
    for pattern in low_value_patterns:
        if re.search(pattern, text, flags=re.IGNORECASE):
            return True
    
    # Require terminal punctuation for sentences longer than 6 words
    words = text.split()
    if len(words) > 6 and not text.endswith(('.', '!', '?')):
        return True
    
    # Very short descriptions are often not meaningful
    if len(text) < 5:
        return True
    
    return False


def should_overwrite_alt(existing_alt: str) -> bool:
    """
    Determine if existing ALT text should be overwritten with generated text.
    
    Args:
        existing_alt: The current ALT text
        
    Returns:
        True if the existing ALT should be overwritten
    """
    return is_low_value_alt(existing_alt)


def is_decorative_shape(shape) -> bool:
    """
    Determine if a shape is purely decorative (lines, borders, empty rectangles).
    
    Args:
        shape: python-pptx shape object
        
    Returns:
        True if shape should be marked as decorative
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return False
    
    shape_type = getattr(shape, 'shape_type', None)
    
    # Lines and connectors are typically decorative
    if shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.CONNECTOR):
        return True
    
    # Check for empty rectangles/shapes (no fill, just outline)
    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            # Check if shape has no meaningful fill
            if hasattr(shape, 'fill'):
                # No fill or transparent fill suggests decorative use
                if not hasattr(shape.fill, 'fore_color') or getattr(shape.fill, 'transparency', 0) > 0.8:
                    # Also check if it's very thin (like a border)
                    if hasattr(shape, 'width') and hasattr(shape, 'height'):
                        width_px = getattr(shape, 'width', 0)
                        height_px = getattr(shape, 'height', 0) 
                        if width_px > 0 and height_px > 0:
                            # Very thin shapes (aspect ratio > 10:1) are likely decorative lines
                            aspect_ratio = max(width_px / height_px, height_px / width_px)
                            if aspect_ratio > 10:
                                return True
        except Exception as e:
            logger.debug(f"Error checking decorative properties: {e}")
    
    return False


def _safe_xpath(element, xpath_expr, namespaces=None):
    """Execute an XPath on python-pptx BaseOxmlElement or plain lxml element."""
    try:
        from pptx.oxml.ns import nsmap as PPTX_NSMAP
    except Exception:
        PPTX_NSMAP = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        }
    
    el = getattr(element, "_element", element)
    try:
        return el.xpath(xpath_expr)  # python-pptx injects namespaces
    except Exception:
        ns = namespaces or getattr(el, "nsmap", None) or PPTX_NSMAP
        return el.xpath(xpath_expr, namespaces=ns)