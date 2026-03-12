"""
Disclosure slide insertion for VisualText Pro.

Ensures every processed PPTX has a standard "Disclosures" slide at position 2
(index 1), or at position 1 (index 0) if the deck is empty. Uses the
presentation's theme/layouts; does not insert if a disclosure slide already
exists anywhere in the deck.

Expected behavior when a disclosure exists elsewhere: we do not move it and do
not insert a second one. "Already present anywhere" means no insertion.
"""

from __future__ import annotations

import logging
from typing import Any

logger = logging.getLogger(__name__)

# Disclosure content (final text per plan)
DISCLOSURE_TITLE = "Disclosures"
DISCLOSURE_BODY_LINES = [
    ("Financial:", 0),
    ("No financial disclosures to report.", 1),
    ("AI Disclosures:", 0),
    (
        "An AI agent was utilized to add alt text to this presentation to meet "
        "Federal Accessibility Requirements.",
        1,
    ),
]

# Substrings (case-insensitive) that identify an existing disclosure slide
DISCLOSURE_TITLE_HINT = "disclosure"  # "disclosure" or "disclosures"
DISCLOSURE_BODY_HINTS = (
    "no financial disclosures",
    "federal accessibility requirements",
)


def _get_all_slide_text(slide: Any) -> str:
    """Extract all text from a slide by concatenating text from every shape."""
    parts = []
    try:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
    except Exception as e:
        logger.debug("Error extracting slide text: %s", e)
    return " ".join(parts)


def _slide_matches_disclosure(slide: Any) -> bool:
    """
    Return True if the slide looks like a disclosure slide (title + body hints).
    Uses all-slide text, not only shapes.title.
    """
    text = _get_all_slide_text(slide).lower()
    if not text:
        return False
    has_title_hint = DISCLOSURE_TITLE_HINT in text
    has_body_hint = any(hint in text for hint in DISCLOSURE_BODY_HINTS)
    return bool(has_title_hint and has_body_hint)


def _has_disclosure_slide_position_aware(presentation: Any) -> bool:
    """
    Check if a disclosure slide already exists (position-aware).
    First check slide at index 1; if it matches, return True. Else scan all
    slides; if any match, return True. Otherwise return False.
    """
    slides = presentation.slides
    if len(slides) > 1:
        if _slide_matches_disclosure(slides[1]):
            return True
    for slide in slides:
        if _slide_matches_disclosure(slide):
            return True
    return False


def _layout_has_title_and_body(layout: Any) -> bool:
    """Return True if layout has both title and body placeholders."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        has_title = False
        has_body = False
        for shape in layout.placeholders:
            if not getattr(shape, "is_placeholder", True):
                continue
            ph = getattr(shape, "placeholder_format", None)
            if ph is None:
                continue
            ph_type = getattr(ph, "type", None)
            if ph_type is None:
                continue
            if ph_type == PP_PLACEHOLDER.TITLE:
                has_title = True
            elif ph_type == PP_PLACEHOLDER.BODY:
                has_body = True
            if has_title and has_body:
                return True
        return has_title and has_body
    except Exception as e:
        logger.debug("Error checking layout placeholders: %s", e)
        return False


def _layout_has_title_only(layout: Any) -> bool:
    """Return True if layout has at least a title placeholder."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        for shape in layout.placeholders:
            if not getattr(shape, "is_placeholder", True):
                continue
            ph = getattr(shape, "placeholder_format", None)
            if ph is None:
                continue
            if getattr(ph, "type", None) == PP_PLACEHOLDER.TITLE:
                return True
        return False
    except Exception as e:
        logger.debug("Error checking layout for title: %s", e)
        return False


def _get_title_and_body_layout(presentation: Any) -> Any | None:
    """
    Pick a slide layout that has both title and body placeholders.
    Fallback: first layout with at least a title placeholder.
    """
    for layout in presentation.slide_layouts:
        if _layout_has_title_and_body(layout):
            return layout
    for layout in presentation.slide_layouts:
        if _layout_has_title_only(layout):
            return layout
    return None


def _move_slide_to_index(presentation: Any, from_index: int, to_index: int) -> None:
    """
    Move the slide at from_index to to_index by reordering _sldIdLst.
    Uses python-pptx internal API; regression tests verify slide order.
    """
    sld_id_lst = presentation.slides._sldIdLst
    slides_list = list(sld_id_lst)
    el = slides_list[from_index]
    sld_id_lst.remove(el)
    sld_id_lst.insert(to_index, el)


def _set_disclosure_content(slide: Any) -> None:
    """Set title and body placeholders on the disclosure slide."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        for shape in slide.shapes:
            if not getattr(shape, "is_placeholder", True):
                continue
            ph = getattr(shape, "placeholder_format", None)
            if ph is None:
                continue
            ph_type = getattr(ph, "type", None)
            if ph_type == PP_PLACEHOLDER.TITLE:
                shape.text = DISCLOSURE_TITLE
            elif ph_type == PP_PLACEHOLDER.BODY:
                tf = getattr(shape, "text_frame", None)
                if tf is not None:
                    paras = tf.paragraphs
                    for i, (line, level) in enumerate(DISCLOSURE_BODY_LINES):
                        if i == 0 and paras:
                            p = paras[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = line
                        p.level = level
                break
    except Exception as e:
        logger.warning("Error setting disclosure content: %s", e)


def ensure_disclosure_slide(presentation: Any) -> None:
    """
    Ensure the presentation has a Disclosures slide at the required position.

    - If the deck has 0 slides: create the disclosure as slide 0 (index 0).
    - If the deck has 1+ slides: insert the disclosure at index 1 (second slide).

    If a disclosure slide already exists anywhere in the deck (identified by
    title/body text), we do not move it and do not insert a second one.

    Uses the presentation's theme/layouts; selects a layout by presence of
    title and body placeholders, not by fixed index.
    """
    if _has_disclosure_slide_position_aware(presentation):
        logger.debug("Disclosure slide already present; skipping insertion")
        return

    layout = _get_title_and_body_layout(presentation)
    if layout is None:
        logger.warning("No suitable layout for disclosure slide; skipping")
        return

    slides = presentation.slides
    n = len(slides)

    if n == 0:
        target_index = 0
    else:
        target_index = 1

    new_slide = slides.add_slide(layout)
    _set_disclosure_content(new_slide)

    if n > 0:
        new_slide_index = n
        _move_slide_to_index(presentation, new_slide_index, target_index)
        logger.debug("Disclosure slide inserted at index %s", target_index)
    else:
        logger.debug("Disclosure slide created at index 0 (empty deck)")
