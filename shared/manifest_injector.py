#!/usr/bin/env python3
"""
Manifest-Based PPTX ALT Text Injector
=====================================

Injects ALT text into PPTX files reading only from the manifest.
No LLaVA calls - all ALT text comes from the manifest SSOT.
"""

from __future__ import annotations
import logging
from pathlib import Path
from typing import Dict, Any, Literal

from alt_manifest import AltManifest

logger = logging.getLogger(__name__)


def inject_from_manifest(
    pptx_path: str,
    manifest_path: str,
    output_path: str = None,
    mode: Literal["preserve", "replace"] = "preserve",
    run_id: str | None = None,
) -> Dict[str, Any]:
    """
    Inject ALT text from manifest into PPTX file.
    
    Args:
        pptx_path: Path to input PPTX file
        manifest_path: Path to manifest JSONL file  
        output_path: Path for output PPTX (default: overwrite input)
        mode: "preserve" (keep existing) or "replace" (overwrite existing)
        run_id: Optional identifier for tracking pipeline runs
        
    Returns:
        Injection results with statistics and decision logging
    """
    logger.info("RUN_ID=%s manifest=%s", run_id, manifest_path)
    logger.info("Injecting ALT text from manifest into %s (mode: %s)", pptx_path, mode)
    
    if output_path is None:
        output_path = pptx_path
    
    try:
        # Load manifest
        manifest = AltManifest(Path(manifest_path))
        entries = manifest.get_all_entries()
        
        if not entries:
            logger.warning("No entries found in manifest")
            return {
                'success': True,
                'total_entries': 0,
                'injected_successfully': 0,
                'skipped_existing': 0,
                'errors': []
            }
        
        logger.info(f"Found {len(entries)} entries in manifest")
        
        # Perform injection using existing robust injector
        result = _inject_using_robust_injector(pptx_path, entries, output_path, mode, manifest)
        
        return result
        
    except Exception as e:
        logger.error(f"Manifest-based injection failed: {e}", exc_info=True)
        return {
            'success': False,
            'error': str(e),
            'total_entries': 0,
            'injected_successfully': 0
        }


def _inject_using_robust_injector(pptx_path: str, entries, output_path: str, 
                                 mode: str, manifest: AltManifest) -> Dict[str, Any]:
    """Use the existing robust injector with manifest data."""
    try:
        # Convert manifest entries to the mapping format expected by injector
        # Use final_alt as primary source, with suggested_alt as fallback
        alt_text_mapping = {}
        decision_log = {}
        
        for entry in entries:
            # Skip entries marked as preserve_existing in decision engine
            if mode == "preserve" and entry.decision_reason == "preserve_existing":
                logger.debug(f"Skipping {entry.instance_key}: preserve_existing decision")
                decision_log[entry.instance_key] = {
                    'alt_used': entry.existing_alt,
                    'decision_reason': 'skipped_preserve_existing',
                    'shape_type': entry.shape_type,
                    'is_group_child': entry.is_group_child,
                    'had_existing_alt': entry.had_existing_alt,
                    'llava_called': entry.llava_called,
                    'injected': False
                }
                continue
            
            # Determine final ALT text to inject using new schema
            alt_to_inject = ""
            decision_reason = ""
            
            if entry.final_alt and entry.final_alt.strip():
                # Use final_alt (our normalized, processed text) - NEW SCHEMA 2.0
                # GREMLIN 1 FIX: Don't strip final_alt - preserve exact punctuation
                alt_to_inject = entry.final_alt
                decision_reason = f"injected_final_alt_from_{entry.decision_reason}"
            elif entry.suggested_alt and entry.suggested_alt.strip():
                # Fallback to legacy suggested_alt field
                # GREMLIN 1 FIX: Don't strip suggested_alt - preserve exact punctuation
                alt_to_inject = entry.suggested_alt
                decision_reason = f"fallback_to_suggested_alt_from_{entry.source}"
            elif mode == "replace" and entry.existing_alt and entry.existing_alt.strip():
                # Replace mode but no new ALT generated - keep existing
                # GREMLIN 1 FIX: Don't strip existing_alt - preserve exact punctuation
                alt_to_inject = entry.existing_alt
                decision_reason = "replace_mode_but_no_new_alt"
            else:
                # No ALT text available - skip this entry
                logger.debug(f"No ALT text available for {entry.instance_key}, skipping")
                decision_log[entry.instance_key] = {
                    'alt_used': '',
                    'decision_reason': 'no_alt_available',
                    'shape_type': entry.shape_type,
                    'is_group_child': entry.is_group_child,
                    'had_existing_alt': entry.had_existing_alt,
                    'llava_called': entry.llava_called,
                    'injected': False
                }
                continue
            
            if alt_to_inject:
                # Use instance_key for mapping (NEW SCHEMA 2.0)
                lookup_key = entry.instance_key if entry.instance_key else entry.key
                alt_text_mapping[lookup_key] = alt_to_inject
                decision_log[lookup_key] = {
                    'alt_used': alt_to_inject,
                    'decision_reason': decision_reason,
                    'shape_type': entry.shape_type,
                    'is_group_child': entry.is_group_child,
                    'had_existing_alt': entry.had_existing_alt,
                    'llava_called': entry.llava_called,
                    'injected': True
                }
        
        logger.info(
            f"Prepared {len(alt_text_mapping)} ALT text mappings for injection"
        )

        if logger.isEnabledFor(logging.DEBUG):
            sample_keys = list(alt_text_mapping.keys())[:5]
            logger.debug("First 5 injection keys: %s", sample_keys)

        # Use existing robust injector
        from core.pptx_alt_injector import PPTXAltTextInjector
        from shared.config_manager import ConfigManager
        
        config_manager = ConfigManager()
        injector = PPTXAltTextInjector(config_manager)
        
        # Perform injection with decision logging
        result = injector.inject_alt_text_from_mapping(
            pptx_path,
            alt_text_mapping,
            output_path,
            mode=mode
        )
        
        # Add decision logging and post-inject verification for traceability
        if result['success']:
            stats = result.get('statistics', {})
            
            # Post-inject verification: read PPT descr and assert equality with final_alt
            logger.info("Starting post-inject verification...")
            verification_result = _post_inject_verification(output_path, alt_text_mapping, manifest)
            result['verification'] = verification_result
            
            # GREMLIN FIXES: Log detailed decisions for each element with punctuation info
            for key, decision_info in decision_log.items():
                alt_used = decision_info['alt_used']
                
                # Log punctuation info for injected ALT text
                punct_info = ""
                if alt_used and decision_info['injected']:
                    ends_with = alt_used[-1] if alt_used else ''
                    has_punct = ends_with in '.!?'
                    punct_info = f" | endswith: '{ends_with}' | has_punct: {has_punct} | len: {len(alt_used)}"
                
                manifest.log_decision(
                    key, 
                    mode, 
                    alt_used, 
                    f"{decision_info['decision_reason']} | shape: {decision_info['shape_type']} | "
                    f"group_child: {decision_info['is_group_child']} | "
                    f"had_existing: {decision_info['had_existing_alt']} | "
                    f"llava_called: {decision_info['llava_called']} | "
                    f"injected: {decision_info['injected']}{punct_info}"
                )
            
            # Log summary statistics
            shape_types = {}
            injected_count = 0
            skipped_count = 0
            
            for decision_info in decision_log.values():
                shape_type = decision_info['shape_type']
                shape_types[shape_type] = shape_types.get(shape_type, 0) + 1
                
                if decision_info['injected']:
                    injected_count += 1
                else:
                    skipped_count += 1
            
            logger.info(f"Injection completed: {injected_count} elements injected, {skipped_count} skipped")
            logger.info(f"Shape type distribution: {shape_types}")
            logger.info(f"Verification: {verification_result.get('verified_count', 0)} verified, "
                       f"{len(verification_result.get('mismatches', []))} mismatches")
            
        return result
        
    except Exception as e:
        logger.error(f"Robust injector failed: {e}", exc_info=True)
        return {
            'success': False,
            'error': str(e),
            'total_entries': len(entries),
            'injected_successfully': 0
        }


def _post_inject_verification(pptx_path: str, alt_text_mapping: Dict[str, str], 
                            manifest: AltManifest) -> Dict[str, Any]:
    """
    Post-inject verification: read PPT descr and assert equality with final_alt.
    
    Args:
        pptx_path: Path to the injected PPTX file
        alt_text_mapping: Mapping of instance_key -> injected ALT text
        manifest: The manifest for context
        
    Returns:
        Verification results with detailed mismatches
    """
    try:
        from pptx import Presentation
        from shared.alt_manifest import create_instance_key
        
        prs = Presentation(pptx_path)
        
        verified_count = 0
        mismatches = []
        verification_errors = []
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                try:
                    # Get shape ID and create instance key
                    shape_id = getattr(shape, 'shape_id', shape_idx)
                    instance_key = create_instance_key(slide_idx, shape_id)
                    
                    # Check if we injected ALT text for this shape
                    expected_alt = alt_text_mapping.get(instance_key)
                    if expected_alt is None:
                        continue  # We didn't inject anything for this shape
                    
                    # Read actual ALT text from PowerPoint
                    actual_alt = ""
                    try:
                        # Extract ALT text from shape
                        if hasattr(shape, '_element'):
                            pic_element = shape._element
                            nvpicpr = pic_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
                            if nvpicpr is not None:
                                actual_alt = nvpicpr.get('descr', '') or nvpicpr.get('title', '')
                    except Exception as e:
                        verification_errors.append(f"Could not read ALT from {instance_key}: {e}")
                        continue
                    
                    # Compare expected vs actual
                    expected_alt = expected_alt.strip()
                    actual_alt = (actual_alt or "").strip()
                    
                    if expected_alt == actual_alt:
                        verified_count += 1
                    else:
                        mismatches.append({
                            'instance_key': instance_key,
                            'slide_idx': slide_idx,
                            'shape_id': shape_id,
                            'expected_alt': expected_alt,
                            'actual_alt': actual_alt,
                            'length_diff': len(actual_alt) - len(expected_alt)
                        })
                        
                except Exception as e:
                    verification_errors.append(f"Error verifying shape at slide {slide_idx}, shape {shape_idx}: {e}")
        
        success = len(mismatches) == 0
        
        result = {
            'success': success,
            'verified_count': verified_count,
            'mismatch_count': len(mismatches),
            'mismatches': mismatches,
            'verification_errors': verification_errors[:5],  # Limit error list
            'total_checked': verified_count + len(mismatches)
        }
        
        if success:
            logger.info(f"✅ Post-inject verification passed: {verified_count} elements verified")
        else:
            logger.warning(f"⚠️  Post-inject verification found {len(mismatches)} mismatches:")
            for mismatch in mismatches[:3]:  # Show first 3 mismatches
                logger.warning(f"  {mismatch['instance_key']}: expected '{mismatch['expected_alt'][:50]}...', "
                             f"got '{mismatch['actual_alt'][:50]}...'")
        
        return result
        
    except Exception as e:
        logger.error(f"Post-inject verification failed: {e}", exc_info=True)
        return {
            'success': False,
            'error': str(e),
            'verified_count': 0,
            'mismatch_count': 0,
            'mismatches': []
        }


def validate_manifest_for_injection(manifest_path: str) -> Dict[str, Any]:
    """
    Validate that manifest has the data needed for injection.
    
    Returns:
        Validation results with statistics
    """
    try:
        manifest = AltManifest(Path(manifest_path))
        entries = manifest.get_all_entries()
        
        if not entries:
            return {
                'valid': False,
                'error': 'No entries found in manifest',
                'total_entries': 0
            }
        
        injectable_count = 0
        missing_alt_text = 0
        by_shape_type = {}
        by_decision_reason = {}
        
        for entry in entries:
            # Count by shape type
            shape_type = entry.shape_type
            by_shape_type[shape_type] = by_shape_type.get(shape_type, 0) + 1
            
            # Count by decision reason
            decision_reason = entry.decision_reason or entry.source
            by_decision_reason[decision_reason] = by_decision_reason.get(decision_reason, 0) + 1
            
            # Check if entry has injectable ALT text
            has_alt = (entry.final_alt and entry.final_alt.strip()) or \
                     (entry.suggested_alt and entry.suggested_alt.strip()) or \
                     (entry.existing_alt and entry.existing_alt.strip())
                     
            if has_alt:
                injectable_count += 1
            else:
                missing_alt_text += 1
        
        stats = manifest.get_statistics()
        
        return {
            'valid': True,
            'total_entries': len(entries),
            'injectable_entries': injectable_count,
            'missing_alt_text': missing_alt_text,
            'by_shape_type': by_shape_type,
            'by_decision_reason': by_decision_reason,
            'llava_calls_in_manifest': stats['llava_calls_made'],
            'manifest_statistics': stats
        }
        
    except Exception as e:
        return {
            'valid': False,
            'error': str(e),
            'total_entries': 0
        }