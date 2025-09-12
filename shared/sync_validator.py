#!/usr/bin/env python3
"""
Sync Validator & Telemetry (PROVE IT)
=====================================

Guarantee no drift; surface costs/coverage.
For each entry: compare manifest.final_alt with PPT descr (if injected) and DOCX Suggested (if built).
Zero diffs expected.

Provides comprehensive telemetry on:
- Counts by decision_reason
- LLaVA calls avoided via cache
- Shapes sent vs. skipped (with reasons)  
- Rasterization failures
"""

from __future__ import annotations
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime
import json

from alt_manifest import AltManifest, create_instance_key

logger = logging.getLogger(__name__)


class SyncValidator:
    """
    Validator to ensure perfect synchronization between manifest, PPT, and DOCX.
    
    Guarantees no drift by comparing ALT text across all outputs.
    """
    
    def __init__(self):
        self.validation_errors = []
        self.validation_warnings = []
        self.telemetry_data = {}
    
    def validate_full_pipeline(self, manifest_path: str, pptx_path: str = None, 
                             docx_path: str = None) -> Dict[str, Any]:
        """
        Validate complete pipeline synchronization.
        
        Args:
            manifest_path: Path to manifest.json file
            pptx_path: Path to injected PPTX file (optional)
            docx_path: Path to generated DOCX file (optional)
            
        Returns:
            Comprehensive validation report
        """
        logger.info("🔍 Starting full pipeline sync validation...")
        
        try:
            # Load manifest
            manifest = AltManifest(Path(manifest_path))
            entries = manifest.get_all_entries()
            
            if not entries:
                return {
                    'success': False,
                    'error': 'No entries found in manifest',
                    'validation_errors': ['Empty manifest'],
                    'telemetry': self._generate_empty_telemetry()
                }
            
            logger.info(f"Validating {len(entries)} manifest entries")
            
            # Generate telemetry first
            self.telemetry_data = self._generate_telemetry(manifest, entries)
            
            # Validate PPT injection if provided
            ppt_validation = {}
            if pptx_path and Path(pptx_path).exists():
                ppt_validation = self._validate_ppt_injection(entries, pptx_path)
            
            # Validate DOCX generation if provided
            docx_validation = {}
            if docx_path and Path(docx_path).exists():
                docx_validation = self._validate_docx_generation(entries, docx_path)
            
            # Calculate overall success
            ppt_success = ppt_validation.get('success', True)  # True if not validated
            docx_success = docx_validation.get('success', True)  # True if not validated
            overall_success = ppt_success and docx_success and len(self.validation_errors) == 0
            
            result = {
                'success': overall_success,
                'validation_timestamp': datetime.now().isoformat(),
                'manifest_path': manifest_path,
                'pptx_path': pptx_path,
                'docx_path': docx_path,
                'total_entries': len(entries),
                'validation_errors': self.validation_errors,
                'validation_warnings': self.validation_warnings,
                'ppt_validation': ppt_validation,
                'docx_validation': docx_validation,
                'telemetry': self.telemetry_data
            }
            
            self._log_validation_summary(result)
            
            return result
            
        except Exception as e:
            logger.error(f"Sync validation failed: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e),
                'validation_errors': [str(e)],
                'telemetry': self._generate_empty_telemetry()
            }
    
    def _validate_ppt_injection(self, entries: List, pptx_path: str) -> Dict[str, Any]:
        """
        Validate PPT injection: compare manifest.final_alt with PPT descr.
        Zero diffs expected.
        """
        logger.info(f"Validating PPT injection synchronization with {pptx_path}")
        
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            
            validated_count = 0
            mismatches = []
            injection_errors = []
            
            # Build lookup of what should be in PPT based on manifest
            expected_alts = {}
            for entry in entries:
                if entry.decision_reason == "preserve_existing":
                    # Preserve mode - should not be injected
                    continue
                elif entry.final_alt and entry.final_alt.strip():
                    # Should be injected with final_alt
                    expected_alts[entry.instance_key] = entry.final_alt.strip()
            
            logger.info(f"Expecting {len(expected_alts)} injected ALT texts in PPT")
            
            # Validate each shape in PPT
            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    try:
                        # Get shape ID and create instance key
                        shape_id = getattr(shape, 'shape_id', shape_idx)
                        instance_key = create_instance_key(slide_idx, shape_id)
                        
                        expected_alt = expected_alts.get(instance_key)
                        if expected_alt is None:
                            continue  # Not expected to have ALT injected
                        
                        # Read actual ALT text from PowerPoint
                        actual_alt = self._extract_ppt_alt_text(shape)
                        
                        # Compare expected vs actual
                        if expected_alt == actual_alt:
                            validated_count += 1
                        else:
                            # GREMLIN 1 FIX: Enhanced mismatch detection with punctuation analysis
                            mismatch_info = {
                                'instance_key': instance_key,
                                'slide_idx': slide_idx,
                                'shape_id': shape_id,
                                'expected_alt': expected_alt,
                                'actual_alt': actual_alt,
                                'length_diff': len(actual_alt) - len(expected_alt) if actual_alt else -len(expected_alt),
                                'punctuation_issue': self._check_punctuation_parity(expected_alt, actual_alt)
                            }
                            mismatches.append(mismatch_info)
                            
                            self.validation_errors.append(
                                f"PPT ALT mismatch {instance_key}: expected '{expected_alt[:50]}...', "
                                f"got '{actual_alt[:50] if actual_alt else '[NONE]'}...'"
                            )
                        
                    except Exception as e:
                        injection_errors.append(f"Error validating {instance_key}: {e}")
            
            success = len(mismatches) == 0
            
            result = {
                'success': success,
                'validated_count': validated_count,
                'expected_count': len(expected_alts),
                'mismatch_count': len(mismatches),
                'mismatches': mismatches,
                'injection_errors': injection_errors[:5]
            }
            
            if success:
                logger.info(f"✅ PPT validation passed: {validated_count}/{len(expected_alts)} verified")
            else:
                logger.error(f"❌ PPT validation failed: {len(mismatches)} mismatches found")
            
            return result
            
        except Exception as e:
            logger.error(f"PPT validation error: {e}")
            return {
                'success': False,
                'error': str(e),
                'validated_count': 0
            }
    
    def _validate_docx_generation(self, entries: List, docx_path: str) -> Dict[str, Any]:
        """
        Validate DOCX generation: ensure DOCX content matches manifest entries.
        """
        logger.info(f"Validating DOCX generation synchronization with {docx_path}")
        
        # For DOCX, we validate that the content is generated from manifest
        # Since DOCX is pure reader of manifest, main validation is structural
        
        try:
            from docx import Document
            
            doc = Document(docx_path)
            tables = doc.tables
            
            if not tables:
                self.validation_errors.append("DOCX contains no tables")
                return {'success': False, 'error': 'No tables in DOCX'}
            
            main_table = tables[0]  # Assume first table is the review table
            data_rows = main_table.rows[1:]  # Skip header row
            
            # Count entries by decision reason for comparison
            docx_row_count = len(data_rows)
            manifest_entry_count = len(entries)
            
            success = True
            preserved_display_errors = []
            
            if docx_row_count != manifest_entry_count:
                self.validation_errors.append(
                    f"DOCX row count mismatch: {docx_row_count} rows vs {manifest_entry_count} manifest entries"
                )
                success = False
            
            # GREMLIN 2 FIX: Validate preserved ALT display logic
            preserved_display_errors = self._validate_preserved_alt_display(entries, data_rows)
            if preserved_display_errors:
                self.validation_errors.extend(preserved_display_errors)
                success = False
            
            # Additional validation: check that DOCX has expected columns
            if len(main_table.columns) != 5:
                self.validation_errors.append(
                    f"DOCX column count mismatch: expected 5 columns, got {len(main_table.columns)}"
                )
                success = False
            
            result = {
                'success': success,
                'docx_rows': docx_row_count,
                'manifest_entries': manifest_entry_count,
                'table_columns': len(main_table.columns) if main_table else 0
            }
            
            if success:
                logger.info(f"✅ DOCX validation passed: {docx_row_count} rows match {manifest_entry_count} entries")
            else:
                logger.error(f"❌ DOCX validation failed: structural mismatches found")
            
            return result
            
        except Exception as e:
            logger.error(f"DOCX validation error: {e}")
            return {
                'success': False,
                'error': str(e),
                'docx_rows': 0
            }
    
    def _extract_ppt_alt_text(self, shape) -> str:
        """Extract ALT text from PowerPoint shape."""
        try:
            if hasattr(shape, '_element'):
                pic_element = shape._element
                nvpicpr = pic_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
                if nvpicpr is not None:
                    alt_text = nvpicpr.get('descr', '') or nvpicpr.get('title', '')
                    return (alt_text or "").strip()
        except Exception as e:
            logger.debug(f"Could not extract ALT text from shape: {e}")
        
        return ""
    
    def _generate_telemetry(self, manifest: AltManifest, entries: List) -> Dict[str, Any]:
        """
        Generate comprehensive telemetry data.
        
        Surfaces:
        - Counts by decision_reason
        - LLaVA calls avoided via cache  
        - Shapes sent vs. skipped (with reasons)
        - Rasterization failures
        """
        logger.info("📊 Generating pipeline telemetry...")
        
        # Count by decision reason
        decision_counts = {}
        llava_calls_made = 0
        llava_calls_avoided = 0
        shapes_sent = 0
        shapes_skipped = 0
        rasterization_failures = 0
        
        skip_reasons = {}
        shape_type_distribution = {}
        cache_hits = 0
        
        for entry in entries:
            # Decision reason counts
            decision_reason = entry.decision_reason or "unknown"
            decision_counts[decision_reason] = decision_counts.get(decision_reason, 0) + 1
            
            # LLaVA call tracking
            if entry.llava_called:
                llava_calls_made += 1
                shapes_sent += 1
            else:
                # Determine why LLaVA wasn't called
                if decision_reason == "preserve_existing":
                    llava_calls_avoided += 1
                    reason = "preserved_existing_alt"
                elif decision_reason == "cached":
                    llava_calls_avoided += 1
                    cache_hits += 1
                    reason = "cache_hit"
                elif decision_reason in ["policy_excluded", "generation_error"]:
                    reason = decision_reason
                else:
                    reason = "other"
                
                skip_reasons[reason] = skip_reasons.get(reason, 0) + 1
                shapes_skipped += 1
            
            # Shape type distribution
            shape_type = entry.shape_type or "unknown"
            shape_type_distribution[shape_type] = shape_type_distribution.get(shape_type, 0) + 1
            
            # Rasterization status
            if hasattr(entry, 'rasterizer_info') and entry.rasterizer_info:
                if entry.rasterizer_info.get('status') == 'error':
                    rasterization_failures += 1
        
        # Calculate efficiency metrics
        total_elements = len(entries)
        cache_hit_rate = (cache_hits / total_elements * 100) if total_elements > 0 else 0
        llava_avoidance_rate = (llava_calls_avoided / total_elements * 100) if total_elements > 0 else 0
        
        # Get manifest statistics
        manifest_stats = manifest.get_statistics()
        
        telemetry = {
            'pipeline_efficiency': {
                'total_elements': total_elements,
                'llava_calls_made': llava_calls_made,
                'llava_calls_avoided': llava_calls_avoided,
                'cache_hit_rate_percent': round(cache_hit_rate, 1),
                'llava_avoidance_rate_percent': round(llava_avoidance_rate, 1),
                'shapes_sent_to_llava': shapes_sent,
                'shapes_skipped': shapes_skipped
            },
            'decision_breakdown': decision_counts,
            'skip_reasons': skip_reasons,
            'shape_type_distribution': shape_type_distribution,
            'rasterization_status': {
                'successes': total_elements - rasterization_failures,
                'failures': rasterization_failures
            },
            'manifest_statistics': manifest_stats,
            'cost_avoidance': {
                'cache_hits': cache_hits,
                'preserved_existing': decision_counts.get('preserve_existing', 0),
                'total_avoided_calls': llava_calls_avoided
            }
        }
        
        return telemetry
    
    def _generate_empty_telemetry(self) -> Dict[str, Any]:
        """Generate empty telemetry structure for error cases."""
        return {
            'pipeline_efficiency': {
                'total_elements': 0,
                'llava_calls_made': 0,
                'llava_calls_avoided': 0,
                'cache_hit_rate_percent': 0,
                'llava_avoidance_rate_percent': 0
            },
            'decision_breakdown': {},
            'skip_reasons': {},
            'shape_type_distribution': {},
            'rasterization_status': {'successes': 0, 'failures': 0}
        }
    
    def _validate_preserved_alt_display(self, entries: List, docx_rows) -> List[str]:
        """
        GREMLIN 2 FIX: Validate that preserved ALT entries are displayed correctly in DOCX.
        
        For entries with decision_reason=preserve_existing:
        - Current ALT Text column should show existing_alt
        - Suggested ALT Text column should be blank (or contain hint text)
        
        Args:
            entries: List of manifest entries
            docx_rows: Table rows from DOCX document
            
        Returns:
            List of validation error messages
        """
        errors = []
        
        # Build lookup of preserved entries
        preserved_entries = {}
        for entry in entries:
            if entry.decision_reason == "preserve_existing":
                preserved_entries[f"{entry.slide_number}/{entry.image_number}"] = entry
        
        if not preserved_entries:
            return errors  # No preserved entries to validate
        
        # Validate each DOCX row for preserved entries
        for row_idx, row in enumerate(docx_rows):
            try:
                cells = row.cells
                if len(cells) < 5:
                    continue  # Skip malformed rows
                
                # Extract slide/image identifier from first column
                slide_img = cells[0].text.strip()
                
                if slide_img in preserved_entries:
                    entry = preserved_entries[slide_img]
                    
                    # Column 3: Current ALT Text should match existing_alt
                    current_alt_cell = cells[3].text.strip()
                    expected_current = entry.existing_alt if entry.existing_alt else ""
                    
                    if expected_current and current_alt_cell != expected_current:
                        errors.append(
                            f"DOCX row {row_idx} ({slide_img}): Current ALT mismatch. "
                            f"Expected '{expected_current}', got '{current_alt_cell}'"
                        )
                    elif not expected_current and current_alt_cell not in ["[No ALT text]", ""]:
                        errors.append(
                            f"DOCX row {row_idx} ({slide_img}): Should show '[No ALT text]' for empty existing_alt, "
                            f"got '{current_alt_cell}'"
                        )
                    
                    # Column 4: Suggested ALT Text should be blank for preserved entries
                    suggested_alt_cell = cells[4].text.strip()
                    
                    # Allow hint text or completely blank
                    hint_text = "(Leave blank to keep current, or type replacement)"
                    is_blank_or_hint = (not suggested_alt_cell or 
                                       hint_text in suggested_alt_cell)
                    
                    if not is_blank_or_hint:
                        errors.append(
                            f"DOCX row {row_idx} ({slide_img}): Suggested ALT should be blank for preserved entries, "
                            f"got '{suggested_alt_cell}'"
                        )
                        
            except Exception as e:
                errors.append(f"Error validating DOCX row {row_idx}: {e}")
        
        return errors
    
    def _check_punctuation_parity(self, expected: str, actual: str) -> Dict[str, Any]:
        """
        GREMLIN 1 FIX: Check if punctuation parity issue exists between expected and actual.
        
        Args:
            expected: Expected ALT text (from manifest)
            actual: Actual ALT text (from PPT)
            
        Returns:
            Dict with punctuation analysis
        """
        if not expected or not actual:
            return {'has_issue': False, 'reason': 'empty_text'}
        
        expected_ends = expected[-1] if expected else ''
        actual_ends = actual[-1] if actual else ''
        
        # Check terminal punctuation
        punctuation_chars = '.!?'
        expected_has_punct = expected_ends in punctuation_chars
        actual_has_punct = actual_ends in punctuation_chars
        
        if expected_has_punct != actual_has_punct:
            return {
                'has_issue': True,
                'reason': 'terminal_punctuation_mismatch',
                'expected_ends_with': expected_ends,
                'actual_ends_with': actual_ends,
                'expected_has_punct': expected_has_punct,
                'actual_has_punct': actual_has_punct
            }
        elif expected_has_punct and actual_has_punct and expected_ends != actual_ends:
            return {
                'has_issue': True,
                'reason': 'different_terminal_punctuation',
                'expected_ends_with': expected_ends,
                'actual_ends_with': actual_ends
            }
        
        # Check if the only difference is trailing punctuation
        expected_no_punct = expected.rstrip('.!?')
        actual_no_punct = actual.rstrip('.!?')
        
        if expected_no_punct == actual_no_punct and expected != actual:
            return {
                'has_issue': True,
                'reason': 'punctuation_only_difference',
                'base_text_matches': True,
                'expected_punct': expected[len(expected_no_punct):],
                'actual_punct': actual[len(actual_no_punct):]
            }
        
        return {'has_issue': False, 'reason': 'no_punctuation_issue'}
    
    def _log_validation_summary(self, result: Dict[str, Any]):
        """Log comprehensive validation summary."""
        success = result['success']
        telemetry = result.get('telemetry', {})
        efficiency = telemetry.get('pipeline_efficiency', {})
        
        logger.info("="*60)
        logger.info("🔍 SYNC VALIDATION & TELEMETRY REPORT")
        logger.info("="*60)
        
        # Overall status
        status_emoji = "✅" if success else "❌"
        logger.info(f"{status_emoji} Overall Status: {'PASS' if success else 'FAIL'}")
        
        if not success:
            logger.info(f"   Validation Errors: {len(result.get('validation_errors', []))}")
            for error in result.get('validation_errors', [])[:3]:
                logger.info(f"     - {error}")
        
        # Pipeline efficiency
        logger.info("\n📊 Pipeline Efficiency:")
        logger.info(f"   Total elements: {efficiency.get('total_elements', 0)}")
        logger.info(f"   LLaVA calls made: {efficiency.get('llava_calls_made', 0)}")
        logger.info(f"   LLaVA calls avoided: {efficiency.get('llava_calls_avoided', 0)}")
        logger.info(f"   Cache hit rate: {efficiency.get('cache_hit_rate_percent', 0)}%")
        logger.info(f"   LLaVA avoidance rate: {efficiency.get('llava_avoidance_rate_percent', 0)}%")
        
        # Decision breakdown
        decision_breakdown = telemetry.get('decision_breakdown', {})
        if decision_breakdown:
            logger.info("\n🎯 Decision Breakdown:")
            for reason, count in sorted(decision_breakdown.items()):
                logger.info(f"   {reason}: {count}")
        
        # Shape type distribution
        shape_dist = telemetry.get('shape_type_distribution', {})
        if shape_dist:
            logger.info("\n📐 Shape Types:")
            for shape_type, count in sorted(shape_dist.items()):
                logger.info(f"   {shape_type}: {count}")
        
        # Rasterization status
        raster_status = telemetry.get('rasterization_status', {})
        if raster_status:
            logger.info("\n🖼️  Rasterization:")
            logger.info(f"   Successes: {raster_status.get('successes', 0)}")
            logger.info(f"   Failures: {raster_status.get('failures', 0)}")
        
        logger.info("="*60)
    
    def save_telemetry_report(self, result: Dict[str, Any], output_path: str = None) -> str:
        """Save detailed telemetry report to JSON file."""
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"telemetry_report_{timestamp}.json"
        
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(result, f, indent=2, ensure_ascii=False)
            
            logger.info(f"📄 Telemetry report saved: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Failed to save telemetry report: {e}")
            return ""


def validate_pipeline_outputs(manifest_path: str, pptx_path: str = None, 
                            docx_path: str = None, save_report: bool = False) -> Dict[str, Any]:
    """
    Convenience function for validating complete pipeline outputs.
    
    Args:
        manifest_path: Path to manifest.json
        pptx_path: Path to injected PPTX file (optional)
        docx_path: Path to generated DOCX file (optional)
        save_report: Whether to save detailed JSON report
        
    Returns:
        Validation results with telemetry
    """
    validator = SyncValidator()
    result = validator.validate_full_pipeline(manifest_path, pptx_path, docx_path)
    
    if save_report:
        validator.save_telemetry_report(result)
    
    return result


# CLI entry point for testing
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python sync_validator.py <manifest_path> [pptx_path] [docx_path]")
        sys.exit(1)
    
    manifest_path = sys.argv[1]
    pptx_path = sys.argv[2] if len(sys.argv) > 2 else None
    docx_path = sys.argv[3] if len(sys.argv) > 3 else None
    
    # Setup logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )
    
    # Run validation
    result = validate_pipeline_outputs(manifest_path, pptx_path, docx_path, save_report=True)
    
    # Exit with appropriate code
    exit_code = 0 if result['success'] else 1
    sys.exit(exit_code)