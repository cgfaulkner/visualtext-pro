"""Smart Selector - Contract-aligned implementation.

Consumes Phase 1 manifest/visual_index as primary source of truth for candidate
identity and element_id. Produces deterministic selector manifest; no pixels,
rendering, or LLM. See docs/smart-selector-contract.md and
docs/IMPLEMENTATION_PLAN_smart_selector_contract_alignment.md.
"""

import json
import re
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional, Set, Tuple

from .types import (
    SelectorManifestRecord,
    SelectorManifest,
)

logger = logging.getLogger(__name__)

SELECTOR_VERSION = "1.0.0"

# Shape-built diagram clustering (connector-or-label rule)
PROXIMITY_PX = 40
REQUIRED_ALIGNMENT_FRACTION = 0.6
CONFIDENCE_THRESHOLD = 0.6
CONNECTOR_WEIGHT = 0.5
NEARBY_LABEL_WEIGHT = 0.25
ADJACENCY_WEIGHT = 0.1
DEFAULT_LABEL_SKIP_TOKENS = frozenset({"the", "a", "an", "and", "or", "of", "to", "in", "on"})

SELECTOR_VERSION_NOTES = (
    "- Connector-or-label rule: virtual_group only when (connector between members) OR (>=2 "
    "meaningful nearby labels). Proximity 40px, alignment fraction 0.6. Confidence threshold 0.6; "
    "below triggers render_and_assist. Anchor selection: smallest element_id lexicographically."
)


def _create_instance_key(slide_idx: int, shape_id: int) -> str:
    """Same formula as Phase 1 (alt_manifest.create_instance_key)."""
    return f"slide_{slide_idx}_shape_{shape_id}"


def _bbox_centroid(bbox: Dict[str, float]) -> Tuple[float, float]:
    """Return (cx, cy) center of bbox (left, top, width, height)."""
    left = bbox.get("left") or 0
    top = bbox.get("top") or 0
    w = bbox.get("width") or 0
    h = bbox.get("height") or 0
    return (left + w / 2, top + h / 2)


def _bbox_distance(bbox1: Dict[str, float], bbox2: Dict[str, float]) -> float:
    """Euclidean distance between bbox centroids (points)."""
    c1 = _bbox_centroid(bbox1)
    c2 = _bbox_centroid(bbox2)
    return ((c1[0] - c2[0]) ** 2 + (c1[1] - c2[1]) ** 2) ** 0.5


def _bboxes_overlap_or_near(
    bbox1: Dict[str, float], bbox2: Dict[str, float], proximity_px: float
) -> bool:
    """True if bboxes overlap or centroids are within proximity_px."""
    if _bbox_distance(bbox1, bbox2) <= proximity_px:
        return True
    l1, t1 = bbox1.get("left", 0), bbox1.get("top", 0)
    w1, h1 = bbox1.get("width", 0), bbox1.get("height", 0)
    l2, t2 = bbox2.get("left", 0), bbox2.get("top", 0)
    w2, h2 = bbox2.get("width", 0), bbox2.get("height", 0)
    return not (l1 + w1 < l2 or l2 + w2 < l1 or t1 + h1 < t2 or t2 + h2 < t1)


def _get_connectors_from_visual_index(visual_index: Dict[str, Any]) -> Set[str]:
    """Element IDs that are connectors (LINE or CONNECTOR shape type). Binary, high weight."""
    out: Set[str] = set()
    for eid, entry in visual_index.items():
        if entry.get("shape_type") in ("LINE", "CONNECTOR"):
            out.add(eid)
    return out


def _get_label_shapes_from_visual_index(visual_index: Dict[str, Any]) -> List[Tuple[str, Dict[str, float], str]]:
    """List of (element_id, bbox, text) for shapes that can act as labels (have text_content)."""
    out: List[Tuple[str, Dict[str, float], str]] = []
    for eid, entry in visual_index.items():
        text = (entry.get("text_content") or entry.get("text") or "").strip()
        if not text:
            continue
        bbox = entry.get("bbox") or {"left": 0, "top": 0, "width": 0, "height": 0}
        out.append((eid, bbox, text))
    return out


def _is_meaningful_label_text(
    text: str,
    min_words: int = 3,
    skip_tokens: Optional[Set[str]] = None,
) -> bool:
    """True if text has at least min_words and is not just skip tokens."""
    if not text or not text.strip():
        return False
    tokens = re.findall(r"\w+", text.strip().lower())
    if len(tokens) < min_words:
        return False
    skip = skip_tokens or DEFAULT_LABEL_SKIP_TOKENS
    meaningful = [t for t in tokens if t not in skip]
    return len(meaningful) >= 1


def nearby_label_text(
    element_keys: Set[str],
    label_shapes: List[Tuple[str, Dict[str, float], str]],
    visual_index: Dict[str, Any],
    proximity_px: float = PROXIMITY_PX,
    required_alignment_fraction: float = REQUIRED_ALIGNMENT_FRACTION,
    min_words: int = 3,
    skip_tokens: Optional[Set[str]] = None,
) -> Tuple[bool, int]:
    """
    True if >=2 meaningful labels are within proximity_px of candidate centroid for >=
    required_alignment_fraction of cluster members. Meaningful = >= min_words, not just
    skip-tokens. Returns (satisfied, count_of_meaningful_nearby_labels).
    """
    if not element_keys or not label_shapes:
        return False, 0
    skip = skip_tokens or DEFAULT_LABEL_SKIP_TOKENS
    count_nearby = 0
    for _lid, lbbox, ltext in label_shapes:
        if not _is_meaningful_label_text(ltext, min_words=min_words, skip_tokens=skip):
            continue
        label_center = _bbox_centroid(lbbox)
        aligned = 0
        for eid in element_keys:
            entry = visual_index.get(eid, {})
            mbbox = entry.get("bbox") or {"left": 0, "top": 0, "width": 0, "height": 0}
            member_center = _bbox_centroid(mbbox)
            dist = ((label_center[0] - member_center[0]) ** 2 + (label_center[1] - member_center[1]) ** 2) ** 0.5
            if dist <= proximity_px:
                aligned += 1
        if aligned >= required_alignment_fraction * len(element_keys):
            count_nearby += 1
    return count_nearby >= 2, count_nearby


def _build_adjacency(
    visual_index: Dict[str, Any],
    top_level_ids: Set[str],
    proximity_px: float,
) -> Dict[str, Set[str]]:
    """For each top-level element_id, set of other top-level ids that overlap or are near."""
    adj: Dict[str, Set[str]] = {eid: set() for eid in top_level_ids}
    ids_list = list(top_level_ids)
    for i, eid1 in enumerate(ids_list):
        b1 = (visual_index.get(eid1) or {}).get("bbox") or {}
        for eid2 in ids_list[i + 1 :]:
            b2 = (visual_index.get(eid2) or {}).get("bbox") or {}
            if _bboxes_overlap_or_near(b1, b2, proximity_px):
                adj[eid1].add(eid2)
                adj[eid2].add(eid1)
    return adj


def _connected_components(adj: Dict[str, Set[str]], nodes: Set[str]) -> List[Set[str]]:
    """Return list of connected component sets (BFS)."""
    seen: Set[str] = set()
    components: List[Set[str]] = []
    for n in nodes:
        if n in seen:
            continue
        comp: Set[str] = set()
        stack = [n]
        while stack:
            cur = stack.pop()
            if cur in seen:
                continue
            seen.add(cur)
            comp.add(cur)
            for nb in adj.get(cur, set()):
                if nb not in seen:
                    stack.append(nb)
        if comp:
            components.append(comp)
    return components


def detect_shape_clusters(
    visual_index: Dict[str, Any],
    top_level_ids: Set[str],
    connectors: Set[str],
    label_shapes: List[Tuple[str, Dict[str, float], str]],
    config: Dict[str, Any],
) -> List[Dict[str, Any]]:
    """
    Form candidate clusters by adjacency/overlap; permit cluster only if
    (connector present between members) OR (>=2 meaningful nearby text labels).
    Anchor = smallest element_id lexicographically. Returns list of cluster dicts:
    {anchor, members, confidence, tie_breaker_signals, reason, escalation}.
    """
    sel = config.get("selector", {})
    syn = sel.get("synthetic_diagram", {})
    enabled = syn.get("enabled", True)
    proximity_px = int(syn.get("proximity_px", PROXIMITY_PX))
    alignment_frac = float(syn.get("required_alignment_fraction", REQUIRED_ALIGNMENT_FRACTION))
    confidence_threshold = float(syn.get("confidence_threshold", CONFIDENCE_THRESHOLD))

    clusters_out: List[Dict[str, Any]] = []

    if not enabled or not top_level_ids:
        return clusters_out

    adjacency = _build_adjacency(visual_index, top_level_ids, proximity_px)
    components = _connected_components(adjacency, top_level_ids)

    for comp in components:
        if len(comp) < 2:
            continue
        has_connector = bool(comp & connectors)
        satisfied, label_count = nearby_label_text(
            comp,
            label_shapes,
            visual_index,
            proximity_px=proximity_px,
            required_alignment_fraction=alignment_frac,
        )
        if not has_connector and not satisfied:
            logger.debug(
                "Cluster rejected: no connector and insufficient nearby labels; "
                "element_ids=%s",
                sorted(comp),
            )
            continue
        anchor = min(comp)
        members = comp - {anchor}
        confidence = 0.0
        if has_connector:
            confidence += CONNECTOR_WEIGHT
        confidence += min(label_count * NEARBY_LABEL_WEIGHT, 0.5)
        confidence += ADJACENCY_WEIGHT
        tie_breaker_signals = {
            "anchor": anchor,
            "connector_present": has_connector,
            "nearby_label_count": label_count,
            "member_count": len(comp),
        }
        reason = (
            "Synthetic diagram: connector present"
            if has_connector
            else f"Synthetic diagram: {label_count} nearby meaningful labels"
        )
        escalation = "render_and_assist" if confidence < confidence_threshold else "none"
        clusters_out.append({
            "anchor": anchor,
            "members": members,
            "confidence": round(confidence, 2),
            "tie_breaker_signals": tie_breaker_signals,
            "reason": reason,
            "escalation": escalation,
        })
        logger.debug(
            "Cluster formed: anchor=%s, members=%s, confidence=%s, %s",
            anchor, sorted(members), confidence, tie_breaker_signals,
        )
    return clusters_out


def _build_parent_map(pptx_path: Path, candidate_ids: Set[str]) -> Dict[str, Optional[str]]:
    """Build element_id -> parent_group element_id from PPTX shape hierarchy (structural only)."""
    parent_map: Dict[str, Optional[str]] = {}
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(pptx_path))

        def walk(shape: Any, slide_idx: int, parent_element_id: Optional[str]) -> None:
            shape_id = getattr(shape, "shape_id", None)
            if shape_id is None:
                return
            element_id = _create_instance_key(slide_idx, shape_id)
            if element_id in candidate_ids and parent_element_id is not None:
                parent_map[element_id] = parent_element_id
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(
                shape, "shapes"
            ):
                next_parent = element_id if element_id in candidate_ids else parent_element_id
                for child in shape.shapes:
                    walk(child, slide_idx, next_parent)

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                walk(shape, slide_idx, None)
    except Exception as e:
        logger.warning(f"Could not build parent map from PPTX: {e}")
    return parent_map


def _is_meaningful_alt(
    alt: str,
    placeholder_patterns: List[str],
    min_chars: int,
) -> bool:
    """True if ALT is non-placeholder and long enough (contract: Preserve Existing ALT)."""
    if not alt or not alt.strip():
        return False
    if len(alt.strip()) < min_chars:
        return False
    for pattern in placeholder_patterns:
        try:
            if re.search(pattern, alt.strip(), re.IGNORECASE):
                return False
        except re.error:
            continue
    return True


def _is_decorative(shape_type: str, bbox: Dict[str, float], config: Dict[str, Any]) -> bool:
    """Structural-only decorative hint (e.g. LINE with negligible area)."""
    if shape_type == "LINE":
        w = bbox.get("width") or 0
        h = bbox.get("height") or 0
        area = w * h
        if area < 1.0:
            return True
    sel = config.get("selector", {})
    threshold = sel.get("decorative_area_threshold")
    if threshold is not None and isinstance(bbox.get("width"), (int, float)) and isinstance(
        bbox.get("height"), (int, float)
    ):
        if (bbox["width"] * bbox["height"]) < threshold:
            return True
    return False


def _decide_record(
    element_id: str,
    shape_type: str,
    existing_alt: str,
    bbox: Dict[str, float],
    parent_group_id: Optional[str],
    inclusion_policy: str,
    placeholder_patterns: List[str],
    min_meaningful_chars: int,
    config: Dict[str, Any],
    groups_selected: Set[str],
) -> Tuple[SelectorManifestRecord, Optional[SelectorManifestRecord]]:
    """Return (main record, optional preserve_conflict child record)."""
    metadata: Dict[str, Any] = {"original_shape_type": shape_type}
    meaningful_alt = _is_meaningful_alt(existing_alt, placeholder_patterns, min_meaningful_chars)

    # 1) Preserve-existing-ALT: policy preserve + meaningful ALT -> include with preserve reason
    if inclusion_policy == "preserve" and meaningful_alt:
        return (
            {
                "selector_version": SELECTOR_VERSION,
                "element_id": element_id,
                "parent_group_id": None,
                "selector_decision": "include_atomic",
                "content_scope": "image",
                "reason_code": "preserve_existing_alt",
                "human_reason": "Existing meaningful ALT preserved per policy",
                "escalation_strategy": "none",
                "metadata": metadata,
            },
            None,
        )

    # 2) Child of a group we're including -> exclude_redundant
    if parent_group_id is not None and parent_group_id in groups_selected:
        # If preserve and child has meaningful ALT, also emit preserve_conflict record
        conflict_record = None
        if inclusion_policy == "preserve" and meaningful_alt:
            conflict_record = {
                "selector_version": SELECTOR_VERSION,
                "element_id": element_id,
                "parent_group_id": parent_group_id,
                "selector_decision": "preserve_conflict",
                "content_scope": "group",
                "reason_code": "preserve_conflict_child_has_alt",
                "human_reason": "Group selected; child has existing ALT (preserve); conflict for review",
                "escalation_strategy": "none",
                "metadata": metadata,
            }
        return (
            {
                "selector_version": SELECTOR_VERSION,
                "element_id": element_id,
                "parent_group_id": parent_group_id,
                "selector_decision": "exclude_redundant",
                "content_scope": "group",
                "reason_code": "exclude_redundant_group_child",
                "human_reason": "Excluded as child of selected group",
                "escalation_strategy": "none",
                "metadata": metadata,
            },
            conflict_record,
        )

    # 3) Decorative
    if _is_decorative(shape_type, bbox, config):
        return (
            {
                "selector_version": SELECTOR_VERSION,
                "element_id": element_id,
                "parent_group_id": None,
                "selector_decision": "exclude_decorative",
                "content_scope": "image",
                "reason_code": "exclude_decorative_structure",
                "human_reason": "Decorative element (structural)",
                "escalation_strategy": "none",
                "metadata": metadata,
            },
            None,
        )

    # 4) Atomic inclusion (default)
    return (
        {
            "selector_version": SELECTOR_VERSION,
            "element_id": element_id,
            "parent_group_id": None,
            "selector_decision": "include_atomic",
            "content_scope": "image",
            "reason_code": "include_atomic_default",
            "human_reason": "Included as atomic visual element",
            "escalation_strategy": "none",
            "metadata": metadata,
        },
        None,
    )


def run_selector(
    pptx_path: Path,
    artifacts: Any,
    config: Dict[str, Any],
    output_path: Optional[Path] = None,
) -> Path:
    """
    Run the Smart Selector using Phase 1 artifacts as primary source of truth.

    Candidates and element_id come from visual_index (instance_key). PPTX is used
    only for shape hierarchy (parent_map) to apply group suppression. Deterministic;
    no pixels or LLM.

    Args:
        pptx_path: Path to PPTX (used for group hierarchy only).
        artifacts: RunArtifacts with visual_index_path, current_alt_by_key, selector_manifest_path.
        config: Configuration (selector.*, alt_text_handling.mode).
        output_path: Where to write selector_manifest.json. Defaults to artifacts.selector_manifest_path.

    Returns:
        Path to written selector_manifest.json.
    """
    logger.info(f"Running Smart Selector v{SELECTOR_VERSION} (Phase 1 artifacts as source)")

    visual_index = artifacts.load_visual_index()
    current_alt_by_key = artifacts.load_current_alt_by_key()
    if not visual_index:
        logger.warning("No visual_index; writing empty selector manifest")
        manifest: SelectorManifest = []
    else:
        candidate_ids = set(visual_index.keys())
        parent_map = _build_parent_map(pptx_path, candidate_ids)

        sel_config = config.get("selector", {})
        inclusion_policy = (
            sel_config.get("inclusion_policy")
            or config.get("alt_text_handling", {}).get("mode", "preserve")
        )
        placeholder_patterns = sel_config.get("placeholder_alt_patterns") or []
        min_meaningful_chars = sel_config.get("min_meaningful_alt_chars", 15)

        # Top-level candidates (no parent in our set)
        top_level_ids: Set[str] = {
            eid for eid in candidate_ids if parent_map.get(eid) is None
        }

        # Explicit GROUP shapes
        group_ids: Set[str] = set()
        for eid in top_level_ids:
            if visual_index.get(eid, {}).get("shape_type") == "GROUP":
                group_ids.add(eid)

        # Synthetic (shape-built) diagram clusters: connector or >=2 nearby labels
        connectors = _get_connectors_from_visual_index(visual_index)
        label_shapes = _get_label_shapes_from_visual_index(visual_index)
        synthetic_clusters = detect_shape_clusters(
            visual_index,
            top_level_ids,
            connectors,
            label_shapes,
            config,
        )
        synthetic_anchor_ids: Set[str] = set()
        synthetic_parent_map: Dict[str, str] = {}
        synthetic_cluster_meta: Dict[str, Dict[str, Any]] = {}
        for clust in synthetic_clusters:
            anchor = clust["anchor"]
            synthetic_anchor_ids.add(anchor)
            synthetic_cluster_meta[anchor] = {
                "confidence": clust.get("confidence"),
                "tie_breaker_signals": clust.get("tie_breaker_signals"),
                "selector_version_notes": SELECTOR_VERSION_NOTES,
                "reason": clust.get("reason", "Included as composite (shape-built diagram)"),
            }
            for mid in clust["members"]:
                synthetic_parent_map[mid] = anchor
            if clust.get("escalation") == "render_and_assist":
                synthetic_cluster_meta[anchor]["escalation"] = "render_and_assist"

        # Combined: explicit groups + synthetic anchors get include_group; their children/members suppressed
        groups_selected = group_ids | synthetic_anchor_ids

        manifest = []
        for element_id in sorted(visual_index.keys()):
            entry = visual_index[element_id]
            shape_type = entry.get("shape_type", "UNKNOWN")
            existing_alt = (
                current_alt_by_key.get(element_id) or entry.get("existing_alt") or ""
            )
            bbox = entry.get("bbox") or {"left": 0, "top": 0, "width": 0, "height": 0}
            parent_group_id = parent_map.get(element_id) or synthetic_parent_map.get(
                element_id
            )

            # Explicit GROUP: include_group
            if element_id in group_ids:
                record = {
                    "selector_version": SELECTOR_VERSION,
                    "element_id": element_id,
                    "parent_group_id": None,
                    "selector_decision": "include_group",
                    "content_scope": "group",
                    "reason_code": "include_group_composite",
                    "human_reason": "Included as composite group",
                    "escalation_strategy": "none",
                    "metadata": {"original_shape_type": shape_type},
                }
                manifest.append(record)
                continue

            # Synthetic diagram anchor: include_group
            if element_id in synthetic_anchor_ids:
                meta = dict(synthetic_cluster_meta.get(element_id) or {})
                meta["original_shape_type"] = shape_type
                escalation = meta.pop("escalation", "none")
                reason_str = meta.pop("reason", "Included as composite (shape-built diagram)")
                record = {
                    "selector_version": SELECTOR_VERSION,
                    "element_id": element_id,
                    "parent_group_id": None,
                    "selector_decision": "include_group",
                    "content_scope": "group",
                    "reason_code": "include_group_synthetic_diagram",
                    "human_reason": reason_str,
                    "escalation_strategy": escalation,
                    "metadata": meta,
                }
                manifest.append(record)
                continue

            # Child of selected group or member of synthetic cluster: exclude_redundant (or preserve_conflict)
            if parent_group_id is not None and parent_group_id in groups_selected:
                is_synthetic = element_id in synthetic_parent_map
                reason_code = (
                    "exclude_redundant_synthetic_diagram_member"
                    if is_synthetic
                    else "exclude_redundant_group_child"
                )
                human_reason = (
                    "Excluded as member of shape-built diagram"
                    if is_synthetic
                    else "Excluded as child of selected group"
                )
                conflict_record = None
                meaningful_alt = _is_meaningful_alt(
                    existing_alt, placeholder_patterns, min_meaningful_chars
                )
                if inclusion_policy == "preserve" and meaningful_alt:
                    conflict_record = {
                        "selector_version": SELECTOR_VERSION,
                        "element_id": element_id,
                        "parent_group_id": parent_group_id,
                        "selector_decision": "preserve_conflict",
                        "content_scope": "group",
                        "reason_code": "preserve_conflict_child_has_alt",
                        "human_reason": "Group selected; child has existing ALT (preserve); conflict for review",
                        "escalation_strategy": "none",
                        "metadata": {"original_shape_type": shape_type},
                    }
                main_record = {
                    "selector_version": SELECTOR_VERSION,
                    "element_id": element_id,
                    "parent_group_id": parent_group_id,
                    "selector_decision": "exclude_redundant",
                    "content_scope": "group",
                    "reason_code": reason_code,
                    "human_reason": human_reason,
                    "escalation_strategy": "none",
                    "metadata": {"original_shape_type": shape_type},
                }
                manifest.append(main_record)
                if conflict_record is not None:
                    manifest.append(conflict_record)
                continue

            main_record, conflict_record = _decide_record(
                element_id,
                shape_type,
                existing_alt,
                bbox,
                parent_group_id,
                inclusion_policy,
                placeholder_patterns,
                min_meaningful_chars,
                config,
                groups_selected,
            )
            manifest.append(main_record)
            if conflict_record is not None:
                manifest.append(conflict_record)

    manifest.sort(key=lambda r: r["element_id"])

    out = output_path if output_path is not None else getattr(
        artifacts, "selector_manifest_path", None
    )
    if out is None:
        sel = config.get("selector", {})
        out = Path(sel.get("output_dir", pptx_path.parent)) / "selector_manifest.json"
    out = Path(out)
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)
    logger.info(f"Selector manifest written to {out} ({len(manifest)} records)")
    return out
