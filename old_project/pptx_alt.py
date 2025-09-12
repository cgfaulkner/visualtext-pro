#!/usr/bin/env python3
"""
PowerPoint Accessibility Auditor - Phase 1
Scans PPTX files for image ALT text compliance and generates reports.
"""

import os
import csv
import logging
import hashlib
import json
import asyncio
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from datetime import datetime
from collections import defaultdict
from io import BytesIO
import subprocess
import shutil
from concept_detector import ConceptDetector
from unified_alt_generator import FlexibleAltGenerator
from image_type_detector import detect_image_type, set_alt_generator
from decorative_filter import (
    is_decorative_image,
    get_image_hash,
    is_force_decorative_by_filename_or_name,
)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
from PIL import Image

# Caches keyed by image hash
logger = logging.getLogger("pptx_alt")
alt_cache_path = Path("alt_cache.json")
thumbnail_cache: dict[str, tuple[str, str]] = {}


def load_alt_cache(path: Path) -> dict:
    """Load ALT text cache from disk.

    Args:
        path: Path to the cache file.

    Returns:
        dict: The loaded cache or an empty dictionary on failure.
    """
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except FileNotFoundError:
        logger.info(f"No cache file found at {path}, starting with empty cache")
    except json.JSONDecodeError:
        logger.warning(f"Cache file at {path} is corrupt. Starting with empty cache")
    except Exception as e:
        logger.warning(f"Could not load cache file {path}: {e}. Starting with empty cache")
    return {}


alt_cache = load_alt_cache(alt_cache_path)


class PowerPointAccessibilityAuditor:
    def __init__(self, config_manager, input_folder=None, output_folder=None, thumbnail_max_width=200):
        self.config_manager = config_manager
        self.config = config_manager.config
        
        # Use config paths or fall back to provided/default values
        paths = self.config.get('paths', {})
        self.input_folder = Path(input_folder or paths.get('input_folder', 'test_pptx'))
        self.output_folder = Path(output_folder or paths.get('output_folder', 'output'))
        self.thumbnail_max_width = thumbnail_max_width
        
        # Create output directory structure using config paths
        self.csv_dir = self.output_folder / 'csv_reports'
        self.thumbnail_dir = Path(paths.get('thumbnail_folder', str(self.output_folder / 'thumbnails')))
        self.log_dir = self.output_folder / 'logs'
        
        for directory in [self.csv_dir, self.thumbnail_dir, self.log_dir]:
            directory.mkdir(parents=True, exist_ok=True)

        # Setup logging
        self._setup_logging()

        # Load ALT text cache from configured path
        global alt_cache, alt_cache_path
        self.alt_cache_path = Path(paths.get('alt_cache', 'alt_cache.json'))
        alt_cache_path = self.alt_cache_path
        alt_cache = load_alt_cache(alt_cache_path)

        # Check external tools
        self.inkscape_path = self._check_inkscape()

        # Setup concept detector if enabled
        alt_cfg = self.config.get('alt_text_handling', {})
        self.max_workers = alt_cfg.get('max_workers', 4)
        concepts_file = alt_cfg.get('concepts_file', 'concepts.yaml')
        self.concept_detector = ConceptDetector(concepts_file)

        # Initialize unified ALT text generator
        self.alt_generator = FlexibleAltGenerator(self.config_manager)
        set_alt_generator(self.alt_generator)

        # Track image occurrences across slides
        self.image_tracker = defaultdict(list)

    def run(self, pptx_path):
        return self._process_presentation(pptx_path)
    
    def _setup_logging(self):
        """Configure logging with both file and console output."""
        log_filename = self.log_dir / f"audit_log_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log"

        # Configure dedicated logger
        self.logger = logging.getLogger("pptx_alt")
        self.logger.setLevel(logging.DEBUG)
        self.logger.handlers.clear()
        self.logger.propagate = False

        formatter = logging.Formatter("%(asctime)s - %(levelname)s - %(message)s")

        file_handler = logging.FileHandler(log_filename, mode="w")
        file_handler.setFormatter(formatter)

        stream_handler = logging.StreamHandler()
        stream_handler.setFormatter(formatter)

        self.logger.addHandler(file_handler)
        self.logger.addHandler(stream_handler)

        self.logger.info(f"Accessibility audit started. Log file: {log_filename}")

    def _check_inkscape(self):
        """Ensure Inkscape is available if conversion is enabled."""
        if not self.config.get('conversion', {}).get('wmf_to_png', True):
            return None

        inkscape_cmd = self.config.get('tools', {}).get('inkscape', 'inkscape')
        path = shutil.which(inkscape_cmd)
        if not path:
            raise RuntimeError(
                f"Inkscape executable '{inkscape_cmd}' not found in PATH"
            )
        return path
    
    def _extract_image_data(self, shape):
        """Extract image data from a shape."""
        try:
            image_part = shape.image
            return image_part.blob
        except Exception as e:
            self.logger.warning(f"Could not extract image data: {e}")
            return None
    
    def _get_shape_position(self, shape):
        """Extract shape position in pixels."""
        try:
            # Convert from EMUs to pixels (1 pixel = 9525 EMUs)
            left_px = int(shape.left / 9525) if hasattr(shape, 'left') else 0
            top_px = int(shape.top / 9525) if hasattr(shape, 'top') else 0
            return (left_px, top_px)
        except Exception as e:
            self.logger.warning(f"Could not extract shape position: {e}")
            return (0, 0)

    def _extract_embedded_image_filename(self, shape):
        """Return the filename of the embedded image if available."""
        try:
            image = shape.image
            return Path(image.filename).name.lower()
        except Exception:
            return ""

    def _convert_vector_to_png(self, image_data, original_name):
        """Convert WMF/EMF bytes to PNG using Inkscape."""
        if not self.inkscape_path:
            return None

        temp_dir = Path(self.config.get('paths', {}).get('temp_folder', 'temp'))
        temp_dir.mkdir(parents=True, exist_ok=True)

        input_path = temp_dir / original_name
        with open(input_path, 'wb') as f:
            f.write(image_data)

        output_name = f"{Path(original_name).stem}.png"
        output_path = self.thumbnail_dir / output_name

        cmd = [self.inkscape_path, str(input_path), "--export-type=png", f"--export-filename={output_path}"]
        try:
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            return str(output_path)
        except subprocess.CalledProcessError as e:
            self.logger.error(f"Inkscape conversion failed for {original_name}: {e}")
            if output_path.exists():
                output_path.unlink()
            return None
        finally:
            if input_path.exists():
                input_path.unlink()
    
    def _save_thumbnail(self, image, pptx_name, slide_num, img_index):
        """Save image thumbnail for a pre-opened PIL image."""
        try:
            # Calculate thumbnail size maintaining aspect ratio
            aspect_ratio = image.height / image.width if image.width else 0
            new_width = min(image.width, self.thumbnail_max_width)
            new_height = int(new_width * aspect_ratio) if aspect_ratio else 0

            img_thumbnail = image.resize((new_width, new_height), Image.Resampling.LANCZOS)

            # Create filename
            filename = f"{pptx_name}_slide{slide_num}_img{img_index}.png"
            thumbnail_path = self.thumbnail_dir / filename

            # Convert to RGB if necessary (for transparency handling)
            if img_thumbnail.mode in ('RGBA', 'LA', 'P'):
                rgb_img = Image.new('RGB', img_thumbnail.size, (255, 255, 255))
                if img_thumbnail.mode == 'P':
                    img_thumbnail = img_thumbnail.convert('RGBA')
                mask = img_thumbnail.split()[-1] if 'A' in img_thumbnail.mode else None
                rgb_img.paste(img_thumbnail, mask=mask)
                img_thumbnail = rgb_img

            img_thumbnail.save(thumbnail_path, 'PNG')
            return str(thumbnail_path), filename  # Return both path and filename

        except Exception as e:
            self.logger.error(f"Error creating thumbnail: {e}")
            return "error_creating_thumbnail", "error.png"

    def _check_decorative_heuristics(
        self,
        shape,
        slide,
        image_data,
        alt_text,
        slide_num,
        img_index,
        dimensions,
        image_hash=None,
    ):
        """Check if an image might be decorative.

        Args:
            shape: The PowerPoint shape containing the image.
            slide: The slide object.
            image_data: Raw image bytes.
            alt_text: ALT text extracted from the image shape.
            slide_num: Slide number for logging context.
            img_index: Image index on the slide.
            dimensions: Tuple of image width and height.
            image_hash: Optional precomputed hash of ``image_data``.

        Returns:
            Tuple[bool, list, str, tuple]: Decorative flag, notes, image hash,
            and (width, height).
        """
        try:
            width, height = dimensions

            # Get shape position
            position = self._get_shape_position(shape)

            # Get shape name (used for keyword matching in decorative filter)
            shape_name = shape.name if hasattr(shape, 'name') else f"Image_{img_index}"

            # Get image hash
            if image_hash is None:
                image_hash = get_image_hash(image_data)

            # Use the decorative_filter module's logic
            is_decorative, notes = is_decorative_image(
                image_data,
                shape_name,
                position,
                (width, height),
                list(slide.shapes),
                image_hash,
                self.image_tracker,
            )

            self.logger.debug(
                f"Decorative check for {shape_name}: decorative={is_decorative}, notes={notes}"
            )

            return is_decorative, notes, image_hash, (width, height)

        except Exception as e:
            self.logger.warning(f"Error in decorative heuristics: {e}")
            return False, ["Error checking decorative status"], None, None
    
    def _get_alt_text_and_decorative(self, shape):
        """
        Extract ALT text and decorative status from shape.
        Returns tuple: (alt_text, is_decorative)
        """
        alt_text = ''
        is_decorative = False

        # Debug helper: inspect raw XML to troubleshoot missing ALT text or decorative flags
        self.logger.debug("Checking shape for ALT text and decorative status")
        
        try:
            # Try to get the raw XML from cNvPr
            if hasattr(shape, '_element') and hasattr(shape._element, 'nvPicPr'):
                nvPicPr = shape._element.nvPicPr
                if hasattr(nvPicPr, 'cNvPr'):
                    cNvPr = nvPicPr.cNvPr
                    
                    # Get the raw XML
                    if hasattr(cNvPr, 'xml'):
                        xml_str = cNvPr.xml
                        self.logger.debug("cNvPr XML: %s", xml_str)
                        
                        # Parse the XML to look for descr attribute
                        import xml.etree.ElementTree as ET
                        try:
                            root = ET.fromstring(xml_str)
                            
                            # Check for descr attribute
                            descr = root.get('descr')
                            if descr:
                                alt_text = descr
                                self.logger.debug("Found descr in XML: '%s'", alt_text)
                            
                            # Check for title attribute
                            title = root.get('title')
                            if title and not alt_text:
                                alt_text = title
                                self.logger.debug("Found title in XML: '%s'", title)
                                
                            # Look for extLst in the XML
                            for child in root:
                                if 'extLst' in child.tag:
                                    ext_xml = ET.tostring(child, encoding='unicode')
                                    self.logger.debug("Found extLst: %s", ext_xml[:200])
                                    
                                    # Check for decorative flag
                                    if 'decorative' in ext_xml and 'val="1"' in ext_xml:
                                        is_decorative = True
                                        self.logger.debug("Found decorative flag in XML")
                                        
                        except Exception as e:
                            self.logger.debug("XML parsing error: %s", e)
            
            # Also check if there's a descr property directly on the shape
            if hasattr(shape, 'descr'):
                self.logger.debug("Shape has descr property: %s", shape.descr)
                if shape.descr and not alt_text:
                    alt_text = shape.descr
                    
            # Check for title property on shape
            if hasattr(shape, 'title'):
                self.logger.debug("Shape has title property: %s", shape.title)
                if shape.title and not alt_text:
                    alt_text = shape.title
                    
            # Try the alternative_text property
            if hasattr(shape, 'alternative_text'):
                self.logger.debug("Shape has alternative_text property: %s", shape.alternative_text)
                if shape.alternative_text and not alt_text:
                    alt_text = shape.alternative_text
                    
        except Exception as e:
            import traceback
            self.logger.debug("Error while extracting ALT text: %s", e)
            self.logger.debug(traceback.format_exc())

        self.logger.info(f"Final result: alt_text='{alt_text}', is_decorative={is_decorative}")
        
        return alt_text, is_decorative
        
    def _get_alt_text(self, shape):
        """Legacy method for compatibility - just returns ALT text."""
        alt_text, _ = self._get_alt_text_and_decorative(shape)
        return alt_text
    
    def _extract_from_raw_xml(self, pptx_path, slide_num, shape_index):
        """
        Direct XML extraction as fallback.
        """
        import zipfile
        import xml.etree.ElementTree as ET
        
        alt_text = ''
        is_decorative = False
        
        try:
            with zipfile.ZipFile(pptx_path, 'r') as zip_file:
                # Read the slide XML
                slide_xml_path = f'ppt/slides/slide{slide_num}.xml'
                with zip_file.open(slide_xml_path) as slide_file:
                    slide_content = slide_file.read()
                    root = ET.fromstring(slide_content)
                    
                    # Define namespaces
                    ns = {
                        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        'a14': 'http://schemas.microsoft.com/office/drawing/2010/main',
                        'adec': 'http://schemas.microsoft.com/office/drawing/2017/decorative'
                    }
                    
                    # Find all pictures
                    pics = root.findall('.//p:pic', ns)
                    
                    if shape_index < len(pics):
                        pic = pics[shape_index]
                        
                        # Look for ALT text
                        descr = pic.find('.//p:cNvPr[@descr]', ns)
                        if descr is not None:
                            alt_text = descr.get('descr', '')
                        
                        # Check decorative in extensions
                        for ext in pic.findall('.//a:ext', ns):
                            if 'a14:decorative' in ET.tostring(ext, encoding='unicode'):
                                is_decorative = True
                                break
                                
        except Exception as e:
            self.logger.debug(f"Raw XML extraction failed: {e}")
        
        return alt_text, is_decorative

    def _collect_slide_context(self, presentation):
        """Collect speaker notes and visible text for each slide."""
        all_context = []
        for slide in presentation.slides:
            notes = ""
            text_content = []

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text or ""

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    text_content.append(shape.text.strip())

            slide_text = "\n".join(text_content).strip()
            all_context.append({"notes": notes.strip(), "text": slide_text})

        return all_context
    
    def _process_presentation(self, pptx_path):
        """Process a single PowerPoint presentation."""
        self.logger.info(f"Processing: {pptx_path.name}")
        
        try:
            presentation = Presentation(str(pptx_path))
            image_data = []
            pending_generation = []
            self.image_tracker.clear()

            all_context = self._collect_slide_context(presentation)
            
            # First pass: collect all images for duplicate detection
            for slide_num, slide in enumerate(presentation.slides, 1):
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_bytes = self._extract_image_data(shape)
                        if image_bytes:
                            image_hash = get_image_hash(image_bytes)
                            self.image_tracker[image_hash].append(slide_num)
            
            # Second pass: analyze images
            img_counter = 0
            for slide_num, slide in enumerate(presentation.slides, 1):
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_counter += 1
                        
                        # Extract image data
                        image_bytes = self._extract_image_data(shape)
                        if not image_bytes:
                            self.logger.warning(f"Skipping image on slide {slide_num} - no data")
                            continue
                        
                        # Determine file extension
                        try:
                            ext = f".{shape.image.ext.lower()}"
                        except Exception:
                            ext = None

                        image_filename = self._extract_embedded_image_filename(shape) or f"slide{slide_num}_picture{img_counter}{ext or ''}"

                        image_hash = get_image_hash(image_bytes)

                        if image_hash in thumbnail_cache:
                            thumbnail_path, thumbnail_filename = thumbnail_cache[image_hash]
                            if ext in ('.wmf', '.emf'):
                                width, height = 0, 0
                            else:
                                img = Image.open(BytesIO(image_bytes))
                                width, height = img.size
                                img.close()
                        else:
                            if ext in ('.wmf', '.emf'):
                                thumbnail_path = self._convert_vector_to_png(image_bytes, image_filename)
                                if not thumbnail_path:
                                    self.logger.error(
                                        f"Conversion failed for {image_filename}, skipping image"
                                    )
                                    continue
                                thumbnail_filename = Path(thumbnail_path).name
                                width, height = 0, 0
                            else:
                                img = Image.open(BytesIO(image_bytes))
                                width, height = img.size
                                # Always generate thumbnail early
                                thumbnail_path, thumbnail_filename = self._save_thumbnail(
                                    img,
                                    pptx_path.stem,
                                    slide_num,
                                    img_counter,
                                )
                                img.close()
                            thumbnail_cache[image_hash] = (thumbnail_path, thumbnail_filename)
                        
                        # Get ALT text using improved extraction
                        alt_text, is_pptx_decorative = self._get_alt_text_and_decorative(shape)

                        skip_alt_values = set(
                            val.strip().lower()
                            for val in self.config.get("reinjection", {}).get("skip_alt_text_if", [])
                        )
                        alt_text_clean = (alt_text or "").strip().lower()
                        alt_text_is_valid = alt_text_clean not in skip_alt_values

                        self.logger.debug(
                            f"ALT text raw: '{alt_text}', cleaned: '{alt_text_clean}', is_valid: {alt_text_is_valid}"
                        )

                        is_missing_alt = not alt_text_is_valid

                        # Extract shape name and run early decorative check
                        shape_name = shape.name.lower() if hasattr(shape, "name") else ""
                        is_force_decorative = is_force_decorative_by_filename_or_name(
                            image_filename,
                            shape_name,
                            self.config,
                        )
                        # Debugging aid: verify filename/shape-based decorative overrides
                        self.logger.debug(
                            "Decorative check: shape_name = %s, result = %s",
                            shape_name,
                            is_force_decorative,
                        )

                        # Check decorative heuristics using the decorative_filter module
                        is_heuristic_decorative, notes, image_hash, dimensions = self._check_decorative_heuristics(
                            shape,
                            slide,
                            image_bytes,
                            alt_text,
                            slide_num,
                            img_counter,
                            (width, height),
                            image_hash,
                        )

                        # Deterministic image identifier for round-trip matching
                        image_id = f"slide{slide_num}_shape{shape_name or img_counter}_hash{image_hash}"

                        # Combine all decorative checks
                        is_decorative = (
                            is_force_decorative
                            or is_heuristic_decorative
                            or is_pptx_decorative
                        )
                        suggested_alt_text = ""
                        tokens_used = None
                        generation_time = None
                        defer_generation = False
                        task_args = None
                        vague_warn = self.config.get("alt_text_handling", {}).get(
                            "vague_output_warning", False
                        )
                        
                        # Update notes if forced decorative
                        if is_force_decorative:
                            notes.insert(0, "Filename matches decorative rule")
                            self.logger.info(f"Image {img_counter} forced decorative by filename: {image_filename}")
                        
                        # Override with PowerPoint's decorative flag if set
                        if is_pptx_decorative:
                            if "Marked as decorative in PowerPoint" not in notes:
                                notes.append("Marked as decorative in PowerPoint")
                            self.logger.info(f"Image {img_counter} marked as decorative in PowerPoint")

                        # Generate suggested ALT text based on decorative status
                        if is_decorative:
                            suggested_alt_text = "[Decorative image - no ALT text needed]"
                            self.logger.info(
                                f"Slide {slide_num}, Image {img_counter}: Marked as decorative"
                            )
                        else:
                            slide_index = slide_num - 1
                            slide_context = all_context[slide_index]
                            # Use single generic prompt to eliminate complexity
                            base_prompt = "Describe this medical/educational image in one sentence for accessibility."
                            image_type = "default"
                            alt_handling = self.config.get("alt_text_handling", {}).get("mode", "preserve")
                            if alt_text_is_valid and alt_handling == "preserve":
                                suggested_alt_text = alt_text
                                self.logger.debug(
                                    f"Preserving existing ALT text: {alt_text}"
                                )
                            else:
                                alt_text = ""
                                self.logger.info(
                                    f"Slide {slide_num}, Image {img_counter}: Generating ALT text with LLaVA"
                                )
                                # Skip complex image type detection and use simple prompt
                                self.alt_generator.set_prompt_type(image_type)

                            # Use simple prompt without complex context building
                            full_prompt = base_prompt.strip()

                            if self.config.get("logging", {}).get("show_prompts"):
                                self.logger.debug(f"LLaVA Prompt:\n{full_prompt}")

                            # Simplify configuration - reduce complexity during debugging
                            reuse_enabled = self.config.get("alt_text_handling", {}).get("reuse_for_identical_images", True)
                            clean_enabled = False  # Disable cleaning to reduce complexity
                            vague_warn = False     # Disable vague warnings to reduce complexity
                            want_metadata = False  # Disable metadata to reduce complexity

                            if alt_text_is_valid and alt_handling == "preserve":
                                pass
                            elif reuse_enabled and image_hash in alt_cache:
                                suggested_alt_text = alt_cache[image_hash]
                                self.logger.debug(
                                    f"Reusing cached ALT for image hash: {image_hash}"
                                )
                            else:
                                defer_generation = True
                                task_args = (
                                    thumbnail_path,
                                    full_prompt,
                                    want_metadata,
                                    clean_enabled,
                                    vague_warn,
                                    image_hash,
                                )
                        
                        # Debug summary of generated ALT text and early decorative status
                        self.logger.debug(
                            "Image %d - alt_text: '%s...', suggested_alt: '%s...'",
                            img_counter,
                            alt_text[:50],
                            suggested_alt_text[:50],
                        )
                        self.logger.debug("Early decorative check → %s", is_force_decorative)
                        
                        # Build the complete image info dictionary
                        img_info = {
                            'filename': pptx_path.name,
                            'slide_number': slide_num,
                            'image_index': img_counter,
                            'image_number': img_counter,
                            'thumbnail_path': str(thumbnail_path),
                            'image_path': str(thumbnail_path),
                            'alt_text': alt_text,
                            'current_alt': alt_text,
                            'is_missing_alt': is_missing_alt,
                            'is_decorative': is_force_decorative,
                            'is_possible_decorative': is_decorative,
                            'suggested_alt': suggested_alt_text,
                            'notes': notes if isinstance(notes, list) else [str(notes)],
                            'decorative_notes': "; ".join(notes) if notes else "",
                            'token_count': tokens_used,
                            'generation_time': generation_time,
                            'image_hash': image_hash,
                            'image_dimensions': dimensions,
                            'image_type': image_type,
                            'position': self._get_shape_position(shape),
                            'original_filename': image_filename,
                            'image_filename': image_filename,
                            'shape_name': shape_name,
                            'image_id': image_id
                        }

                        self.logger.info(f"Processed image {img_counter}: decorative={is_decorative}, has_alt={not is_missing_alt}")
                        image_data.append(img_info)
                        if defer_generation and task_args:
                            pending_generation.append((len(image_data) - 1, task_args))
                        else:
                            if vague_warn and any(
                                p in suggested_alt_text.lower()
                                for p in ["slide from a presentation", "appears to be"]
                            ):
                                self.logger.warning(
                                    f"Vague ALT: '{suggested_alt_text}' for image hash {image_hash}"
                                )
                            self.logger.debug(
                                f"Generated ALT text via unified generator: {suggested_alt_text}"
                            )
            
            if pending_generation:
                asyncio.run(self._run_alt_generation(pending_generation, image_data))

            # Summary logging
            total_images = len(image_data)
            decorative_count = sum(1 for img in image_data if img['is_decorative'])
            missing_alt_count = sum(1 for img in image_data if img['is_missing_alt'] and not img['is_decorative'])

            self.logger.info(
                f"Presentation summary: {total_images} images, {decorative_count} decorative, {missing_alt_count} missing ALT text"
            )

            return image_data
            
        except Exception as e:
            self.logger.error(f"Error processing {pptx_path.name}: {e}")
            import traceback
            self.logger.error(traceback.format_exc())
            return []

    def _generate_alt_text_sync(
        self,
        thumbnail_path,
        full_prompt,
        want_metadata,
        clean_enabled,
        vague_warn,
        image_hash,
    ):
        """Generate ALT text synchronously in a worker thread."""
        if want_metadata:
            result = self.alt_generator.generate_alt_text(
                thumbnail_path,
                custom_prompt=full_prompt,
                return_metadata=True,
            )
            suggested_alt_text, metadata = (
                result if isinstance(result, tuple) else (result, {})
            )
            tokens_used = metadata.get("tokens_used")
            generation_time = metadata.get("generation_time")
        else:
            suggested_alt_text = self.alt_generator.generate_alt_text(
                thumbnail_path,
                custom_prompt=full_prompt,
            )
            tokens_used = None
            generation_time = None

        if not suggested_alt_text:
            suggested_alt_text = "[ALT text generation failed]"
        if clean_enabled:
            from alt_cleaner import clean_alt_text

            suggested_alt_text = clean_alt_text(suggested_alt_text)
        alt_cache[image_hash] = suggested_alt_text
        if vague_warn and any(
            p in suggested_alt_text.lower()
            for p in ["slide from a presentation", "appears to be"]
        ):
            self.logger.warning(
                f"Vague ALT: '{suggested_alt_text}' for image hash {image_hash}"
            )
        self.logger.debug(
            f"Generated ALT text via unified generator: {suggested_alt_text}"
        )
        return suggested_alt_text, tokens_used, generation_time

    async def _run_alt_generation(self, tasks, image_data):
        """Run ALT text generation tasks concurrently."""
        loop = asyncio.get_running_loop()
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            futures = [
                loop.run_in_executor(executor, self._generate_alt_text_sync, *args)
                for _, args in tasks
            ]
            results = await asyncio.gather(*futures)
        for (index, _), (alt, tokens, gen_time) in zip(tasks, results):
            image_data[index]["suggested_alt"] = alt
            image_data[index]["token_count"] = tokens
            image_data[index]["generation_time"] = gen_time

    def _write_csv_report(self, results, pptx_path):
        """Write results to CSV file."""
        if not results:
            self.logger.warning(f"No images found in {pptx_path.name}")
            return
        
        csv_filename = self.csv_dir / f"{pptx_path.stem}_accessibility_report.csv"
        
        try:
            with open(csv_filename, 'w', newline='', encoding='utf-8') as csvfile:
                fieldnames = [
                    'filename', 'slide_number', 'image_index', 'alt_text',
                    'is_missing_alt', 'thumbnail_path', 'original_filename',
                    'is_possible_decorative', 'notes',
                    'image_filename', 'shape_name', 'image_id'
                ]
                writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
                writer.writeheader()
                
                # Convert notes list to string for CSV
                for result in results:
                    csv_row = result.copy()
                    if isinstance(csv_row.get('notes'), list):
                        csv_row['notes'] = "; ".join(csv_row['notes'])
                    writer.writerow(csv_row)
            
            self.logger.info(f"Report saved: {csv_filename.name}")
            
            # Log summary statistics
            missing_alt_count = sum(1 for r in results if r['is_missing_alt'])
            decorative_count = sum(1 for r in results if r['is_possible_decorative'])
            self.logger.info(
                f"  - Total images: {len(results)}, "
                f"Missing ALT: {missing_alt_count}, "
                f"Possibly decorative: {decorative_count}"
            )
            
        except Exception as e:
            self.logger.error(f"Error writing CSV report: {e}")
    
    def audit_folder(self):
        """Audit all PowerPoint files in the input folder."""
        pptx_files = list(self.input_folder.glob("*.pptx"))
        
        if not pptx_files:
            self.logger.warning(f"No .pptx files found in {self.input_folder}")
            return
        
        self.logger.info(f"Found {len(pptx_files)} PowerPoint files to process")
        
        for pptx_path in pptx_files:
            # Skip temporary files
            if pptx_path.name.startswith('~$'):
                continue

            results = self._process_presentation(pptx_path)
            if results:
                self._write_csv_report(results, pptx_path)

        # Optionally persist ALT text cache
        if alt_cache:
            try:
                with open(self.alt_cache_path, "w", encoding="utf-8") as f:
                    json.dump(alt_cache, f, indent=2)
            except Exception as e:
                self.logger.warning(
                    f"Could not save ALT cache to {self.alt_cache_path}: {e}"
                )

        self.logger.info("Accessibility audit completed")


def main():
    """Main entry point."""
    import sys
    from config_manager import ConfigManager
    
    if len(sys.argv) < 2:
        print("Usage: python pptx_accessibility_audit.py <input_folder> [output_folder]")
        sys.exit(1)
    
    input_folder = sys.argv[1]
    output_folder = sys.argv[2] if len(sys.argv) > 2 else None  # Let config handle default
    
    # Initialize config manager with validation
    try:
        config_manager = ConfigManager()
    except ValueError as e:
        print(f"❌ Configuration Error: {e}")
        sys.exit(1)
    
    if not Path(input_folder).exists():
        print(f"Error: Input folder '{input_folder}' does not exist")
        sys.exit(1)

    auditor = PowerPointAccessibilityAuditor(config_manager, input_folder, output_folder)
    auditor.audit_folder()


if __name__ == "__main__":
    main()
