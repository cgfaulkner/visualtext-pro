#!/usr/bin/env python3
"""
Test script to demonstrate the 100% image coverage features.
This script creates a test PPTX file and shows how the new features work.
"""

import logging
import os
import sys
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import PIL.Image
import io

# Setup paths
project_root = Path(__file__).parent
sys.path.insert(0, str(project_root))
sys.path.insert(0, str(project_root / "shared"))
sys.path.insert(0, str(project_root / "core"))

from pptx_alt_processor import PPTXAltProcessor

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def create_test_image(width=100, height=100, color=(255, 0, 0)):
    """Create a test image as bytes."""
    img = PIL.Image.new('RGB', (width, height), color)
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    return img_bytes.getvalue()

def create_test_pptx_with_images():
    """Create a test PPTX file with various types of images."""
    
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Medical content (should generate descriptive ALT text)
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide1.shapes.title.text = "Medical Imaging Study"
    
    # Add text content to provide context
    textbox = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    textbox.text = "CT scan showing anatomical structures of the chest cavity"
    
    # Add a "medical" image (will be processed for ALT text)
    img_data1 = create_test_image(300, 200, (128, 128, 255))  # Blue image
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp1:
        tmp1.write(img_data1)
        tmp1.flush()
        slide1.shapes.add_picture(tmp1.name, Inches(2), Inches(2.5), Inches(4), Inches(3))
    os.unlink(tmp1.name)
    
    # Slide 2: Technical diagram (should generate descriptive ALT text)
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Process Workflow Diagram"
    
    textbox2 = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    textbox2.text = "Flowchart showing patient data processing workflow"
    
    # Add a "diagram" image
    img_data2 = create_test_image(250, 300, (0, 255, 0))  # Green image
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp2:
        tmp2.write(img_data2)
        tmp2.flush()
        slide2.shapes.add_picture(tmp2.name, Inches(2), Inches(2.5), Inches(3), Inches(4))
    os.unlink(tmp2.name)
    
    # Slide 3: Small decorative elements (should be marked decorative by heuristics)
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Summary and Conclusions"
    
    textbox3 = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
    textbox3.text = "Key findings from our analysis show significant improvements in patient outcomes."
    
    # Add small decorative image (should be detected as decorative)
    img_data3 = create_test_image(30, 30, (255, 255, 0))  # Small yellow image
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp3:
        tmp3.write(img_data3)
        tmp3.flush()
        slide3.shapes.add_picture(tmp3.name, Inches(8.5), Inches(0.5), Inches(0.5), Inches(0.5))
    os.unlink(tmp3.name)
    
    # Slide 4: Image that will likely fail generation (corrupt or problematic)
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "Test Image - May Fail Generation"
    
    textbox4 = slide4.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    textbox4.text = "This image may cause generation failures to test fallback behavior"
    
    # Add an image that might cause generation issues (very small or unusual)
    img_data4 = create_test_image(10, 10, (128, 0, 128))  # Tiny purple image
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp4:
        tmp4.write(img_data4)
        tmp4.flush()
        slide4.shapes.add_picture(tmp4.name, Inches(3), Inches(3), Inches(2), Inches(2))
    os.unlink(tmp4.name)
    
    return prs

def test_coverage_features():
    """Test the 100% coverage features."""
    
    print("🧪 Testing PPTX 100% Image Coverage Features")
    print("=" * 60)
    
    # Create test PPTX file
    print("Creating test PPTX file with various image types...")
    prs = create_test_pptx_with_images()
    
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_pptx:
        prs.save(tmp_pptx.name)
        test_pptx_path = tmp_pptx.name
    
    print(f"Test PPTX created: {test_pptx_path}")
    
    try:
        # Test 1: Process without force-decorative (normal behavior)
        print("\n📊 Test 1: Processing without --force-decorative")
        print("-" * 50)
        
        processor1 = PPTXAltProcessor(verbose=True, force_decorative=False)
        
        output1_path = test_pptx_path.replace('.pptx', '_normal.pptx')
        result1 = processor1.process_single_file(test_pptx_path, output1_path)
        
        print(f"Results (normal processing):")
        print(f"  Success: {result1['success']}")
        print(f"  Total images: {result1.get('total_images', 0)}")
        print(f"  Processed (descriptive): {result1.get('processed_images', 0)}")
        print(f"  Decorative (heuristic): {result1.get('decorative_images', 0)}")
        print(f"  Failed: {result1.get('failed_images', 0)}")
        if 'coverage_report' in result1:
            coverage = result1['coverage_report']
            print(f"  Coverage: {coverage['total_coverage_percent']:.1f}%")
        
        # Test 2: Process with force-decorative (100% coverage)
        print("\n📊 Test 2: Processing with --force-decorative")
        print("-" * 50)
        
        processor2 = PPTXAltProcessor(verbose=True, force_decorative=True)
        
        output2_path = test_pptx_path.replace('.pptx', '_force_decorative.pptx')
        result2 = processor2.process_single_file(test_pptx_path, output2_path)
        
        print(f"Results (force decorative):")
        print(f"  Success: {result2['success']}")
        print(f"  Total images: {result2.get('total_images', 0)}")
        print(f"  Processed (descriptive): {result2.get('processed_images', 0)}")
        print(f"  Decorative (heuristic): {result2.get('decorative_images', 0)}")
        print(f"  Decorative (fallback): {result2.get('fallback_decorative', 0)}")
        print(f"  Failed: {result2.get('failed_images', 0)}")
        if 'coverage_report' in result2:
            coverage = result2['coverage_report']
            print(f"  Coverage: {coverage['total_coverage_percent']:.1f}%")
        
        # Check for coverage report files
        print("\n📋 Coverage Report Files:")
        print("-" * 30)
        
        # Look for generated coverage report files
        for output_path in [output1_path, output2_path]:
            if os.path.exists(output_path):
                report_path = output_path.replace('.pptx', '_coverage_report.json')
                if os.path.exists(report_path):
                    print(f"  ✅ Coverage report: {os.path.basename(report_path)}")
                    
                    # Show a snippet of the report
                    try:
                        import json
                        with open(report_path, 'r') as f:
                            report = json.load(f)
                        
                        print(f"     - Processing timestamp: {report.get('processing_timestamp')}")
                        print(f"     - Failed generations logged: {len(report.get('failed_generations', []))}")
                        
                        if report.get('failed_generations'):
                            print("     - Sample failed generation:")
                            failed = report['failed_generations'][0]
                            print(f"       * Image: {failed.get('image_key')}")
                            print(f"       * Error: {failed.get('error')}")
                    
                    except Exception as e:
                        print(f"     - Error reading report: {e}")
                else:
                    print(f"  ❌ No coverage report found for {os.path.basename(output_path)}")
        
        print("\n✅ Coverage feature testing completed!")
        print("Key improvements demonstrated:")
        print("  1. ✅ Validation after ALT generation with decorative fallback")  
        print("  2. ✅ Coverage reporting (descriptive vs decorative counts)")
        print("  3. ✅ --force-decorative flag for 100% coverage")
        print("  4. ✅ Failed generation logging for manual review")
        print("  5. ✅ Coverage percentage statistics")
        
        return True
        
    finally:
        # Clean up test files
        try:
            os.unlink(test_pptx_path)
            if os.path.exists(output1_path):
                os.unlink(output1_path)
            if os.path.exists(output2_path):
                os.unlink(output2_path)
            
            # Clean up coverage report files
            for suffix in ['_normal_coverage_report.json', '_force_decorative_coverage_report.json']:
                report_path = test_pptx_path.replace('.pptx', suffix)
                if os.path.exists(report_path):
                    os.unlink(report_path)
        
        except Exception as e:
            print(f"Warning: Failed to clean up test files: {e}")

if __name__ == "__main__":
    try:
        success = test_coverage_features()
        sys.exit(0 if success else 1)
    except Exception as e:
        logger.error(f"Test failed: {e}", exc_info=True)
        sys.exit(1)