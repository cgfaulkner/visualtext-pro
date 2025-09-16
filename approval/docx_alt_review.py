# approval/docx_alt_review.py
from pathlib import Path
from datetime import datetime
import logging
import re

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_ORIENT
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

from shared.pipeline_artifacts import normalize_final_alt_map

def set_cell_margins(cell, left=None, right=None, top=None, bottom=None):
    """Set cell margins in twips (1/20 pt). 180 ≈ 0.125"."""
    tcPr = cell._tc.get_or_add_tcPr()
    tcMar = OxmlElement('w:tcMar')
    for side, val in (('top', top), ('start', left), ('bottom', bottom), ('end', right)):
        if val is None:
            continue
        el = OxmlElement(f'w:{side}')
        el.set(qn('w:w'), str(val))
        el.set(qn('w:type'), 'dxa')
        tcMar.append(el)
    tcPr.append(tcMar)

def force_portrait_with_sane_geometry(doc):
    """Force portrait and set sane page geometry."""
    section = doc.sections[0]
    section.orientation = WD_ORIENT.PORTRAIT
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.left_margin = Inches(0.5)
    section.right_margin = Inches(0.5) 
    section.top_margin = Inches(0.5)
    section.bottom_margin = Inches(0.5)

def collect_original_alts(pptx_path) -> dict:
    """Collect original ALT text from PPTX before injection."""
    from pptx import Presentation
    try:
        prs = Presentation(pptx_path)
        cache = {}
        for s_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                try:
                    # Create same stable key as used in processing
                    if hasattr(shape, 'shape_id'):
                        shape_id = shape.shape_id
                    else:
                        continue  # Skip shapes without stable ID
                    
                    # Try to get image hash for complete key
                    image_hash = ""
                    try:
                        if hasattr(shape, 'image') and hasattr(shape.image, 'blob'):
                            from hashlib import md5
                            image_hash = md5(shape.image.blob).hexdigest()[:8]
                    except:
                        pass
                    
                    key = f"slide_{s_idx}_shapeid_{shape_id}_hash_{image_hash}"
                    
                    # Get ALT text - try alternative_text first, then title
                    alt = ""
                    try:
                        alt = getattr(shape, "alternative_text", "") or ""
                    except:
                        alt = ""
                    if not alt:
                        try:
                            alt = getattr(shape, "title", "") or ""
                        except:
                            pass
                    
                    cache[key] = alt.strip()
                except Exception:
                    continue
        return cache
    except Exception as e:
        print(f"Warning: Could not collect original ALT text: {e}")
        return {}

def force_print_layout(doc, zoom_percent=110):
    """Force Word to open in Print Layout mode with specified zoom."""
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    
    settings = doc._part.document.settings._element
    view = settings.find(qn('w:view'))
    if view is None:
        view = OxmlElement('w:view')
        settings.append(view)
    view.set(qn('w:val'), 'print')
    
    zoom = settings.find(qn('w:zoom'))
    if zoom is None:
        zoom = OxmlElement('w:zoom')
        settings.append(zoom)
    zoom.set(qn('w:percent'), str(zoom_percent))

def repeat_header_on_each_page(row):
    """Make table header row repeat on every page."""
    tr = row._tr
    trPr = tr.get_or_add_trPr()
    tblHeader = OxmlElement('w:tblHeader')
    tblHeader.set(qn('w:val'), "true")
    trPr.append(tblHeader)

def set_row_no_break(row):
    """Prevent row from breaking across pages."""
    trPr = row._tr.get_or_add_trPr()
    cantSplit = OxmlElement('w:cantSplit')
    trPr.append(cantSplit)

def shade_cell(cell, hex_color):
    """Apply background shading to a cell."""
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def set_font_properties(run, size=11, bold=False, italic=False, color=None):
    """Set font properties for a run."""
    run.font.name = 'Calibri'
    run.font.size = Pt(size)
    run.bold = bold
    run.italic = italic
    if color:
        run.font.color.rgb = RGBColor.from_string(color)

def add_header_footer(doc, title, input_filename):
    """Add professional header and footer."""
    for section in doc.sections:
        # Header
        header = section.header
        header_para = header.paragraphs[0]
        header_para.clear()
        
        # Left-aligned title
        title_run = header_para.add_run(f"{title} – ALT Review")
        set_font_properties(title_run, size=11, bold=True)
        
        # Right-aligned date/time
        header_para.add_run('\t' * 6)  # Tab to right side for portrait
        date_run = header_para.add_run(datetime.now().strftime("%Y-%m-%d %H:%M"))
        set_font_properties(date_run, size=10)
        
        header_para.paragraph_format.tab_stops.add_tab_stop(Inches(6.5))  # Portrait tab stop
        
        # Footer
        footer = section.footer
        footer_para = footer.paragraphs[0]
        footer_para.clear()
        
        # Left-aligned filename
        filename_run = footer_para.add_run(input_filename)
        set_font_properties(filename_run, size=10, color="666666")
        
        # Right-aligned page numbers
        footer_para.add_run('\t' * 6)
        page_run = footer_para.add_run("Page ")
        set_font_properties(page_run, size=10, color="666666")
        
        footer_para.paragraph_format.tab_stops.add_tab_stop(Inches(6.5))  # Portrait tab stop

# ALT text skip set from injector preserve mode
SKIP_ALT_VALUES = {"", "undefined", "(None)", "N/A", "Not reviewed", "n/a"}

def get_current_alt_from_shape(shape):
    """Extract current ALT text from shape's cNvPr element."""
    try:
        if hasattr(shape, '_element'):
            # Look for cNvPr in picture elements
            cNvPr_elements = shape._element.xpath('.//a:cNvPr', namespaces={
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
            })
            if cNvPr_elements:
                cNvPr = cNvPr_elements[0]
                # descr is the actual ALT text; title is a fallback
                alt_text = cNvPr.get('descr') or cNvPr.get('title') or ''
                return alt_text.strip()
    except Exception:
        pass
    return ""

def process_alt_for_report(current_alt, generate_alt_func=None):
    """
    Process ALT text for report based on current ALT and skip logic.
    
    Args:
        current_alt: Current ALT text from shape
        generate_alt_func: Function to generate new ALT text (called only if needed)
    
    Returns:
        Tuple of (current_alt_display, suggested_alt_display)
    """
    current_cleaned = (current_alt or "").strip()
    has_current = current_cleaned and current_cleaned not in SKIP_ALT_VALUES
    
    if has_current:
        # Current ALT exists and is valid - don't generate, just mirror
        return current_cleaned, current_cleaned
    else:
        # Current ALT is missing or in skip set - generate suggested
        suggested = ""
        if generate_alt_func:
            try:
                suggested = generate_alt_func() or ""
            except Exception:
                suggested = ""
        
        return "[No ALT text]", normalize_alt_text(suggested)

def normalize_alt_text(text):
    """Ensure ALT text ends with proper punctuation."""
    if not text or not text.strip():
        return text
    text = text.strip()
    if text and text[-1] not in ('.', '!', '?'):
        text += '.'
    return text

def pick_alts_for_report(key, final_alt_map, existing_alt_from_ppt):
    """
    - If the PPT had ALT already, use it for BOTH Current and Suggested
      (so they match exactly).
    - If the PPT had none, Suggested = generated canonical ALT; Current = "".
    """
    existing = (existing_alt_from_ppt or "").strip()
    record = final_alt_map.get(key, {}) if isinstance(final_alt_map, dict) else {}

    if not existing and isinstance(record, dict):
        existing = (record.get('existing_alt') or "").strip()

    final_alt = (record.get('final_alt') or "").strip() if isinstance(record, dict) else ""
    generated = (record.get('generated_alt') or "").strip() if isinstance(record, dict) else ""

    suggested = final_alt or generated

    if existing:
        return existing, existing  # Current, Suggested
    return "", suggested

def generate_alt_review_doc(processed_images, lecture_title: str, output_path: str, original_pptx_path: str = None, final_alt_map: dict = None, status_map: dict = None):
    """
    Generate a professional matrix-style accessibility review document.
    
    processed_images: list[dict] with keys:
      image_path, slide_number, image_number, current_alt, suggested_alt,
      is_decorative, image_key, (optional) slide_title, slide_notes, status_info
    """
    if final_alt_map:
        final_alt_map = normalize_final_alt_map(final_alt_map)

    doc = Document()
    
    # Force portrait orientation with sane geometry and print layout
    force_portrait_with_sane_geometry(doc)
    force_print_layout(doc, zoom_percent=110)
    
    # Collect original ALT text if we have the source PPTX path
    original_alt_cache = {}
    if original_pptx_path:
        original_alt_cache = collect_original_alts(original_pptx_path)
    
    # Set default style to Calibri 11pt
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    
    # Add header and footer
    input_filename = Path(output_path).stem.replace('_ALT_Review', '')
    add_header_footer(doc, lecture_title, input_filename)
    
    # Title
    title_para = doc.add_paragraph()
    title_run = title_para.add_run(f"{lecture_title} – Accessibility ALT Review")
    set_font_properties(title_run, size=16, bold=True)
    title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    title_para.paragraph_format.space_after = Pt(8)
    
    # Add fallback policy info to header
    try:
        # Try to get policy from config or status_map
        fallback_policy = "none"  # default
        if status_map:
            # Look for policy info in status_map
            for status_info in status_map.values():
                if isinstance(status_info, dict) and 'policy' in status_info:
                    fallback_policy = status_info['policy']
                    break
        
        policy_para = doc.add_paragraph()
        policy_run = policy_para.add_run(f"Fallback Policy: {fallback_policy}")
        set_font_properties(policy_run, size=10, italic=True, color="666666")
        policy_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        policy_para.paragraph_format.space_after = Pt(12)
    except Exception:
        # Just add a small spacer if policy info not available
        spacer_para = doc.add_paragraph()
        spacer_para.paragraph_format.space_after = Pt(4)
    
    # Intro paragraph
    intro_para = doc.add_paragraph()
    intro_run = intro_para.add_run(
        "Review Suggested ALT, check Decorative as needed, and add Review Notes. "
        "Leave Approved ALT blank if Decorative is checked."
    )
    set_font_properties(intro_run, size=11)
    intro_para.paragraph_format.space_after = Pt(8)
    
    # Process ALT text for all images with corrected logic
    processed_alt_data = []
    identical_count = 0
    missing_current = 0
    
    # Debug: Log key mappings to help with alignment
    logger = logging.getLogger(__name__)
    logger.info(f"Processing {len(processed_images)} images for review document")
    if final_alt_map:
        logger.info(f"final_alt_map contains {len(final_alt_map)} entries")
        logger.debug(f"final_alt_map keys: {list(final_alt_map.keys())[:5]}...")  # Show first 5 keys
    if status_map:
        logger.info(f"status_map contains {len(status_map)} entries")
        logger.debug(f"status_map keys: {list(status_map.keys())[:5]}...")  # Show first 5 keys
    
    for item in processed_images:
        # Get the image key and look up original ALT text
        image_key = item.get('image_key', '')
        original_alt = original_alt_cache.get(image_key, '').strip() if image_key else ''
        
        # Debug: Log key for troubleshooting
        logger.debug(f"Processing image_key: {image_key}")
        
        # Use canonical final_alt_map if provided, otherwise fall back to processed_images
        if final_alt_map:
            current_alt, suggested_alt = pick_alts_for_report(image_key, final_alt_map, original_alt)
        else:
            # Fallback to old logic for backward compatibility
            generated_alt = item.get('suggested_alt', '').strip()
            current_alt = original_alt if original_alt else ""
            suggested_alt = original_alt if original_alt else (generated_alt if generated_alt else "")
        
        # For display
        current_display = current_alt if current_alt else "[No ALT text]"
        suggested_display = suggested_alt if suggested_alt else "[No ALT text]"
        
        # Track statistics
        if current_display == "[No ALT text]":
            missing_current += 1
        if (current_display != "[No ALT text]" and 
            suggested_display != "[No ALT text]" and 
            current_display == suggested_display):
            identical_count += 1
            
        # Store processed data for table building
        processed_item = item.copy()
        processed_item['current_alt_display'] = current_display
        processed_item['suggested_alt_display'] = suggested_display
        processed_alt_data.append(processed_item)
    
    total_images = len(processed_images)
    
    # Calculate bucket totals for honest reporting
    bucket_counts = {
        'preserved': 0, 'generated': 0, 'needs_alt': 0, 
        'fallback_injected': 0, 'decorative': 0
    }
    
    for item in processed_alt_data:
        current_alt = item.get('current_alt_display', '')
        suggested_alt = item.get('suggested_alt_display', '')
        
        # Derive status for counting based on same logic as status derivation
        if item.get('is_decorative', False):
            bucket_counts['decorative'] += 1
        elif current_alt != "[No ALT text]":
            if suggested_alt not in ["[No suggestion]", "[No ALT text]", ""] and current_alt == suggested_alt:
                bucket_counts['generated'] += 1
            else:
                bucket_counts['preserved'] += 1
        elif suggested_alt.startswith("FALLBACK:"):
            bucket_counts['fallback_injected'] += 1
        else:
            bucket_counts['needs_alt'] += 1
    
    # Summary with bucket totals
    summary_para = doc.add_paragraph()
    bucket_text = (f"Preserved: {bucket_counts['preserved']} • Generated: {bucket_counts['generated']} • "
                  f"Fallback: {bucket_counts['fallback_injected']} • Needs ALT: {bucket_counts['needs_alt']} • "
                  f"Decorative: {bucket_counts['decorative']}")
    summary_run = summary_para.add_run(
        f"Total elements: {total_images}   |   {bucket_text}"
    )
    set_font_properties(summary_run, size=10, color="666666")
    summary_para.paragraph_format.space_after = Pt(12)
    
    # Fixed table with optimized columns for portrait layout
    table = doc.add_table(rows=1, cols=6)
    table.autofit = False  # Critical: let our fixed widths win
    
    # Columns: [Slide / ID, Thumbnail, Current, Suggested, Status, Decor.]  
    # 8.5" page - 1.0" margins = 7.5" usable width, optimized for shorter headers
    col_widths = [Inches(0.8), Inches(0.8), Inches(2.2), Inches(2.2), Inches(1.0), Inches(0.5)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
        for cell in table.columns[i].cells:
            cell.width = w
    
    hdr = table.rows[0].cells
    # Shortened labels to prevent wrap
    header_names = ["Slide / ID", "Thumbnail", "Current", "Suggested", "Status", "Decor."]
    
    for i, name in enumerate(header_names):
        hdr[i].text = name
        # Header formatting to prevent letter-wrapping
        p = hdr[i].paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        hdr[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        # Tighten header font to prevent wrapping
        for run in p.runs:
            run.font.size = Pt(10)
        p.paragraph_format.keep_together = True
        p.paragraph_format.keep_with_next = True
        
        # Header styling
        for r in p.runs:
            set_font_properties(r, size=9.5, bold=True, color="FFFFFF")
        
        # Dark gray header background
        shade_cell(hdr[i], "4F4F4F")
    
    # Make header repeat on every page
    repeat_header_on_each_page(table.rows[0])
    
    # Data rows with improved cell formatting
    for row_idx, item in enumerate(processed_alt_data):
        row = table.add_row()
        set_row_no_break(row)
        cells = row.cells
        
        # Apply zebra striping (light gray on even rows)
        if row_idx % 2 == 0:
            for cell in cells:
                shade_cell(cell, "F7F7F7")
        
        # Set all cell widths and basic formatting
        for i, cell in enumerate(cells):
            cell.width = col_widths[i]
            cell.vertical_alignment = WD_ALIGN_VERTICAL.TOP  # Top-left align body cells
        
        # Column 0: Slide / Img (combined)
        cell = cells[0]
        para = cell.paragraphs[0]
        para.clear()
        slide_num = item.get('slide_number', '')
        image_num = item.get('image_number', '')
        run = para.add_run(f"S{slide_num} / I{image_num}")
        set_font_properties(run, size=11, bold=True)
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        para.paragraph_format.space_after = Pt(6)
        para.paragraph_format.line_spacing = 1.0
        
        # Column 1: Thumbnail with improved sizing and right padding
        cell = cells[1]
        cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER  # Center thumbnail vertically
        set_cell_margins(cell, right=180)  # ~0.125" gutter on the right
        para = cell.paragraphs[0]
        para.clear()
        
        thumb_path = item.get('image_path')
        if thumb_path and Path(thumb_path).exists():
            try:
                # Add thumbnail with constrained sizing
                run = para.add_run()
                pic = run.add_picture(thumb_path, width=Inches(1.40))  # a hair smaller than the column
                    
                # Add ALT text to picture for accessibility
                pic._inline.graphic.graphicData.pic.nvPicPr.cNvPr.set('descr', 
                    f"Thumbnail of slide {item.get('slide_number', '')} image {item.get('image_number', '')}")
            except Exception:
                run = para.add_run("(thumbnail unavailable)")
                set_font_properties(run, size=10, italic=True, color="999999")
        else:
            run = para.add_run("(no thumbnail)")
            set_font_properties(run, size=10, italic=True, color="999999")
        
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        para.paragraph_format.space_after = Pt(6)
        para.paragraph_format.line_spacing = 1.0
        
        # Column 2: Current ALT Text (using processed data)
        cell = cells[2]
        para = cell.paragraphs[0]
        para.clear()
        para.alignment = WD_ALIGN_PARAGRAPH.LEFT  # Top-left align
        para.paragraph_format.space_after = Pt(6)
        para.paragraph_format.line_spacing = 1.0
        
        current_alt_display = item.get('current_alt_display', '[No ALT text]')
        
        if current_alt_display == "[No ALT text]":
            run = para.add_run(current_alt_display)
            set_font_properties(run, size=10.5, color="999999")
            # Highlight missing current ALT with yellow background
            shade_cell(cell, "FFF2CC")
        else:
            run = para.add_run(current_alt_display)
            set_font_properties(run, size=10.5)
        
        # Column 3: Suggested ALT Text (using processed data)
        cell = cells[3]
        para = cell.paragraphs[0]
        para.clear()
        para.alignment = WD_ALIGN_PARAGRAPH.LEFT  # Top-left align
        para.paragraph_format.space_after = Pt(6)
        para.paragraph_format.line_spacing = 1.0
        
        suggested_alt_display = item.get('suggested_alt_display', '[No suggestion]')
        
        if not suggested_alt_display or suggested_alt_display == "[No suggestion]":
            run = para.add_run("[No suggestion]")
            set_font_properties(run, size=10.5, color="999999")
        else:
            run = para.add_run(suggested_alt_display)
            set_font_properties(run, size=10.5)
            
            # Highlight identical current/suggested with green background
            if (current_alt_display != "[No ALT text]" and 
                current_alt_display == suggested_alt_display):
                shade_cell(cell, "E2F0D9")
        
        # Column 4: Status information
        cell = cells[4]
        para = cell.paragraphs[0]
        para.clear()
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        para.paragraph_format.space_after = Pt(6)
        para.paragraph_format.line_spacing = 1.0
        
        # Get status information from item, status_map, or derive it
        status_info = item.get('status_info', {})
        status_display = status_info.get('status', '')
        
        # Try to get status from status_map using image_key
        if not status_display and status_map and image_key:
            status_info = status_map.get(image_key, {})
            status_display = status_info.get('status', '')
        
        # If no status provided, derive it from current state
        if not status_display:
            status_display, status_text = _derive_status_from_current_state(
                current_alt_display, suggested_alt_display, item
            )
        else:
            # Map provided status to user-friendly text
            status_text = _map_status_to_display(status_display)
        
        run = para.add_run(status_text)
        
        # Color code the status based on mapped display text
        status_lower = status_text.lower()
        if 'needs alt' in status_lower:
            set_font_properties(run, size=9, color="CC0000")  # Red for needs attention
            shade_cell(cell, "FFE6E6")
        elif 'fallback' in status_lower:
            set_font_properties(run, size=9, color="FF8C00")  # Orange for fallback
            shade_cell(cell, "FFF2E6")
        elif 'decorative' in status_lower:
            set_font_properties(run, size=9, color="888888")  # Gray for decorative
            shade_cell(cell, "F0F0F0")
        elif 'preserved' in status_lower:
            set_font_properties(run, size=9, color="0066CC")  # Blue for preserved
        elif 'generated' in status_lower:
            set_font_properties(run, size=9, color="006600")  # Green for generated
        else:
            set_font_properties(run, size=9, color="000000")  # Black for unknown
        
        # Column 5: Decorative checkbox (single line, centered)
        cell = cells[5]
        para = cell.paragraphs[0]
        para.clear()
        run = para.add_run("Yes" if item.get('is_decorative', False) else "No")
        set_font_properties(run, size=11)
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        para.paragraph_format.space_after = Pt(6)
        para.paragraph_format.line_spacing = 1.0
    
    # Final portrait enforcement (catches any stray landscape sections)
    force_portrait_with_sane_geometry(doc)
    
    # Save document
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    doc.save(output_path)
    return output_path


def _map_status_to_display(status_display: str) -> str:
    """Map internal status codes to user-friendly display text."""
    status_mapping = {
        'preserved': 'Preserved',
        'generated': 'Generated', 
        'fallback_injected': 'Fallback (PPT-gated)',
        'needs_alt': 'Needs ALT',
        'decorative_skipped': 'Decorative'
    }
    
    # Handle various status formats
    if status_display in status_mapping:
        return status_mapping[status_display]
    elif status_display.startswith('NEEDS ALT'):
        return 'Needs ALT'
    elif status_display.startswith('AUTO-LOWCONF'):
        return 'Fallback (doc-only)'
    elif status_display == 'FALLBACK_INJECTED':
        return 'Fallback (PPT-gated)'
    else:
        return status_display.replace('_', ' ').title()


def _derive_status_from_current_state(current_alt_display: str, suggested_alt_display: str, item: dict) -> tuple[str, str]:
    """
    Derive status when not provided by examining current ALT state.
    Returns (internal_status, display_text)
    """
    # Check if decorative
    if item.get('is_decorative', False):
        return 'decorative_skipped', 'Decorative'
    
    # Check if current ALT exists
    has_current = current_alt_display != "[No ALT text]"
    has_suggested = suggested_alt_display not in ["[No suggestion]", "[No ALT text]", ""]
    
    if has_current:
        # Check if it's a generic PowerPoint fallback
        current_text = current_alt_display.lower()
        if any(phrase in current_text for phrase in [
            "this is a powerpoint shape", 
            "image of", 
            "picture", 
            "graphic",
            "unknown"
        ]):
            return 'fallback_injected', 'Fallback (PPT-gated)'
        
        # Check if suggested matches current (generated) or different (preserved)
        if has_suggested and current_alt_display == suggested_alt_display:
            return 'generated', 'Generated'
        else:
            return 'preserved', 'Preserved'
    
    elif has_suggested:
        # Current is empty but suggestion exists
        if suggested_alt_display.startswith("FALLBACK:"):
            return 'needs_alt', 'Fallback (doc-only)'
        else:
            return 'needs_alt', 'Needs ALT'
    
    else:
        # No current ALT and no suggestion
        return 'needs_alt', 'Needs ALT'