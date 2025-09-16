"""Tests for preserve-mode overrides when final ALT text is supplied."""

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from core.pptx_alt_injector import PPTXAltTextInjector


def test_final_alt_overrides_existing_in_preserve_mode(tmp_path) -> None:
    """Editing final_alt should rewrite ALT text even if one already exists."""

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(2),
    )

    existing_text = "Generated description of a rectangle."
    cnvpr = shape._element.find(
        './/{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr'
    )
    assert cnvpr is not None
    cnvpr.set('descr', existing_text)

    pptx_path = tmp_path / "source.pptx"
    prs.save(pptx_path)

    injector = PPTXAltTextInjector()

    key_presentation = Presentation(str(pptx_path))
    identifier_map = injector._build_image_identifier_mapping(key_presentation)
    del key_presentation
    assert identifier_map, "Expected the presentation to expose a shape identifier"
    image_key = next(iter(identifier_map.keys()))

    updated_text = "Updated reviewer-provided description."
    alt_mapping = {
        image_key: {
            'existing_alt': existing_text,
            'generated_alt': existing_text,
            'final_alt': updated_text,
            'decision': None,
            'existing_meaningful': True,
            'source_existing': 'pptx',
            'source_generated': 'llava',
        }
    }

    output_path = tmp_path / "output.pptx"
    result = injector.inject_alt_text_from_mapping(
        pptx_path=str(pptx_path),
        alt_text_mapping=alt_mapping,
        output_path=str(output_path),
    )

    updated_presentation = Presentation(str(output_path))
    updated_shape = updated_presentation.slides[0].shapes[0]
    updated_alt = injector._read_current_alt(updated_shape)

    assert updated_alt == updated_text
    assert result['statistics']['written_final'] == 1
