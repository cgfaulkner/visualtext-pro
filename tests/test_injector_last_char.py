"""Tests for verifying ALT text injection writes final character correctly."""

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from core.pptx_alt_injector import PPTXAltTextInjector


def test_injected_text_last_character_matches() -> None:
    """Ensure the last character of written ALT matches the final ALT text."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1), Inches(1), Inches(1), Inches(1)
    )

    injector = PPTXAltTextInjector()
    raw_text = "Example ALT text."
    final_alt = injector._apply_final_normalization_gate(raw_text, "test_key", "unit_test")
    injector._inject_alt(shape, raw_text, "test_key", "unit_test")
    written_descr = injector._read_current_alt(shape)

    assert written_descr[-1] == final_alt[-1]
