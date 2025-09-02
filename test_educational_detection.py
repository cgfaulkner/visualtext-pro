#!/usr/bin/env python3
"""
Test script to verify that title areas, headers, and educational content are properly protected from decorative marking.
"""

import logging
import sys
from pathlib import Path

# Add the core directory to the Python path
sys.path.insert(0, str(Path(__file__).parent / "core"))

from pptx_processor import PPTXAccessibilityProcessor, PPTXShapeInfo

def setup_logging():
    """Setup detailed debug logging."""
    logging.basicConfig(
        level=logging.DEBUG,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        handlers=[
            logging.StreamHandler(sys.stdout),
            logging.FileHandler('educational_detection_test.log')
        ]
    )

def test_title_protection():
    """Test that title areas and headers are protected."""
    setup_logging()
    logger = logging.getLogger(__name__)
    
    # Find the test PPTX file
    test_file = Path("Documents to Review") / "test1_llava_latest_backup test names.pptx"
    
    if not test_file.exists():
        logger.error(f"Test file not found: {test_file}")
        return False
    
    logger.info("🧪 Testing title area and educational content protection")
    logger.info(f"📁 Test file: {test_file}")
    
    try:
        processor = PPTXAccessibilityProcessor(debug=True)
        
        # Load the presentation
        from pptx import Presentation
        presentation = Presentation(str(test_file))
        
        logger.info(f"📊 Loaded presentation with {len(presentation.slides)} slides")
        
        # Test title/header detection directly
        logger.info("🔍 Testing title/header detection on all shapes...")
        
        title_count = 0
        text_box_count = 0
        
        for slide_idx, slide in enumerate(presentation.slides):
            logger.info(f"📄 Analyzing slide {slide_idx + 1}...")
            
            slide_text = processor._extract_slide_text(slide) if processor.include_slide_text else ""
            
            for shape_idx, shape in enumerate(slide.shapes):
                try:
                    shape_info = PPTXShapeInfo(shape, slide_idx, shape_idx, slide_text)
                    
                    # Test title detection
                    is_title = processor._is_title_or_header_area(shape_info, debug=True, indent="  ")
                    if is_title:
                        title_count += 1
                        logger.info(f"  ✅ TITLE DETECTED: {shape_info.shape_type_name} '{shape_info.shape_name}' ({shape_info.width_px}x{shape_info.height_px}px)")
                    
                    # Test text box detection
                    is_text_box = processor._is_text_box_with_content(shape_info, debug=True, indent="  ")
                    if is_text_box:
                        text_box_count += 1
                        logger.info(f"  ✅ TEXT BOX DETECTED: {shape_info.shape_type_name} with text: '{shape_info.text_content[:50]}...'")
                        
                    # Test educational content detection
                    is_educational = processor._is_educational_shape(shape_info, debug=True, indent="  ")
                    if is_educational:
                        logger.info(f"  ✅ EDUCATIONAL CONTENT: {shape_info.shape_type_name}")
                    
                    # For shapes with text, verify they're protected
                    if shape_info.has_text:
                        is_decorative = processor._is_shape_decorative(shape_info, debug=True, indent="  ")
                        if is_decorative:
                            logger.warning(f"  ⚠️ TEXT SHAPE MARKED DECORATIVE: {shape_info.shape_type_name} with text: '{shape_info.text_content[:30]}...'")
                        else:
                            logger.info(f"  ✅ TEXT SHAPE PROTECTED: {shape_info.shape_type_name} with text: '{shape_info.text_content[:30]}...'")
                    
                except Exception as e:
                    logger.warning(f"Error analyzing shape {shape_idx} on slide {slide_idx + 1}: {e}")
                    continue
        
        logger.info(f"📊 SUMMARY:")
        logger.info(f"  Titles/Headers detected: {title_count}")
        logger.info(f"  Text boxes detected: {text_box_count}")
        
        return True
        
    except Exception as e:
        logger.error(f"💥 Test failed with exception: {e}", exc_info=True)
        return False

if __name__ == "__main__":
    print("🧪 Testing educational content and title protection...")
    success = test_title_protection()
    if success:
        print("✅ Test completed successfully!")
        print("📄 Check 'educational_detection_test.log' for detailed output")
    else:
        print("❌ Test failed!")
        sys.exit(1)