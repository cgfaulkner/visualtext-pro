#!/usr/bin/env python3
"""
Create a sample medical PPTX presentation with realistic images for end-to-end testing.
This creates actual embedded images that can be processed by the ALT text system.
"""

import logging
import sys
from pathlib import Path
from io import BytesIO
import base64

# Third-party imports for PPTX creation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from PIL import Image, ImageDraw, ImageFont
    PPTX_AVAILABLE = True
    PIL_AVAILABLE = True
except ImportError as e:
    PPTX_AVAILABLE = False
    PIL_AVAILABLE = False
    ERROR_MSG = str(e)

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def create_medical_image(width: int, height: int, image_type: str, color_scheme: tuple = (255, 255, 255)) -> BytesIO:
    """
    Create a realistic medical image using PIL.
    
    Args:
        width: Image width in pixels
        height: Image height in pixels
        image_type: Type of medical image ('heart', 'xray', 'brain', 'logo')
        color_scheme: RGB color scheme
        
    Returns:
        BytesIO: Image data as bytes
    """
    # Create image with background
    if image_type == 'xray':
        # Dark background for X-ray
        img = Image.new('RGB', (width, height), (25, 25, 35))
    elif image_type == 'logo':
        # White background for logo
        img = Image.new('RGB', (width, height), (255, 255, 255))
    else:
        # Light background for diagrams
        img = Image.new('RGB', (width, height), (240, 240, 250))
    
    draw = ImageDraw.Draw(img)
    
    try:
        # Try to use a system font, fallback to default
        font_large = ImageFont.truetype("Arial", 24)
        font_medium = ImageFont.truetype("Arial", 16) 
        font_small = ImageFont.truetype("Arial", 12)
    except:
        font_large = ImageFont.load_default()
        font_medium = ImageFont.load_default()
        font_small = ImageFont.load_default()
    
    if image_type == 'heart':
        # Create anatomical heart diagram
        center_x, center_y = width // 2, height // 2
        
        # Heart outline (simplified)
        heart_points = [
            (center_x, center_y - 50),  # Top
            (center_x - 40, center_y - 30),  # Left curve
            (center_x - 30, center_y + 20),  # Left bottom
            (center_x, center_y + 60),  # Bottom point
            (center_x + 30, center_y + 20),  # Right bottom
            (center_x + 40, center_y - 30),  # Right curve
            (center_x, center_y - 50)   # Back to top
        ]
        
        # Draw heart shape
        draw.polygon(heart_points, fill=(200, 100, 100), outline=(150, 50, 50))
        
        # Add chambers
        draw.ellipse((center_x - 25, center_y - 25, center_x - 5, center_y + 5), 
                    fill=(180, 80, 80), outline=(120, 40, 40))
        draw.ellipse((center_x + 5, center_y - 25, center_x + 25, center_y + 5),
                    fill=(180, 80, 80), outline=(120, 40, 40))
        
        # Labels
        draw.text((center_x - 60, center_y - 40), "LA", fill=(0, 0, 0), font=font_small)
        draw.text((center_x + 40, center_y - 40), "RA", fill=(0, 0, 0), font=font_small)
        draw.text((center_x - 60, center_y + 10), "LV", fill=(0, 0, 0), font=font_small)
        draw.text((center_x + 40, center_y + 10), "RV", fill=(0, 0, 0), font=font_small)
        
        # Title
        draw.text((10, 10), "Human Heart Anatomy", fill=(0, 0, 0), font=font_medium)
        
    elif image_type == 'xray':
        # Create chest X-ray simulation
        center_x, center_y = width // 2, height // 2
        
        # Chest outline
        draw.ellipse((center_x - 80, center_y - 60, center_x + 80, center_y + 80), 
                    fill=(80, 80, 90), outline=(120, 120, 140))
        
        # Ribs (simplified)
        for i in range(6):
            y_offset = center_y - 40 + (i * 15)
            draw.arc((center_x - 60, y_offset - 10, center_x + 60, y_offset + 10),
                    start=0, end=180, fill=(150, 150, 170), width=2)
        
        # Spine
        draw.rectangle((center_x - 3, center_y - 50, center_x + 3, center_y + 60),
                      fill=(200, 200, 220))
        
        # Heart shadow
        heart_points = [
            (center_x - 20, center_y - 10),
            (center_x - 35, center_y + 5),
            (center_x - 25, center_y + 30),
            (center_x + 10, center_y + 25),
            (center_x + 25, center_y + 10),
            (center_x + 15, center_y - 15)
        ]
        draw.polygon(heart_points, fill=(120, 120, 140))
        
        # Lung fields with infiltrates (pneumonia simulation)
        # Left lung with infiltrate
        draw.ellipse((center_x - 70, center_y - 20, center_x - 10, center_y + 40),
                    fill=(90, 90, 100), outline=(110, 110, 130))
        # Infiltrate pattern (cloudy appearance)
        for i in range(10):
            x = center_x - 60 + (i * 6)
            y = center_y + 15 + (i % 3 * 8)
            draw.ellipse((x, y, x + 8, y + 8), fill=(140, 140, 160))
        
        # Right lung
        draw.ellipse((center_x + 10, center_y - 20, center_x + 70, center_y + 40),
                    fill=(90, 90, 100), outline=(110, 110, 130))
        
        # Labels
        draw.text((10, 10), "Chest X-Ray - PA View", fill=(200, 200, 220), font=font_medium)
        draw.text((10, height - 30), "Bilateral infiltrates visible", fill=(200, 200, 220), font=font_small)
        
    elif image_type == 'brain':
        # Create brain MRI simulation
        center_x, center_y = width // 2, height // 2
        
        # Brain outline
        draw.ellipse((center_x - 70, center_y - 80, center_x + 70, center_y + 60),
                    fill=(180, 180, 200), outline=(120, 120, 140))
        
        # Cerebral hemispheres division
        draw.line((center_x, center_y - 70, center_x, center_y + 50),
                 fill=(100, 100, 120), width=2)
        
        # Ventricles
        draw.ellipse((center_x - 15, center_y - 20, center_x + 15, center_y + 10),
                    fill=(220, 220, 240), outline=(160, 160, 180))
        
        # Brain structures
        # Cerebellum
        draw.ellipse((center_x - 50, center_y + 30, center_x + 50, center_y + 55),
                    fill=(160, 160, 180), outline=(120, 120, 140))
        
        # Brainstem
        draw.rectangle((center_x - 8, center_y + 20, center_x + 8, center_y + 45),
                      fill=(140, 140, 160))
        
        # Labels
        draw.text((10, 10), "Brain MRI - Axial T1", fill=(0, 0, 0), font=font_medium)
        draw.text((center_x - 80, center_y - 40), "L", fill=(0, 0, 0), font=font_small)
        draw.text((center_x + 65, center_y - 40), "R", fill=(0, 0, 0), font=font_small)
        
    elif image_type == 'logo':
        # Create medical institution logo
        center_x, center_y = width // 2, height // 2
        
        # Medical cross
        cross_size = 20
        draw.rectangle((center_x - 3, center_y - cross_size, center_x + 3, center_y + cross_size),
                      fill=(200, 50, 50))
        draw.rectangle((center_x - cross_size, center_y - 3, center_x + cross_size, center_y + 3),
                      fill=(200, 50, 50))
        
        # Institution name
        draw.text((center_x - 40, center_y + 35), "MEDICAL CENTER", fill=(0, 0, 0), font=font_small)
        
        # Border
        draw.rectangle((5, 5, width - 5, height - 5), outline=(100, 100, 100), width=2)
        
    elif image_type == 'chart':
        # Create medical data chart
        # Chart background
        chart_x, chart_y = 50, 50
        chart_w, chart_h = width - 100, height - 100
        draw.rectangle((chart_x, chart_y, chart_x + chart_w, chart_y + chart_h),
                      fill=(250, 250, 250), outline=(0, 0, 0))
        
        # Grid lines
        for i in range(5):
            y = chart_y + (i * chart_h // 4)
            draw.line((chart_x, y, chart_x + chart_w, y), fill=(200, 200, 200))
        
        for i in range(6):
            x = chart_x + (i * chart_w // 5)
            draw.line((x, chart_y, x, chart_y + chart_h), fill=(200, 200, 200))
        
        # Data points (heart rate over time)
        data_points = [
            (chart_x + 20, chart_y + chart_h - 60),
            (chart_x + 60, chart_y + chart_h - 80),
            (chart_x + 100, chart_y + chart_h - 70),
            (chart_x + 140, chart_y + chart_h - 90),
            (chart_x + 180, chart_y + chart_h - 85),
            (chart_x + 220, chart_y + chart_h - 75)
        ]
        
        # Draw line graph
        for i in range(len(data_points) - 1):
            draw.line((data_points[i], data_points[i + 1]), fill=(200, 100, 100), width=3)
        
        # Draw points
        for point in data_points:
            draw.ellipse((point[0] - 3, point[1] - 3, point[0] + 3, point[1] + 3),
                        fill=(150, 50, 50))
        
        # Labels
        draw.text((10, 10), "Heart Rate Monitoring", fill=(0, 0, 0), font=font_medium)
        draw.text((chart_x, chart_y + chart_h + 10), "Time (hours)", fill=(0, 0, 0), font=font_small)
        draw.text((10, chart_y + chart_h // 2), "BPM", fill=(0, 0, 0), font=font_small)
        
    # Convert to BytesIO
    img_bytes = BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    
    return img_bytes


def create_medical_presentation() -> str:
    """
    Create a comprehensive medical presentation with realistic images.
    
    Returns:
        str: Path to created PPTX file
    """
    if not PPTX_AVAILABLE or not PIL_AVAILABLE:
        raise ImportError(f"Required libraries not available: {ERROR_MSG}")
    
    # Create presentation
    prs = Presentation()
    
    # Remove default slide layout and create custom
    slide_layouts = prs.slide_layouts
    
    # Slide 1: Title Slide with Logo
    logger.info("Creating slide 1: Title slide with medical logo")
    slide1 = prs.slides.add_slide(slide_layouts[0])  # Title layout
    slide1.shapes.title.text = "Cardiovascular Disease: Diagnosis and Treatment"
    slide1.shapes.placeholders[1].text = "Medical Grand Rounds\nDepartment of Cardiology\nPresented by: Dr. Smith"
    
    # Add medical logo
    logo_img = create_medical_image(120, 80, 'logo')
    logo_pic = slide1.shapes.add_picture(logo_img, Inches(8), Inches(0.5), Inches(1.2), Inches(0.8))
    # Note: No ALT text added - this will be generated by the system
    
    # Slide 2: Heart Anatomy
    logger.info("Creating slide 2: Heart anatomy diagram")
    slide2 = prs.slides.add_slide(slide_layouts[1])  # Content layout
    slide2.shapes.title.text = "Cardiac Anatomy and Physiology"
    
    content = slide2.shapes.placeholders[1]
    content.text = """Key Anatomical Structures:
• Left Atrium (LA) - Receives oxygenated blood
• Right Atrium (RA) - Receives deoxygenated blood  
• Left Ventricle (LV) - Pumps blood to systemic circulation
• Right Ventricle (RV) - Pumps blood to pulmonary circulation
• Valves regulate unidirectional blood flow"""
    
    # Add anatomical heart diagram
    heart_img = create_medical_image(400, 300, 'heart')
    heart_pic = slide2.shapes.add_picture(heart_img, Inches(0.5), Inches(2), Inches(4), Inches(3))
    
    # Add slide notes for context
    notes_slide = slide2.notes_slide
    notes_slide.notes_text_frame.text = """This anatomical diagram shows the four-chamber structure of the human heart. The diagram illustrates the left and right atria (upper chambers) and left and right ventricles (lower chambers). This is essential for understanding cardiac physiology and pathology."""
    
    # Slide 3: Chest X-Ray Analysis
    logger.info("Creating slide 3: Chest X-ray with pathology")
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = "Chest Radiography: Pneumonia Case Study"
    
    content3 = slide3.shapes.placeholders[1]
    content3.text = """Patient Presentation:
• 65-year-old male with fever and cough
• Shortness of breath for 3 days
• Physical exam: Crackles in bilateral lower lobes

Radiographic Findings:
• Bilateral lower lobe infiltrates
• No pleural effusion
• Normal cardiac silhouette"""
    
    # Add chest X-ray
    xray_img = create_medical_image(350, 400, 'xray')
    xray_pic = slide3.shapes.add_picture(xray_img, Inches(5), Inches(1.5), Inches(3.5), Inches(4))
    
    # Add clinical notes
    notes_slide3 = slide3.notes_slide
    notes_slide3.notes_text_frame.text = """This posterior-anterior (PA) chest X-ray demonstrates bilateral lower lobe infiltrates consistent with community-acquired pneumonia. The infiltrates appear as increased opacity in the lung bases. The cardiac silhouette is normal in size and contour."""
    
    # Slide 4: Brain MRI
    logger.info("Creating slide 4: Brain MRI scan")
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = "Neuroimaging: Brain MRI Interpretation"
    
    content4 = slide4.shapes.placeholders[1]
    content4.text = """MRI Technique:
• T1-weighted axial sequences
• 1.5 Tesla field strength
• Contrast: Gadolinium-enhanced

Normal Structures Visible:
• Cerebral hemispheres
• Lateral ventricles
• Brainstem and cerebellum
• Gray and white matter differentiation"""
    
    # Add brain MRI
    brain_img = create_medical_image(320, 350, 'brain')
    brain_pic = slide4.shapes.add_picture(brain_img, Inches(5.5), Inches(1.8), Inches(3.2), Inches(3.5))
    
    notes_slide4 = slide4.notes_slide
    notes_slide4.notes_text_frame.text = """Axial T1-weighted MRI image of the brain showing normal anatomical structures. The image demonstrates good gray-white matter differentiation with normal ventricular system. No acute abnormalities are identified."""
    
    # Slide 5: Clinical Data Chart
    logger.info("Creating slide 5: Patient monitoring data")
    slide5 = prs.slides.add_slide(slide_layouts[1])
    slide5.shapes.title.text = "Patient Monitoring: Vital Signs Trending"
    
    content5 = slide5.shapes.placeholders[1]
    content5.text = """24-Hour Monitoring Results:
• Continuous cardiac monitoring
• Hourly vital sign assessment
• Blood pressure: Stable 120/80
• Temperature: Resolved fever
• Oxygen saturation: 98% on room air

Clinical Interpretation:
• Improving hemodynamic status
• Response to treatment evident"""
    
    # Add monitoring chart
    chart_img = create_medical_image(400, 280, 'chart')
    chart_pic = slide5.shapes.add_picture(chart_img, Inches(0.8), Inches(2.3), Inches(4), Inches(2.8))
    
    notes_slide5 = slide5.notes_slide
    notes_slide5.notes_text_frame.text = """This graph shows heart rate monitoring over a 24-hour period. The data demonstrates gradual stabilization of heart rate following treatment initiation. The trend indicates positive response to therapeutic intervention."""
    
    # Slide 6: Treatment Summary
    logger.info("Creating slide 6: Treatment summary")
    slide6 = prs.slides.add_slide(slide_layouts[1])
    slide6.shapes.title.text = "Treatment Protocol and Outcomes"
    
    content6 = slide6.shapes.placeholders[1]
    content6.text = """Evidence-Based Treatment:
• Antibiotic therapy: Azithromycin 500mg daily
• Supportive care: Oxygen therapy PRN
• Monitoring: Serial chest X-rays
• Duration: 7-day course

Clinical Outcomes:
• Fever resolution within 48 hours
• Improved respiratory symptoms
• Radiographic improvement on day 5
• Successful outpatient management"""
    
    # Add small logo for consistency
    footer_logo = create_medical_image(80, 40, 'logo')
    footer_pic = slide6.shapes.add_picture(footer_logo, Inches(8.5), Inches(6.8), Inches(0.8), Inches(0.4))
    
    # Save presentation
    output_file = Path("medical_sample_presentation.pptx")
    prs.save(str(output_file))
    
    logger.info(f"Created medical presentation: {output_file}")
    logger.info(f"Total slides: {len(prs.slides)}")
    
    # Count images for verification
    total_images = 0
    for i, slide in enumerate(prs.slides):
        slide_images = 0
        for shape in slide.shapes:
            if hasattr(shape, 'image') and shape.image:
                slide_images += 1
                total_images += 1
        logger.info(f"  Slide {i+1}: {slide_images} images")
    
    logger.info(f"Total embedded images: {total_images}")
    
    return str(output_file)


def main():
    """Create sample medical PPTX presentation."""
    try:
        print("Medical Sample PPTX Creator")
        print("=" * 40)
        
        if not PPTX_AVAILABLE or not PIL_AVAILABLE:
            print(f"❌ Required libraries not available:")
            print(f"   {ERROR_MSG}")
            print("\nInstall requirements:")
            print("   pip install python-pptx Pillow")
            return 1
        
        # Create presentation
        output_file = create_medical_presentation()
        
        print(f"\n✅ Successfully created medical sample PPTX:")
        print(f"   File: {output_file}")
        print(f"   Size: {Path(output_file).stat().st_size / 1024:.1f} KB")
        print(f"\nThis presentation contains:")
        print(f"   📊 6 slides with medical content")
        print(f"   🖼️  6 embedded images (no ALT text)")
        print(f"   📝 Clinical context and slide notes")
        print(f"   🏥 Realistic medical scenarios")
        print(f"\nReady for end-to-end ALT text processing!")
        
        return 0
        
    except Exception as e:
        logger.error(f"Failed to create medical presentation: {e}")
        print(f"❌ Error: {e}")
        return 1


if __name__ == "__main__":
    exit(main())