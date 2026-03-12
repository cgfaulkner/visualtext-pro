"""
Offline placeholder ALT injection (no AI provider calls).

Reuses core/pptx_alt_injector shape traversal (iter_shapes_flattened) and ALT writing.
Used when --offline-mode is fill-missing or overwrite-all and provider is offline.
"""

from pathlib import Path
from typing import Any, Dict, Optional

from pptx import Presentation

from shared.config_manager import ConfigManager
from pptx_alt_injector import PPTXAltTextInjector


PLACEHOLDER_PREFIX = "[ALT pending]"
MAX_SHAPE_NAME_LEN = 60


def _is_image_shape(shape: Any) -> bool:
    """
    Treat as image if shape_type is PICTURE or shape has an .image relationship.
    Group containers are not images; their children may be.
    """
    try:
        if hasattr(shape, "image") and shape.image:
            return True
        if hasattr(shape, "shape_type"):
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
    except Exception:
        pass
    return False


def _has_meaningful_text(shape: Any) -> bool:
    """
    True if shape has a text frame with non-empty stripped text.
    All access guarded to avoid missing/empty attributes.
    """
    if not getattr(shape, "has_text_frame", False):
        return False
    tf = getattr(shape, "text_frame", None)
    if not tf:
        return False
    text = (getattr(tf, "text", None) or "").strip()
    return bool(text)


def _is_in_scope(
    shape: Any,
    scope: str,
    injector: "PPTXAltTextInjector",
) -> bool:
    """
    True if shape is in the given placeholder scope.
    - images: picture shapes only (_is_image_shape).
    - visuals: injector visual elements, excluding TEXT_BOX and shapes with meaningful text.
    """
    if scope == "images":
        return _is_image_shape(shape)
    if scope == "visuals":
        if not injector._is_visual_element_for_injection(shape):
            return False
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TEXT_BOX:
                return False
        except Exception:
            pass
        if _has_meaningful_text(shape):
            return False
        return True
    return False


def _alt_present(
    shape: Any,
    injector: PPTXAltTextInjector,
    treat_pending_as_missing: bool = False,
) -> bool:
    """
    ALT is present if either descr or title is non-empty (same fields injector writes).
    Do not treat shape name as ALT.
    If treat_pending_as_missing is True, text starting with "[ALT pending]" counts as missing.
    """
    descr, title = injector._get_existing_descr_and_title(shape)
    combined = (descr or "").strip() or (title or "").strip()
    if not combined:
        return False
    if treat_pending_as_missing and combined.startswith(PLACEHOLDER_PREFIX):
        return False
    return True


def _placeholder_text(shape: Any) -> str:
    """Build placeholder: '[ALT pending] {shape_name}' (truncate name; fallback '[ALT pending]')."""
    name = getattr(shape, "name", "") or ""
    if name:
        s = name[:MAX_SHAPE_NAME_LEN] if len(name) > MAX_SHAPE_NAME_LEN else name
        return f"{PLACEHOLDER_PREFIX} {s}"
    return PLACEHOLDER_PREFIX


def run_placeholder_injection(
    pptx_path: str,
    offline_mode: str,
    config_path: Optional[str] = None,
    debug_offline_placeholders: bool = False,
    treat_pending_as_missing: bool = False,
    placeholder_scope: str = "images",
) -> Dict[str, Any]:
    """
    Inject placeholder ALT only (no AI calls). Uses injector.iter_shapes_flattened
    so all targets (images or visuals) including inside groups are considered.

    Args:
        pptx_path: Path to PPTX file (modified in place).
        offline_mode: 'fill-missing' (only set where ALT missing/empty) or 'overwrite-all'.
        config_path: Optional config path for ConfigManager.
        debug_offline_placeholders: If True, print per-target decisions.
        treat_pending_as_missing: If True, treat "[ALT pending]..." as missing (default False).
        placeholder_scope: 'images' (pictures only) or 'visuals' (all visual elements, excl. text).

    Returns:
        Dict with: targets_found, targets_missing_alt_found (always); when scope==images
        also images_found, images_missing_alt_found; plus placeholders_applied,
        existing_alt_preserved, injection_attempted, injection_failed, success, error.
    """
    result: Dict[str, Any] = {
        "targets_found": 0,
        "targets_missing_alt_found": 0,
        "placeholders_applied": 0,
        "existing_alt_preserved": 0,
        "injection_attempted": 0,
        "injection_failed": 0,
        "success": False,
        "error": None,
    }
    path = Path(pptx_path)
    if not path.exists():
        result["error"] = f"File not found: {pptx_path}"
        return result

    try:
        config_manager = ConfigManager(config_path)
        injector = PPTXAltTextInjector(config_manager)
        presentation = Presentation(str(path))

        for slide_idx, shape, in_group in injector.iter_shapes_flattened(presentation):
            if not _is_in_scope(shape, placeholder_scope, injector):
                continue
            result["targets_found"] += 1
            shape_name = getattr(shape, "name", "") or ""
            is_picture = _is_image_shape(shape)
            kind = "picture" if is_picture else "shape"
            shape_type_str = str(getattr(shape, "shape_type", ""))
            alt_present = _alt_present(shape, injector, treat_pending_as_missing)
            if not alt_present:
                result["targets_missing_alt_found"] += 1

            if offline_mode == "fill-missing" and alt_present:
                result["existing_alt_preserved"] += 1
                if debug_offline_placeholders:
                    print(
                        f"  slide={slide_idx} shape={shape_name!r} type={shape_type_str} "
                        f"kind={kind} grouped={in_group} alt_present=True action=preserved"
                    )
                continue

            # Intend to inject: record attempt, read before, inject, read after, verify
            result["injection_attempted"] += 1
            before_descr, before_title = injector._get_existing_descr_and_title(shape)
            text = _placeholder_text(shape)
            try:
                injector._inject_alt_text_robust(shape, text)
            except Exception as exc:
                result["injection_failed"] += 1
                if debug_offline_placeholders:
                    print(
                        f"  WRITE_EXCEPTION slide={slide_idx} name={shape_name!r} "
                        f"kind={kind} in_group={in_group}: {exc!r}"
                    )
                continue
            after_descr, after_title = injector._get_existing_descr_and_title(shape)
            after_combined = (after_descr or "").strip() or (after_title or "").strip()
            if not after_combined:
                result["injection_failed"] += 1
                if debug_offline_placeholders:
                    print(
                        f"  WRITE_FAILED slide={slide_idx} name={shape_name!r} in_group={in_group} "
                        f"before=({before_descr!r}, {before_title!r}) "
                        f"after=({after_descr!r}, {after_title!r})"
                    )
            else:
                result["placeholders_applied"] += 1
                if debug_offline_placeholders:
                    print(
                        f"  slide={slide_idx} shape={shape_name!r} type={shape_type_str} "
                        f"kind={kind} grouped={in_group} alt_present={alt_present} action=applied (verified)"
                    )

        if placeholder_scope == "images":
            result["images_found"] = result["targets_found"]
            result["images_missing_alt_found"] = result["targets_missing_alt_found"]
        from shared.disclosure_slide import ensure_disclosure_slide
        ensure_disclosure_slide(presentation)
        presentation.save(str(path))
        result["success"] = True
    except Exception as e:
        result["error"] = str(e)
    return result
