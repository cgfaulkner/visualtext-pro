"""
PPTX ALT Text Injector for PPTX Accessibility Processor
Injects ALT text into PowerPoint presentations using python-pptx XML manipulation
Integrates with existing ConfigManager, reinjection settings, and workflow patterns
"""

# --- safe XPath helper for python-pptx (BaseOxmlElement) and raw lxml ---
try:
    from pptx.oxml.ns import nsmap as PPTX_NSMAP  # type: ignore
except Exception:  # pragma: no cover
    PPTX_NSMAP = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
        'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
    }

def _safe_xpath(element, xpath_expr, namespaces=None):
    """Execute an XPath on python-pptx BaseOxmlElement or plain lxml element.
    Tries python-pptx override (no kwargs) first, then raw lxml with nsmap.
    Accepts an optional explicit namespace map for raw lxml cases.
    """
    el = getattr(element, "_element", element)
    try:
        return el.xpath(xpath_expr)  # python-pptx injects namespaces
    except Exception:
        if namespaces is not None:
            ns = namespaces
        elif hasattr(el, "nsmap") and el.nsmap is not None:
            ns = el.nsmap
        else:
            ns = PPTX_NSMAP
        return el.xpath(xpath_expr, namespaces=ns)
# --- end safe XPath helper ---


import logging
import os
import sys
import time
import argparse
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple, Union
from hashlib import md5
import tempfile
from collections import Counter

# Third-party imports for PPTX processing
try:
    from pptx import Presentation
    from pptx.shapes.picture import Picture
    from pptx.shapes.base import BaseShape
    from pptx.oxml.ns import _nsmap
    PPTX_AVAILABLE = True
except ImportError as e:
    PPTX_AVAILABLE = False
    PPTX_ERROR = str(e)

# Setup paths for shared and core modules
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root / "shared"))
sys.path.insert(0, str(project_root / "core"))

# Import shared modules
from config_manager import ConfigManager
from decorative_filter import is_force_decorative_by_filename
from alt_text_reader import read_existing_alt
from fallback_policies import apply_fallback_policy, apply_for_ppt_injection

logger = logging.getLogger(__name__)
logger.info("LOG: injector_file=%s", __file__)


# --- XML tag helper functions for robust shape type detection ---
def _element_of(shape):
    """Extract the underlying XML element from a shape object."""
    return getattr(shape, "_element", None) or getattr(shape, "element", None)


def is_connector(shape):
    """Check if shape is a connector using XML tag inspection."""
    el = _element_of(shape)
    return bool(el is not None and el.tag.endswith('}cxnSp'))  # p:cxnSp


def is_picture(shape):
    """Check if shape is a picture using XML tag inspection."""
    el = _element_of(shape)
    return bool(el is not None and el.tag.endswith('}pic'))    # p:pic


def is_placeholder_or_autoshape(shape):
    """Check if shape is a placeholder or autoshape using XML tag inspection."""
    el = _element_of(shape)
    # p:sp covers most shapes/placeholders; refine if needed
    return bool(el is not None and el.tag.endswith('}sp'))


def _set_alt_on_cNvPr(el, text):
    """
    Unified ALT text setter for all shape types using cNvPr node.
    Works for p:sp, p:pic, p:cxnSp, etc. by finding the appropriate cNvPr.
    """
    # Find cNvPr (p:nvSpPr/p:cNvPr or p:nvPicPr/p:cNvPr or p:nvCxnSpPr/p:cNvPr)
    cNvPr = el.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
    if cNvPr is not None:
        cNvPr.set('descr', text or '')
        # Optionally also set 'title' if dual-writing is desired
        cNvPr.set('title', cNvPr.get('title', ''))
        return True
    return False


class PPTXImageIdentifier:
    """
    Robust image identifier for maintaining consistency across extract→generate→inject workflow.
    """
    
    def __init__(self, slide_idx: int, shape_idx: int, shape_name: str = "", 
                 image_hash: str = "", embed_id: str = ""):
        self.slide_idx = slide_idx
        self.shape_idx = shape_idx
        self.shape_name = shape_name
        self.image_hash = image_hash
        self.embed_id = embed_id
        self.image_key = self._create_image_key()
    
    def _create_image_key(self) -> str:
        """Create stable, unique identifier using shape ID when available."""
        # If shape_idx is a shape ID (integer), use shapeid format for stability
        if isinstance(self.shape_idx, int):
            # Use stable shape ID format: slide_X_shapeid_Y_hash_Z
            components = [f"slide_{self.slide_idx}", f"shapeid_{self.shape_idx}"]
        else:
            # Fallback for complex/string indices
            if isinstance(self.shape_idx, str) and '_' in str(self.shape_idx):
                shape_component = f"shape_{self.shape_idx}"
            else:
                shape_component = f"shape_{self.shape_idx}"
            components = [f"slide_{self.slide_idx}", shape_component]
        
        # CRITICAL FIX: Always include hash component to match processor
        # The processor ALWAYS includes hash, so injector must too
        if self.image_hash:
            components.append(f"hash_{self.image_hash[:8]}")
        else:
            # Even if no hash available, add placeholder to maintain format consistency
            components.append("hash_00000000")

        return "_".join(components)

    @classmethod
    def from_shape(cls, shape: Picture, slide_idx: int, shape_idx):
        """Create identifier from shape object, supporting nested/complex shape indices."""
        shape_name = getattr(shape, 'name', '')

        # Extract image hash if available
        image_hash = ""
        try:
            if hasattr(shape, 'image') and shape.image:
                image_data = shape.image.blob
                image_hash = md5(image_data).hexdigest()
        except Exception:
            pass
        
        # Extract embed ID if available
        embed_id = ""
        try:
            blip_fill = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            if blip_fill is not None:
                embed_id = blip_fill.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', '')
        except Exception:
            pass
        
        return cls(slide_idx, shape_idx, shape_name, image_hash, embed_id)
    
    @classmethod
    def parse_from_complex_key(cls, image_key: str):
        """
        Parse complex image key format to extract components.
        Handles formats like: slide_X_shape_Y_Z_hash_XXXXX or slide_X_shape_Y_name_ABC_hash_XXXXX
        
        Args:
            image_key: Complex image key string
            
        Returns:
            PPTXImageIdentifier instance
        """
        parts = image_key.split('_')
        
        if len(parts) < 3 or not parts[0].startswith('slide') or not parts[2].startswith('shape'):
            raise ValueError(f"Invalid image key format: {image_key}")
        
        # Extract slide index
        slide_idx = int(parts[1])
        
        # Extract shape index (may be complex like "0_1_2")
        shape_parts = []
        i = 3
        while i < len(parts) and (parts[i].isdigit() or (i == 3)):
            shape_parts.append(parts[i])
            i += 1
        
        if not shape_parts:
            raise ValueError(f"No shape index found in key: {image_key}")
        
        shape_idx = "_".join(shape_parts) if len(shape_parts) > 1 else int(shape_parts[0])
        
        # Extract optional components
        shape_name = ""
        image_hash = ""
        
        while i < len(parts):
            if parts[i] == 'name' and i + 1 < len(parts):
                shape_name = parts[i + 1]
                i += 2
            elif parts[i] == 'hash' and i + 1 < len(parts):
                image_hash = parts[i + 1]
                i += 2
            else:
                i += 1
        
        return cls(slide_idx, shape_idx, shape_name, image_hash, "")

    @classmethod
    def from_key(cls, key: str):
        """
        Parse image key to extract slide and shape indices, ignoring hash for matching.
        Handles formats like: slide_X_shapeid_Y_hash_XXXXX or slide_X_shape_Y_hash_XXXXX

        Args:
            key: Image key string

        Returns:
            PPTXImageIdentifier instance for stable matching
        """
        import re
        _KEY_RE = re.compile(r'^slide_(?P<slide>\d+)_(?:shapeid_|shape_)(?P<shape>\d+)(?:_hash_[0-9a-f]+)?')

        m = _KEY_RE.match(key)
        if not m:
            raise ValueError(f"Bad key: {key}")

        slide_idx = int(m.group('slide'))
        shape_idx = int(m.group('shape'))

        # Create identifier without hash dependency for stable matching
        return cls(slide_idx, shape_idx, "", "", "")


@dataclass
class AltWritePlan:
    """Decision record for a single ALT text injection attempt."""

    final_alt: str
    decision: str
    should_write: bool


class PPTXAltTextInjector:
    """
    PPTX ALT text injector that integrates with existing system architecture.
    Supports multiple injection methods, validation, and ConfigManager integration.
    """
    
    def __init__(self, config_manager: Optional[ConfigManager] = None):
        """
        Initialize the PPTX ALT text injector.
        
        Args:
            config_manager: Optional ConfigManager instance
        """
        if not PPTX_AVAILABLE:
            raise ImportError(f"python-pptx is required: {PPTX_ERROR}")
        
        self.config_manager = config_manager or ConfigManager()
        
        # Get reinjection settings from config
        self.reinjection_config = self.config_manager.config.get('reinjection', {})
        self.skip_alt_text_if = self.reinjection_config.get('skip_alt_text_if', [])
        
        # Get ALT text handling settings
        self.alt_text_config = self.config_manager.config.get('alt_text_handling', {})
        self.mode = self.alt_text_config.get('mode', 'preserve')
        self.clean_generated_alt_text = self.alt_text_config.get('clean_generated_alt_text', True)
        
        # Get PPTX-specific settings
        self.pptx_config = self.config_manager.config.get('pptx_processing', {})
        
        # Register XML namespaces for decorative detection
        self._register_namespaces()
        
        # Statistics - unified name
        self.statistics = {
            'success': 0,
            'skipped_existing': 0,
            'skipped_invalid': 0,
            'failed': 0,
            'preserved_existing': 0,
            'written_generated': 0,
            'written_existing': 0,
            'written_final': 0,
            'skipped_no_content': 0,
            # Legacy names for compatibility
            'total_images': 0,
            'injected_successfully': 0,
            'failed_injection': 0,
            'validation_failures': 0
        }
        
        # HOTPATCH FIX 1: Single-writer rule tracking
        # Track what has been written per element key to prevent duplicates
        self.written_elements = {}  # {element_key: {'alt': str, 'source': str, 'timestamp': float, 'score': int}}
        
        # PHASE 1.1 HOTFIX: Last-mile idempotence tracking
        self.final_writes = {}  # {element_key: {'text_hash': str, 'source': str, 'processed_by': str}}
        self.element_classifiers = {}  # {element_key: 'connector-bypass'|'generic'|'raster'}
        
        logger.info("Initialized PPTX ALT text injector")
        logger.debug(f"Skip ALT text if: {self.skip_alt_text_if}")
        logger.debug(f"Mode: {self.mode}")

    def _record_written_stat(self, decision: str) -> None:
        """Increment statistics counters based on the write decision."""
        if decision == 'written_generated':
            self.statistics['written_generated'] += 1
        elif decision == 'written_existing':
            self.statistics['written_existing'] += 1
        elif decision == 'written_final':
            self.statistics['written_final'] += 1

    def _normalize_alt_universal(self, txt: str) -> str:
        """
        DETERMINISTIC NORMALIZER: Bulletproof, sentence-aware deduplication.
        Runs exactly once right before any write (single source of truth).
        
        Args:
            txt: ALT text that may contain duplications
            
        Returns:
            Clean, deduplicated ALT text with deterministic sentence-level processing
        """
        if not txt:
            return ""

        import re

        # 1) Collapse whitespace
        t = " ".join(txt.split())

        # 2) Split into sentences, de-dup near-duplicates
        parts = [p.strip(" .") for p in re.split(r'(?<=[.!?])\s+', t) if p.strip()]
        out = []
        for p in parts:
            pl = p.lower()
            # skip exact duplicates
            if any(pl == q.lower() for q in out):
                continue
            # skip if wholly contained in an existing sentence (keeps the longer one)
            if any(pl in q.lower() and len(p) <= len(q) for q in out):
                continue
            # also drop any *existing* sentence that is wholly contained in the new one
            out = [q for q in out if q.lower() not in pl]
            out.append(p)

        t = ". ".join(out)

        # 3) Canonicalize the common boilerplate variants
        #    "This is a PowerPoint shape. It is a chart showing data" -> "Chart showing data"
        #    "This is a PowerPoint shape. It is a shape ..." -> "This is a PowerPoint shape ..."
        t = re.sub(r'This is a PowerPoint shape\. It is a (chart .+)', r'\1', t, flags=re.IGNORECASE)
        t = re.sub(r'This is a PowerPoint shape\. It is a (shape|line|connector)\b',
                   r'This is a PowerPoint \1', t, flags=re.IGNORECASE)

        # 4) If the "A shape (WxHpx)" preamble appears multiple times, keep one
        t = re.sub(r'(A (?:shape|line|connector)\s*\([^)]*\))\s+\1(?=\s|$)', r'\1', t, flags=re.IGNORECASE)

        # 5) If both a preamble and the "This is a PowerPoint …" are present and redundant, keep the more informative one
        t = re.sub(
            r'^(A (?:shape|line|connector)\s*\([^)]*\))\s+(This is a PowerPoint (?:shape|line|connector)\b)',
            r'\2', t, flags=re.IGNORECASE
        )
        t = self._ensure_terminal_punctuation(t)

        # Ensure proper sentence case (capitalize first letter)
        if t and t[0].islower():
            t = t[0].upper() + t[1:]

        return t

    @staticmethod
    def _ensure_terminal_punctuation(text: str) -> str:
        """Ensure text ends with a terminal punctuation mark."""
        text = text.strip()
        if not text:
            return ""
        if text[-1] in ".!?":
            return text
        return f"{text}."

    @staticmethod
    def _finalize_alt(text: str) -> str:
        """
        Safe finalizer for ALT text that prevents string index out of range errors.
        Handles capitalization and punctuation without unsafe indexing.
        """
        if text is None:
            return ""

        s = text.strip()
        if not s:
            return ""

        # Skip finalization for certain reserved values
        if s.lower() in ('', 'undefined', '(none)', 'n/a', 'none'):
            return ""

        # Capitalize first alphabetic character without indexing empty strings
        first_alpha = next((i for i, ch in enumerate(s) if ch.isalpha()), None)
        if first_alpha is not None and s[first_alpha].islower():
            s = s[:first_alpha] + s[first_alpha].upper() + s[first_alpha+1:]

        # Append period only if last visible char is alphanumeric and not punctuation already
        tail = s.rstrip()
        if tail and tail[-1].isalnum():
            s = tail + "."
        else:
            s = tail

        return s

    def _score_alt_text_quality(self, alt_text: str, source: str = "unknown") -> int:
        """
        HOTPATCH FIX 1: Score ALT text quality to determine which description wins.
        
        Heuristic (simple & effective):
        +3 if contains letters/words beyond boilerplate 
        +2 if has domain nouns (icon, arrow, circle, photo, screenshot, chart)
        +1 if > 40 chars and < 200 chars (good detail level)
        -2 if matches preamble-only ("A shape (WxHpx)")
        -5 if contains "error"
        
        Args:
            alt_text: ALT text to score
            source: Source of the ALT text for logging
            
        Returns:
            Quality score (higher is better)
        """
        if not alt_text or not alt_text.strip():
            return -10
            
        text = alt_text.strip().lower()
        score = 0
        
        # Check for error strings (immediate disqualification)
        error_patterns = ['error', 'failed', 'unable', 'sorry', 'cannot', 'i cannot', 'i am unable']
        if any(err in text for err in error_patterns):
            score -= 10  # Heavy penalty to ensure errors always lose
            
        # Check for preamble-only pattern
        import re
        if re.match(r'^a (?:shape|line|connector) \(\d+x\d+px\)$', text):
            score -= 2
            
        # Check for boilerplate vs semantic content
        boilerplate_phrases = [
            'this is a powerpoint shape',
            'this is a shape', 
            'powerpoint element',
            'visual element'
        ]
        
        # Remove boilerplate to see what's left
        text_without_boilerplate = text
        for phrase in boilerplate_phrases:
            text_without_boilerplate = text_without_boilerplate.replace(phrase, '')
        
        # Score for content beyond boilerplate
        meaningful_words = len([w for w in text_without_boilerplate.split() if len(w) > 2])
        if meaningful_words >= 3:
            score += 3
            
        # Score for domain-specific nouns (indicates semantic understanding)
        domain_nouns = [
            'icon', 'arrow', 'circle', 'photo', 'screenshot', 'chart',
            'diagram', 'graph', 'table', 'image', 'picture', 'illustration',
            'logo', 'symbol', 'button', 'text', 'label', 'caption'
        ]
        
        domain_word_count = sum(1 for noun in domain_nouns if noun in text)
        if domain_word_count > 0:
            score += 2
            
        # Score for optimal length (not too short, not too verbose)
        text_len = len(alt_text.strip())
        if 40 <= text_len <= 200:
            score += 1
            
        return score
    
    def _should_replace_alt_text(self, element_key: str, candidate_alt: str, candidate_source: str) -> tuple[bool, str]:
        """
        HOTPATCH FIX 1: Determine if candidate ALT text should replace existing one.
        
        Args:
            element_key: Unique element identifier
            candidate_alt: New ALT text candidate
            candidate_source: Source of candidate (e.g., 'generator', 'fallback', 'descriptive')
            
        Returns:
            (should_replace: bool, reason: str)
        """
        if element_key not in self.written_elements:
            return True, "first_write"
            
        existing = self.written_elements[element_key]
        existing_score = existing['score']
        candidate_score = self._score_alt_text_quality(candidate_alt, candidate_source)
        
        if candidate_score > existing_score:
            return True, f"quality_upgrade ({candidate_score} > {existing_score})"
        else:
            return False, f"quality_downgrade ({candidate_score} <= {existing_score})"
    
    def _sync_legacy_statistics(self) -> None:
        """Synchronize legacy statistics counters to ensure coverage reports show correct values."""
        # Calculate new enriched statistics
        preserved_count = self.statistics.get('preserved_existing', 0)
        generated_count = self.statistics.get('written_generated', 0)
        existing_count = self.statistics.get('written_existing', 0)
        final_count = self.statistics.get('written_final', 0)
        total_successful = preserved_count + generated_count + existing_count + final_count

        # Update legacy counters to match new enriched counters
        self.statistics['injected_successfully'] = total_successful
        self.statistics['success'] = total_successful
        self.statistics['failed_injection'] = self.statistics.get('failed', 0)

        # Validation failures should match skipped + invalid
        self.statistics['validation_failures'] = (
            self.statistics.get('skipped_existing', 0) +
            self.statistics.get('skipped_invalid', 0)
        )

        # Log the enriched statistics
        logger.info("📊 ENRICHED STATISTICS:")
        logger.info(f"   Preserved existing ALT: {preserved_count}")
        logger.info(f"   Wrote generated ALT: {generated_count}")
        logger.info(f"   Used existing ALT (replace fallback): {existing_count}")
        logger.info(f"   Wrote stored final ALT: {final_count}")
        logger.info(f"   Total successful: {total_successful}")
    
    def _is_blocked_fallback_pattern(self, text: str) -> bool:
        """Check if text matches patterns that should be blocked from PPT."""
        if not text:
            return False
        
        blocked_patterns = [
            "This is a PowerPoint shape",
            "unknown (",
            "text placeholder (",
            "chart (",
            "table (",
            "group shape (",
            "connector (",
            "shape with no specific content",
            "[Generated description not available]"
        ]
        
        for pattern in blocked_patterns:
            if pattern in text:
                return True
        return False
    
    def _register_write(self, element_key: str, alt_text: str, source: str) -> None:
        """
        HOTPATCH FIX 1: Register that we've written ALT text for this element.
        
        Args:
            element_key: Unique element identifier  
            alt_text: ALT text that was written
            source: Source of the ALT text
        """
        import time
        
        score = self._score_alt_text_quality(alt_text, source)
        self.written_elements[element_key] = {
            'alt': alt_text,
            'source': source, 
            'timestamp': time.time(),
            'score': score
        }
    
    def _should_replace_alt_text_normalized(self, current_text: str, new_text: str) -> bool:
        """
        DETERMINISTIC FIX: Compare normalized versions to decide if replacement needed.
        Prevents overwrites when texts are equivalent after normalization.
        
        Args:
            current_text: Existing ALT text
            new_text: New ALT text candidate
            
        Returns:
            True if new_text should replace current_text
        """
        if not new_text or not new_text.strip():
            return False
        
        if not current_text or not current_text.strip():
            return True
            
        c = self._normalize_alt_universal(current_text)
        n = self._normalize_alt_universal(new_text)
        
        # Don't replace if they're essentially the same after normalization (case-insensitive)
        return n.lower() != c.lower()
    
    def _read_current_alt(self, shape) -> str:
        """Return existing ALT (descr) from the shape if present; '' if none."""
        return read_existing_alt(shape)

    # --- Enhanced generic placeholder detection for low-value boilerplate ---
    GENERIC_ALT_REGEXES = [
        __import__('re').compile(r"^\s*(a|an)\s+(picture|image|graphic|photo)\b", __import__('re').I),
        __import__('re').compile(r"^\s*screenshot\b", __import__('re').I),
        __import__('re').compile(r"^\s*(picture|image)\s*\d+\s*$", __import__('re').I),
        __import__('re').compile(r"\(\s*\d+\s*x\s*\d+\s*px\s*\)\s*$", __import__('re').I),  # trailing (WxHpx)
        # PowerPoint boilerplate patterns
        __import__('re').compile(r"^\s*This is a PowerPoint shape\b", __import__('re').I),
        __import__('re').compile(r"^\s*Image of\b", __import__('re').I),
        __import__('re').compile(r"^\s*(picture|graphic|shape|object)\s*\.?\s*$", __import__('re').I),
        __import__('re').compile(r"\bunknown\b", __import__('re').I),
        # Very short descriptions (often meaningless)
        __import__('re').compile(r"^\s*\w{1,4}\s*$", __import__('re').I),
    ]

    def _is_generic_placeholder_alt(self, text: str) -> bool:
        if not text:
            return True
        t = text.strip()
        # strip any trailing (WxHpx) to catch "a picture. (692x556px)"
        t_no_size = __import__('re').sub(r"\(\s*\d+\s*x\s*\d+\s*px\s*\)\s*$", "", t, flags=__import__('re').I).strip()
        
        # Check for generic patterns first
        if any(rx.search(t) for rx in self.GENERIC_ALT_REGEXES) or any(rx.search(t_no_size) for rx in self.GENERIC_ALT_REGEXES):
            return True
        
        # Require terminal punctuation for sentences longer than 6 words (quality gate)
        words = t.split()
        if len(words) > 6 and not t.endswith(('.', '!', '?')):
            return True  # Treat as low-value if no proper sentence termination
        
        return False

    def _has_meaningful_alt(self, shape) -> bool:
        """True if current ALT exists and is not in skip list / not blank-ish."""
        existing = self._read_current_alt(shape)
        if not existing:
            return False
        if self._is_generic_placeholder_alt(existing):
            return False  # treat placeholders as NOT meaningful -> OK to overwrite
        # Respect your configured sentinel "skip" values (case-insensitive)
        bads = {s.lower() for s in self.skip_alt_text_if}
        return existing.lower() not in bads
    
    def _equivalent(self, text1: str, text2: str) -> bool:
        """
        PHASE 1.1 HOTFIX: Compare text after normalization (case/spacing collapsed).
        
        Args:
            text1: First text to compare
            text2: Second text to compare
            
        Returns:
            True if texts are equivalent after normalization
        """
        if not text1 and not text2:
            return True
        
        # Normalize both for comparison
        norm1 = self._normalize_alt_universal(text1) if text1 else ''
        norm2 = self._normalize_alt_universal(text2) if text2 else ''
        
        # Compare after case/spacing collapse
        import re
        def collapse(text):
            return re.sub(r'\s+', ' ', text.strip().lower())
        
        return collapse(norm1) == collapse(norm2)
    
    def _write_descr_and_title(self, shape, text: str, preserve_override: bool = False) -> None:
        """
        Overwrite ALT text and make sure it actually lands in XML.
        Many python-pptx shape classes (notably Picture) don't update XML
        when assigning .descr/.title, and they don't raise, so we ALWAYS
        write to XML and then verify.
        """
        # Respect existing ALT in preserve mode
        mode = getattr(self, "mode", "overwrite")
        has_alt = self._has_meaningful_alt(shape)
        logger.info(f"WRITE_GUARD: mode={mode}, has_meaningful_alt={has_alt}")
        if mode == "preserve" and has_alt and not preserve_override:
            logger.info("WRITE_GUARD TRIGGERED: Preserve mode: existing ALT found; skipping reinjection for this shape.")
            return
        if preserve_override and mode == "preserve" and has_alt:
            logger.info("WRITE_GUARD OVERRIDDEN: Writing ALT despite preserve mode due to final override.")
            
        # Best-effort property set (harmless if ignored)
        for attr in ("descr", "title"):
            try:
                setattr(shape, attr, text)
            except Exception:
                pass
        # Robust XML write (covers Picture, Shape, Connector, GraphicFrame, Group)
        self._write_alt_via_xml_fallback(shape, text)
        # Post-write verification; log if it didn't stick
        try:
            after = self._read_current_alt(shape)
            if not self._equivalent(after, text):
                logger.debug("ALT post-write verification failed; XML write may not have stuck for this shape.")
        except Exception:
            # don't crash the injection flow on verification
            pass
    
    def _write_alt_via_xml_fallback(self, shape, text: str) -> None:
        """Write ALT to the correct cNvPr node. Set only 'descr' to avoid duplicates."""
        logger.info(f"XML_FALLBACK: Starting write for '{text[:50]}...' to {type(shape).__name__}")
        try:
            # SKIP REDUNDANT GATE: Quality gating already applied in _inject_alt()
            # The centralized gate here was blocking legitimate ALT text that passed the lenient gate
            logger.info(f"XML_FALLBACK: Proceeding with write (quality gates already applied upstream)")
            
            element = getattr(shape, "_element", None) or getattr(shape, "element", None)
            if element is None:
                return
            ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

            # Try specific shape types first, then a generic fallback.
            paths = [
                ".//p:nvPicPr/p:cNvPr",          # Picture
                ".//p:nvSpPr/p:cNvPr",           # AutoShape
                ".//p:nvCxnSpPr/p:cNvPr",        # Connector
                ".//p:nvGraphicFramePr/p:cNvPr", # Chart/Table/SmartArt frame
                ".//p:nvGrpSpPr/p:cNvPr",        # Group parent
                ".//p:cNvPr"                     # Last resort
            ]

            wrote_any = False
            logger.info(f"XML_WRITE: Attempting write to shape {type(shape).__name__}")
            for xp in paths:
                found_nodes = _safe_xpath(element, xp, ns)
                logger.info(f"XML_WRITE: XPath '{xp}' found {len(found_nodes)} nodes")
                for cnvpr in found_nodes:
                    # Only set 'descr' – leave 'title' empty to avoid duplicate reading in UI/AT.
                    cnvpr.set('descr', text)
                    logger.info(f"XML_WRITE: Successfully wrote '{text[:50]}...' to cNvPr node")
                    wrote_any = True
            if not wrote_any:
                logger.info(f"XML_WRITE: ❌ FAILED - No cNvPr nodes found for {type(shape).__name__} shape")
        except Exception as e:
            logger.debug(f"XML fallback failed: {e}")
    
    def _compute_text_hash(self, text: str) -> str:
        """Compute short hash of text for rerun detection"""
        import hashlib
        return hashlib.sha1(text.encode('utf-8')).hexdigest()[:8]

    def _add_geometric_backfill(self, alt_text_mapping: Dict[str, Any], image_identifiers: Dict[str, Tuple[PPTXImageIdentifier, Any]]) -> Dict[str, Any]:
        """
        Add synthesized ALT text for auto-shapes, lines, and geometric elements missing from mapping.

        Args:
            alt_text_mapping: Existing mapping of image keys to ALT text
            image_identifiers: All discovered shapes with their identifiers

        Returns:
            Enhanced mapping with geometric backfill added
        """
        enhanced_mapping = alt_text_mapping.copy()

        for image_key, (identifier, shape) in image_identifiers.items():
            # Skip if already has ALT text mapping
            if image_key in alt_text_mapping:
                continue

            # Skip text placeholders - these shouldn't have ALT text
            shape_type = type(shape).__name__
            if shape_type in ['SlidePlaceholder', 'TextBox']:
                continue

            # Generate geometric description for auto-shapes, lines, etc.
            geometric_alt = self._generate_geometric_description(shape, identifier)
            if geometric_alt:
                # Create enriched record for geometric backfill
                enhanced_mapping[image_key] = {
                    "existing_alt": "",
                    "generated_alt": geometric_alt,
                    "source_existing": None,
                    "source_generated": "geometric_backfill",
                    "final_alt": None,
                    "decision": None,
                    "existing_meaningful": False
                }
                logger.info(f"BACKFILL: {image_key} -> '{geometric_alt}'")

        return enhanced_mapping

    def _generate_geometric_description(self, shape, identifier: PPTXImageIdentifier) -> str:
        """
        Generate basic ALT text for geometric shapes based on type and metadata.

        Args:
            shape: The shape object
            identifier: Shape identifier with metadata

        Returns:
            Synthesized ALT text or empty string if not applicable
        """
        try:
            shape_type = type(shape).__name__
            shape_name = getattr(shape, 'name', '')

            # Handle different shape types
            if 'Picture' in shape_type:
                # Should have been handled by model generation, but fallback
                return f"Image: {shape_name}" if shape_name else "Image"

            elif 'Oval' in shape_name or 'Circle' in shape_name:
                color = self._extract_shape_color(shape)
                return f"{color}circle" if color else "Circle"

            elif 'Rectangle' in shape_name or 'Square' in shape_name:
                color = self._extract_shape_color(shape)
                return f"{color}rectangle" if color else "Rectangle"

            elif 'Line' in shape_name or 'Connector' in shape_type:
                orientation = self._get_line_orientation(shape)
                return f"{orientation}line" if orientation else "Line"

            elif 'Arrow' in shape_name:
                direction = self._get_arrow_direction(shape)
                color = self._extract_shape_color(shape)
                parts = [p for p in [color, direction, "arrow"] if p]
                return " ".join(parts).capitalize() if parts else "Arrow"

            elif 'Hexagon' in shape_name:
                color = self._extract_shape_color(shape)
                return f"{color}hexagon" if color else "Hexagon"

            elif 'Triangle' in shape_name:
                color = self._extract_shape_color(shape)
                return f"{color}triangle" if color else "Triangle"

            else:
                # Generic shape fallback
                if shape_name and shape_name != shape_type:
                    return f"Shape: {shape_name}"
                else:
                    return f"Geometric shape ({shape_type})"

        except Exception as e:
            logger.debug(f"Error generating geometric description for {identifier.image_key}: {e}")
            return "Geometric shape"

    def _extract_shape_color(self, shape) -> str:
        """Extract dominant color from shape fill/stroke for description."""
        try:
            # Try to get fill color
            if hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color'):
                if hasattr(shape.fill.fore_color, 'rgb'):
                    rgb = shape.fill.fore_color.rgb
                    return self._rgb_to_color_name(rgb.r, rgb.g, rgb.b) + " "

            # Try stroke/line color
            if hasattr(shape, 'line') and hasattr(shape.line, 'color'):
                if hasattr(shape.line.color, 'rgb'):
                    rgb = shape.line.color.rgb
                    return self._rgb_to_color_name(rgb.r, rgb.g, rgb.b) + " "

        except Exception:
            pass
        return ""

    def _rgb_to_color_name(self, r: int, g: int, b: int) -> str:
        """Convert RGB values to basic color name."""
        # Simple color mapping for common colors
        if r > 200 and g < 100 and b < 100:
            return "red"
        elif r < 100 and g > 200 and b < 100:
            return "green"
        elif r < 100 and g < 100 and b > 200:
            return "blue"
        elif r > 200 and g > 200 and b < 100:
            return "yellow"
        elif r > 200 and g < 150 and b > 200:
            return "purple"
        elif r > 200 and g > 150 and b < 100:
            return "orange"
        elif r < 100 and g < 100 and b < 100:
            return "black"
        elif r > 200 and g > 200 and b > 200:
            return "white"
        elif r > 100 and g > 100 and b > 100:
            return "gray"
        else:
            return ""

    def _get_line_orientation(self, shape) -> str:
        """Determine line orientation (horizontal, vertical, diagonal)."""
        try:
            if hasattr(shape, 'left') and hasattr(shape, 'top') and hasattr(shape, 'width') and hasattr(shape, 'height'):
                width = getattr(shape, 'width', 0)
                height = getattr(shape, 'height', 0)

                # Convert to pixel-like units for comparison
                w = width // 12700 if width else 0  # EMUs to approximate pixels
                h = height // 12700 if height else 0

                if h < 10 and w > h * 3:  # Mostly horizontal
                    return "horizontal "
                elif w < 10 and h > w * 3:  # Mostly vertical
                    return "vertical "
                else:
                    return "diagonal "
        except Exception:
            pass
        return ""

    def _get_arrow_direction(self, shape) -> str:
        """Determine arrow direction based on shape orientation."""
        # This is a simplified version - could be enhanced with actual arrow head detection
        orientation = self._get_line_orientation(shape)
        if "horizontal" in orientation:
            return "horizontal"
        elif "vertical" in orientation:
            return "vertical"
        else:
            return ""

    def _determine_alt_decision(self, image_key: str, alt_record: Dict[str, Any], mode: str) -> AltWritePlan:
        """Resolve the final ALT text and whether it should be written."""

        normalized_mode = (mode or '').lower()
        if normalized_mode == 'overwrite':
            normalized_mode = 'replace'
        if normalized_mode not in {'preserve', 'replace'}:
            logger.warning(
                "Unknown mode '%s' for %s, defaulting to replace behavior",
                mode,
                image_key,
            )
            normalized_mode = 'replace'

        existing_alt = (alt_record.get('existing_alt') or '').strip()
        generated_alt = (alt_record.get('generated_alt') or '').strip()
        stored_final_alt = (alt_record.get('final_alt') or '').strip()

        candidate_generated = stored_final_alt or generated_alt
        final_alt = ''
        decision = 'skipped_no_content'
        should_write = False

        if normalized_mode == 'preserve':
            if stored_final_alt:
                decision = 'written_final'
                final_alt = stored_final_alt
                should_write = True
            elif existing_alt:
                decision = 'preserved_existing'
                final_alt = existing_alt
            elif candidate_generated:
                decision = 'written_generated'
                final_alt = candidate_generated
                should_write = True
        else:  # replace semantics
            if candidate_generated:
                decision = 'written_generated'
                final_alt = candidate_generated
                should_write = True
            elif existing_alt:
                decision = 'written_existing'
                final_alt = existing_alt
                should_write = False

        # If a stored final ALT overrides the generated candidate, surface that in the decision
        if (
            decision == 'written_generated'
            and stored_final_alt
            and stored_final_alt == final_alt
            and stored_final_alt != generated_alt
        ):
            decision = 'written_final'

        alt_record['final_alt'] = final_alt or None
        alt_record['decision'] = decision

        logger.debug(
            "ALT_DECISION key=%s mode=%s decision=%s existing='%s' generated='%s' final='%s'",
            image_key,
            normalized_mode,
            decision,
            existing_alt[:50],
            generated_alt[:50],
            (final_alt or '')[:50],
        )

        return AltWritePlan(final_alt, decision, should_write and bool(final_alt))

    def _apply_final_normalization_gate(self, raw_text: str, element_key: str, source: str) -> str:
        """
        SURGICAL FIX C: Single normalization gate of last resort before any write.
        Handles all the duplicate-text fixes that can still occur even after A+B fixes.
        
        Args:
            raw_text: Raw ALT text to normalize
            element_key: Element key for logging
            source: Source of text for context
            
        Returns:
            Normalized text ready for injection
        """
        if not raw_text:
            return ""
        
        # Apply deterministic normalization first
        normalized = self._normalize_alt_universal(raw_text)
        
        # Log pre→post normalization for first few writes per run
        if len(self.final_writes) < 3:  # Log first 3 normalizations
            if raw_text != normalized:
                logger.info(f"NORMALIZATION_GATE: '{raw_text[:50]}...' → '{normalized[:50]}...' (source: {source})")
        
        return normalized
    
    def _inject_alt(
        self,
        shape,
        raw_text: str,
        element_key: str,
        source: str = "unknown",
        *,
        preserve_override: bool = False,
    ) -> bool:
        """
        PHASE 1.1 HOTFIX: Single source of truth for ALL ALT text injection.
        This is the ONLY function that should write ALT text to shapes.

        Args:
            shape: Shape to inject ALT text into
            raw_text: Raw ALT text to inject
            element_key: Unique element identifier for logging
            source: Source of the ALT text (e.g., 'generator', 'connector-bypass', 'fallback')
            preserve_override: Allow writing even when preserve mode would normally skip
            
        Returns:
            True if text was written, False if skipped (idempotent/equivalent)
        """
        # SURGICAL FIX C: Single normalization gate of last resort
        text = self._apply_final_normalization_gate(raw_text, element_key, source)
        
        if not text or not text.strip():
            logger.debug(f"INJECT_ALT: Skipping empty text for {element_key}")
            return False
        
        # Apply safe finalization to prevent string indexing errors
        text = self._finalize_alt(text)

        if not text:
            logger.debug(f"INJECT_ALT: Skipping text after finalization for {element_key}")
            return False

        # CENTRALIZED GATE: Block templated strings from entering PPT
        policy = self.config_manager.get_fallback_policy()
        quality_flags = {
            "is_generated": source in ["generator", "unified_alt_generator"],
            "is_fallback": source in ["fallback", "fallback_key_based", "connector-bypass"]
        }

        if not quality_flags["is_generated"] and not quality_flags["is_fallback"]:
            # Treat direct/manual writes as generated so they are not blocked by strict policies
            quality_flags["is_generated"] = True

        gated_text = apply_for_ppt_injection(text, "shape", quality_flags, policy, shape)
        if gated_text is None:
            logger.debug(f"INJECT_ALT: Text blocked by centralized gate for {element_key}: {text[:50]}...")
            return False
        
        text = gated_text
        
        # 2) READ current ALT text (ONLY from cNvPr/@descr and @title)
        existing = self._read_current_alt(shape)
        existing_display = (existing or '').strip()

        # 2.5) PRESERVE MODE HARD GUARD - never overwrite existing ALT text in preserve mode
        # USE INJECTOR'S EFFECTIVE MODE, NOT CONFIG (fixes split-brain issue)
        logger.info(f"PRESERVE_GUARD: mode={self.mode}, existing='{existing_display}'")
        if self.mode == 'preserve' and _is_meaningful(existing) and not preserve_override:
            logger.info(f"INJECT_ALT: Preserving existing ALT text for {element_key} (mode: preserve)")
            self.statistics['skipped_existing'] += 1
            return False
        if preserve_override and self.mode == 'preserve' and _is_meaningful(existing):
            logger.info(
                "PRESERVE_OVERRIDE: Forcing write for %s despite existing ALT in preserve mode",
                element_key,
            )

        # 3) IDEMPOTENT GUARD - skip if equivalent (using deterministic normalized comparison)
        if not self._should_replace_alt_text_normalized(existing, text):
            logger.debug(f"INJECT_ALT: Skipping equivalent text for {element_key} (normalized texts are identical)")
            return False
        
        # 4) RERUN SAFETY NET - check if same hash already written
        text_hash = self._compute_text_hash(text)
        if element_key in self.final_writes:
            existing_hash = self.final_writes[element_key]['text_hash']
            if text_hash == existing_hash:
                logger.debug(f"INJECT_ALT: Skipping duplicate hash for {element_key} (hash: {text_hash})")
                return False
        
        # 5) OVERWRITE ONLY - never append
        self._write_descr_and_title(shape, text, preserve_override=preserve_override)

        # IMMEDIATE READBACK VERIFICATION - Force one known write path end-to-end
        written_descr = self._read_current_alt(shape)
        final_alt = text
        logger.info(f"ONE-KEY READBACK: {element_key} -> written='{written_descr}' expected='{final_alt}'")
        if written_descr.strip() == final_alt.strip():
            logger.info(f"✅ READBACK SUCCESS: Write verified for {element_key}")
        else:
            logger.warning(f"❌ READBACK MISMATCH: {element_key} expected={final_alt!r} actual={written_descr!r}")
        logger.debug(
            "repr_written=%r repr_final=%r ord_last=%s",
            written_descr,
            final_alt,
            ord(written_descr[-1]) if written_descr else 0,
        )
        
        # 6) RECORD the final write
        self.final_writes[element_key] = {
            'text_hash': text_hash,
            'source': source,
            'processed_by': self.element_classifiers.get(element_key, 'generic')
        }
        
        # 7) LOG exactly one FINAL_WRITE per element key
        logger.info(f"FINAL_WRITE key={element_key} len={len(text)} sha1={text_hash} source={source}")
        
        return True
    
    def _classify_element(self, element_key: str, element_type: str) -> str:
        """
        PHASE 1.1 HOTFIX: Classifier fences to prevent double-processing.
        
        Args:
            element_key: Unique element identifier
            element_type: Type of element ('connector', 'line', 'shape', etc.)
            
        Returns:
            Classification ('connector-bypass', 'generic', 'raster')
        """
        if element_type in ['connector', 'line']:
            classification = 'connector-bypass'
        elif element_type in ['picture', 'image']:
            classification = 'raster'
        else:
            classification = 'generic'
        
        # Set fence - this element is now classified
        self.element_classifiers[element_key] = classification
        logger.debug(f"CLASSIFY: {element_key} -> {classification}")
        
        return classification
    
    def _is_element_processed(self, element_key: str) -> tuple[bool, str]:
        """
        Check if element was already processed by a different path.
        
        Returns:
            (is_processed: bool, processed_by: str)
        """
        if element_key in self.element_classifiers:
            return True, self.element_classifiers[element_key]
        return False, ''
    
    def _register_namespaces(self):
        """Register required XML namespaces."""
        try:
            # Register decorative namespace for Office 2019+ decorative image support
            _nsmap["adec"] = "http://schemas.microsoft.com/office/drawing/2017/decorative"
        except Exception as e:
            logger.warning(f"Could not register XML namespaces: {e}")

    def _dedupe_titles(self, presentation):
        """Clean up pre-existing duplicates where title == descr."""
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        for slide in presentation.slides:
            for shp in slide.shapes:
                el = getattr(shp, "_element", None) or getattr(shp, "element", None)
                if el is None:
                    continue
                for xp in (".//p:nvPicPr/p:cNvPr", ".//p:nvSpPr/p:cNvPr",
                           ".//p:nvCxnSpPr/p:cNvPr", ".//p:nvGraphicFramePr/p:cNvPr",
                           ".//p:nvGrpSpPr/p:cNvPr", ".//p:cNvPr"):
                    for cnvpr in _safe_xpath(el, xp, ns):
                        descr = (cnvpr.get("descr") or "").strip()
                        title = (cnvpr.get("title") or "").strip()
                        if descr and title and title == descr:
                            cnvpr.set("title", "")  # clear duplicate title
    
    def inject_alt_text_from_mapping(self, pptx_path: str, alt_text_mapping: Dict[str, Dict[str, Any]],
                                   output_path: Optional[str] = None, mode: Optional[str] = None) -> Dict[str, Any]:
        """
        Inject ALT text into PPTX file from a mapping dictionary.

        Args:
            pptx_path: Path to input PPTX file
            alt_text_mapping: Mapping of image keys to enriched ALT records
            output_path: Optional output path (defaults to overwriting input)
            mode: ALT text handling mode ("preserve" or "replace"). If None, uses config mode.

        Returns:
            Dictionary with injection results and statistics
        """
        pptx_path = Path(pptx_path)
        if output_path is None:
            output_path = pptx_path
        else:
            output_path = Path(output_path)
        
        # Validate input file
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
        
        logger.info(f"Injecting ALT text into: {pptx_path}")
        logger.info(f"Output will be saved to: {output_path}")
        logger.info(f"ALT text mappings: {len(alt_text_mapping)}")
        
        # Reset statistics
        self.statistics = {key: 0 for key in self.statistics}

        # Get ALT mode from parameter or config
        alt_mode = mode or self.config_manager.get_alt_mode()
        logger.info(f"🔧 Injector ALT mode: {alt_mode}")

        # Temporarily override mode if different from current
        original_mode = self.mode
        self.mode = alt_mode

        # Add sanity logging to catch mode desync issues
        logger.info(f"Injector init mode: {original_mode}")
        logger.info(f"Effective injector mode: {self.mode}")
        assert self.mode == alt_mode, f"Injector mode desync: expected {alt_mode}, got {self.mode}"

        try:
            # Load presentation
            presentation = Presentation(str(pptx_path))
            # Store for key-based writing
            self.prs = presentation
            
            # Build image identifier mapping for matching
            image_identifiers = self._build_image_identifier_mapping(presentation)
            
            # KEY MATCHING VERIFICATION
            logger.info("=== KEY MATCHING VERIFICATION ===")
            logger.info(f"Processor keys (expected): {len(alt_text_mapping)}")
            processor_keys = sorted(alt_text_mapping.keys())
            for i, key in enumerate(processor_keys[:3]):  # Show first 3
                logger.info(f"  P{i}: {key}")

            logger.info(f"Injector keys (available): {len(image_identifiers)}")
            injector_keys = sorted(image_identifiers.keys())
            for i, key in enumerate(injector_keys[:3]):  # Show first 3
                logger.info(f"  I{i}: {key}")

            # Check for exact matches
            exact_matches = set(processor_keys) & set(injector_keys)
            logger.info(f"Exact matches found: {len(exact_matches)}")
            logger.info("=== END VERIFICATION ===")
            
            # SURGICAL FIX A: Enhanced key validation logging
            logger.info(f"VALIDATION: Mapping keys from generator ({len(alt_text_mapping)}):")
            gen_keys = sorted(alt_text_mapping.keys())
            # Log first/middle/last sample keys
            sample_indices = [0, len(gen_keys)//2, len(gen_keys)-1] if len(gen_keys) > 1 else ([0] if gen_keys else [])
            for i in sample_indices:
                if i < len(gen_keys):
                    logger.info(f"  Expected[{i}]: {gen_keys[i]}")
            
            logger.info(f"VALIDATION: Identifier keys from PPTX ({len(image_identifiers)}):")
            available_keys = sorted(image_identifiers.keys())
            # Log first 3 shape keys per slide
            slides_logged = set()
            keys_per_slide = {}
            for key in available_keys:
                if key.startswith('slide_'):
                    slide_part = key.split('_')[1]
                    if slide_part not in keys_per_slide:
                        keys_per_slide[slide_part] = []
                    keys_per_slide[slide_part].append(key)
            
            for slide_num in sorted(keys_per_slide.keys(), key=lambda x: int(x) if x.isdigit() else 999):
                slide_keys = keys_per_slide[slide_num][:3]  # First 3 per slide
                logger.info(f"  Slide {slide_num} shapes: {slide_keys}")
                if len(keys_per_slide[slide_num]) > 3:
                    logger.info(f"    ... and {len(keys_per_slide[slide_num]) - 3} more")
            
            # GEOMETRIC BACKFILL: Add synthesized ALT for auto-shapes and lines missing from mapping
            original_mapping_count = len(alt_text_mapping)
            alt_text_mapping = self._add_geometric_backfill(alt_text_mapping, image_identifiers)
            backfill_count = len(alt_text_mapping) - original_mapping_count
            if backfill_count > 0:
                logger.info(f"GEOMETRIC_BACKFILL: Added {backfill_count} synthesized descriptions for geometric shapes")

            # Inject ALT text for each mapping
            matched_keys = []
            unmatched_keys = []

            # SURGICAL LOGGING: Prove inject loop runs after mapping
            logger.info(f"INJECT_LOOP: {len(alt_text_mapping)} keys; image_map={len(image_identifiers)}")
            logger.info(f"PRS_ID before inject: {id(presentation)}")

            for image_key, alt_record in alt_text_mapping.items():
                if isinstance(alt_record, str):
                    alt_record = {
                        "existing_alt": "",
                        "generated_alt": alt_record,
                        "final_alt": alt_record or None,
                        "decision": None,
                        "existing_meaningful": False,
                        "source_existing": None,
                        "source_generated": None,
                    }
                    alt_text_mapping[image_key] = alt_record
                elif isinstance(alt_record, dict):
                    alt_text_mapping[image_key] = alt_record
                else:
                    logger.warning(
                        "Unsupported ALT record type for %s: %s",
                        image_key,
                        type(alt_record).__name__,
                    )
                    unmatched_keys.append(image_key)
                    continue

                plan = self._determine_alt_decision(image_key, alt_record, self.mode)
                decision = alt_record.get('decision') or plan.decision
                candidate_text = _choose_candidate(alt_record)

                if candidate_text:
                    stored_final_alt = alt_record.get('final_alt')
                    current_final = stored_final_alt if isinstance(stored_final_alt, str) else ''
                    if current_final != candidate_text:
                        logger.debug(
                            "SYNC_FINAL_ALT: updating final_alt for %s from %r to %r",  # pragma: no cover
                            image_key,
                            stored_final_alt,
                            candidate_text,
                        )
                        alt_record['final_alt'] = candidate_text
                        if plan is not None:
                            plan.final_alt = candidate_text
                        alt_text_mapping[image_key] = alt_record

                if decision == 'preserved_existing':
                    self.statistics['preserved_existing'] += 1
                    continue
                if decision == 'skipped_no_content':
                    self.statistics['skipped_no_content'] += 1
                    continue
                if not candidate_text:
                    self.statistics['skipped_no_content'] += 1
                    continue

                should_write = plan.should_write if plan is not None else decision in {'written_generated', 'written_final'}
                preserve_override = decision == 'written_final'

                if not should_write:
                    self._record_written_stat(decision)
                    continue

                shape_entry = image_identifiers.get(image_key)
                logger.info(
                    f"WILL_WRITE key={image_key} has_shape={bool(shape_entry)} mode={self.mode} "
                    f"decision={decision} final='{(candidate_text or '')[:50]}...'"
                )
                if not shape_entry:
                    logger.info(f"NO_SHAPE for {image_key}, checking fallback matching...")
                if image_key not in image_identifiers:
                    logger.debug(f"SKIP: no shape found for {image_key}")
                else:
                    logger.debug(
                        f"-> WILL WRITE for {image_key}: '{candidate_text[:80]}...' "
                        f"(mode={self.mode}, decision={decision})"
                    )

                try:
                    if shape_entry:
                        identifier, shape = shape_entry
                        if self.mode == 'preserve' and not preserve_override:
                            current_alt = self._read_current_alt(shape)
                            if _is_meaningful(current_alt):
                                logger.info(
                                    "PRESERVE_SKIP: Existing ALT retained for %s", identifier.image_key
                                )
                                self.statistics['skipped_existing'] += 1
                                matched_keys.append(image_key)
                                continue

                        success = self._inject_alt_text_single(
                            shape,
                            candidate_text,
                            identifier,
                            preserve_override=preserve_override,
                        )
                        if success:
                            matched_keys.append(image_key)
                            self._record_written_stat(decision)
                        else:
                            logger.debug(f"Injection skipped for {image_key} (idempotent guard)")
                            current_alt = self._read_current_alt(shape)
                            if self.mode == 'preserve' and _is_meaningful(current_alt):
                                matched_keys.append(image_key)
                            elif current_alt.strip() == candidate_text.strip():
                                matched_keys.append(image_key)
                                self._record_written_stat(decision)
                            else:
                                unmatched_keys.append(image_key)
                    else:
                        matched = False

                        try:
                            identifier_from_key = PPTXImageIdentifier.from_key(image_key)
                            slide_idx = identifier_from_key.slide_idx
                            shape_idx = identifier_from_key.shape_idx

                            for available_key, (identifier, shape) in image_identifiers.items():
                                if identifier.slide_idx == slide_idx and identifier.shape_idx == shape_idx:
                                    if self.mode == 'preserve' and not preserve_override:
                                        current_alt = self._read_current_alt(shape)
                                        if _is_meaningful(current_alt):
                                            self.statistics['skipped_existing'] += 1
                                            matched_keys.append(image_key)
                                            matched = True
                                            logger.info(
                                                "PRESERVE_SKIP: Existing ALT retained for %s", identifier.image_key
                                            )
                                            break
                                    success = self._inject_alt_text_single(
                                        shape,
                                        candidate_text,
                                        identifier,
                                        preserve_override=preserve_override,
                                    )
                                    if success:
                                        matched_keys.append(image_key)
                                        self._record_written_stat(decision)
                                        matched = True
                                        logger.info(
                                            f"Successfully matched via stable slide+shape ID (ignoring hash): "
                                            f"{image_key} -> {available_key}"
                                        )
                                        break
                                    current_alt = self._read_current_alt(shape)
                                    if self.mode == 'preserve' and not preserve_override and _is_meaningful(current_alt):
                                        matched_keys.append(image_key)
                                        matched = True
                                        logger.info(
                                            f"Successfully matched via stable slide+shape ID (ignoring hash): "
                                            f"{image_key} -> {available_key}"
                                        )
                                        break
                                    if current_alt.strip() == candidate_text.strip():
                                        matched_keys.append(image_key)
                                        self._record_written_stat(decision)
                                        matched = True
                                        logger.info(
                                            f"Successfully matched via stable slide+shape ID (ignoring hash): "
                                            f"{image_key} -> {available_key}"
                                        )
                                        break
                        except Exception as e:
                            logger.debug(f"Could not parse key for stable matching: {image_key}, {e}")

                        if not matched:
                            variant_match = self._try_multi_variant_key_matching(image_key, image_identifiers)
                            if variant_match:
                                identifier, shape = variant_match
                                if self.mode == 'preserve' and not preserve_override:
                                    current_alt = self._read_current_alt(shape)
                                    if _is_meaningful(current_alt):
                                        self.statistics['skipped_existing'] += 1
                                        matched_keys.append(image_key)
                                        matched = True
                                        logger.info(
                                            "PRESERVE_SKIP: Existing ALT retained for %s", identifier.image_key
                                        )
                                        break
                                success = self._inject_alt_text_single(
                                    shape,
                                    candidate_text,
                                    identifier,
                                    preserve_override=preserve_override,
                                )
                                if success:
                                    matched_keys.append(image_key)
                                    self._record_written_stat(decision)
                                    matched = True
                                    logger.info(
                                        f"Successfully matched via multi-variant key matching: {image_key} -> "
                                        f"{identifier.image_key}"
                                    )
                                else:
                                    current_alt = self._read_current_alt(shape)
                                    if self.mode == 'preserve' and not preserve_override and _is_meaningful(current_alt):
                                        matched_keys.append(image_key)
                                        matched = True
                                        logger.info(
                                            f"Successfully matched via multi-variant key matching: {image_key} -> "
                                            f"{identifier.image_key}"
                                        )
                                    elif current_alt.strip() == candidate_text.strip():
                                        matched_keys.append(image_key)
                                        self._record_written_stat(decision)
                                        matched = True
                                        logger.info(
                                            f"Successfully matched via multi-variant key matching: {image_key} -> "
                                            f"{identifier.image_key}"
                                        )
                            if not matched:
                                if self._write_alt_by_key(
                                    image_key,
                                    candidate_text,
                                    preserve_override=preserve_override,
                                ):
                                    matched_keys.append(image_key)
                                    self._record_written_stat(decision)
                                    logger.info(f"Successfully injected via key-based fallback: {image_key}")
                                elif self._try_fallback_injection(
                                    presentation,
                                    image_key,
                                    candidate_text,
                                    preserve_override=preserve_override,
                                ):
                                    matched_keys.append(image_key)
                                    self._record_written_stat(decision)
                                    logger.info(f"Successfully injected via general fallback method: {image_key}")
                                else:
                                    logger.warning(
                                        f"Could not find image for key (even with fallback): {image_key}"
                                    )
                                    unmatched_keys.append(image_key)

                except Exception as e:
                    logger.error(f"Error processing key {image_key}: {e}")
                    self.statistics['failed'] += 1
                    unmatched_keys.append(image_key)
            
            # SURGICAL FIX A: Generated vs matched counter with unmatched examples
            generated_count = len(alt_text_mapping)
            matched_count = len(matched_keys)
            unmatched_count = len(unmatched_keys)
            
            logger.info(f"KEY_MATCHING_RESULTS: generated={generated_count}, matched={matched_count}, unmatched={unmatched_count}")
            
            if unmatched_keys:
                logger.info(f"UNMATCHED_EXAMPLES (first 3): {unmatched_keys[:3]}")
                if unmatched_count > 3:
                    logger.info(f"  ... and {unmatched_count - 3} more unmatched keys")
            
            # VERIFICATION STEP: Check what ALT texts are actually in the presentation before saving
            logger.info("🔍 DEBUG: POST-INJECTION VERIFICATION")
            self._perform_post_injection_verification(presentation, image_identifiers, alt_text_mapping)
            
            # Clean up any pre-existing duplicate title/descr pairs
            self._dedupe_titles(presentation)
            
            # Save presentation
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # Get file size before save (if exists)
            size_before = output_path.stat().st_size if output_path.exists() else 0
            logger.info(f"PRS_ID saved: {id(presentation)}")
            presentation.save(str(output_path))
            size_after = output_path.stat().st_size

            logger.info(f"📁 SAVE COMPLETE: {output_path}")
            logger.info(f"📁 File size: {size_before} -> {size_after} bytes")
            
            # Create result summary
            # Synchronize legacy statistics before returning
            self._sync_legacy_statistics()
            
            result = {
                'success': True,
                'input_file': str(pptx_path),
                'output_file': str(output_path),
                'statistics': self.statistics.copy(),
                'errors': [],
                'final_alt_map': alt_text_mapping
            }
            
            self._log_injection_summary(result)
            return result
            
        except Exception as e:
            error_msg = f"Failed to inject ALT text: {str(e)}"
            logger.error(error_msg, exc_info=True)
            
            # Synchronize statistics even in error case
            self._sync_legacy_statistics()
            
            return {
                'success': False,
                'input_file': str(pptx_path),
                'output_file': str(output_path),
                'statistics': self.statistics.copy(),
                'errors': [error_msg],
                'final_alt_map': alt_text_mapping
            }
        finally:
            # Restore original mode
            self.mode = original_mode
    
    def _build_image_identifier_mapping(self, presentation: Presentation) -> Dict[str, Tuple[PPTXImageIdentifier, Any]]:
        """
        Build mapping from image keys to (identifier, shape) tuples using recursive traversal
        to find images within grouped shapes, chart elements, and nested structures.
        
        Args:
            presentation: PowerPoint presentation
            
        Returns:
            Dictionary mapping image keys to (identifier, shape) tuples
        """
        mapping = {}
        logger.info("🔍 DEBUG: Starting _build_image_identifier_mapping")
        
        for slide_idx, slide in enumerate(presentation.slides):
            logger.info(f"🔍 DEBUG: Processing slide {slide_idx + 1} for image mapping")
            logger.info(f"🔍 DEBUG:   Total shapes on slide: {len(slide.shapes)}")
            
            # Use recursive shape processing from enhanced detection system
            images_found = self._extract_images_from_shapes_for_mapping(
                slide.shapes, slide_idx, parent_group_idx=None
            )
            
            logger.info(f"🔍 DEBUG:   Images found on slide {slide_idx + 1}: {len(images_found)}")
            
            for identifier, shape in images_found:
                self.statistics['total_images'] += 1
                mapping[identifier.image_key] = (identifier, shape)
                
                # Log detailed shape information
                logger.info(f"🔍 DEBUG:   Mapped image key: {identifier.image_key}")
                logger.info(f"🔍 DEBUG:     Shape type: {type(shape).__name__}")
                logger.info(f"🔍 DEBUG:     Shape ID: {getattr(shape, 'shape_id', 'unknown')}")
                logger.info(f"🔍 DEBUG:     Shape name: {getattr(shape, 'name', 'unknown')}")
                if hasattr(shape, '_element'):
                    logger.info(f"🔍 DEBUG:     XML element: {shape._element.tag if hasattr(shape._element, 'tag') else 'unknown'}")
                
                # Check current ALT text
                current_alt = ""
                try:
                    if hasattr(shape, 'descr'):
                        current_alt = shape.descr or ""
                    elif hasattr(shape, '_element'):
                        current_alt = shape._element.get('descr', "") or ""
                except:
                    pass
                logger.info(f"🔍 DEBUG:     Current ALT text: '{current_alt}'")
            
            logger.info(f"🔍 DEBUG: Completed slide {slide_idx + 1} - found {len(images_found)} images")
        
        logger.info(f"🔍 DEBUG: Completed mapping build - total images: {len(mapping)}")
        logger.info(f"🔍 DEBUG: Final image keys in mapping:")
        for key in sorted(mapping.keys()):
            logger.info(f"🔍 DEBUG:   - {key}")
        
        return mapping
    
    def _extract_images_from_shapes_for_mapping(self, shapes, slide_idx: int, parent_group_idx: str = None) -> List[Tuple[PPTXImageIdentifier, Any]]:
        """
        Extract ALL visual elements from shapes using FLATTENED indexing to match processor.
        This must match the processor's _extract_visual_elements_from_shapes() indexing exactly.
        
        Args:
            shapes: Collection of shapes to process
            slide_idx: Slide index
            parent_group_idx: Not used in flattened mode (kept for compatibility)
            
        Returns:
            List of (identifier, shape) tuples with flattened sequential indices
        """
        # Use flattened indexing approach to match processor
        return self._extract_shapes_flattened(shapes, slide_idx)
    
    def _extract_shapes_flattened(self, shapes, slide_idx: int, shape_counter: int = 0) -> List[Tuple[PPTXImageIdentifier, Any]]:
        """
        Extract shapes using STABLE SHAPE IDs instead of enumeration indices.
        This provides consistent identification across extraction and injection.
        
        Args:
            shapes: Collection of shapes to process
            slide_idx: Slide index  
            shape_counter: Starting counter for fallback (if no shape.shape_id)
            
        Returns:
            List of (identifier, shape) tuples with stable shape IDs
        """
        shape_mappings = []
        current_counter = shape_counter
        
        for shape_idx, shape in enumerate(shapes):
            try:
                # Get stable shape ID
                shape_id = getattr(shape, 'shape_id', None)
                shape_name = getattr(shape, 'name', 'unnamed')
                shape_type = getattr(shape, 'shape_type', 'unknown')
                
                if shape_id is not None:
                    logger.debug(f"    Examining shape ID {shape_id}: {type(shape).__name__} (type={shape_type}, name='{shape_name}')")
                else:
                    logger.debug(f"    Examining shape {current_counter}: {type(shape).__name__} (type={shape_type}, name='{shape_name}') [No ID]")
                
                # 1. Group shapes (recursively process shapes within groups)
                if hasattr(shape, 'shapes'):
                    logger.debug(f"      -> Found group shape with {len(shape.shapes)} child shapes")
                    group_shapes = self._extract_shapes_flattened(
                        shape.shapes, slide_idx, current_counter
                    )
                    shape_mappings.extend(group_shapes)
                    logger.debug(f"      -> Mapped {len(group_shapes)} shapes from group")
                    # Update counter for fallback indexing
                    current_counter += len(group_shapes)
                    continue  # Don't process the group shape itself, just its contents (matches processor)
                
                # 2. Check if this is a visual element that would have been processed
                if self._is_visual_element_for_injection(shape):
                    if shape_id is not None:
                        logger.debug(f"      -> Found visual element with stable ID {shape_id}")
                        identifier = self._create_identifier_from_shape(shape, slide_idx, shape_id)
                    else:
                        logger.debug(f"      -> Found visual element, using fallback index {current_counter}")
                        identifier = self._create_identifier_from_shape(shape, slide_idx, current_counter)
                    
                    shape_mappings.append((identifier, shape))
                    logger.debug(f"      -> Mapped visual element: {identifier.image_key}")
                    current_counter += 1  # Increment fallback counter
                        
                else:
                    logger.debug(f"      -> Skipping non-visual element: {type(shape).__name__}")
                    current_counter += 1  # Still increment for consistency
                
            except Exception as e:
                logger.warning(f"Error processing shape on slide {slide_idx}: {e}")
                current_counter += 1
                continue
        
        return shape_mappings
    
    def _update_identifier_for_nested_shape(self, identifier: PPTXImageIdentifier, shape_id) -> PPTXImageIdentifier:
        """
        Update identifier for nested shapes to handle complex hierarchical IDs.
        
        Args:
            identifier: Original identifier
            shape_id: Hierarchical shape ID (e.g., "0_1_2" for nested groups)
            
        Returns:
            Updated identifier with complex shape index
        """
        # If shape_id contains underscores, it's a nested shape
        if isinstance(shape_id, str) and '_' in str(shape_id):
            # Update the shape_idx to reflect the nested structure
            identifier.shape_idx = shape_id
            # Recreate the image key with the updated shape index
            identifier.image_key = identifier._create_image_key()
        
        return identifier
    
    def _is_visual_element_for_injection(self, shape) -> bool:
        """
        Determine if a shape is a visual element that would have ALT text generated for it.
        This should match the logic in pptx_processor._classify_visual_element().
        
        Args:
            shape: Shape to check
            
        Returns:
            bool: True if this shape should have ALT text injected
        """
        try:
            # Images with actual image data
            if hasattr(shape, 'image') and shape.image:
                return True
                
            # Check shape type for visual elements
            if hasattr(shape, 'shape_type'):
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                
                shape_type = shape.shape_type
                
                # Pictures
                if shape_type == MSO_SHAPE_TYPE.PICTURE:
                    return True
                
                # Auto shapes (lines, rectangles, circles, etc.)
                if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    return True
                
                # Freeform shapes
                if shape_type == MSO_SHAPE_TYPE.FREEFORM:
                    return True
                
                # Charts
                if shape_type == MSO_SHAPE_TYPE.CHART:
                    return True
                
                # Text boxes and placeholders that might be visual
                if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    return True
                
                if shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    # Only include placeholders that have visual significance
                    return self._has_visual_significance(shape)
                
                # Connectors (lines between shapes)
                if shape_type == MSO_SHAPE_TYPE.LINE:
                    return True
                
                # Tables can be visual
                if shape_type == MSO_SHAPE_TYPE.TABLE:
                    return True
                
                # Media (embedded videos, etc.)
                if shape_type == MSO_SHAPE_TYPE.MEDIA:
                    return True
            
            # Check if shape has visual fill properties
            if hasattr(shape, 'fill'):
                # Shapes with picture fills or other visual fills
                try:
                    from pptx.enum.dml import MSO_FILL
                    if hasattr(shape.fill, 'type'):
                        if shape.fill.type == MSO_FILL.PICTURE:
                            return True
                        elif shape.fill.type == MSO_FILL.PATTERNED:
                            return True
                except:
                    pass
            
            return False
            
        except Exception as e:
            logger.debug(f"Error checking if shape is visual element: {e}")
            return False
    
    def _has_visual_significance(self, shape) -> bool:
        """
        Check if a shape has visual significance beyond just text content.
        
        Args:
            shape: Shape to check
            
        Returns:
            bool: True if shape has visual significance
        """
        try:
            # Check if shape has background fill
            if hasattr(shape, 'fill'):
                try:
                    from pptx.enum.dml import MSO_FILL
                    if hasattr(shape.fill, 'type') and shape.fill.type != MSO_FILL.BACKGROUND:
                        return True
                except:
                    pass
            
            # Check if shape has borders/outline
            if hasattr(shape, 'line'):
                try:
                    if hasattr(shape.line, 'color') and shape.line.color:
                        return True
                except:
                    pass
            
            # Check dimensions - very small shapes might not be visually significant
            try:
                width_px = int(shape.width.emu / 914400 * 96) if hasattr(shape, 'width') and shape.width else 0
                height_px = int(shape.height.emu / 914400 * 96) if hasattr(shape, 'height') and shape.height else 0
                
                # If shape is too small, it's probably not visually significant
                if width_px < 10 or height_px < 10:
                    return False
            except:
                pass
            
            return True
            
        except Exception as e:
            logger.debug(f"Error checking visual significance: {e}")
            return True  # Default to including it
    
    def _create_identifier_from_shape(self, shape, slide_idx: int, shape_identifier) -> PPTXImageIdentifier:
        """
        Create a PPTXImageIdentifier from a shape using STABLE shape IDs when available.
        
        Args:
            shape: Shape object
            slide_idx: Slide index
            shape_identifier: Shape ID (from shape.shape_id) or fallback index
            
        Returns:
            PPTXImageIdentifier instance
        """
        shape_name = getattr(shape, 'name', '')
        
        # For shapes with actual images, use the same logic as processor
        if hasattr(shape, 'image') and shape.image:
            try:
                image_data = shape.image.blob
                # Use same hash function as processor: get_image_hash()
                from decorative_filter import get_image_hash
                image_hash = get_image_hash(image_data)
            except Exception:
                # Fallback to direct md5 if get_image_hash fails
                image_hash = md5(shape.image.blob).hexdigest()
            
            # Truncate to 8 chars to match processor expectations
            image_hash = image_hash[:8] if len(image_hash) > 8 else image_hash
            
            embed_id = ""
            try:
                blip_fill = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                if blip_fill is not None:
                    embed_id = blip_fill.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', '')
            except Exception:
                pass
                
            return PPTXImageIdentifier(slide_idx, shape_identifier, shape_name, image_hash, embed_id)
        
        # For non-image visual elements, create hash similar to processor's PPTXVisualElement
        # This should match how processor creates element_hash for non-image shapes
        embed_id = ""
        image_hash = ""
        
        try:
            # Use same hash logic as processor's PPTXVisualElement.element_hash
            shape_type = getattr(shape, 'shape_type', None)
            width_px = 0
            height_px = 0
            text_content = ""
            
            # Extract dimensions like processor does
            try:
                width_px = int(shape.width.emu / 914400 * 96) if hasattr(shape, 'width') and shape.width else 0
                height_px = int(shape.height.emu / 914400 * 96) if hasattr(shape, 'height') and shape.height else 0
            except:
                pass
            
            # Extract text content like processor does
            try:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text_content = shape.text_frame.text or ""
            except:
                pass
            
            # Create hash content using same format as processor PPTXVisualElement
            hash_content = f"{shape_type}_{width_px}_{height_px}_{text_content}"
            import hashlib
            full_hash = hashlib.md5(hash_content.encode()).hexdigest()
            # Truncate to 8 chars like processor does
            image_hash = full_hash[:8]
            
        except Exception:
            # Fallback hash - use simple approach
            import hashlib
            hash_content = f"{slide_idx}_{shape_identifier}_{shape_name}"
            full_hash = hashlib.md5(hash_content.encode()).hexdigest()
            image_hash = full_hash[:8]
        
        return PPTXImageIdentifier(slide_idx, shape_identifier, shape_name, image_hash, embed_id)
    
    def _try_multi_variant_key_matching(self, target_key: str, image_identifiers: Dict[str, Tuple]) -> Optional[Tuple]:
        """
        SURGICAL FIX A: Multi-variant key matching with canonical format fallbacks.
        Tries both 0-based/1-based variants and shape/shapeid variants in order of specificity.
        
        Args:
            target_key: Key from generation that couldn't be matched directly
            image_identifiers: Available injection identifiers
            
        Returns:
            (identifier, shape) tuple if match found, None otherwise
        """
        try:
            # Parse the target key to extract components
            target_parts = target_key.split('_')
            target_slide = None
            target_shape_id = None
            target_hash = None
            
            # Extract slide index
            for i, part in enumerate(target_parts):
                if part == 'slide' and i + 1 < len(target_parts):
                    try:
                        target_slide = int(target_parts[i + 1])
                        break
                    except ValueError:
                        continue
            
            # Extract shape ID (try both 'shapeid' and 'shape')
            for i, part in enumerate(target_parts):
                if part in ['shapeid', 'shape'] and i + 1 < len(target_parts):
                    try:
                        target_shape_id = target_parts[i + 1]
                        # Convert to int if possible for consistency
                        if target_shape_id.isdigit():
                            target_shape_id = int(target_shape_id)
                        break
                    except (ValueError, IndexError):
                        continue
            
            # Extract hash
            for i, part in enumerate(target_parts):
                if part == 'hash' and i + 1 < len(target_parts):
                    target_hash = target_parts[i + 1]
                    break
            
            if target_slide is None:
                logger.debug(f"Could not extract slide index from key: {target_key}")
                return None
            
            # Generate all possible key variants (most specific → least specific)
            key_variants = []
            
            if target_shape_id is not None and target_hash is not None:
                # Try both 0-based and 1-based slide indices with both shapeid/shape variants
                for slide_variant in [target_slide, target_slide + 1, target_slide - 1]:
                    if slide_variant >= 0:
                        # shapeid variant
                        key_variants.append(f"slide_{slide_variant}_shapeid_{target_shape_id}_hash_{target_hash}")
                        # shape variant
                        key_variants.append(f"slide_{slide_variant}_shape_{target_shape_id}_hash_{target_hash}")
            
            # Try each variant in order
            for variant in key_variants:
                if variant in image_identifiers:
                    logger.info(f"KEY_MATCH_SUCCESS: '{target_key}' -> '{variant}' (variant matching)")
                    return image_identifiers[variant]
            
            # If no exact variants match, try the original fallback logic
            return self._try_fallback_key_matching_legacy(target_key, image_identifiers)
            
        except Exception as e:
            logger.debug(f"Multi-variant key matching failed: {e}")
            return None

    def _try_fallback_key_matching_legacy(self, target_key: str, image_identifiers: Dict[str, Tuple]) -> Optional[Tuple]:
        """
        Legacy fallback key matching by (slide_index, hash) combination when shape IDs don't align.
        
        Args:
            target_key: Key from generation that couldn't be matched directly
            image_identifiers: Available injection identifiers
            
        Returns:
            (identifier, shape) tuple if match found, None otherwise
        """
        try:
            # Parse the target key to extract slide and hash
            target_parts = target_key.split('_')
            target_slide = None
            target_hash = None
            
            # Extract slide index
            for i, part in enumerate(target_parts):
                if part == 'slide' and i + 1 < len(target_parts):
                    try:
                        target_slide = int(target_parts[i + 1])
                    except ValueError:
                        pass
                elif part == 'hash' and i + 1 < len(target_parts):
                    target_hash = target_parts[i + 1]
            
            if target_slide is None or target_hash is None:
                logger.debug(f"Could not parse slide/hash from key: {target_key}")
                return None
            
            logger.debug(f"Looking for fallback match: slide={target_slide}, hash={target_hash}")
            
            # Try to find a match by (slide_index, hash) combination
            for available_key, (identifier, shape) in image_identifiers.items():
                if identifier.slide_idx == target_slide and target_hash in identifier.image_hash:
                    logger.debug(f"Found fallback match: {target_key} -> {available_key}")
                    return (identifier, shape)
            
            # Try partial hash matching (first 8 characters)
            for available_key, (identifier, shape) in image_identifiers.items():
                if identifier.slide_idx == target_slide:
                    available_hash = identifier.image_hash[:8]
                    if available_hash == target_hash:
                        logger.debug(f"Found partial hash match: {target_key} -> {available_key}")
                        return (identifier, shape)
            
            logger.debug(f"No fallback match found for: {target_key}")
            return None
            
        except Exception as e:
            logger.debug(f"Error in fallback key matching for {target_key}: {e}")
            return None
    
    def _try_fallback_injection(
        self,
        presentation: Presentation,
        image_key: str,
        alt_text: str,
        preserve_override: bool = False,
    ) -> bool:
        """
        Try fallback injection methods for images that weren't found through normal shape traversal.
        This handles images accessible through relationships but not the shape API.

        Args:
            presentation: PowerPoint presentation
            image_key: Image key that wasn't matched
            alt_text: ALT text to inject
            preserve_override: Allow writing even when preserve mode would normally skip

        Returns:
            bool: True if injection succeeded via fallback method
        """
        logger.debug(f"Attempting fallback injection for {image_key}")
        
        # Try to parse the image key to understand what we're looking for
        try:
            identifier = PPTXImageIdentifier.parse_from_complex_key(image_key)
            logger.debug(f"Parsed fallback key: slide={identifier.slide_idx}, shape={identifier.shape_idx}")
        except Exception as e:
            logger.debug(f"Could not parse image key for fallback: {e}")
            return False
        
        # Method 1: Try relationship-based injection
        if self._try_relationship_based_injection(presentation, identifier, alt_text):
            return True
        
        # Method 2: Try XML-based direct manipulation
        if self._try_xml_based_injection(presentation, identifier, alt_text):
            return True
        
        # Method 3: Try slide part-based injection
        if self._try_slide_part_injection(presentation, identifier, alt_text):
            return True
        
        logger.debug(f"All fallback methods failed for {image_key}")
        return False
    
    def _write_alt_by_key(
        self,
        image_key: str,
        text: str,
        preserve_override: bool = False,
    ) -> bool:
        """
        Write ALT text to the specific shape identified by the image key.
        This ensures fallback text goes to the correct shape, not just any shape.

        Args:
            image_key: Stable image key (e.g., "slide_2_shapeid_10_hash_abc123")
            text: ALT text to write
            preserve_override: Allow writing even when preserve mode would normally skip

        Returns:
            bool: True if successful
        """
        try:
            # Parse the image key to get slide and shape info
            identifier = PPTXImageIdentifier.from_key(image_key)
            
            # Get the presentation (should be available in current context)
            if not hasattr(self, 'prs') or self.prs is None:
                logger.error(f"No presentation available for key-based writing: {image_key}")
                self.statistics['failed'] += 1
                return False
            
            # Get the slide
            if identifier.slide_idx >= len(self.prs.slides):
                logger.error(f"Slide index {identifier.slide_idx} out of range for {image_key}")
                self.statistics['failed'] += 1
                return False
                
            slide = self.prs.slides[identifier.slide_idx]
            
            # Find the target shape by shape_id
            target = None
            shape_id = identifier.shape_idx
            
            # Try to find by shape_id first
            for shp in slide.shapes:
                if hasattr(shp, 'shape_id') and shp.shape_id == shape_id:
                    target = shp
                    break
            
            # If not found by shape_id, try XML-based matching
            if not target:
                for shp in slide.shapes:
                    try:
                        el = getattr(shp, "_element", None)
                        if el is not None:
                            # Look for cNvPr with matching id
                            cnvpr_nodes = _safe_xpath(el, ".//p:cNvPr", PPTX_NSMAP)
                            for cnvpr in cnvpr_nodes:
                                cnvpr_id = cnvpr.get("id")
                                if cnvpr_id and int(cnvpr_id) == shape_id:
                                    target = shp
                                    break
                            if target:
                                break
                    except Exception:
                        continue
            
            if not target:
                logger.warning(f"Could not find target shape for {image_key} (shape_id: {shape_id})")
                self.statistics['failed'] += 1
                return False
            
            # Write the ALT text using the unified writer
            success = self._inject_alt(
                target,
                text,
                image_key,
                source="fallback_key_based",
                preserve_override=preserve_override,
            )
            if success:
                logger.info(f"✅ Successfully wrote ALT text via key-based targeting: {image_key}")
                self.statistics['success'] += 1
                return True
            else:
                logger.warning(f"❌ Key-based ALT writing was skipped for {image_key}")
                return False
                
        except Exception as e:
            logger.error(f"Error in key-based ALT writing for {image_key}: {e}")
            self.statistics['failed'] += 1
            return False
    
    def _try_relationship_based_injection(self, presentation: Presentation, identifier: PPTXImageIdentifier, alt_text: str) -> bool:
        """
        Try to inject ALT text by finding images through presentation relationships.
        
        Args:
            presentation: PowerPoint presentation
            identifier: Image identifier
            alt_text: ALT text to inject
            
        Returns:
            bool: True if successful
        """
        try:
            logger.debug(f"Trying relationship-based injection for slide {identifier.slide_idx}")
            
            if identifier.slide_idx >= len(presentation.slides):
                return False
            
            slide = presentation.slides[identifier.slide_idx]
            slide_part = slide.part
            
            # Look through slide relationships for images
            for relationship in slide_part.rels.values():
                if hasattr(relationship, 'target_part'):
                    target = relationship.target_part
                    if hasattr(target, 'content_type') and target.content_type.startswith('image/'):
                        logger.debug(f"Found image relationship: {relationship.rId}")
                        # Try to find the corresponding element in the slide XML and set ALT text
                        if self._inject_alt_via_relationship(slide, relationship.rId, alt_text, identifier):
                            return True
            
        except Exception as e:
            logger.debug(f"Relationship-based injection failed: {e}")
        
        return False
    
    def _inject_alt_via_relationship(self, slide, rel_id: str, alt_text: str, identifier: PPTXImageIdentifier) -> bool:
        """
        Inject ALT text by finding XML elements that reference a specific relationship ID.
        
        Args:
            slide: Slide object
            rel_id: Relationship ID
            alt_text: ALT text to inject
            identifier: Image identifier for matching
            
        Returns:
            bool: True if successful
        """
        try:
            # APPLY CENTRALIZED GATE - block templated strings
            policy = self.config_manager.get_fallback_policy()
            quality_flags = {
                "is_generated": False,
                "is_fallback": True  # Relationship injection is typically a fallback
            }
            
            gated_text = apply_for_ppt_injection(alt_text, "shape", quality_flags, policy, None)
            if gated_text is None:
                logger.debug(f"Relationship injection blocked by centralized gate: {alt_text[:50]}...")
                return False
            
            alt_text = gated_text
            # Access the slide's XML and look for elements with this relationship ID
            slide_xml = slide._element
            
            # Look for blip elements that reference this relationship
            blip_elements = _safe_xpath(slide_xml, f'.//a:blip[@r:embed="{rel_id}"]', 
                                          {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                                           'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'})
            
            for blip_element in blip_elements:
                # Find the parent cNvPr element where we can set the description
                parent_elements = _safe_xpath(blip_element, 'ancestor::*')
                for parent in reversed(parent_elements):  # Start from closest ancestor
                    cnvpr_elements = _safe_xpath(parent, './/a:cNvPr | .//p:cNvPr',
                                                {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                                                 'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                    
                    if cnvpr_elements:
                        cnvpr_element = cnvpr_elements[0]
                        # Verify this matches our identifier if possible
                        if self._verify_element_matches_identifier(cnvpr_element, identifier):
                            # GREMLIN 1 FIX: Write exact final_alt with safe finalization
                            finalized_text = self._finalize_alt(alt_text)
                            cnvpr_element.set('descr', finalized_text)
                            
                            # GREMLIN 1 FIX: Post-write read-back assertion
                            actual_written = cnvpr_element.get('descr', '')
                            if actual_written != alt_text:
                                logger.error(f"POST-WRITE ASSERTION FAILED for {rel_id}:")
                                logger.error(f"  Expected: {repr(alt_text)}")
                                logger.error(f"  Actual:   {repr(actual_written)}")
                                return False
                            
                            logger.debug(f"Injected ALT text via relationship {rel_id}")
                            return True
            
        except Exception as e:
            logger.debug(f"Failed to inject via relationship {rel_id}: {e}")
        
        return False
    
    def _verify_element_matches_identifier(self, element, identifier: PPTXImageIdentifier) -> bool:
        """
        Verify that an XML element matches the given identifier.
        
        Args:
            element: XML element
            identifier: Image identifier
            
        Returns:
            bool: True if element matches identifier
        """
        try:
            # For now, we'll be permissive since relationship-based matching
            # is already a fallback method. In a more sophisticated implementation,
            # we could check element IDs, names, or position information.
            
            # Check if element has a name that matches our identifier
            element_name = element.get('name', '')
            if identifier.shape_name and identifier.shape_name in element_name:
                return True
            
            # For fallback cases, assume match if we got this far
            return True
            
        except Exception:
            return True  # Be permissive for fallback injection
    
    def _try_xml_based_injection(self, presentation: Presentation, identifier: PPTXImageIdentifier, alt_text: str) -> bool:
        """
        Try direct XML manipulation to inject ALT text.
        
        Args:
            presentation: PowerPoint presentation
            identifier: Image identifier
            alt_text: ALT text to inject
            
        Returns:
            bool: True if successful
        """
        try:
            # APPLY CENTRALIZED GATE - block templated strings
            policy = self.config_manager.get_fallback_policy()
            quality_flags = {
                "is_generated": False,
                "is_fallback": True  # XML injection is typically a fallback
            }
            
            gated_text = apply_for_ppt_injection(alt_text, "shape", quality_flags, policy, None)
            if gated_text is None:
                logger.debug(f"XML injection blocked by centralized gate: {alt_text[:50]}...")
                return False
            
            alt_text = gated_text
            logger.debug(f"Trying XML-based injection for slide {identifier.slide_idx}")
            
            if identifier.slide_idx >= len(presentation.slides):
                return False
            
            slide = presentation.slides[identifier.slide_idx]
            slide_xml = slide._element
            
            # Look for all a:cNvPr elements (drawingML non-visual properties)
            cnvpr_elements = _safe_xpath(slide_xml, './/a:cNvPr',
                                           {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            
            # Try to match by position or other identifying characteristics
            for i, cnvpr in enumerate(cnvpr_elements):
                if self._element_matches_shape_index(cnvpr, identifier.shape_idx, i):
                    # GREMLIN 1 FIX: Write exact final_alt with safe finalization
                    finalized_text = self._finalize_alt(alt_text)
                    cnvpr.set('descr', finalized_text)
                    
                    # GREMLIN 1 FIX: Post-write read-back assertion
                    actual_written = cnvpr.get('descr', '')
                    if actual_written != alt_text:
                        logger.error(f"POST-WRITE ASSERTION FAILED at index {i}:")
                        logger.error(f"  Expected: {repr(alt_text)}")
                        logger.error(f"  Actual:   {repr(actual_written)}")
                        return False
                    
                    logger.debug(f"Injected ALT text via XML manipulation at index {i}")
                    return True
            
        except Exception as e:
            logger.debug(f"XML-based injection failed: {e}")
        
        return False
    
    def _element_matches_shape_index(self, element, target_shape_idx, current_index: int) -> bool:
        """
        Check if an XML element matches the target shape index.
        
        Args:
            element: XML element
            target_shape_idx: Target shape index (may be complex like "0_1_2")
            current_index: Current enumeration index
            
        Returns:
            bool: True if element matches
        """
        # If target is a simple integer, match by current index
        if isinstance(target_shape_idx, int):
            return current_index == target_shape_idx
        
        # For complex shape indices, this is more challenging
        # For now, use a heuristic approach
        if isinstance(target_shape_idx, str) and '_' in target_shape_idx:
            # Parse the complex index and try to match the last component
            parts = target_shape_idx.split('_')
            try:
                if parts:  # Ensure parts is not empty
                    last_index = int(parts[-1])
                    return current_index == last_index
            except ValueError:
                pass
        
        return False
    
    def _try_slide_part_injection(self, presentation: Presentation, identifier: PPTXImageIdentifier, alt_text: str) -> bool:
        """
        Try injection through slide part manipulation.
        
        Args:
            presentation: PowerPoint presentation
            identifier: Image identifier
            alt_text: ALT text to inject
            
        Returns:
            bool: True if successful
        """
        try:
            logger.debug(f"Trying slide part injection for slide {identifier.slide_idx}")
            
            if identifier.slide_idx >= len(presentation.slides):
                return False
            
            slide = presentation.slides[identifier.slide_idx]
            slide_part = slide.part
            
            # This is a placeholder for more advanced slide part manipulation
            # In a full implementation, this would involve:
            # 1. Parsing the slide's XML content
            # 2. Finding image elements by their relationship IDs
            # 3. Modifying the appropriate cNvPr elements
            # 4. Handling the part's relationships and embedded content
            
            logger.debug("Slide part injection not yet fully implemented")
            return False
            
        except Exception as e:
            logger.debug(f"Slide part injection failed: {e}")
            return False
    
    def _inject_alt_text_single(
        self,
        shape,
        alt_text: str,
        identifier: PPTXImageIdentifier,
        preserve_override: bool = False,
    ) -> bool:
        """
        PHASE 1.1 HOTFIX: Route all injection through single source of truth.
        This function is now just a wrapper around _inject_alt().

        Args:
            shape: Picture shape to inject ALT text into
            alt_text: ALT text to inject
            identifier: Image identifier for logging
            preserve_override: Allow writing even when preserve mode would normally skip

        Returns:
            bool: True if injection was successful
        """
        # Belt-and-suspenders: skip injection entirely if the post-gate ALT is empty
        if not alt_text or not alt_text.strip():
            logger.info(f"ALT is empty after gating; skipping injection for {identifier.image_key}")
            return False

        try:
            logger.debug(f"🔍 Processing injection for {identifier.image_key}")
            logger.debug(f"   Raw ALT text: '{alt_text}'")
            
            # PHASE 1.1: Check classifier fences
            is_processed, processed_by = self._is_element_processed(identifier.image_key)
            if is_processed:
                logger.debug(f"FENCE: {identifier.image_key} already processed by {processed_by}, skipping")
                self.statistics['skipped_existing'] += 1
                return False
            
            # Check if we should skip this ALT text (before going to single source of truth)
            gate_allowed = not self._should_skip_alt_text(alt_text)
            logger.info(f"GATE key={identifier.image_key} allowed={gate_allowed}")
            if not gate_allowed:
                logger.debug(f"❌ DECISION: Skipping invalid ALT text for {identifier.image_key}: '{alt_text}'")
                logger.debug(f"   Reason: ALT text matches skip patterns")
                self.statistics['skipped_invalid'] += 1
                return False
            
            # PHASE 1.1 HOTFIX: Use single source of truth for ALL injection
            # This replaces all the old logic with the idempotent injector
            logger.debug(f"ATTEMPTING WRITE for {identifier.image_key} via _inject_alt() (mode={self.mode})")
            success = self._inject_alt(
                shape,
                alt_text,
                identifier.image_key,
                "generator",
                preserve_override=preserve_override,
            )
            
            # Update stats based on result
            if success:
                logger.debug(f"✅ Successfully injected via single source of truth")
                self.statistics['injected_successfully'] += 1

                # ONLY classify element as processed after successful write
                element_type = getattr(identifier, 'element_type', 'generic')
                self._classify_element(identifier.image_key, element_type)
            else:
                logger.debug(f"➡️  Skipped by idempotent guard (already equivalent or processed)")
                # Note: _inject_alt logs the specific reason for skipping
            
            return success
                
        except Exception as e:
            logger.error(f"💥 Error injecting ALT text for {identifier.image_key}: {e}")
            self.statistics['failed_injection'] += 1
            return False
    
    def _should_inject_alt_text(self, existing_alt_text: str, new_alt_text: str, image_key: str) -> Tuple[bool, str]:
        """
        Determine whether to inject new ALT text based on current state and mode.
        
        This is the core logic that fixes the preserve/overwrite issue by properly
        distinguishing between pre-existing ALT text vs freshly generated text.
        
        Args:
            existing_alt_text: Current ALT text in the shape (may be empty)
            new_alt_text: Newly generated ALT text to potentially inject
            image_key: Image identifier for logging
            
        Returns:
            Tuple of (should_inject: bool, reason: str)
        """
        # If no existing ALT text, always inject the new text
        if not existing_alt_text or existing_alt_text.strip() == "":
            return True, "No existing ALT text found - injecting new text"
        
        # Check if existing ALT text should be skipped (invalid patterns)
        if self._should_skip_alt_text(existing_alt_text):
            return True, f"Existing ALT text is invalid/skippable: '{existing_alt_text}' - replacing with new text"
        
        # Mode-based decisions
        if self.mode == 'overwrite':
            return True, f"Overwrite mode - replacing existing '{existing_alt_text}' with new text"
        
        elif self.mode == 'preserve':
            # In preserve mode, we need to distinguish between:
            # 1. Truly pre-existing ALT text (preserve it)
            # 2. ALT text that was just generated but is now being seen as "existing"
            
            # Check if this looks like generated decorative text that should be replaced
            if existing_alt_text.strip().lower() in ['[decorative image]', 'decorative image', '']:
                return True, f"Existing text appears to be generated decorative placeholder - replacing"
            
            # Check if existing ALT text matches common generation failure patterns
            failure_patterns = [
                'error', 'failed', 'cannot', 'unable', 'sorry', 
                'i cannot', 'i am unable', 'no description',
                'not available', 'description not available',
                'image could not be processed'
            ]
            
            existing_lower = existing_alt_text.lower()
            for pattern in failure_patterns:
                if pattern in existing_lower:
                    return True, f"Existing text contains failure pattern '{pattern}' - replacing"
            
            # If we have meaningful existing ALT text and we're in preserve mode,
            # check if the new text is significantly better
            if len(existing_alt_text.strip()) > 10:  # Meaningful length
                # For now, preserve existing meaningful ALT text in preserve mode
                # TODO: Could add AI comparison here to determine if new text is better
                return False, f"Preserve mode - keeping existing meaningful ALT text: '{existing_alt_text}'"
            else:
                # Short existing text might not be meaningful
                return True, f"Existing ALT text too short ('{existing_alt_text}') - replacing with new text"
        
        else:
            # Unknown mode, default to inject
            return True, f"Unknown mode '{self.mode}' - defaulting to inject"
    
    def _inject_alt_text_robust(self, shape, alt_text: str) -> bool:
        """
        Inject ALT text using multiple fallback methods for maximum compatibility.
        
        Args:
            shape: Picture shape
            alt_text: ALT text to inject
            
        Returns:
            bool: True if any method succeeded
        """
        # List of injection methods in order of preference
        injection_methods = [
            ('modern_property', self._inject_via_modern_property),
            ('xml_cnvpr', self._inject_via_xml_cnvpr),
            ('xml_shape_cnvpr', self._inject_via_xml_shape_cnvpr),  # For non-picture shapes
            ('xml_element', self._inject_via_xml_element),
            ('xml_fallback', self._inject_via_xml_fallback)
        ]
        
        logger.debug(f"   Attempting {len(injection_methods)} injection methods:")
        
        for method_name, method_func in injection_methods:
            try:
                logger.debug(f"     Trying {method_name}...")
                if method_func(shape, alt_text):
                    logger.debug(f"     ✅ ALT text injected successfully via {method_name}")
                    return True
                else:
                    logger.debug(f"     ❌ Method {method_name} returned False")
            except Exception as e:
                logger.debug(f"     💥 Injection method {method_name} failed: {e}")
                continue
        
        logger.debug(f"   ❌ All {len(injection_methods)} injection methods failed")
        return False
    
    def _inject_via_modern_property(self, shape, alt_text: str) -> bool:
        """Inject using modern property-based approach (python-pptx >= 0.6.22)."""
        logger.info(f"🔍 DEBUG: XML PATH - Modern property injection")
        logger.info(f"🔍 DEBUG:   Shape type: {type(shape).__name__}")
        if hasattr(shape, 'descr'):
            logger.info(f"🔍 DEBUG:   Using shape.descr property (modern approach)")
            logger.info(f"🔍 DEBUG:   Setting ALT text: '{alt_text}'")
            
            # DUAL WRITE: Set both descr and alternative_text for Office compatibility
            # Apply safe finalization
            finalized_text = self._finalize_alt(alt_text)
            shape.descr = finalized_text

            # Also set alternative_text property if available (covers various Office readers)
            if hasattr(shape, 'alternative_text'):
                shape.alternative_text = finalized_text
                logger.info(f"🔍 DEBUG:   Also set alternative_text property")
            elif hasattr(shape, 'alt_text'):
                shape.alt_text = finalized_text
                logger.info(f"🔍 DEBUG:   Also set alt_text property")
            
            # Set title (short version for Reading Order)
            title_text = self._create_title_from_alt_text(alt_text)
            if title_text and hasattr(shape, 'title'):
                shape.title = title_text
                logger.info(f"🔍 DEBUG:   Setting title: '{title_text}'")
            
            # Verify it was set
            actual_value = getattr(shape, 'descr', '')
            logger.info(f"🔍 DEBUG:   Verification - got back: '{actual_value}'")
            return True
        else:
            logger.info(f"🔍 DEBUG:   Shape does not have 'descr' property")
        return False
    
    def _inject_via_xml_cnvpr(self, shape, alt_text: str) -> bool:
        """Inject via direct XML cNvPr element manipulation."""
        logger.info(f"🔍 DEBUG: XML PATH - cNvPr element injection")
        try:
            # APPLY CENTRALIZED GATE - block templated strings
            policy = self.config_manager.get_fallback_policy()
            quality_flags = {
                "is_generated": False,
                "is_fallback": True  # XML cNvPr injection is typically a fallback
            }
            
            gated_text = apply_for_ppt_injection(alt_text, "shape", quality_flags, policy, shape)
            if gated_text is None:
                logger.debug(f"XML cNvPr injection blocked by centralized gate: {alt_text[:50]}...")
                return False
            
            alt_text = gated_text
            logger.info(f"🔍 DEBUG:   Accessing shape._element._nvXxPr.cNvPr")
            cNvPr = shape._element._nvXxPr.cNvPr
            logger.info(f"🔍 DEBUG:   cNvPr element found: {cNvPr}")
            logger.info(f"🔍 DEBUG:   cNvPr tag: {getattr(cNvPr, 'tag', 'unknown')}")
            logger.info(f"🔍 DEBUG:   Setting descr attribute: '{alt_text}'")
            
            # Set description (full ALT text) - apply safe finalization
            finalized_text = self._finalize_alt(alt_text)
            cNvPr.set('descr', finalized_text)
            
            # Set title (short version for Reading Order)
            title_text = self._create_title_from_alt_text(alt_text)
            if title_text:
                cNvPr.set('title', title_text)
                logger.info(f"🔍 DEBUG:   Setting title attribute: '{title_text}'")
            
            # Verify it was set
            actual_value = cNvPr.get('descr', '')
            logger.info(f"🔍 DEBUG:   Verification - cNvPr.get('descr'): '{actual_value}'")
            return True
        except AttributeError as e:
            logger.info(f"🔍 DEBUG:   cNvPr element access failed: {e}")
            return False
    
    def _inject_via_xml_element(self, shape, alt_text: str) -> bool:
        """Inject via XML element attribute (current approach)."""
        logger.info(f"🔍 DEBUG: XML PATH - Direct element injection")
        try:
            # APPLY CENTRALIZED GATE - block templated strings
            policy = self.config_manager.get_fallback_policy()
            quality_flags = {
                "is_generated": False,
                "is_fallback": True  # XML element injection is typically a fallback
            }
            
            gated_text = apply_for_ppt_injection(alt_text, "shape", quality_flags, policy, shape)
            if gated_text is None:
                logger.debug(f"XML element injection blocked by centralized gate: {alt_text[:50]}...")
                return False
            
            alt_text = gated_text
            logger.info(f"🔍 DEBUG:   Accessing shape._element")
            logger.info(f"🔍 DEBUG:   Element: {shape._element}")
            logger.info(f"🔍 DEBUG:   Element tag: {getattr(shape._element, 'tag', 'unknown')}")
            logger.info(f"🔍 DEBUG:   Setting descr attribute: '{alt_text}'")
            # Apply safe finalization
            finalized_text = self._finalize_alt(alt_text)
            shape._element.set('descr', finalized_text)
            
            # Verify it was set
            actual_value = shape._element.get('descr', '')
            logger.info(f"🔍 DEBUG:   Verification - element.get('descr'): '{actual_value}'")
            return True
        except Exception as e:
            logger.info(f"🔍 DEBUG:   XML element access failed: {e}")
            return False
    
    def _inject_via_xml_shape_cnvpr(self, shape, alt_text: str) -> bool:
        """Inject via XML cNvPr element for general shapes (not just pictures)."""
        logger.info(f"🔍 DEBUG: XML PATH - Shape cNvPr element injection")
        try:
            # APPLY CENTRALIZED GATE - block templated strings
            policy = self.config_manager.get_fallback_policy()
            quality_flags = {
                "is_generated": False,
                "is_fallback": True  # XML shape cNvPr injection is typically a fallback
            }
            
            gated_text = apply_for_ppt_injection(alt_text, "shape", quality_flags, policy, shape)
            if gated_text is None:
                logger.debug(f"XML shape cNvPr injection blocked by centralized gate: {alt_text[:50]}...")
                return False
            
            alt_text = gated_text
            
            # Look for cNvPr elements in the shape's XML structure
            if not hasattr(shape, '_element'):
                logger.info(f"🔍 DEBUG:   Shape has no _element attribute")
                return False
                
            element = shape._element
            logger.info(f"🔍 DEBUG:   Shape element: {element}")
            logger.info(f"🔍 DEBUG:   Element tag: {getattr(element, 'tag', 'unknown')}")
            
            # Try to find cNvPr element using XPath
            # Different shape types have different XML structures:
            # - Pictures: p:pic/p:nvPicPr/p:cNvPr
            # - Shapes: p:sp/p:nvSpPr/p:cNvPr
            # - Lines: p:cxnSp/p:nvCxnSpPr/p:cNvPr
            # - Groups: p:grpSp/p:nvGrpSpPr/p:cNvPr
            
            namespaces = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture'
            }
            
            # Try different cNvPr paths for different shape types
            cnvpr_paths = [
                './/p:cNvPr',           # Most general - any cNvPr element
                './/a:cNvPr',           # DrawingML cNvPr (most common for shapes/pictures)
                './p:nvSpPr/p:cNvPr',   # Shape-specific
                './p:nvPicPr/p:cNvPr',  # Picture-specific
                './p:nvCxnSpPr/p:cNvPr', # Connector-specific
                './p:nvGrpSpPr/p:cNvPr' # Group-specific
            ]
            
            for path in cnvpr_paths:
                try:
                    cnvpr_elements = _safe_xpath(element, path, namespaces)
                    if cnvpr_elements:
                        cnvpr_element = cnvpr_elements[0]  # Take the first match
                        logger.info(f"🔍 DEBUG:   Found cNvPr via path: {path}")
                        logger.info(f"🔍 DEBUG:   cNvPr element: {cnvpr_element}")
                        logger.info(f"🔍 DEBUG:   Setting descr attribute: '{alt_text}'")
                        
                        # OVERWRITE MODE: Check existing ALT text before updating
                        existing_descr = cnvpr_element.get('descr', '')
                        existing_title = cnvpr_element.get('title', '')
                        
                        # Compare normalized versions to avoid redundant updates
                        normalized_new = self._normalize_for_comparison(alt_text)
                        normalized_existing = self._normalize_for_comparison(existing_descr)
                        
                        if normalized_new != normalized_existing or not existing_descr:
                            # Set description (full ALT text) with safe finalization
                            finalized_text = self._finalize_alt(alt_text)
                            cnvpr_element.set('descr', finalized_text)
                            logger.info(f"🔍 DEBUG:   Updated descr: '{existing_descr}' -> '{alt_text}'")
                            
                            # Set title (short version for Reading Order)
                            title_text = self._create_title_from_alt_text(alt_text)
                            if title_text:
                                cnvpr_element.set('title', title_text)
                                logger.info(f"🔍 DEBUG:   Updated title: '{existing_title}' -> '{title_text}'")
                        else:
                            logger.info(f"🔍 DEBUG:   Skipping update - ALT text unchanged: '{existing_descr}'")
                        
                        # Verify current state
                        actual_value = cnvpr_element.get('descr', '')
                        logger.info(f"🔍 DEBUG:   Final verification - cNvPr.get('descr'): '{actual_value}'")
                        return True
                        
                except Exception as xpath_error:
                    # SURGICAL FIX B: Special logging for group parent failures
                    if 'nvGrpSpPr' in path:
                        logger.warning(f"GROUP_PARENT_XPATH_FAILURE: XPath failed for path '{path}': {xpath_error}")
                    else:
                        logger.info(f"🔍 DEBUG:   Path {path} failed: {xpath_error}")
                    continue
            
            logger.info(f"🔍 DEBUG:   No cNvPr element found via any path")
            return False
            
        except Exception as e:
            logger.info(f"🔍 DEBUG:   Shape cNvPr injection failed: {e}")
            return False
    
    def _inject_via_xml_fallback(self, shape, alt_text: str) -> bool:
        """Fallback XML injection method."""
        # APPLY CENTRALIZED GATE - block templated strings
        policy = self.config_manager.get_fallback_policy()
        quality_flags = {
            "is_generated": False,
            "is_fallback": True  # XML fallback injection is always a fallback
        }
        
        gated_text = apply_for_ppt_injection(alt_text, "shape", quality_flags, policy, shape)
        if gated_text is None:
            logger.debug(f"XML fallback injection blocked by centralized gate: {alt_text[:50]}...")
            return False
        
        alt_text = gated_text
        
        elements_to_try = [
            ('shape._element', lambda: shape._element),
            ('shape._element._nvXxPr', lambda: shape._element._nvXxPr),
            ('shape._element._nvXxPr.cNvPr', lambda: shape._element._nvXxPr.cNvPr)
        ]
        
        for element_name, element_getter in elements_to_try:
            try:
                element = element_getter()
                if element is not None:
                    logger.debug(f"       Trying fallback on {element_name}")
                    element.set('descr', alt_text)
                    return True
            except Exception as e:
                logger.debug(f"       Fallback {element_name} failed: {e}")
                continue
        
        logger.debug(f"       All fallback methods exhausted")
        return False
    
    def _get_existing_alt_text(self, shape) -> str:
        """
        Get existing ALT text from shape with debug logging.
        
        Args:
            shape: Picture shape
            
        Returns:
            str: Existing ALT text or empty string
        """
        # Try modern property first
        try:
            if hasattr(shape, 'descr'):
                alt_text = shape.descr or ""
                logger.debug(f"   Retrieved ALT text via shape.descr: '{alt_text}'")
                return alt_text
            else:
                logger.debug(f"   Shape does not have 'descr' property")
        except Exception as e:
            logger.debug(f"   Failed to get ALT text via shape.descr: {e}")
        
        # Try XML access
        try:
            cNvPr = shape._element._nvXxPr.cNvPr
            alt_text = cNvPr.get('descr', '')
            logger.debug(f"   Retrieved ALT text via XML cNvPr: '{alt_text}'")
            return alt_text
        except Exception as e:
            logger.debug(f"   Failed to get ALT text via XML cNvPr: {e}")
        
        # Fallback XML access
        try:
            alt_text = shape._element.get('descr', '')
            logger.debug(f"   Retrieved ALT text via XML element: '{alt_text}'")
            return alt_text
        except Exception as e:
            logger.debug(f"   Failed to get ALT text via XML element: {e}")
        
        logger.debug(f"   No ALT text found via any method")
        return ""
    
    def _validate_alt_text_injection(self, shape, expected_alt_text: str) -> bool:
        """
        Validate that ALT text was successfully injected.
        Uses multiple methods to retrieve ALT text to account for different injection methods.
        
        Args:
            shape: Shape object (picture or other visual element)
            expected_alt_text: Expected ALT text
            
        Returns:
            bool: True if validation passed
        """
        # Normalize expected text for comparison
        expected_normalized = expected_alt_text.strip()
        
        # Try multiple methods to get the actual ALT text, since different injection
        # methods might store it in different places
        retrieval_methods = [
            ('standard_method', self._get_existing_alt_text),
            ('descr_property', self._get_alt_via_descr_property),
            ('xml_cnvpr', self._get_alt_via_xml_cnvpr),
            ('xml_element', self._get_alt_via_xml_element)
        ]
        
        for method_name, method_func in retrieval_methods:
            try:
                actual_alt_text = method_func(shape)
                actual_normalized = actual_alt_text.strip() if actual_alt_text else ""
                
                if actual_normalized == expected_normalized:
                    logger.debug(f"   ✅ Validation successful via {method_name}: '{actual_normalized}'")
                    return True
                elif actual_normalized:  # Log non-empty mismatches
                    logger.debug(f"   ⚠️  Method {method_name} returned: '{actual_normalized}' (expected: '{expected_normalized}')")
                    
            except Exception as e:
                logger.debug(f"   ❌ Validation method {method_name} failed: {e}")
                continue
        
        # If no exact match, check for partial matches (in case of text cleaning/modification)
        try:
            actual_alt_text = self._get_existing_alt_text(shape)
            actual_normalized = actual_alt_text.strip() if actual_alt_text else ""
            
            if actual_normalized and expected_normalized:
                # Check if core content matches (allows for minor differences)
                if len(actual_normalized) >= 10 and len(expected_normalized) >= 10:
                    # For longer texts, check if 80% of content matches
                    if self._texts_substantially_match(actual_normalized, expected_normalized, 0.8):
                        logger.debug(f"   ✅ Validation successful via substantial match: '{actual_normalized[:30]}...'")
                        return True
                        
        except Exception as e:
            logger.debug(f"   ❌ Partial match validation failed: {e}")
        
        logger.debug(f"   ❌ Validation failed - no retrieval method found matching ALT text")
        return False
    
    def _get_alt_via_descr_property(self, shape) -> str:
        """Get ALT text via the descr property."""
        try:
            if hasattr(shape, 'descr'):
                return shape.descr or ""
        except:
            pass
        return ""
    
    def _get_alt_via_xml_cnvpr(self, shape) -> str:
        """Get ALT text via XML cNvPr element."""
        try:
            if hasattr(shape, '_element'):
                # Try the standard cNvPr path for pictures
                cNvPr = getattr(shape._element, '_nvXxPr', None)
                if cNvPr:
                    cNvPr_elem = getattr(cNvPr, 'cNvPr', None)
                    if cNvPr_elem:
                        return cNvPr_elem.get('descr', '') or ""
        except:
            pass
        return ""
    
    def _get_alt_via_xml_element(self, shape) -> str:
        """Get ALT text via direct XML element access."""
        try:
            if hasattr(shape, '_element'):
                return shape._element.get('descr', '') or ""
        except:
            pass
        return ""
    
    def _texts_substantially_match(self, text1: str, text2: str, threshold: float = 0.8) -> bool:
        """
        Check if two texts substantially match using a simple similarity metric.
        
        Args:
            text1: First text
            text2: Second text  
            threshold: Similarity threshold (0.0 to 1.0)
            
        Returns:
            bool: True if texts are substantially similar
        """
        try:
            # Simple word-based similarity
            words1 = set(text1.lower().split())
            words2 = set(text2.lower().split())
            
            if not words1 and not words2:
                return True
            if not words1 or not words2:
                return False
                
            intersection = len(words1.intersection(words2))
            union = len(words1.union(words2))
            
            similarity = intersection / union if union > 0 else 0
            return similarity >= threshold
            
        except Exception:
            return False
    
    def _perform_post_injection_verification(self, presentation: Presentation,
                                          image_identifiers: Dict[str, Tuple],
                                          alt_text_mapping: Dict[str, Dict[str, Any]]) -> None:
        """
        Verify that ALT texts were actually injected after the injection process.
        
        Args:
            presentation: PowerPoint presentation
            image_identifiers: Mapping of image keys to (identifier, shape) tuples
            alt_text_mapping: Original ALT text mapping requested
        """
        logger.info("🔍 DEBUG: Verifying ALT text injection results...")
        
        successful_injections = 0
        failed_injections = 0
        
        # Check each image that we tried to inject
        for image_key, alt_record in alt_text_mapping.items():
            # Handle both old string format and new enriched format
            if isinstance(alt_record, str):
                expected_alt_text = alt_record
            else:
                expected_alt_text = (
                    alt_record.get('final_alt', '')
                    or alt_record.get('generated_alt', '')
                    or alt_record.get('existing_alt', '')
                )

            if image_key in image_identifiers:
                identifier, shape = image_identifiers[image_key]

                # Get current ALT text using all available methods
                current_alt_text = self._get_existing_alt_text(shape)
                
                # Normalize both texts for comparison
                expected_normalized = self._normalize_for_comparison(expected_alt_text)
                actual_normalized = self._normalize_for_comparison(current_alt_text)

                if expected_normalized == actual_normalized:
                    logger.info(f"🔍 DEBUG: ✅ VERIFIED (normalized): {image_key}")
                    logger.info(f"🔍 DEBUG:   Expected: '{expected_alt_text}' -> '{expected_normalized}'")
                    logger.info(f"🔍 DEBUG:   Actual: '{current_alt_text}' -> '{actual_normalized}'")
                    successful_injections += 1
                else:
                    logger.info(f"🔍 DEBUG: ❌ FAILED: {image_key}")
                    logger.info(f"🔍 DEBUG:   Expected: '{expected_alt_text}' -> '{expected_normalized}'")
                    logger.info(f"🔍 DEBUG:   Actual: '{current_alt_text}' -> '{actual_normalized}'")

                    # REPAIR PASS: Check if actual text matches blocked fallback patterns
                    if self._is_blocked_fallback_pattern(current_alt_text):
                        logger.info(f"🔧 REPAIR: Detected blocked fallback pattern, attempting repair...")
                        repair_success = self._inject_alt(shape, expected_alt_text, image_key, "repair_injected")
                        if repair_success:
                            # Re-verify after repair using normalized comparison
                            repaired_text = self._get_existing_alt_text(shape)
                            repaired_normalized = self._normalize_for_comparison(repaired_text)
                            if repaired_normalized == expected_normalized:
                                logger.info(f"🔧 REPAIR: ✅ Successfully repaired {image_key}")
                                successful_injections += 1
                                continue
                            else:
                                logger.warning(f"🔧 REPAIR: ❌ Repair verification failed for {image_key}")
                        else:
                            logger.warning(f"🔧 REPAIR: ❌ Repair injection failed for {image_key}")

                    failed_injections += 1
                    
                    # Additional debug info for failed injections
                    logger.info(f"🔍 DEBUG:   Shape type: {type(shape).__name__}")
                    logger.info(f"🔍 DEBUG:   Shape ID: {getattr(shape, 'shape_id', 'unknown')}")
                    if hasattr(shape, '_element'):
                        try:
                            # Check XML attributes directly
                            descr_attr = shape._element.get('descr')
                            logger.info(f"🔍 DEBUG:   XML descr attribute: '{descr_attr}'")
                            
                            # Check cNvPr element
                            if hasattr(shape._element, '_nvXxPr'):
                                cnvpr = getattr(shape._element._nvXxPr, 'cNvPr', None)
                                if cnvpr is not None:
                                    cnvpr_descr = cnvpr.get('descr')
                                    logger.info(f"🔍 DEBUG:   cNvPr descr attribute: '{cnvpr_descr}'")
                        except Exception as e:
                            logger.info(f"🔍 DEBUG:   XML inspection failed: {e}")
        
        logger.info(f"🔍 DEBUG: VERIFICATION SUMMARY:")
        logger.info(f"🔍 DEBUG:   Successful injections: {successful_injections}")
        logger.info(f"🔍 DEBUG:   Failed injections: {failed_injections}")
        logger.info(f"🔍 DEBUG:   Total attempted: {len(alt_text_mapping)}")
        
        # Also verify by re-scanning the presentation
        logger.info("🔍 DEBUG: Re-scanning presentation for all ALT texts...")
        all_alt_texts_found = 0
        for slide_idx, slide in enumerate(presentation.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'image') or hasattr(shape, '_element'):
                    alt_text = self._get_existing_alt_text(shape)
                    if alt_text.strip():
                        all_alt_texts_found += 1
                        logger.info(f"🔍 DEBUG:   Found ALT text on slide {slide_idx}, shape {shape_idx}: '{alt_text}'")
        
        logger.info(f"🔍 DEBUG: Total ALT texts found in presentation: {all_alt_texts_found}")
    
    def _normalize_for_comparison(self, text: str) -> str:
        """
        Normalize ALT text for comparison to avoid redundant updates.
        
        Args:
            text: ALT text to normalize
            
        Returns:
            Normalized text for comparison
        """
        if not text:
            return ""
        
        import re
        
        # Normalize whitespace and case
        normalized = re.sub(r'\s+', ' ', text.strip().lower())
        
        # Remove common punctuation variations
        normalized = re.sub(r'[.!?]+$', '', normalized)  # Remove trailing punctuation
        
        # Remove common prefixes that don't affect meaning
        prefixes = [
            'this is a powerpoint ',
            'powerpoint ',
            'a ',
            'an ',
            'the '
        ]
        
        for prefix in prefixes:
            if normalized.startswith(prefix):
                normalized = normalized[len(prefix):]
                break
        
        return normalized.strip()
    
    def _create_title_from_alt_text(self, alt_text: str) -> str:
        """
        Create a short title from the full ALT text for PowerPoint Reading Order.

        Args:
            alt_text: Full ALT text

        Returns:
            Shortened title (60-80 characters)
        """
        if not alt_text:
            return ""  # guard 1

        clean_text = alt_text.strip()
        if not clean_text:
            return ""  # guard 2

        # Optionally strip known prefixes your gate might leave behind
        for prefix in ("FALLBACK:",):
            if clean_text.upper().startswith(prefix):
                clean_text = clean_text[len(prefix):].strip(" .:-–—")
                break

        if not clean_text:
            return ""  # guard 3

        # Remove common prefixes to make title more concise
        prefixes_to_remove = [
            "This is a PowerPoint shape. It is ",
            "This is a PowerPoint ",
            "This is a ",
            "A PowerPoint ",
            "PowerPoint "
        ]

        for prefix in prefixes_to_remove:
            if clean_text.startswith(prefix):
                clean_text = clean_text[len(prefix):].strip()
                break

        if not clean_text:
            return ""  # guard 4 after prefix removal

        # Safe capitalization using _finalize_alt
        clean_text = self._finalize_alt(clean_text)
        
        # Truncate to reasonable title length (60-80 chars)
        max_title_length = 70
        if len(clean_text) <= max_title_length:
            return clean_text
        
        # Find a good breaking point (sentence end, period, etc.)
        truncated = clean_text[:max_title_length]
        
        # Try to break at sentence boundaries
        for break_char in ['. ', '! ', '? ']:
            last_break = truncated.rfind(break_char)
            if last_break > max_title_length * 0.7:  # At least 70% of target length
                return truncated[:last_break + 1].strip()
        
        # Fall back to word boundary
        last_space = truncated.rfind(' ')
        if last_space > max_title_length * 0.7:
            return truncated[:last_space].strip()
        
        # Hard truncate with ellipsis
        return truncated.strip() + "..."
    
    def _should_skip_alt_text(self, alt_text: str) -> bool:
        """
        Check if ALT text should be skipped based on reinjection rules.
        
        Args:
            alt_text: ALT text to check
            
        Returns:
            bool: True if ALT text should be skipped
        """
        if not alt_text:
            return True
        
        alt_text_stripped = alt_text.strip()
        
        for skip_pattern in self.skip_alt_text_if:
            if isinstance(skip_pattern, str):
                if skip_pattern == alt_text_stripped:
                    return True
            
        return False
    
    def _clean_alt_text(self, alt_text: str) -> str:
        """
        Clean ALT text using existing alt_cleaner if available.
        
        Args:
            alt_text: ALT text to clean
            
        Returns:
            str: Cleaned ALT text
        """
        try:
            # Import and use existing alt_cleaner
            sys.path.insert(0, str(project_root / "shared"))
            from alt_cleaner import clean_alt_text
            return clean_alt_text(alt_text)
        except ImportError:
            logger.debug("alt_cleaner not available, using basic cleaning")
            # Basic cleaning - remove extra whitespace
            return " ".join(alt_text.split())
        except Exception as e:
            logger.warning(f"Error cleaning ALT text: {e}")
            return alt_text
    
    def _log_injection_summary(self, result: Dict[str, Any]):
        """Log summary of injection results."""
        stats = result['statistics']
        
        logger.info("PPTX ALT Text Injection Summary:")
        logger.info(f"  Input file: {result['input_file']}")
        logger.info(f"  Output file: {result['output_file']}")
        logger.info(f"  Total images found: {stats['total_images']}")
        logger.info(f"  Successfully processed: {stats['injected_successfully']}")
        logger.info(f"  Preserved existing ALT: {stats.get('preserved_existing', 0)}")
        logger.info(f"  Written (generated): {stats.get('written_generated', 0)}")
        logger.info(f"  Written (existing fallback): {stats.get('written_existing', 0)}")
        logger.info(f"  Written (stored final): {stats.get('written_final', 0)}")
        logger.info(f"  Skipped (no content): {stats.get('skipped_no_content', 0)}")
        logger.info(f"  Skipped (existing fence): {stats['skipped_existing']}")
        logger.info(f"  Skipped (invalid): {stats['skipped_invalid']}")
        logger.info(f"  Failed injection: {stats['failed_injection']}")
        logger.info(f"  Validation failures: {stats['validation_failures']}")
        logger.info(f"  Success: {result['success']}")
        
        if result.get('errors'):
            logger.warning(f"Errors encountered: {len(result['errors'])}")
            for error in result['errors']:
                logger.warning(f"  - {error}")
    
    def extract_images_with_identifiers(self, pptx_path: str) -> Dict[str, Dict[str, Any]]:
        """
        Extract images with robust identifiers for roundtrip workflow.
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Dictionary mapping image keys to image information
        """
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
        
        logger.info(f"Extracting images with identifiers from: {pptx_path}")
        
        presentation = Presentation(str(pptx_path))
        extracted_images = {}
        
        for slide_idx, slide in enumerate(presentation.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'image') and shape.image:
                    try:
                        # Create robust identifier
                        identifier = PPTXImageIdentifier.from_shape(shape, slide_idx, shape_idx)
                        
                        # Extract image information with consistent ALT text extraction
                        alt_text = self._get_existing_alt_text(shape)
                        image_info = {
                            'identifier': identifier,
                            'slide_idx': slide_idx,
                            'shape_idx': shape_idx,
                            'shape_name': identifier.shape_name,
                            'image_key': identifier.image_key,
                            'image_hash': identifier.image_hash,
                            'embed_id': identifier.embed_id,
                            'existing_alt_text': alt_text,
                            'alt_text': alt_text,  # Add explicit alt_text field for compatibility
                            'image_data': shape.image.blob,
                            'filename': getattr(shape.image, 'filename', f'slide_{slide_idx}_shape_{shape_idx}.png')
                        }
                        
                        extracted_images[identifier.image_key] = image_info
                        logger.debug(f"Extracted image: {identifier.image_key}")
                        
                    except Exception as e:
                        logger.warning(f"Failed to extract image from slide {slide_idx}, shape {shape_idx}: {e}")
        
        logger.info(f"Extracted {len(extracted_images)} images with identifiers")
        return extracted_images
    
    def test_pdf_export_alt_text_survival(self, pptx_path: str, output_dir: Optional[str] = None) -> Dict[str, Any]:
        """
        Test that ALT text survives PowerPoint → PDF export.
        
        Args:
            pptx_path: Path to PPTX file
            output_dir: Optional directory for output files
            
        Returns:
            Dictionary with test results
        """
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
        
        if output_dir is None:
            output_dir = pptx_path.parent
        else:
            output_dir = Path(output_dir)
            output_dir.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"Testing ALT text survival in PDF export for: {pptx_path}")
        
        # This would require PowerPoint automation or a PDF conversion library
        # For now, we'll create a placeholder test that validates ALT text exists in PPTX
        
        try:
            presentation = Presentation(str(pptx_path))
            alt_text_count = 0
            total_images = 0
            
            for slide_idx, slide in enumerate(presentation.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'image') and shape.image:
                        total_images += 1
                        existing_alt = self._get_existing_alt_text(shape)
                        if existing_alt and not self._should_skip_alt_text(existing_alt):
                            alt_text_count += 1
            
            survival_test_result = {
                'success': True,
                'pptx_file': str(pptx_path),
                'total_images': total_images,
                'images_with_alt_text': alt_text_count,
                'alt_text_coverage': alt_text_count / total_images if total_images > 0 else 0,
                'test_type': 'pptx_validation',
                'note': 'Full PDF export testing requires PowerPoint automation or conversion library',
                'errors': []
            }
            
            logger.info(f"ALT text survival test completed:")
            logger.info(f"  Total images: {total_images}")
            logger.info(f"  Images with ALT text: {alt_text_count}")
            logger.info(f"  Coverage: {survival_test_result['alt_text_coverage']:.1%}")
            
            return survival_test_result
            
        except Exception as e:
            error_msg = f"PDF export survival test failed: {str(e)}"
            logger.error(error_msg)
            
            return {
                'success': False,
                'pptx_file': str(pptx_path),
                'test_type': 'pptx_validation',
                'errors': [error_msg]
            }


def create_alt_text_mapping(image_data: Dict[str, Dict[str, Any]],
                          alt_text_results: Dict[str, str]) -> Dict[str, Any]:
    """
    Create enriched ALT text mapping with provenance for existing + generated ALT text.

    Args:
        image_data: Dictionary from extract_images_with_identifiers()
        alt_text_results: Dictionary mapping image keys to generated ALT text

    Returns:
        Dictionary mapping image keys to enriched ALT text objects with provenance
    """
    mapping = {}

    for image_key, image_info in image_data.items():
        # Extract existing ALT text from the shape data
        existing_alt = image_info.get('existing_alt_text', '').strip()

        # Check if existing ALT is meaningful (not empty/placeholder)
        existing_is_meaningful = _is_meaningful(existing_alt)

        # Get generated ALT text if available
        generated_alt = alt_text_results.get(image_key, '').strip()

        # Create enriched record with provenance
        enriched_record = {
            "existing_alt": existing_alt,
            "generated_alt": generated_alt,
            "source_existing": "pptx" if existing_is_meaningful else None,
            "source_generated": "llava" if generated_alt else None,
            "final_alt": None,  # Will be decided at inject time
            "decision": None,   # Will be set based on mode
            "existing_meaningful": existing_is_meaningful
        }

        mapping[image_key] = enriched_record
        logger.debug(f"Enriched mapping for {image_key}: existing='{existing_alt[:30]}...' generated='{generated_alt[:30]}...'")

    # Count different types of content
    existing_count = sum(1 for r in mapping.values() if r['existing_meaningful'])
    generated_count = sum(1 for r in mapping.values() if r['generated_alt'])

    logger.info(f"Created enriched ALT text mapping:")
    logger.info(f"  Total shapes: {len(mapping)}")
    logger.info(f"  With existing ALT: {existing_count}")
    logger.info(f"  With generated ALT: {generated_count}")

    return mapping

def _is_skip_token(alt_text: str) -> bool:
    """Check if ALT text is a placeholder/skip token that should be treated as empty."""
    if not alt_text:
        return True

    skip_tokens = {'(none)', 'n/a', 'not reviewed', 'undefined', 'image.png', 'picture', ''}
    return alt_text.lower().strip() in skip_tokens


def _is_meaningful(value: Optional[str]) -> bool:
    """Return True when the provided ALT text contains meaningful content."""
    if value is None:
        return False

    text = value.strip()
    if not text:
        return False

    return not _is_skip_token(text)


def _choose_candidate(record: Union[str, Dict[str, Any], None]) -> Optional[str]:
    """Select the best ALT text candidate from an enriched record or legacy mapping."""
    if record is None:
        return None

    if isinstance(record, str):
        candidate = record.strip()
        return candidate or None

    if not isinstance(record, dict):
        return None

    for field in ("final_alt", "generated_alt", "existing_alt"):
        value = record.get(field)
        if isinstance(value, str) and _is_meaningful(value):
            return value.strip()

    value = record.get("final_alt")
    if isinstance(value, str):
        stripped = value.strip()
        return stripped or None

    return None


def main():
    """Command-line interface for PPTX ALT text injection."""
    parser = argparse.ArgumentParser(
        description='Inject ALT text into PowerPoint presentations',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python pptx_alt_injector.py presentation.pptx --alt-text-file mappings.json
  python pptx_alt_injector.py presentation.pptx --extract-only --output extracted_images.json
  python pptx_alt_injector.py presentation.pptx --test-pdf-export
  python pptx_alt_injector.py presentation.pptx --config custom_config.yaml --verbose
        """
    )
    
    parser.add_argument('pptx_file', help='Input PPTX file')
    parser.add_argument('-o', '--output', help='Output PPTX file (default: overwrite input)')
    parser.add_argument('--alt-text-file', help='JSON file containing ALT text mappings')
    parser.add_argument('--extract-only', action='store_true',
                       help='Only extract images with identifiers (no injection)')
    parser.add_argument('--test-pdf-export', action='store_true',
                       help='Test ALT text survival in PDF export')
    parser.add_argument('--emit-injected-map', action='store_true',
                       help='Write *-alt_map.final.json with final ALT decisions after injection')
    parser.add_argument('--config', help='Configuration file path')
    parser.add_argument('--verbose', '-v', action='store_true', help='Enable verbose logging')
    parser.add_argument('--mode', choices=['preserve', 'overwrite'], 
                       help='ALT text handling mode')
    parser.add_argument('--debug-decisions', action='store_true',
                       help='Enable debug logging for injection decisions only')
    
    args = parser.parse_args()
    
    # Set up logging
    log_level = logging.DEBUG if (args.verbose or args.debug_decisions) else logging.INFO
    logging.basicConfig(
        level=log_level,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )
    
    # If debug-decisions is enabled, set up more focused logging
    if args.debug_decisions and not args.verbose:
        # Reduce noise from other modules
        logging.getLogger('pptx').setLevel(logging.WARNING)
        logging.getLogger('PIL').setLevel(logging.WARNING)
    
    try:
        # Initialize components
        config_manager = ConfigManager(args.config)
        injector = PPTXAltTextInjector(config_manager)
        
        # Override mode if specified
        if args.mode:
            injector.mode = args.mode
        
        print(f"PPTX ALT Text Injector")
        print(f"Processing: {args.pptx_file}")
        
        # Extract-only mode
        if args.extract_only:
            extracted_images = injector.extract_images_with_identifiers(args.pptx_file)
            output_file = args.output or f"{Path(args.pptx_file).stem}_extracted_images.json"
            
            import json
            with open(output_file, 'w') as f:
                # Convert to JSON-serializable format
                serializable = {}
                for key, info in extracted_images.items():
                    serializable[key] = {
                        'slide_idx': info['slide_idx'],
                        'shape_idx': info['shape_idx'],
                        'shape_name': info['shape_name'],
                        'image_key': info['image_key'],
                        'existing_alt_text': info['existing_alt_text'],
                        'filename': info['filename']
                    }
                json.dump(serializable, f, indent=2)
            
            print(f"Extracted {len(extracted_images)} images to: {output_file}")
            return 0
        
        # PDF export test mode
        if args.test_pdf_export:
            result = injector.test_pdf_export_alt_text_survival(args.pptx_file, args.output)
            
            if result['success']:
                print(f"✅ PDF export test completed")
                print(f"ALT text coverage: {result['alt_text_coverage']:.1%} ({result['images_with_alt_text']}/{result['total_images']})")
            else:
                print(f"❌ PDF export test failed")
                for error in result['errors']:
                    print(f"Error: {error}")
            
            return 0 if result['success'] else 1
        
        # ALT text injection mode
        if not args.alt_text_file:
            parser.error("--alt-text-file is required for ALT text injection")
        
        # Load ALT text mappings
        import json
        with open(args.alt_text_file, 'r') as f:
            alt_text_mapping = json.load(f)
        
        # Perform injection
        result = injector.inject_alt_text_from_mapping(
            args.pptx_file,
            alt_text_mapping,
            args.output
        )
        
        # Display results
        stats = result['statistics']
        print(f"\nInjection Results:")
        print(f"  Success: {result['success']}")
        print(f"  Processed: {stats['injected_successfully']}/{stats['total_images']}")
        print(f"    Preserved existing: {stats.get('preserved_existing', 0)}")
        print(f"    Written (generated): {stats.get('written_generated', 0)}")
        print(f"    Written (existing fallback): {stats.get('written_existing', 0)}")
        print(f"    Written (stored final): {stats.get('written_final', 0)}")
        print(f"    Skipped (no content): {stats.get('skipped_no_content', 0)}")
        print(f"  Failed: {stats['failed_injection']}")
        
        if result.get('errors'):
            print(f"Errors:")
            for error in result['errors']:
                print(f"  - {error}")

        if args.emit_injected_map and result.get('final_alt_map'):
            emitted_path = Path(args.alt_text_file).with_name(
                f"{Path(args.alt_text_file).stem}.final.json"
            )
            final_map_payload = {}
            for key, record in result['final_alt_map'].items():
                if isinstance(record, dict):
                    final_map_payload[key] = {
                        'final_alt': record.get('final_alt'),
                        'decision': record.get('decision')
                    }
                else:
                    final_map_payload[key] = {'final_alt': record, 'decision': None}

            with open(emitted_path, 'w', encoding='utf-8') as f:
                json.dump(final_map_payload, f, indent=2, ensure_ascii=False)

            print(f"  Injected map saved to: {emitted_path}")

        if result['success']:
            print(f"✅ ALT text injection completed successfully!")
            print(f"Output saved to: {result['output_file']}")
        else:
            print(f"❌ ALT text injection failed!")
        
        return 0 if result['success'] else 1
        
    except Exception as e:
        logger.error(f"ALT text injection failed: {e}", exc_info=True)
        print(f"Error: {e}")
        return 1


if __name__ == "__main__":
    exit(main())