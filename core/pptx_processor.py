"""
PPTX Accessibility Processor - Adapted from PDF processor to work with PowerPoint files.
Integrates with existing ConfigManager, FlexibleAltGenerator, medical prompts, and decorative detection.
"""

# --- safe XPath helper for python-pptx (BaseOxmlElement) and raw lxml ---
try:
    from pptx.oxml.ns import nsmap as PPTX_NSMAP  # type: ignore
except Exception:  # pragma: no cover
    PPTX_NSMAP = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
        'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
    }

def _safe_xpath(element, xpath_expr, namespaces=None):
    """Execute an XPath on python-pptx BaseOxmlElement or plain lxml element.
    Tries python-pptx override (no kwargs) first, then raw lxml with nsmap.
    Accepts an optional explicit namespace map for raw lxml cases.
    """
    el = getattr(element, "_element", element)
    try:
        return el.xpath(xpath_expr)  # python-pptx injects namespaces
    except Exception:
        ns = namespaces or getattr(el, "nsmap", None) or PPTX_NSMAP
        return el.xpath(xpath_expr, namespaces=ns)
# --- end safe XPath helper ---


import logging
import os
import sys
import tempfile
import time
import base64
import re
import io
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple
from collections import defaultdict
from hashlib import md5

# Third-party imports for PPTX processing
try:
    from pptx import Presentation
    from pptx.shapes.picture import Picture
    from pptx.shapes.base import BaseShape
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.shapes.connector import Connector
    from pptx.shapes.freeform import FreeformBuilder
    from pptx.shapes.autoshape import Shape as AutoShape
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    from pptx.oxml.ns import _nsmap
except ImportError as e:
    raise ImportError(
        "python-pptx is required for PPTX processing. Install with: pip install python-pptx"
    ) from e

# Setup paths for shared modules
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root / "shared"))

# Import shared modules
from config_manager import ConfigManager
from unified_alt_generator import FlexibleAltGenerator
from decorative_filter import (
    is_force_decorative_by_filename, 
    is_decorative_image,
    get_image_hash,
    validate_decorative_config
)

def describe_shape_with_details(shape) -> str:
    """Return a grammatically correct phrase like:
       'a text box. (921x195px)' or 'an unknown. (921x195px)'.
       Caller will prepend 'This is a PowerPoint shape. It is '."""
    import re
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    def emu_to_px(emu_value):
        """Convert EMU (English Metric Units) to pixels. 1 px ≈ 9525 EMUs"""
        return round(emu_value / 9525)
    
    # Get shape type description and sanitize enum noise like 'picture (13)'
    shape_type_desc = "unknown"
    try:
        if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if hasattr(shape, 'auto_shape_type'):
                shape_type_desc = str(shape.auto_shape_type).split('.')[-1]
        elif hasattr(shape, 'shape_type'):
            shape_type_desc = str(shape.shape_type).split('.')[-1]
    except Exception:
        pass
    
    shape_type_desc = shape_type_desc.replace('_', ' ').lower()
    shape_type_desc = re.sub(r'\s*\(\d+\)', '', shape_type_desc).strip()  # drop '(13)' etc.
    if not shape_type_desc:
        shape_type_desc = "unknown"
    
    # Get dimensions
    try:
        width_px = emu_to_px(shape.width) if hasattr(shape, 'width') else 0
        height_px = emu_to_px(shape.height) if hasattr(shape, 'height') else 0
        size_info = f" ({width_px}x{height_px}px)" if width_px > 0 and height_px > 0 else ""
    except Exception:
        size_info = ""
    
    # Choose article; treat words that start with vowel sounds — and 'unknown' — as 'an'
    needs_an = shape_type_desc[:1] in "aeiou" or shape_type_desc.startswith("unknown")
    article = "an" if needs_an else "a"
    # Return phrase to be appended after 'It is '
    return f"{article} {shape_type_desc}.{size_info}"

# Import PIL for shape-to-image rendering
try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    logging.warning("PIL not available - shape-to-image rendering will be disabled")

logger = logging.getLogger(__name__)


class PPTXImageInfo:
    """Container for PPTX image information."""
    
    def __init__(self, shape: BaseShape, slide_idx: int, shape_idx: int, 
                 image_data: bytes, filename: str, slide_text: str = "", is_rendered: bool = False):
        self.shape = shape
        self.slide_idx = slide_idx
        self.shape_idx = shape_idx
        self.image_data = image_data
        self.filename = filename
        self.slide_text = slide_text
        self.is_rendered = is_rendered  # Flag to track if this was rendered from a shape
        self.image_hash = get_image_hash(image_data)
        
        # Extract shape properties
        self.width = shape.width.emu if shape.width else 0
        self.height = shape.height.emu if shape.height else 0
        self.left = shape.left.emu if shape.left else 0
        self.top = shape.top.emu if shape.top else 0
        
        # Convert EMU to pixels (1 EMU = 1/914400 inch, assume 96 DPI)
        self.width_px = int(self.width / 914400 * 96) if self.width else 0
        self.height_px = int(self.height / 914400 * 96) if self.height else 0
        self.left_px = int(self.left / 914400 * 96) if self.left else 0
        self.top_px = int(self.top / 914400 * 96) if self.top else 0
        
        # Unique identifier consistent with PPTXAltTextInjector
        self.image_key = self._create_consistent_image_key(slide_idx, shape_idx, shape)
    
    def _create_consistent_image_key(self, slide_idx: int, shape_idx, shape: BaseShape) -> str:
        """Create stable image key using shape ID instead of enumeration index."""
        # Use shape.shape_id for stable identification (consistent across runs)
        shape_id = getattr(shape, 'shape_id', None)
        if shape_id is not None:
            # Use stable shape ID format: slide_X_shapeid_Y_hash_Z
            hash_value = self.image_hash[:8] if self.image_hash else f"{slide_idx}{shape_id}img"[:8]
            return f"slide_{slide_idx}_shapeid_{shape_id}_hash_{hash_value}"
        else:
            # Fallback to index-based key for shapes without IDs
            hash_value = self.image_hash[:8] if self.image_hash else f"{slide_idx}{shape_idx}img"[:8]
            return f"slide_{slide_idx}_shape_{shape_idx}_hash_{hash_value}"


class PPTXVisualElement:
    """Container for PPTX visual element information (images, shapes, charts)."""
    
    def __init__(self, shape: BaseShape, slide_idx: int, shape_idx: int, 
                 slide_text: str = "", element_type: str = "unknown"):
        self.shape = shape
        self.slide_idx = slide_idx
        self.shape_idx = shape_idx
        self.slide_text = slide_text
        self.element_type = element_type  # 'image', 'shape', 'chart', etc.
        
        # For images, store the image data
        self.image_data = None
        self.filename = None
        self.image_hash = None
        
        # Try to extract image data if this is a picture
        if hasattr(shape, 'image') and shape.image:
            try:
                self.image_data = shape.image.blob
                self.filename = f"image_{slide_idx}_{shape_idx}.{shape.image.ext}"
                self.image_hash = get_image_hash(self.image_data)
                self.element_type = "image"
            except:
                pass
        
        # Extract shape properties
        self.width = getattr(shape, 'width', None)
        self.height = getattr(shape, 'height', None) 
        self.left = getattr(shape, 'left', None)
        self.top = getattr(shape, 'top', None)
        
        # Convert EMU to pixels safely
        try:
            self.width_px = int(self.width.emu / 914400 * 96) if self.width else 0
            self.height_px = int(self.height.emu / 914400 * 96) if self.height else 0
            self.left_px = int(self.left.emu / 914400 * 96) if self.left else 0
            self.top_px = int(self.top.emu / 914400 * 96) if self.top else 0
        except:
            self.width_px = self.height_px = self.left_px = self.top_px = 0
        
        # Extract text content (needed for hash generation)
        self.has_text = False
        self.text_content = ""
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                self.text_content = shape.text_frame.text
                self.has_text = bool(self.text_content.strip())
        except:
            pass
            
        # Shape type information (needed for hash generation)
        self.shape_type = getattr(shape, 'shape_type', None)
        self.shape_name = getattr(shape, 'name', '') or ''
        
        # Generate element hash for duplicate detection
        if self.image_hash:
            self.element_hash = self.image_hash[:8]  # Truncate to 8 chars to match injector expectations
        else:
            # For non-images, create hash based on properties
            try:
                hash_content = f"{self.shape_type}_{self.width_px}_{self.height_px}_{self.text_content}"
                import hashlib
                self.element_hash = hashlib.md5(hash_content.encode()).hexdigest()[:8]  # 8 chars to match injector
            except:
                # Fallback to simple string-based hash
                self.element_hash = f"{slide_idx}{shape_idx}{element_type}"[:8]  # Keep under 8 chars
        
        # Generate stable element key using shape ID when available
        shape_id = getattr(shape, 'shape_id', None)
        if shape_id is not None:
            # Use stable shape ID format: slide_X_shapeid_Y_hash_Z
            self.element_key = f"slide_{slide_idx}_shapeid_{shape_id}_hash_{self.element_hash}"
        else:
            # Fallback to index-based key for shapes without IDs
            self.element_key = f"slide_{slide_idx}_shape_{shape_idx}_hash_{self.element_hash}"

class PPTXShapeInfo:
    """Container for PPTX shape information for decorative detection."""
    
    def __init__(self, shape: BaseShape, slide_idx: int, shape_idx: int, slide_text: str = ""):
        self.shape = shape
        self.slide_idx = slide_idx
        self.shape_idx = shape_idx
        self.slide_text = slide_text
        
        # Extract shape properties
        self.shape_name = getattr(shape, 'name', 'unnamed')
        self.shape_type = getattr(shape, 'shape_type', None)
        self.shape_type_name = self._get_shape_type_name()
        
        # Dimensions
        self.width = shape.width.emu if hasattr(shape, 'width') and shape.width else 0
        self.height = shape.height.emu if hasattr(shape, 'height') and shape.height else 0
        self.left = shape.left.emu if hasattr(shape, 'left') and shape.left else 0
        self.top = shape.top.emu if hasattr(shape, 'top') and shape.top else 0
        
        # Convert EMU to pixels
        self.width_px = int(self.width / 914400 * 96) if self.width else 0
        self.height_px = int(self.height / 914400 * 96) if self.height else 0
        self.left_px = int(self.left / 914400 * 96) if self.left else 0
        self.top_px = int(self.top / 914400 * 96) if self.top else 0
        
        # Check for text content
        self.has_text = hasattr(shape, 'text') and bool(shape.text and shape.text.strip())
        self.text_content = shape.text.strip() if self.has_text else ""
        
        # Create unique identifier
        self.shape_key = f"slide_{slide_idx}_shape_{shape_idx}_{self.shape_name}"
    
    def _get_shape_type_name(self) -> str:
        """Get human-readable shape type name with enhanced detection."""
        if self.shape_type is None:
            return "unknown"
        
        try:
            # Find the name of the shape type enum
            for attr_name in dir(MSO_SHAPE_TYPE):
                if not attr_name.startswith('_') and not callable(getattr(MSO_SHAPE_TYPE, attr_name, None)):
                    try:
                        attr_value = getattr(MSO_SHAPE_TYPE, attr_name)
                        if attr_value == self.shape_type:
                            # Add additional subtype information for AUTO_SHAPE
                            if attr_name == "AUTO_SHAPE" and hasattr(self.shape, 'auto_shape_type'):
                                try:
                                    auto_shape_type = self.shape.auto_shape_type
                                    return f"{attr_name}({auto_shape_type})"
                                except:
                                    pass
                            return attr_name
                    except (AttributeError, TypeError):
                        continue
            return f"MSO_SHAPE_TYPE({self.shape_type})"
        except Exception as e:
            return f"error_getting_type({self.shape_type})"


class PPTXAccessibilityProcessor:
    """
    PPTX accessibility processor that integrates with the existing PDF ALT text system.
    Reuses ConfigManager, FlexibleAltGenerator, medical prompts, and decorative detection.
    """
    
    def __init__(self, config_manager: Optional[ConfigManager] = None, debug: bool = False):
        """
        Initialize the PPTX accessibility processor.
        
        Args:
            config_manager: Optional ConfigManager instance. If None, creates a new one.
        """
        self.config_manager = config_manager or ConfigManager()
        self.debug = debug
        
        # Validate decorative configuration
        if not validate_decorative_config(self.config_manager.config):
            logger.warning("Decorative configuration validation failed")
        
        # Initialize ALT text generator
        try:
            self.alt_generator = FlexibleAltGenerator(self.config_manager)
            logger.info("Initialized PPTX accessibility processor with ALT text generator")
        except Exception as e:
            logger.error(f"Failed to initialize ALT text generator: {e}")
            raise
        
        # Get processing configuration
        self.processing_config = self.config_manager.config.get('pptx_processing', {})
        
        # Enhanced decorative detection settings with educational content bias
        # Reduced threshold as we now have better educational content detection
        self.decorative_size_threshold = self.processing_config.get('decorative_size_threshold', 30)
        self.skip_decorative = self.processing_config.get('skip_decorative_images', True)
        
        # Context extraction settings
        self.include_slide_notes = self.processing_config.get('include_slide_notes', True)
        self.include_slide_text = self.processing_config.get('include_slide_text', True)
        self.max_context_length = self.processing_config.get('max_context_length', 200)
        
        # Semantic analysis settings
        self.enable_semantic_icon_labels = self.processing_config.get('enable_semantic_icon_labels', False)
        
        logger.debug(f"Decorative size threshold: {self.decorative_size_threshold}px")
        logger.debug(f"Skip decorative images: {self.skip_decorative}")
        logger.debug(f"Include slide notes: {self.include_slide_notes}")
        logger.debug(f"Include slide text: {self.include_slide_text}")
    
    def process_pptx(self, pptx_path: str, output_path: Optional[str] = None, 
                    failed_generation_callback=None, debug: bool = False) -> Dict[str, Any]:
        """
        Process a PPTX file to add ALT text to all visual elements.
        
        Args:
            pptx_path: Path to the input PPTX file
            output_path: Optional path for output file. If None, overwrites original.
            failed_generation_callback: Callback function for failed generations
            
        Returns:
            Dictionary with processing statistics
        """
        start_time = time.time()
        pptx_path = Path(pptx_path)
        debug = debug or self.debug
        
        # Initialize result structure
        result = {
            'success': False,
            'input_file': str(pptx_path),
            'output_file': '',
            'total_slides': 0,
            'total_visual_elements': 0,
            'processed_visual_elements': 0,
            'failed_visual_elements': 0,
            'generation_time': 0.0,
            'injection_time': 0.0,
            'total_time': 0.0,
            'errors': []
        }
        
        # Validate input file
        if not pptx_path.exists():
            error_msg = f"PPTX file not found: {pptx_path}"
            logger.error(error_msg)
            result['errors'].append(error_msg)
            return result
        
        # Determine output path
        if output_path is None:
            output_path = pptx_path  # Overwrite original
        else:
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
        
        result['output_file'] = str(output_path)
        
        logger.info(f"Processing PPTX: {pptx_path.name}")
        logger.info(f"Output will be saved to: {output_path}")
        
        try:
            # Step 1: Extract all visual elements from PPTX
            logger.info("Step 1: Extracting all visual elements from PPTX...")
            extraction_start = time.time()
            
            presentation, visual_elements = self._extract_all_visual_elements(str(pptx_path))
            
            extraction_time = time.time() - extraction_start
            logger.info(f"Visual element extraction completed in {extraction_time:.2f}s")
            
            result['total_slides'] = len(presentation.slides)
            result['total_visual_elements'] = len(visual_elements)
            
            if not visual_elements:
                logger.warning(f"No visual elements found in PPTX: {pptx_path.name}")
                result['success'] = True  # Not an error, just no visual elements to process
                result['total_time'] = time.time() - start_time
                return result
            
            logger.info(f"Found {len(visual_elements)} visual elements across {result['total_slides']} slides")
            
            # Step 2: Generate ALT text for all visual elements
            logger.info("Step 2: Generating ALT text for all visual elements...")
            generation_start = time.time()
            
            alt_text_mapping = {}
            element_tracker = defaultdict(list)  # Track duplicate elements
            
            for visual_element in visual_elements:
                # Track element occurrences for duplicate detection
                element_key = getattr(visual_element, 'element_hash', str(visual_element.element_key))
                element_tracker[element_key].append(visual_element)
            
            for visual_element in visual_elements:
                generation_failure_reason = None
                
                try:
                    # Generate ALT text for ALL visual elements - let LLaVa decide if decorative
                    if debug:
                        logger.info(f"🔍 DEBUG: Processing {visual_element.element_type}: {visual_element.element_key}")
                        logger.info(f"🔍 DEBUG: Size: {visual_element.width_px}x{visual_element.height_px}px")
                        if visual_element.filename:
                            logger.info(f"🔍 DEBUG: Filename: {visual_element.filename}")
                        logger.info(f"🔍 DEBUG: Slide text: {visual_element.slide_text[:100]}...")
                    
                    alt_text, failure_reason = self._generate_alt_text_for_visual_element(visual_element, debug)
                    
                    if alt_text and alt_text.strip() and alt_text.strip() != "":
                        # Check for LLaVA errors and handle them
                        if self._is_llava_error(alt_text.strip()):
                            # LLaVA returned error - try fallback
                            fallback_description = self._handle_llava_error_with_fallback(visual_element, debug)
                            if fallback_description:
                                alt_text = fallback_description
                                if debug:
                                    logger.info(f"🔄 DEBUG: LLaVA error handled with fallback for {visual_element.element_key}")
                        
                        # Successfully generated valid ALT text - normalize to remove duplications
                        normalized_alt_text = self._normalize_alt(alt_text.strip())
                        alt_text_mapping[visual_element.element_key] = {
                            'alt_text': normalized_alt_text,
                            'shape': visual_element.shape,
                            'slide_idx': visual_element.slide_idx,
                            'shape_idx': visual_element.shape_idx
                        }
                        result['processed_visual_elements'] += 1
                        if debug:
                            logger.info(f"✅ DEBUG: Generated ALT text for {visual_element.element_key}: {alt_text[:50]}...")
                        else:
                            logger.info(f"Generated ALT text for {visual_element.element_key}: {alt_text[:50]}...")
                    else:
                        # Generation failed - try creating fallback descriptive ALT text
                        generation_failure_reason = failure_reason or "Empty or invalid ALT text returned"
                        
                        # Instead of generic "PowerPoint shape element", use descriptive text
                        if visual_element.element_type in ['shape', 'text_placeholder', 'text_box', 'line', 'connector']:
                            fallback_description = self._create_enhanced_fallback_description(visual_element)
                            
                            # Add bypass annotation for session data visibility
                            bypass_reason = self._check_element_bypass(visual_element)
                            if bypass_reason:
                                # Mark as bypassed for session data
                                # HOTPATCH FIX 2: Use compose_alt at join point instead of simple concatenation
                                bypass_annotation = f"[BYPASS: {bypass_reason}]"
                                composed_description = self._compose_alt([bypass_annotation, fallback_description])
                                normalized_description = composed_description
                                alt_text_mapping[visual_element.element_key] = {
                                    'alt_text': normalized_description,
                                    'shape': visual_element.shape,
                                    'slide_idx': visual_element.slide_idx,
                                    'shape_idx': visual_element.shape_idx,
                                    'bypass_reason': bypass_reason,
                                    'bypassed': True,
                                    'fallback_used': True
                                }
                                if debug:
                                    logger.info(f"🚧 DEBUG: Used bypassed fallback for {visual_element.element_key}: {bypass_reason}")
                                else:
                                    logger.info(f"Used bypassed fallback for {visual_element.element_key}: {bypass_reason}")
                            else:
                                # Normal fallback without bypass - HOTPATCH FIX 2: Use compose_alt for consistency
                                normalized_description = self._compose_alt([fallback_description])
                                alt_text_mapping[visual_element.element_key] = {
                                    'alt_text': normalized_description,
                                    'shape': visual_element.shape,
                                    'slide_idx': visual_element.slide_idx,
                                    'shape_idx': visual_element.shape_idx,
                                    'fallback_used': True
                                }
                                if debug:
                                    logger.info(f"✅ DEBUG: Used fallback description for {visual_element.element_key}: {fallback_description}")
                                else:
                                    logger.info(f"Used fallback description for {visual_element.element_key}: {fallback_description}")
                            
                            result['processed_visual_elements'] += 1
                        else:
                            # For other element types, still count as failed
                            result['failed_visual_elements'] += 1
                            
                            if debug:
                                logger.warning(f"❌ DEBUG: Generation failed for {visual_element.element_key}: {generation_failure_reason}")
                                
                            # Log failed generation for manual review
                            if failed_generation_callback:
                                failed_generation_callback(
                                    visual_element.element_key,
                                    {
                                        'slide_idx': visual_element.slide_idx,
                                        'shape_idx': visual_element.shape_idx,
                                        'element_type': visual_element.element_type,
                                        'filename': visual_element.filename,
                                        'width_px': visual_element.width_px,
                                        'height_px': visual_element.height_px,
                                        'slide_text': visual_element.slide_text
                                    },
                                    f"ALT text generation failed: {generation_failure_reason}"
                                )
                
                except Exception as e:
                    generation_failure_reason = f"Exception during generation: {str(e)}"
                    
                    # Try fallback description even for exceptions on shapes
                    if visual_element.element_type in ['shape', 'text_placeholder', 'text_box', 'line', 'connector']:
                        try:
                            fallback_description = self._create_enhanced_fallback_description(visual_element)
                            
                            # Add bypass annotation for session data visibility
                            bypass_reason = self._check_element_bypass(visual_element)
                            if bypass_reason:
                                # Mark as bypassed for session data
                                # HOTPATCH FIX 2: Use compose_alt at join point instead of simple concatenation
                                bypass_annotation = f"[BYPASS: {bypass_reason}]"
                                composed_description = self._compose_alt([bypass_annotation, fallback_description])
                                normalized_description = composed_description
                                alt_text_mapping[visual_element.element_key] = {
                                    'alt_text': normalized_description,
                                    'shape': visual_element.shape,
                                    'slide_idx': visual_element.slide_idx,
                                    'shape_idx': visual_element.shape_idx,
                                    'bypass_reason': bypass_reason,
                                    'bypassed': True,
                                    'fallback_used': True,
                                    'exception_fallback': True
                                }
                                if debug:
                                    logger.info(f"🚧 DEBUG: Used bypassed fallback after exception for {visual_element.element_key}: {bypass_reason}")
                                else:
                                    logger.info(f"Used bypassed fallback after exception for {visual_element.element_key}: {bypass_reason}")
                            else:
                                # Normal fallback without bypass - HOTPATCH FIX 2: Use compose_alt for consistency
                                normalized_description = self._compose_alt([fallback_description])
                                alt_text_mapping[visual_element.element_key] = {
                                    'alt_text': normalized_description,
                                    'shape': visual_element.shape,
                                    'slide_idx': visual_element.slide_idx,
                                    'shape_idx': visual_element.shape_idx,
                                    'fallback_used': True,
                                    'exception_fallback': True
                                }
                                if debug:
                                    logger.info(f"✅ DEBUG: Used fallback description after exception for {visual_element.element_key}: {fallback_description}")
                                else:
                                    logger.info(f"Used fallback description after exception for {visual_element.element_key}: {fallback_description}")
                            
                            result['processed_visual_elements'] += 1
                        except Exception as fallback_e:
                            # Fallback failed too
                            result['failed_visual_elements'] += 1
                            if debug:
                                logger.error(f"💥 DEBUG: Exception processing {visual_element.element_key} and fallback failed: {e}, fallback: {fallback_e}", exc_info=True)
                            else:
                                error_msg = f"Error processing {visual_element.element_key}: {str(e)}"
                                logger.error(error_msg)
                                result['errors'].append(error_msg)
                    else:
                        result['failed_visual_elements'] += 1
                        if debug:
                            logger.error(f"💥 DEBUG: Exception processing {visual_element.element_key}: {e}", exc_info=True)
                        else:
                            error_msg = f"Error processing {visual_element.element_key}: {str(e)}"
                            logger.error(error_msg)
                            result['errors'].append(error_msg)
                        
                    # Log failed generation for manual review (only if no fallback was used)
                    if visual_element.element_key not in alt_text_mapping and failed_generation_callback:
                        failed_generation_callback(
                            visual_element.element_key,
                            {
                                'slide_idx': visual_element.slide_idx,
                                'shape_idx': visual_element.shape_idx,
                                'element_type': visual_element.element_type,
                                'filename': visual_element.filename,
                                'width_px': visual_element.width_px,
                                'height_px': visual_element.height_px,
                                'slide_text': visual_element.slide_text
                            },
                            f"Exception during generation: {str(e)}"
                        )
            
            result['generation_time'] = time.time() - generation_start
            logger.info(f"ALT text generation completed in {result['generation_time']:.2f}s")
            
            # Step 3: Validate ALT text coverage before injection
            logger.info("Step 3: Validating visual element ALT text coverage...")
            validation_result = self._validate_visual_element_coverage(visual_elements, alt_text_mapping, debug)
            
            if not validation_result['complete_coverage']:
                missing_count = validation_result['missing_count']
                error_msg = f"Incomplete ALT text coverage: {missing_count} visual elements missing ALT text"
                logger.error(error_msg)
                result['errors'].append(error_msg)
                
                if debug:
                    logger.error("❌ DEBUG: Visual elements missing ALT text:")
                    for missing_key in validation_result['missing_elements']:
                        logger.error(f"   - {missing_key}")
            
            # Step 4: Inject ALT text into PPTX
            if alt_text_mapping:
                logger.info("Step 4: Adding ALT text to PPTX...")
                injection_start = time.time()
                
                if debug:
                    logger.info(f"🔍 DEBUG: Injecting {len(alt_text_mapping)} ALT text mappings")
                    for key, info in list(alt_text_mapping.items())[:3]:  # Show first 3
                        logger.info(f"🔍 DEBUG: {key} -> '{info['alt_text'][:30]}...'")
                
                injection_success, final_alt_map = self._inject_alt_text_to_pptx(
                    presentation, alt_text_mapping, str(output_path), debug
                )
                
                result['injection_time'] = time.time() - injection_start
                result['final_alt_map'] = final_alt_map  # Store canonical mapping for approval docs
                logger.info(f"ALT text injection completed in {result['injection_time']:.2f}s")
                
                if injection_success:
                    result['success'] = True
                    logger.info("✅ PPTX processing completed successfully!")
                    
                    # Report visual element coverage
                    if result['total_visual_elements'] > 0:
                        coverage = (result['processed_visual_elements'] / result['total_visual_elements']) * 100
                        logger.info(f"📊 Visual element ALT text coverage: {result['processed_visual_elements']}/{result['total_visual_elements']} ({coverage:.1f}%)")
                        
                        if coverage == 100.0:
                            logger.info("🎯 100% visual element ALT text coverage achieved!")
                    else:
                        logger.info("📊 No visual elements found to process")
                else:
                    error_msg = "ALT text injection failed"
                    logger.error(error_msg)
                    result['errors'].append(error_msg)
            else:
                logger.warning("No ALT text to inject - this should not happen with proper fallback")
                result['success'] = False
                result['errors'].append("No ALT text mappings generated - fallback system failed")
            
        except Exception as e:
            error_msg = f"Unexpected error during PPTX processing: {str(e)}"
            logger.error(error_msg, exc_info=True)
            result['errors'].append(error_msg)
        
        # Calculate total processing time
        result['total_time'] = time.time() - start_time
        
        # Log final statistics
        self._log_processing_summary(result)
        
        return result
    
    def _render_shape_to_image(self, shape: BaseShape, slide_idx: int, shape_idx: int, slide_context: str = "") -> Optional[PPTXImageInfo]:
        """
        Render a PowerPoint shape to an image for LLaVa processing.
        
        Args:
            shape: Shape to render
            slide_idx: Slide index
            shape_idx: Shape index
            slide_context: Slide text context
            
        Returns:
            PPTXImageInfo with rendered image data or None if rendering failed
        """
        if not PIL_AVAILABLE:
            logger.warning("PIL not available - cannot render shapes to images")
            return None
        
        try:
            # Get shape dimensions and position
            width_emu = getattr(shape, 'width', 0)
            height_emu = getattr(shape, 'height', 0)
            
            if width_emu <= 0 or height_emu <= 0:
                return None
            
            # Convert EMU to pixels (914400 EMU = 1 inch, 96 DPI)
            width_px = max(int(width_emu / 914400 * 96), 50)
            height_px = max(int(height_emu / 914400 * 96), 50)
            
            # Create image canvas with white background
            img = Image.new('RGB', (width_px, height_px), 'white')
            draw = ImageDraw.Draw(img)
            
            # Render shape based on type
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            shape_type = getattr(shape, 'shape_type', None)
            
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                self._render_auto_shape(draw, shape, width_px, height_px)
            elif shape_type == MSO_SHAPE_TYPE.LINE:
                self._render_line_shape(draw, shape, width_px, height_px)
            elif shape_type == MSO_SHAPE_TYPE.FREEFORM:
                self._render_freeform_shape(draw, shape, width_px, height_px)
            elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                self._render_text_box(draw, shape, width_px, height_px)
            else:
                # Generic shape rendering
                self._render_generic_shape(draw, shape, width_px, height_px)
            
            # Convert PIL image to bytes
            img_bytes = io.BytesIO()
            img.save(img_bytes, format='PNG')
            image_data = img_bytes.getvalue()
            
            # Create PPTXImageInfo
            filename = f"rendered_shape_{slide_idx}_{shape_idx}.png"
            image_info = PPTXImageInfo(
                shape=shape,
                slide_idx=slide_idx,
                shape_idx=shape_idx,
                image_data=image_data,
                filename=filename,
                slide_text=slide_context[:self.max_context_length] if slide_context else "",
                is_rendered=True  # Flag to indicate this was rendered
            )
            
            logger.debug(f"Rendered shape to image: {filename} ({width_px}x{height_px}px)")
            return image_info
            
        except Exception as e:
            logger.warning(f"Failed to render shape {shape_idx} on slide {slide_idx}: {e}")
            return None
    
    def _render_auto_shape(self, draw: ImageDraw.Draw, shape: AutoShape, width: int, height: int):
        """Render AutoShape (circles, rectangles, etc.)"""
        try:
            # Get shape fill color
            fill_color = self._get_shape_fill_color(shape)
            line_color = self._get_shape_line_color(shape)
            line_width = self._get_shape_line_width(shape)
            
            # Get shape type for specific rendering
            auto_shape_type = getattr(shape, 'auto_shape_type', None)
            
            if auto_shape_type:
                from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                
                if auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                    # Draw circle/ellipse
                    draw.ellipse([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
                elif auto_shape_type in [MSO_AUTO_SHAPE_TYPE.RECTANGLE, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE]:
                    # Draw rectangle
                    draw.rectangle([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
                elif auto_shape_type == MSO_AUTO_SHAPE_TYPE.HEXAGON:
                    # Draw hexagon
                    points = self._get_hexagon_points(width, height)
                    draw.polygon(points, fill=fill_color, outline=line_color)
                else:
                    # Generic rectangle for unknown shapes
                    draw.rectangle([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
            else:
                # Default rectangle
                draw.rectangle([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
                
        except Exception as e:
            logger.warning(f"Failed to render auto shape: {e}")
            # Fallback: draw a simple rectangle
            draw.rectangle([0, 0, width-1, height-1], fill='lightgray', outline='black')
    
    def _render_line_shape(self, draw: ImageDraw.Draw, shape: BaseShape, width: int, height: int):
        """Render line shape"""
        try:
            line_color = self._get_shape_line_color(shape)
            line_width = self._get_shape_line_width(shape)
            
            # Draw line from top-left to bottom-right (simplified)
            draw.line([0, 0, width-1, height-1], fill=line_color, width=line_width)
            
        except Exception as e:
            logger.warning(f"Failed to render line shape: {e}")
            draw.line([0, 0, width-1, height-1], fill='black', width=2)
    
    def _render_freeform_shape(self, draw: ImageDraw.Draw, shape: BaseShape, width: int, height: int):
        """Render freeform shape"""
        try:
            fill_color = self._get_shape_fill_color(shape)
            line_color = self._get_shape_line_color(shape)
            
            # For freeform, draw a polygon (simplified)
            # This is a basic implementation - complex freeforms would need path parsing
            points = [(0, height//2), (width//4, 0), (3*width//4, 0), (width-1, height//2), (width-1, height-1), (0, height-1)]
            draw.polygon(points, fill=fill_color, outline=line_color)
            
        except Exception as e:
            logger.warning(f"Failed to render freeform shape: {e}")
            draw.rectangle([0, 0, width-1, height-1], fill='lightgray', outline='black')
    
    def _render_text_box(self, draw: ImageDraw.Draw, shape: BaseShape, width: int, height: int):
        """Render text box with background"""
        try:
            fill_color = self._get_shape_fill_color(shape)
            line_color = self._get_shape_line_color(shape)
            line_width = self._get_shape_line_width(shape)
            
            # Draw background
            draw.rectangle([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
            
            # Add text if available
            if hasattr(shape, 'text') and shape.text:
                try:
                    font = ImageFont.load_default()
                    text_color = 'black'
                    
                    # Simple text positioning (centered)
                    text_width = len(shape.text) * 8  # Rough estimate
                    text_height = 12
                    x = max(0, (width - text_width) // 2)
                    y = max(0, (height - text_height) // 2)
                    
                    draw.text((x, y), shape.text[:50], fill=text_color, font=font)  # Limit text length
                except:
                    pass  # Text rendering is optional
                    
        except Exception as e:
            logger.warning(f"Failed to render text box: {e}")
            draw.rectangle([0, 0, width-1, height-1], fill='white', outline='black')
    
    def _render_generic_shape(self, draw: ImageDraw.Draw, shape: BaseShape, width: int, height: int):
        """Generic shape rendering fallback"""
        try:
            fill_color = self._get_shape_fill_color(shape)
            line_color = self._get_shape_line_color(shape)
            line_width = self._get_shape_line_width(shape)
            
            # Draw as rectangle with visual indication
            draw.rectangle([0, 0, width-1, height-1], fill=fill_color, outline=line_color, width=line_width)
            
            # Add X marks to indicate it's a generic shape
            draw.line([0, 0, width-1, height-1], fill=line_color, width=1)
            draw.line([0, height-1, width-1, 0], fill=line_color, width=1)
            
        except Exception as e:
            logger.warning(f"Failed to render generic shape: {e}")
            draw.rectangle([0, 0, width-1, height-1], fill='lightgray', outline='black')
    
    def _get_shape_fill_color(self, shape: BaseShape) -> str:
        """Get shape fill color"""
        try:
            if hasattr(shape, 'fill') and shape.fill:
                fill = shape.fill
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    # Try to get RGB color
                    try:
                        rgb = fill.fore_color.rgb
                        return f"#{rgb:06x}"
                    except:
                        pass
                
                # Check for solid fill
                from pptx.enum.dml import MSO_FILL_TYPE
                if hasattr(fill, 'type') and fill.type == MSO_FILL_TYPE.SOLID:
                    return 'lightblue'  # Default solid color
                    
            # Shape-specific defaults
            shape_type = getattr(shape, 'shape_type', None)
            auto_shape_type = getattr(shape, 'auto_shape_type', None)
            
            if auto_shape_type:
                from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                if auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                    return 'lightblue'  # Blue circle default
                elif auto_shape_type == MSO_AUTO_SHAPE_TYPE.HEXAGON:
                    return 'purple'  # Purple hexagon default
                    
            return 'lightgray'  # Generic default
            
        except Exception:
            return 'lightgray'
    
    def _get_shape_line_color(self, shape: BaseShape) -> str:
        """Get shape line color"""
        try:
            if hasattr(shape, 'line') and shape.line:
                line = shape.line
                if hasattr(line, 'color') and line.color:
                    try:
                        rgb = line.color.rgb
                        return f"#{rgb:06x}"
                    except:
                        pass
            return 'black'  # Default line color
        except Exception:
            return 'black'
    
    def _get_shape_line_width(self, shape: BaseShape) -> int:
        """Get shape line width"""
        try:
            if hasattr(shape, 'line') and shape.line:
                line = shape.line
                if hasattr(line, 'width') and line.width:
                    # Convert EMU to pixels (rough approximation)
                    width_emu = line.width
                    width_px = max(1, int(width_emu / 914400 * 96 / 72))  # Convert to reasonable pixel width
                    return min(width_px, 10)  # Cap at 10px
            return 2  # Default line width
        except Exception:
            return 2
    
    def _get_hexagon_points(self, width: int, height: int) -> List[Tuple[int, int]]:
        """Generate hexagon points"""
        cx, cy = width // 2, height // 2
        radius_x, radius_y = width // 2 - 2, height // 2 - 2
        
        import math
        points = []
        for i in range(6):
            angle = i * math.pi / 3
            x = cx + radius_x * math.cos(angle)
            y = cy + radius_y * math.sin(angle)
            points.append((int(x), int(y)))
        
        return points
    
    def _should_render_shape_to_image(self, shape: BaseShape) -> bool:
        """
        Determine if a shape should be rendered to an image for LLaVa processing.
        
        Args:
            shape: Shape to evaluate
            
        Returns:
            True if shape should be rendered to image
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            shape_type = getattr(shape, 'shape_type', None)
            
            # Don't render shapes that already contain images
            if hasattr(shape, 'image') and shape.image:
                return False
            
            # Don't render group shapes (they're processed recursively)
            if hasattr(shape, 'shapes'):
                return False
            
            # If the shape has a picture fill, treat it as visual so we render it to an image
            try:
                from pptx.enum.dml import MSO_FILL_TYPE
                if hasattr(shape, 'fill') and hasattr(shape.fill, 'type'):
                    if shape.fill.type == MSO_FILL_TYPE.PICTURE:
                        return True
            except Exception:
                pass
            
            # Render visual shape types that should be processed by LLaVa
            visual_shape_types = [
                MSO_SHAPE_TYPE.AUTO_SHAPE,      # Circles, rectangles, hexagons, etc.
                MSO_SHAPE_TYPE.LINE,            # Lines and connectors
                MSO_SHAPE_TYPE.FREEFORM,        # Custom drawn shapes
                MSO_SHAPE_TYPE.TEXT_BOX,        # Text boxes with visual styling
            ]
            
            if shape_type in visual_shape_types:
                # Additional checks for shapes with visual content
                
                # Check if shape has visual significance (fill, border, etc.)
                if self._has_visual_significance_for_rendering(shape):
                    return True
                
                # Check if text box has substantial text content
                if shape_type == MSO_SHAPE_TYPE.TEXT_BOX and hasattr(shape, 'text') and shape.text:
                    # Render text boxes with meaningful content
                    text_content = shape.text.strip()
                    if len(text_content) > 10:  # More than just a few characters
                        return True
            
            return False
            
        except Exception as e:
            logger.debug(f"Error determining if shape should be rendered: {e}")
            return False
    
    def _has_visual_significance_for_rendering(self, shape: BaseShape) -> bool:
        """
        Check if shape has visual significance that warrants rendering for LLaVa.
        More permissive than the decorative detection logic.
        
        Args:
            shape: Shape to evaluate
            
        Returns:
            True if shape has visual significance
        """
        try:
            # Check for fill colors/patterns
            if hasattr(shape, 'fill') and shape.fill:
                from pptx.enum.dml import MSO_FILL_TYPE
                if hasattr(shape.fill, 'type') and shape.fill.type == MSO_FILL_TYPE.SOLID:
                    return True
                if hasattr(shape.fill, 'type') and shape.fill.type != MSO_FILL_TYPE.NO_FILL:
                    return True
            
            # Check for borders/outlines
            if hasattr(shape, 'line') and shape.line:
                from pptx.enum.dml import MSO_LINE_STYLE
                if hasattr(shape.line, 'style') and hasattr(shape.line, 'width'):
                    if shape.line.width and shape.line.width > 0:
                        return True
            
            # Check size - larger shapes are more likely to be content
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                width_emu = getattr(shape, 'width', 0)
                height_emu = getattr(shape, 'height', 0)
                
                if width_emu and height_emu:
                    # Convert to pixels for evaluation
                    width_px = int(width_emu / 914400 * 96)
                    height_px = int(height_emu / 914400 * 96)
                    
                    # Consider shapes larger than 30x30 pixels as potentially significant
                    if width_px >= 30 and height_px >= 30:
                        return True
            
            return False
            
        except Exception:
            return False
    
    def _extract_images_from_pptx(self, pptx_path: str) -> Tuple[Presentation, List[PPTXImageInfo]]:
        """
        Extract all images from PPTX with their context, including grouped shapes, 
        chart elements, embedded objects, and images in text boxes.
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Tuple of (Presentation object, List of PPTXImageInfo objects)
        """
        presentation = Presentation(pptx_path)
        image_infos = []
        
        for slide_idx, slide in enumerate(presentation.slides):
            logger.debug(f"Processing slide {slide_idx + 1}")
            
            # Extract slide text for context
            slide_text = self._extract_slide_text(slide) if self.include_slide_text else ""
            
            # Extract slide notes for context
            slide_notes = self._extract_slide_notes(slide) if self.include_slide_notes else ""
            
            # Combine slide context
            slide_context = []
            if slide_text:
                slide_context.append(slide_text)
            if slide_notes:
                slide_context.append(f"Notes: {slide_notes}")
            slide_context_str = " ".join(slide_context)
            
            # Debug: Log all shapes found on this slide with comprehensive enumeration
            logger.debug(f"📄 Processing slide {slide_idx + 1} with {len(slide.shapes)} top-level shapes")
            self._enumerate_all_shapes(slide.shapes, indent="  ", slide_idx=slide_idx)
            
            # Process all shapes recursively to find images
            images_found_on_slide = self._extract_images_from_shapes(
                slide.shapes, slide_idx, slide_context_str
            )
            
            image_infos.extend(images_found_on_slide)
            logger.debug(f"Found {len(images_found_on_slide)} images on slide {slide_idx + 1}")
        
        # Also attempt to find images through presentation relationships
        # This can catch images that aren't accessible through the shape API
        logger.debug("Attempting to find additional images through presentation relationships...")
        relationship_images = self._extract_images_from_relationships(presentation)
        
        logger.info(f"Extracted {len(image_infos)} images via shapes, {len(relationship_images)} via relationships from {len(presentation.slides)} slides")
        logger.info(f"Total unique images found: {len(image_infos)}")
        return presentation, image_infos
    
    def _extract_images_from_shapes(self, shapes, slide_idx: int, slide_context: str, parent_group_idx: int = None) -> List[PPTXImageInfo]:
        """
        Recursively extract images from shapes, including grouped shapes, charts, and embedded objects.
        
        Args:
            shapes: Collection of shapes to process
            slide_idx: Slide index
            slide_context: Slide context text
            parent_group_idx: Index of parent group if this is a nested call
            
        Returns:
            List of PPTXImageInfo objects
        """
        image_infos = []
        
        for shape_idx, shape in enumerate(shapes):
            try:
                # Create unique shape identifier
                if parent_group_idx is not None:
                    shape_id = f"{parent_group_idx}_{shape_idx}"
                else:
                    shape_id = shape_idx
                
                # Debug: Log detailed shape information
                shape_name = getattr(shape, 'name', 'unnamed')
                shape_type = getattr(shape, 'shape_type', 'unknown')
                logger.debug(f"    Examining shape {shape_id}: {type(shape).__name__} (type={shape_type}, name='{shape_name}')")
                
                # Check for various types of images
                images_from_shape = []
                
                # 1. Direct picture shapes (original logic)
                if hasattr(shape, 'image') and shape.image:
                    logger.debug(f"      -> Found direct picture shape")
                    try:
                        image_data = shape.image.blob
                        filename = getattr(shape.image, 'filename', f'slide_{slide_idx}_shape_{shape_id}.png')
                        
                        image_info = PPTXImageInfo(
                            shape=shape,
                            slide_idx=slide_idx,
                            shape_idx=shape_id,
                            image_data=image_data,
                            filename=filename,
                            slide_text=slide_context[:self.max_context_length] if slide_context else ""
                        )
                        images_from_shape.append(image_info)
                        logger.debug(f"      -> Extracted direct image: {filename}")
                        
                    except Exception as e:
                        logger.warning(f"Failed to extract direct image from shape {shape_id}: {e}")
                
                # 2. Group shapes (recursively process shapes within groups)
                if hasattr(shape, 'shapes'):
                    logger.debug(f"      -> Found group shape with {len(shape.shapes)} child shapes")
                    group_images = self._extract_images_from_shapes(
                        shape.shapes, slide_idx, slide_context, shape_id
                    )
                    images_from_shape.extend(group_images)
                    logger.debug(f"      -> Extracted {len(group_images)} images from group")
                
                # 3. Chart shapes (may contain images)
                if hasattr(shape, 'chart'):
                    logger.debug(f"      -> Found chart shape")
                    chart_images = self._extract_images_from_chart(shape.chart, slide_idx, shape_id, slide_context)
                    images_from_shape.extend(chart_images)
                    logger.debug(f"      -> Extracted {len(chart_images)} images from chart")
                
                # 4. Text boxes with image fills
                if hasattr(shape, 'text_frame') and hasattr(shape, 'fill'):
                    logger.debug(f"      -> Examining text box for image fill")
                    fill_images = self._extract_images_from_fill(shape.fill, slide_idx, shape_id, slide_context, shape_name)
                    images_from_shape.extend(fill_images)
                    if fill_images:
                        logger.debug(f"      -> Extracted {len(fill_images)} images from text box fill")
                
                # 5. Shape fills (any shape can have an image fill)
                elif hasattr(shape, 'fill'):
                    logger.debug(f"      -> Examining shape fill")
                    fill_images = self._extract_images_from_fill(shape.fill, slide_idx, shape_id, slide_context, shape_name)
                    images_from_shape.extend(fill_images)
                    if fill_images:
                        logger.debug(f"      -> Extracted {len(fill_images)} images from shape fill")
                
                # 6. OLE objects and embedded content
                if hasattr(shape, '_element'):
                    logger.debug(f"      -> Examining XML element for embedded objects")
                    ole_images = self._extract_images_from_ole(shape._element, slide_idx, shape_id, slide_context, shape_name)
                    images_from_shape.extend(ole_images)
                    if ole_images:
                        logger.debug(f"      -> Extracted {len(ole_images)} images from OLE objects")
                
                # 7. Render visual shapes to images for LLaVa processing
                # If no images were found from this shape but it's a visual element, render it
                if not images_from_shape and self._should_render_shape_to_image(shape):
                    logger.debug(f"      -> Rendering shape to image for LLaVa processing")
                    rendered_image = self._render_shape_to_image(shape, slide_idx, shape_id, slide_context)
                    if rendered_image:
                        images_from_shape.append(rendered_image)
                        logger.debug(f"      -> Successfully rendered shape: {rendered_image.filename}")
                
                image_infos.extend(images_from_shape)
                
                if images_from_shape:
                    logger.debug(f"    Total images from shape {shape_id}: {len(images_from_shape)}")
                
            except Exception as e:
                logger.warning(f"Error processing shape {shape_idx} on slide {slide_idx}: {e}")
                continue
        
        return image_infos
    
    def _extract_images_from_chart(self, chart, slide_idx: int, shape_id: str, slide_context: str) -> List[PPTXImageInfo]:
        """
        Extract images from chart elements (chart backgrounds, data point images, etc.).
        
        Args:
            chart: Chart object
            slide_idx: Slide index
            shape_id: Shape identifier
            slide_context: Slide context text
            
        Returns:
            List of PPTXImageInfo objects
        """
        images = []
        
        try:
            # Check chart plot area fill
            if hasattr(chart, 'plot_area') and hasattr(chart.plot_area, 'fill'):
                fill_images = self._extract_images_from_fill(
                    chart.plot_area.fill, slide_idx, f"{shape_id}_plot", slide_context, "chart_plot_area"
                )
                images.extend(fill_images)
            
            # Check chart area fill
            if hasattr(chart, 'chart_area') and hasattr(chart.chart_area, 'fill'):
                fill_images = self._extract_images_from_fill(
                    chart.chart_area.fill, slide_idx, f"{shape_id}_area", slide_context, "chart_area"
                )
                images.extend(fill_images)
            
            # Check series fills (data points might have image fills)
            if hasattr(chart, 'series'):
                for series_idx, series in enumerate(chart.series):
                    if hasattr(series, 'fill'):
                        fill_images = self._extract_images_from_fill(
                            series.fill, slide_idx, f"{shape_id}_series_{series_idx}", 
                            slide_context, f"chart_series_{series_idx}"
                        )
                        images.extend(fill_images)
            
        except Exception as e:
            logger.debug(f"Error extracting images from chart: {e}")
        
        return images
    
    def _extract_images_from_fill(self, fill, slide_idx: int, shape_id: str, slide_context: str, source_name: str) -> List[PPTXImageInfo]:
        """
        Extract images from fill objects (picture fills, texture fills, etc.).
        
        Args:
            fill: Fill object
            slide_idx: Slide index
            shape_id: Shape identifier
            slide_context: Slide context text
            source_name: Name describing the source of this fill
            
        Returns:
            List of PPTXImageInfo objects
        """
        images = []
        
        try:
            # Check if this is a picture fill
            if hasattr(fill, 'type') and fill.type is not None:
                # Import fill type constants
                from pptx.dml.fill import MSO_FILL_TYPE
                
                if fill.type == MSO_FILL_TYPE.PICTURE:
                    logger.debug(f"        -> Found picture fill in {source_name}")
                    
                    # Try to extract image data from picture fill
                    if hasattr(fill, '_fill') and hasattr(fill._fill, 'blipFill'):
                        blip_fill = fill._fill.blipFill
                        if hasattr(blip_fill, 'blip') and hasattr(blip_fill.blip, 'rId'):
                            # This indicates there's an image, but we need to get the actual data
                            # For now, create a placeholder entry - the actual extraction would need
                            # access to the presentation's part relationships
                            logger.debug(f"        -> Picture fill found but data extraction needs relationship resolution")
                            
                            # Create a minimal image info to track this image
                            # Note: We can't create a full PPTXImageInfo without actual image data
                            # This is a limitation that would need deeper PPTX internals access
                            pass
                
        except Exception as e:
            logger.debug(f"Error extracting images from fill: {e}")
        
        return images
    
    def _extract_images_from_ole(self, element, slide_idx: int, shape_id: str, slide_context: str, source_name: str) -> List[PPTXImageInfo]:
        """
        Extract images from OLE objects and other embedded content.
        
        Args:
            element: XML element
            slide_idx: Slide index
            shape_id: Shape identifier
            slide_context: Slide context text
            source_name: Name describing the source
            
        Returns:
            List of PPTXImageInfo objects
        """
        images = []
        
        try:
            # Look for embedded objects in the XML
            # This would require parsing the XML structure for oleObj elements
            # and embedded image relationships
            
            # Check for oleObj elements
            ole_objects = _safe_xpath(element, './/p:oleObj', namespaces=element.nsmap) if element.nsmap else []
            if ole_objects:
                logger.debug(f"        -> Found {len(ole_objects)} OLE objects in {source_name}")
                # OLE object image extraction would require access to embedded parts
            
            # Check for embedded pictures in alternative locations
            embedded_pics = _safe_xpath(element, './/pic:pic', namespaces=element.nsmap) if element.nsmap else []
            if embedded_pics:
                logger.debug(f"        -> Found {len(embedded_pics)} embedded pictures in {source_name}")
            
        except Exception as e:
            logger.debug(f"Error extracting images from OLE: {e}")
        
        return images
    
    def _enumerate_all_shapes(self, shapes, indent: str = "", slide_idx: int = None):
        """
        Recursively enumerate and log comprehensive information about all shapes with enhanced detection.
        
        Args:
            shapes: Collection of shapes to enumerate
            indent: Indentation string for nested shapes
            slide_idx: Current slide index for context
        """
        total_shapes = len(shapes) if hasattr(shapes, '__len__') else 0
        if slide_idx is not None and not indent:
            logger.debug(f"🔍 Slide {slide_idx + 1}: Found {total_shapes} top-level shapes")
        
        for i, shape in enumerate(shapes):
            try:
                # Get comprehensive shape type information
                shape_type_name = "unknown"
                shape_type_value = getattr(shape, 'shape_type', None)
                shape_class = type(shape).__name__
                
                # Enhanced shape type detection with better subtype analysis
                try:
                    if shape_type_value is not None:
                        # Find the name of the shape type enum
                        for attr_name in dir(MSO_SHAPE_TYPE):
                            if not attr_name.startswith('_') and getattr(MSO_SHAPE_TYPE, attr_name) == shape_type_value:
                                shape_type_name = attr_name
                                break
                        else:
                            shape_type_name = f"MSO_SHAPE_TYPE({shape_type_value})"
                        
                        # Enhanced subtype detection for AUTO_SHAPE
                        if shape_type_name == "AUTO_SHAPE":
                            auto_shape_details = self._get_auto_shape_details(shape, debug, indent)
                            if auto_shape_details:
                                shape_type_name = f"AUTO_SHAPE({auto_shape_details})"
                        
                        # Enhanced CONNECTOR detection
                        elif shape_type_name == "CONNECTOR" or (hasattr(shape, 'connector_type') and shape.connector_type is not None):
                            connector_details = self._get_connector_details(shape, debug, indent)
                            if connector_details:
                                shape_type_name = f"CONNECTOR({connector_details})"
                        
                        # Enhanced LINE detection
                        elif shape_type_name == "LINE":
                            line_details = self._get_line_details(shape, debug, indent)
                            if line_details:
                                shape_type_name = f"LINE({line_details})"
                        
                        # Detect shapes that might be misclassified
                        elif shape_type_name == "unknown" or shape_type_value == 0:
                            # Try alternative detection methods
                            alternative_type = self._detect_shape_by_properties(shape, debug, indent)
                            if alternative_type:
                                shape_type_name = f"DETECTED_AS_{alternative_type}"
                                
                except Exception as type_error:
                    logger.debug(f"{indent}    Error getting type for shape {i}: {type_error}")
                    shape_type_name = str(shape_type_value) if shape_type_value is not None else "unknown"
                    
                    # Even on error, try basic property detection
                    try:
                        alternative_type = self._detect_shape_by_properties(shape, debug, indent)
                        if alternative_type:
                            shape_type_name = f"ERROR_RECOVERY_{alternative_type}"
                    except:
                        pass
                
                # Collect comprehensive shape properties
                properties = []
                
                # Basic identification properties
                shape_name = getattr(shape, 'name', 'unnamed')
                if shape_name and shape_name != 'unnamed':
                    properties.append(f"name='{shape_name}'")
                
                shape_id = getattr(shape, 'shape_id', None)
                if shape_id is not None:
                    properties.append(f"id={shape_id}")
                
                # Dimensions and positioning
                if hasattr(shape, 'width') and hasattr(shape, 'height'):
                    try:
                        width_px = int(shape.width.emu / 914400 * 96) if shape.width else 0
                        height_px = int(shape.height.emu / 914400 * 96) if shape.height else 0
                        properties.append(f"size={width_px}x{height_px}px")
                        
                        if hasattr(shape, 'left') and hasattr(shape, 'top'):
                            left_px = int(shape.left.emu / 914400 * 96) if shape.left else 0
                            top_px = int(shape.top.emu / 914400 * 96) if shape.top else 0
                            properties.append(f"pos=({left_px},{top_px})")
                    except Exception as dim_error:
                        logger.debug(f"{indent}    Error getting dimensions: {dim_error}")
                
                # Content detection
                has_image = hasattr(shape, 'image') and shape.image
                if has_image:
                    try:
                        image_format = getattr(shape.image, 'ext', 'unknown')
                        properties.append(f"IMAGE({image_format})")
                    except:
                        properties.append("IMAGE")
                
                has_chart = hasattr(shape, 'chart')
                if has_chart:
                    try:
                        chart_type = getattr(shape.chart, 'chart_type', 'unknown')
                        properties.append(f"CHART({chart_type})")
                    except:
                        properties.append("CHART")
                
                # Group detection with count
                has_shapes = hasattr(shape, 'shapes')
                if has_shapes:
                    child_count = len(shape.shapes) if shape.shapes else 0
                    properties.append(f"GROUP({child_count})")
                
                # Text content analysis
                has_text = hasattr(shape, 'text') and shape.text
                if has_text:
                    text_content = shape.text.strip()
                    if len(text_content) > 0:
                        text_preview = text_content[:30].replace('\n', ' ')
                        properties.append(f"TEXT='{text_preview}{'...' if len(text_content) > 30 else ''}'")
                        properties.append(f"chars={len(text_content)}")
                    else:
                        properties.append("TEXT=<empty>")
                elif hasattr(shape, 'text_frame'):
                    properties.append("TEXT_FRAME")
                
                # Fill and formatting properties
                if hasattr(shape, 'fill'):
                    try:
                        fill_type = getattr(shape.fill, 'type', None)
                        if fill_type is not None:
                            from pptx.dml.fill import MSO_FILL_TYPE
                            if fill_type == MSO_FILL_TYPE.PICTURE:
                                properties.append("PICTURE_FILL")
                            elif fill_type == MSO_FILL_TYPE.TEXTURED:
                                properties.append("TEXTURE_FILL")
                            elif fill_type == MSO_FILL_TYPE.SOLID:
                                properties.append("SOLID_FILL")
                            elif fill_type == MSO_FILL_TYPE.GRADIENT:
                                properties.append("GRADIENT_FILL")
                            elif fill_type == MSO_FILL_TYPE.BACKGROUND:
                                properties.append("BACKGROUND_FILL")
                            else:
                                properties.append(f"FILL({fill_type})")
                    except Exception as fill_error:
                        properties.append("FILL")
                        logger.debug(f"{indent}    Error getting fill type: {fill_error}")
                
                # Line/border properties for connectors and lines
                if hasattr(shape, 'line'):
                    try:
                        if hasattr(shape.line, 'color') or hasattr(shape.line, 'width'):
                            properties.append("HAS_LINE")
                    except:
                        pass
                
                # Connector-specific properties
                if shape_type_name.startswith("CONNECTOR") or shape_type_name.startswith("LINE"):
                    try:
                        if hasattr(shape, 'begin_connect'):
                            begin_connected = shape.begin_connect is not None
                            properties.append(f"begin_connected={begin_connected}")
                        if hasattr(shape, 'end_connect'):
                            end_connected = shape.end_connect is not None
                            properties.append(f"end_connected={end_connected}")
                    except Exception as conn_error:
                        logger.debug(f"{indent}    Error getting connector info: {conn_error}")
                
                # Rotation and 3D properties
                try:
                    if hasattr(shape, 'rotation'):
                        rotation = shape.rotation
                        if rotation != 0:
                            properties.append(f"rotation={rotation}°")
                except:
                    pass
                
                # Build and log the shape information
                props_str = f" [{', '.join(properties)}]" if properties else ""
                logger.debug(f"{indent}Shape {i:2d}: {shape_class} -> {shape_type_name}{props_str}")
                
                # Special handling for unknown or fallback types
                if shape_type_name == "unknown" or "error_getting_type" in shape_type_name:
                    logger.debug(f"{indent}    ⚠️  FALLBACK DETECTION: {shape_class} with properties: {properties}")
                    
                    # Try alternative detection methods
                    xml_tag = None
                    if hasattr(shape, '_element'):
                        try:
                            xml_tag = shape._element.tag
                            logger.debug(f"{indent}    XML tag: {xml_tag}")
                        except:
                            pass
                
                # Recursively enumerate grouped shapes with enhanced tracking
                if has_shapes and shape.shapes:
                    logger.debug(f"{indent}  📁 Group contents ({len(shape.shapes)} shapes):")
                    self._enumerate_all_shapes(shape.shapes, indent + "    ", slide_idx)
                
            except Exception as e:
                logger.debug(f"{indent}Shape {i}: ❌ Error during enumeration - {e}")
                # Still log what we can
                shape_class = type(shape).__name__ if shape else "None"
                logger.debug(f"{indent}    Class: {shape_class}")
                try:
                    if hasattr(shape, '_element') and shape._element is not None:
                        xml_tag = getattr(shape._element, 'tag', 'no_tag')
                        logger.debug(f"{indent}    XML: {xml_tag}")
                except:
                    pass
    
    def _extract_images_from_relationships(self, presentation: Presentation) -> List[PPTXImageInfo]:
        """
        Extract images by directly parsing presentation relationships and parts.
        This can find images that aren't accessible through the normal shape API.
        
        Args:
            presentation: Presentation object
            
        Returns:
            List of PPTXImageInfo objects
        """
        images = []
        
        try:
            # Access the presentation part and its relationships
            prs_part = presentation.part
            
            # Get all image parts from relationships
            image_parts = []
            for relationship in prs_part.rels.values():
                if hasattr(relationship, 'target_part'):
                    target = relationship.target_part
                    # Check if this is an image part
                    if hasattr(target, 'content_type') and target.content_type.startswith('image/'):
                        image_parts.append((relationship.rId, target))
            
            logger.debug(f"Found {len(image_parts)} image parts in presentation relationships")
            
            # Also check slide-level relationships
            for slide_idx, slide in enumerate(presentation.slides):
                try:
                    slide_part = slide.part
                    for relationship in slide_part.rels.values():
                        if hasattr(relationship, 'target_part'):
                            target = relationship.target_part
                            if hasattr(target, 'content_type') and target.content_type.startswith('image/'):
                                image_parts.append((f"slide_{slide_idx}_{relationship.rId}", target))
                except Exception as e:
                    logger.debug(f"Error checking slide {slide_idx} relationships: {e}")
            
            logger.debug(f"Total image parts found: {len(image_parts)}")
            
            # Create image info objects for relationship-based images
            # Note: These won't have shape context, but they represent actual images in the file
            for rel_id, image_part in image_parts:
                try:
                    image_data = image_part.blob
                    filename = getattr(image_part, 'partname', f'relationship_{rel_id}.png')
                    
                    # Create a minimal image info - we don't have shape context for these
                    logger.debug(f"Found relationship image: {filename} ({len(image_data)} bytes)")
                    
                except Exception as e:
                    logger.debug(f"Error extracting image from relationship {rel_id}: {e}")
        
        except Exception as e:
            logger.debug(f"Error extracting images from relationships: {e}")
        
        return images
    
    def _extract_all_visual_elements(self, pptx_path: str) -> Tuple[Presentation, List[PPTXVisualElement]]:
        """
        Extract ALL visual elements from PPTX (images, shapes, charts, etc.) for ALT text generation.
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Tuple of (Presentation object, List of PPTXVisualElement objects)
        """
        presentation = Presentation(pptx_path)
        visual_elements = []
        
        logger.info(f"Starting comprehensive visual element extraction from {len(presentation.slides)} slides...")
        
        for slide_idx, slide in enumerate(presentation.slides):
            # Extract slide text for context
            slide_text = self._extract_slide_text(slide) if self.include_slide_text else ""
            
            logger.debug(f"Processing slide {slide_idx + 1} with {len(slide.shapes)} shapes")
            
            # Process all shapes recursively
            slide_elements = self._extract_visual_elements_from_shapes(
                slide.shapes, slide_idx, slide_text, self.debug
            )
            
            visual_elements.extend(slide_elements)
            
            logger.debug(f"Found {len(slide_elements)} visual elements on slide {slide_idx + 1}")
        
        logger.info(f"Extracted {len(visual_elements)} total visual elements")
        
        # Step 1.5: Process group ALT text roll-up AFTER children are extracted
        # This creates parent ALT text for group shapes where PowerPoint Reading Order checks
        if visual_elements:
            logger.info("Processing group ALT text roll-up for PowerPoint Reading Order...")
            group_parent_elements = self._process_group_alt_rollup(presentation, visual_elements, self.debug)
            if group_parent_elements:
                logger.info(f"Added {len(group_parent_elements)} group parent elements for ALT roll-up")
                visual_elements.extend(group_parent_elements)
        
        return presentation, visual_elements
    
    def _process_group_alt_rollup(self, presentation: Presentation, visual_elements: List[PPTXVisualElement], debug: bool = False) -> List[PPTXVisualElement]:
        """
        Process group ALT text roll-up for PowerPoint Reading Order compatibility.
        
        PowerPoint Reading Order evaluates group parents (p:grpSp), not children.
        This method analyzes groups that contain visual elements and creates appropriate
        parent ALT text by rolling up from meaningful children.
        
        Args:
            presentation: PowerPoint presentation
            visual_elements: List of extracted visual elements (children)
            debug: Enable debug logging
            
        Returns:
            List of new PPTXVisualElement objects for group parents
        """
        group_parents = []
        groups_processed = 0
        groups_with_alt = 0
        children_marked_decorative = 0
        
        if debug:
            logger.debug("🔍 Starting group ALT roll-up analysis...")
        
        for slide_idx, slide in enumerate(presentation.slides):
            if debug:
                logger.debug(f"🔍 Processing slide {slide_idx + 1} for groups...")
            
            # Enhanced per-slide counters
            slide_found_groups = 0
            slide_found_autoshapes = 0
            slide_injected_ok = 0
            
            # Count elements on this slide for enhanced logging
            for shape in slide.shapes:
                if hasattr(shape, 'shapes') and shape.shapes:
                    slide_found_groups += 1
                elif hasattr(shape, 'shape_type'):
                    try:
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                            slide_found_autoshapes += 1
                    except:
                        pass
            
            slide_groups_found, slide_groups_with_alt, slide_children_marked = self._process_slide_groups_for_rollup(
                slide.shapes, slide_idx, visual_elements, debug
            )
            
            # Enhanced logging with per-slide counters
            if slide_groups_found > 0 or slide_found_autoshapes > 0:
                logger.info(f"📊 Slide {slide_idx + 1} enhanced processing summary:")
                logger.info(f"   🔧 found_groups: {slide_found_groups}")
                logger.info(f"   🔷 found_autoshapes: {slide_found_autoshapes}")
                logger.info(f"   ✅ groups_with_alt: {slide_groups_with_alt}")
                logger.info(f"   🎯 injected_ok: {slide_groups_with_alt}") # Groups that got ALT text
            
            groups_processed += slide_groups_found
            groups_with_alt += slide_groups_with_alt  
            children_marked_decorative += slide_children_marked
            
            # Enhanced per-slide logging
            if slide_groups_found > 0:
                logger.info(f"🔍 Slide {slide_idx + 1} group processing:")
                logger.info(f"   📊 Groups found: {slide_groups_found}")
                logger.info(f"   ✅ Groups with parent ALT: {slide_groups_with_alt}")
                logger.info(f"   👶 Children marked decorative: {slide_children_marked}")
                if debug:
                    success_rate = (slide_groups_with_alt/slide_groups_found*100) if slide_groups_found > 0 else 0
                    logger.debug(f"   🎯 Slide success rate: {success_rate:.1f}%")
        
        # Enhanced logging summary
        logger.info(f"🔄 Group ALT Roll-up Summary:")
        logger.info(f"   📊 Groups processed: {groups_processed}")
        logger.info(f"   ✅ Groups with parent ALT: {groups_with_alt}")
        logger.info(f"   🎯 Success rate: {(groups_with_alt/groups_processed*100):.1f}%" if groups_processed > 0 else "   🎯 Success rate: 0.0%")
        logger.info(f"   👶 Children marked decorative: {children_marked_decorative}")
        
        if debug and groups_processed > 0:
            logger.debug(f"🔍 Detailed group processing:")
            logger.debug(f"   - Groups without ALT roll-up: {groups_processed - groups_with_alt}")
            logger.debug(f"   - Average children marked decorative per group: {children_marked_decorative/groups_processed:.1f}")
        
        return group_parents
    
    def _process_slide_groups_for_rollup(self, shapes, slide_idx: int, visual_elements: List[PPTXVisualElement], 
                                       debug: bool = False, depth: int = 0) -> Tuple[int, int, int]:
        """
        Recursively process shapes on a slide to find groups and roll up ALT text.
        
        Args:
            shapes: Collection of shapes to process
            slide_idx: Slide index
            visual_elements: List of visual elements (for finding children)
            debug: Enable debug logging
            depth: Recursion depth for nested groups
            
        Returns:
            Tuple of (groups_found, groups_with_alt, children_marked_decorative)
        """
        groups_found = 0
        groups_with_alt = 0  
        children_marked_decorative = 0
        indent = "  " * depth
        
        for shape_idx, shape in enumerate(shapes):
            try:
                # Check if this is a group shape
                if hasattr(shape, 'shapes') and shape.shapes:
                    groups_found += 1
                    shape_id = getattr(shape, 'shape_id', f"group_{slide_idx}_{shape_idx}")
                    
                    # Enhanced group processing logging
                    logger.info(f"🔧 Processing group_id={shape_id}, depth={depth}, children={len(shape.shapes)}")
                    if debug:
                        logger.debug(f"🔍 {indent}Group structure analysis:")
                        logger.debug(f"🔍 {indent}  Depth: {depth}")
                        logger.debug(f"🔍 {indent}  Shape ID: {shape_id}")
                        logger.debug(f"🔍 {indent}  Children count: {len(shape.shapes)}")
                        logger.debug(f"🔍 {indent}  Shape type: {type(shape).__name__}")
                        if hasattr(shape, 'name'):
                            logger.debug(f"🔍 {indent}  Shape name: '{shape.name}'")
                        if hasattr(shape, 'width') and hasattr(shape, 'height'):
                            try:
                                width_px = int(shape.width.emu / 914400 * 96) if shape.width else 0
                                height_px = int(shape.height.emu / 914400 * 96) if shape.height else 0
                                logger.debug(f"🔍 {indent}  Dimensions: {width_px}x{height_px}px")
                            except:
                                pass
                    
                    if debug:
                        logger.debug(f"🔍 {indent}Found group (ID: {shape_id}) with {len(shape.shapes)} children")
                    
                    # CRITICAL FIX: Process nested groups FIRST before analyzing current group
                    # This ensures deep nested groups get their ALT text before parent groups analyze them
                    nested_groups_found, nested_groups_with_alt, nested_children_marked = self._process_slide_groups_for_rollup(
                        shape.shapes, slide_idx, visual_elements, debug, depth + 1
                    )
                    groups_found += nested_groups_found
                    groups_with_alt += nested_groups_with_alt
                    children_marked_decorative += nested_children_marked
                    
                    # Now analyze group children (which may now include processed nested groups with ALT text)
                    child_analysis = self._analyze_group_children_for_rollup(
                        shape.shapes, slide_idx, visual_elements, debug, depth + 1
                    )
                    
                    # Decide if this group should get parent ALT text
                    should_create_parent_alt, parent_alt_text, children_to_mark_decorative = self._decide_group_alt_rollup(
                        shape, child_analysis, debug, indent
                    )
                    
                    if should_create_parent_alt:
                        groups_with_alt += 1
                        children_marked_decorative += len(children_to_mark_decorative)
                        
                        # Enhanced roll-up success logging
                        logger.info(f"✅ Group {shape_id} at depth {depth} received ALT text: '{parent_alt_text[:50]}...'")
                        if len(children_to_mark_decorative) > 0:
                            logger.info(f"   🎯 Marked {len(children_to_mark_decorative)} children as decorative")
                        
                        if debug:
                            logger.debug(f"🔍 {indent}✅ Rolling up ALT to group parent: '{parent_alt_text}'")
                            logger.debug(f"🔍 {indent}   Will mark {len(children_to_mark_decorative)} children as decorative")
                            logger.debug(f"🔍 {indent}   Child analysis summary:")
                            logger.debug(f"🔍 {indent}     Meaningful children: {len(child_analysis.get('meaningful_children', []))}")
                            logger.debug(f"🔍 {indent}     Text-only children: {len(child_analysis.get('text_only_children', []))}")
                            logger.debug(f"🔍 {indent}     Decorative children: {len(child_analysis.get('decorative_children', []))}")
                        
                        # ANTI-CONCATENATION GUARD: Check if group would be writing duplicate text
                        try:
                            # Try to get existing group ALT text
                            existing_group_alt = ""
                            if hasattr(shape, '_element'):
                                group_element = shape._element
                                cnvpr_elements = _safe_xpath(group_element, './/p:nvGrpSpPr/p:cNvPr', namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                                if cnvpr_elements:
                                    existing_group_alt = cnvpr_elements[0].get('descr', '') or ''
                            
                            # Only inject if this would actually change the content (after normalization)
                            if existing_group_alt.strip() == parent_alt_text.strip():
                                if debug:
                                    logger.debug(f"🔍 {indent}🔒 Skipping group injection - text already present")
                            else:
                                # Inject ALT text directly into the group parent XML
                                self._inject_group_parent_alt_text(shape, parent_alt_text, debug, indent)
                                
                        except Exception as guard_error:
                            if debug:
                                logger.debug(f"🔍 {indent}⚠️  Group guard failed, proceeding with injection: {guard_error}")
                            # Fallback to normal injection if guard fails
                            self._inject_group_parent_alt_text(shape, parent_alt_text, debug, indent)
                        
                        # Mark redundant children as decorative
                        self._mark_group_children_decorative(children_to_mark_decorative, debug, indent)
                    else:
                        # Enhanced logging for groups without roll-up
                        logger.info(f"❌ Group {shape_id} at depth {depth} did not receive ALT text")
                        if debug:
                            logger.debug(f"🔍 {indent}⏭️ Group doesn't need parent ALT text")
                            logger.debug(f"🔍 {indent}   Child analysis summary:")
                            logger.debug(f"🔍 {indent}     Meaningful children: {len(child_analysis.get('meaningful_children', []))}")
                            logger.debug(f"🔍 {indent}     Text-only children: {len(child_analysis.get('text_only_children', []))}")
                            logger.debug(f"🔍 {indent}     Decorative children: {len(child_analysis.get('decorative_children', []))}")
                            
            except Exception as e:
                if debug:
                    logger.debug(f"🔍 {indent}Error processing shape for group rollup: {e}")
                continue
        
        return groups_found, groups_with_alt, children_marked_decorative
    
    def _analyze_group_children_for_rollup(self, group_shapes, slide_idx: int, visual_elements: List[PPTXVisualElement], 
                                         debug: bool = False, depth: int = 0) -> Dict[str, Any]:
        """
        Analyze children of a group to determine ALT text roll-up strategy.
        
        Args:
            group_shapes: Shapes within the group
            slide_idx: Slide index  
            visual_elements: List of visual elements
            debug: Enable debug logging
            depth: Recursion depth
            
        Returns:
            Dictionary with analysis results
        """
        indent = "  " * depth
        analysis = {
            'total_children': len(group_shapes),
            'meaningful_children': [],
            'decorative_children': [],
            'text_only_children': [],
            'has_single_meaningful_child': False,
            'has_multiple_meaningful_children': False,
            'recommended_parent_alt': "",
            'children_with_alt': []
        }
        
        # Find visual elements that belong to this group
        for child_shape in group_shapes:
            child_id = getattr(child_shape, 'id', None)
            child_name = getattr(child_shape, 'name', 'unnamed')
            
            # Look for visual elements matching this child
            matching_elements = []
            for ve in visual_elements:
                if (hasattr(ve, 'shape') and ve.shape == child_shape) or \
                   (child_id and hasattr(ve, 'shape_idx') and ve.shape_idx == child_id):
                    matching_elements.append(ve)
            
            # Classify child based on what we know
            if matching_elements:
                # This child has generated visual elements
                for element in matching_elements:
                    if hasattr(element, 'element_type') and element.element_type in ['image', 'shape', 'chart']:
                        analysis['meaningful_children'].append({
                            'shape': child_shape,
                            'element': element,
                            'type': element.element_type,
                            'id': child_id,
                            'name': child_name
                        })
                        analysis['children_with_alt'].append(element)
            elif hasattr(child_shape, 'text_frame') and child_shape.text_frame and child_shape.text_frame.text.strip():
                # Text-only child
                analysis['text_only_children'].append({
                    'shape': child_shape,
                    'text': child_shape.text_frame.text.strip()[:100],
                    'id': child_id,
                    'name': child_name
                })
            else:
                # Likely decorative or structural
                analysis['decorative_children'].append({
                    'shape': child_shape,
                    'id': child_id,
                    'name': child_name
                })
        
        # Determine roll-up strategy
        meaningful_count = len(analysis['meaningful_children'])
        analysis['has_single_meaningful_child'] = meaningful_count == 1
        analysis['has_multiple_meaningful_children'] = meaningful_count > 1
        
        if debug:
            logger.debug(f"🔍 {indent}Group analysis: {meaningful_count} meaningful, {len(analysis['text_only_children'])} text-only, {len(analysis['decorative_children'])} decorative")
        
        return analysis
    
    def _decide_group_alt_rollup(self, group_shape, child_analysis: Dict[str, Any], debug: bool = False, indent: str = "") -> Tuple[bool, str, List]:
        """
        Decide whether a group should get parent ALT text and what it should be.
        
        Args:
            group_shape: The group shape
            child_analysis: Analysis results from _analyze_group_children_for_rollup
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            Tuple of (should_create_alt, parent_alt_text, children_to_mark_decorative)
        """
        meaningful_children = child_analysis['meaningful_children']
        meaningful_count = len(meaningful_children)
        
        # DESCRIBE EVERYTHING POLICY - ALL groups get ALT text
        # No meaningful children - create generic group description
        if meaningful_count == 0:
            text_children_count = len(child_analysis['text_only_children'])
            decorative_children_count = len(child_analysis['decorative_children'])
            total_children = child_analysis['total_children']
            
            # Try to determine if this is a semantic icon group
            semantic_type = self._detect_group_semantic_type(group_shape, child_analysis)
            
            if semantic_type:
                parent_alt = f"Group representing {semantic_type} icon"
                policy_applied = "semantic_parent_generic"
            elif text_children_count > 0:
                parent_alt = f"Group containing {text_children_count} text element{'s' if text_children_count != 1 else ''}"
                policy_applied = "text_only_group"
            elif decorative_children_count > 0:
                parent_alt = f"Group containing {total_children} decorative element{'s' if total_children != 1 else ''}"
                policy_applied = "decorative_group"
            else:
                parent_alt = f"Group containing {total_children} element{'s' if total_children != 1 else ''}"
                policy_applied = "generic_group"
            
            # Enhanced roll-up decision logging
            logger.info(f"📋 Applied roll-up policy: {policy_applied} -> '{parent_alt}'")
                
            if debug:
                logger.debug(f"🔍 {indent}Creating group description: '{parent_alt}' (semantic: {semantic_type or 'none'})")
            return True, parent_alt, []
        
        children_to_mark_decorative = []
        
        # Single meaningful child - ROBUST ROLL-UP POLICY
        if meaningful_count == 1:
            child = meaningful_children[0]
            element = child['element']
            
            # Check if group parent has semantic meaning (icon-like)
            parent_semantic_type = self._detect_group_semantic_type(group_shape, child_analysis)
            
            if parent_semantic_type:
                # Parent has semantic meaning - keep parent ALT, mark child decorative
                parent_alt = f"Group representing {parent_semantic_type} icon"
                children_to_mark_decorative = [child['shape']]
                policy_applied = "parent_kept_child_decorative"
                
                # Enhanced roll-up decision logging
                logger.info(f"📋 Applied roll-up policy: {policy_applied} -> '{parent_alt}'")
                
                if debug:
                    logger.debug(f"🔍 {indent}Parent semantic strategy: '{parent_alt}' (child becomes decorative)")
                    
                return True, parent_alt, children_to_mark_decorative
            else:
                # Child carries meaning - create composite description
                parent_alt = f"Group containing {element.element_type}"
                # Don't mark child as decorative - both parent and child keep their ALT
                children_to_mark_decorative = []
                policy_applied = "child_promoted_both_keep"
                
                # Enhanced roll-up decision logging
                logger.info(f"📋 Applied roll-up policy: {policy_applied} -> '{parent_alt}'")
                
                if debug:
                    logger.debug(f"🔍 {indent}Child meaningful strategy: '{parent_alt}' (both keep ALT)")
                
                return True, parent_alt, children_to_mark_decorative
        
        # Multiple meaningful children - ROBUST ROLL-UP POLICY
        if meaningful_count > 1:
            # Check if parent has semantic meaning
            parent_semantic_type = self._detect_group_semantic_type(group_shape, child_analysis)
            
            if parent_semantic_type:
                # Parent is semantic (icon) - container description, children keep content
                parent_alt = f"Group representing {parent_semantic_type} icon with {meaningful_count} elements"
                # Don't mark children decorative - they provide content detail
                children_to_mark_decorative = []
                
                if debug:
                    logger.debug(f"🔍 {indent}Semantic parent + meaningful children: container/content split")
            else:
                # Parent is structural - create composite description
                child_types = []
                for child in meaningful_children:
                    child_type = child['element'].element_type
                    if child_type == 'image':
                        child_types.append('image')
                    elif child_type == 'shape':
                        child_types.append('shape')
                    elif child_type == 'chart':
                        child_types.append('chart')
                
                # Build description based on child types
                if len(set(child_types)) == 1:
                    # All same type
                    child_type = child_types[0]
                    parent_alt = f"Group of {meaningful_count} {child_type}s"
                else:
                    # Mixed types
                    parent_alt = f"Group containing {meaningful_count} visual elements"
                
                # Don't mark children as decorative when multiple meaningful children
                # Let them keep their individual ALT text
                children_to_mark_decorative = []
            
            # Enhanced roll-up decision logging for multiple children
            if parent_semantic_type:
                policy_applied = "composite_semantic_container"
            else:
                policy_applied = "composite_structural"
            logger.info(f"📋 Applied roll-up policy: {policy_applied} -> '{parent_alt}'")
            
            if debug:
                logger.debug(f"🔍 {indent}Multiple meaningful children strategy: '{parent_alt}' (children keep ALT)")
            
            return True, parent_alt, children_to_mark_decorative
        
        return False, "", []
    
    def _detect_group_semantic_type(self, group_shape, child_analysis: Dict[str, Any]) -> Optional[str]:
        """
        Detect if a group represents a semantic icon type (lightbulb, brain, etc.).
        Disabled by default to avoid false positives.
        
        Args:
            group_shape: The group shape
            child_analysis: Analysis results from group children
            
        Returns:
            Semantic type string if detected, None otherwise
        """
        # Check if semantic icon labels are enabled
        if not getattr(self, "enable_semantic_icon_labels", False):
            return None
            
        try:
            # Analyze shape patterns and geometry to detect common icons
            total_children = child_analysis['total_children']
            decorative_count = len(child_analysis['decorative_children'])
            
            # Get group dimensions for pattern analysis
            group_width = getattr(group_shape, 'width', None)
            group_height = getattr(group_shape, 'height', None)
            
            if group_width and group_height:
                width_px = int(group_width.emu / 914400 * 96) if hasattr(group_width, 'emu') else 0
                height_px = int(group_height.emu / 914400 * 96) if hasattr(group_height, 'emu') else 0
                aspect_ratio = width_px / height_px if height_px > 0 else 1
                
                # Icon detection based on patterns:
                
                # Lightbulb pattern: circular top + narrower bottom, multiple decorative elements
                if 0.8 <= aspect_ratio <= 1.4 and total_children >= 3 and decorative_count >= 2:
                    return "lightbulb"
                
                # Brain pattern: organic/complex shape with many small elements
                if 0.9 <= aspect_ratio <= 1.6 and total_children >= 5:
                    return "brain"
                    
                # Lungs pattern: two-part symmetric structure
                if 1.2 <= aspect_ratio <= 2.0 and total_children >= 2:
                    return "lungs"
                
                # Graduation cap pattern: rectangular base + triangular top
                if 0.7 <= aspect_ratio <= 1.5 and total_children >= 2:
                    return "graduation cap"
            
            # Fallback pattern analysis based on child count and complexity
            if total_children >= 5:
                return "complex icon"
            elif total_children >= 3:
                return "composite icon"
                
        except Exception as e:
            # If semantic detection fails, return None
            pass
            
        return None
    
    def _inject_group_parent_alt_text(self, group_shape, alt_text: str, debug: bool = False, indent: str = ""):
        """
        Inject ALT text directly into group parent XML structure.
        
        PowerPoint Reading Order checks p:grpSp/nvGrpSpPr/cNvPr for @descr and @title.
        
        Args:
            group_shape: Group shape object
            alt_text: ALT text to inject
            debug: Enable debug logging
            indent: Logging indentation
        """
        try:
            if not hasattr(group_shape, '_element'):
                if debug:
                    logger.debug(f"🔍 {indent}Group shape has no _element - cannot inject XML ALT text")
                return
            
            group_element = group_shape._element
            
            # Look for the nvGrpSpPr/cNvPr structure
            # XPath for group non-visual properties
            cnvpr_xpath = ".//p:nvGrpSpPr/p:cNvPr"
            namespaces = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
            
            cnvpr_elements = _safe_xpath(group_element, cnvpr_xpath, namespaces=namespaces)
            
            if cnvpr_elements:
                cnvpr = cnvpr_elements[0]
                
                # Set only descr for groups to avoid duplicate reads in some UIs/AT
                cnvpr.set('descr', alt_text)
                # Leave title blank for groups
                # (we still use title for pictures/shapes elsewhere if needed)
                
                if debug:
                    logger.debug(f"🔍 {indent}✅ Injected ALT into group cNvPr: '{alt_text}' (title omitted)")
            else:
                if debug:
                    logger.debug(f"🔍 {indent}⚠️  Could not find group cNvPr element for ALT injection")
                    logger.debug(f"🔍 {indent}   Group element: {group_element.tag if hasattr(group_element, 'tag') else 'unknown'}")
                
        except Exception as e:
            if debug:
                logger.debug(f"🔍 {indent}❌ Error injecting group parent ALT text: {e}")
    
    def _mark_group_children_decorative(self, children_to_mark: List, debug: bool = False, indent: str = ""):
        """
        Mark group children as decorative to avoid redundant ALT text.
        
        Args:
            children_to_mark: List of child shapes to mark as decorative
            debug: Enable debug logging
            indent: Logging indentation
        """
        for child_shape in children_to_mark:
            try:
                if not hasattr(child_shape, '_element'):
                    continue
                
                child_element = child_shape._element
                
                # Look for cNvPr element in child
                # Different child types have different paths:
                # - Pictures: p:pic/p:nvPicPr/p:cNvPr
                # - Shapes: p:sp/p:nvSpPr/p:cNvPr  
                # - Lines: p:cxnSp/p:nvCxnSpPr/p:cNvPr
                
                cnvpr_paths = [
                    ".//p:cNvPr",           # General path
                    ".//pic:cNvPr"          # Picture-specific
                ]
                
                namespaces = {
                    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                    'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture'
                }
                
                cnvpr_element = None
                for xpath in cnvpr_paths:
                    elements = _safe_xpath(child_element, xpath, namespaces=namespaces)
                    if elements:
                        cnvpr_element = elements[0]
                        break
                
                if cnvpr_element is not None:
                    # Create decorative extension element
                    # This follows Office 2019+ standard for decorative images
                    from lxml import etree
                    
                    # Create extLst if it doesn't exist
                    extlst = cnvpr_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}extLst')
                    if extlst is None:
                        extlst = etree.SubElement(cnvpr_element, '{http://schemas.openxmlformats.org/drawingml/2006/main}extLst')
                    
                    # Create decorative extension
                    ext = etree.SubElement(extlst, '{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
                    ext.set('uri', '{C809F854-F9FF-4C6D-B9C3-6A6F6E8C5B8D}')  # Office decorative URI
                    
                    decorative = etree.SubElement(ext, '{http://schemas.microsoft.com/office/drawing/2017/decorative}decorative')
                    decorative.set('val', '1')
                    
                    if debug:
                        child_id = getattr(child_shape, 'id', 'unknown')
                        logger.debug(f"🔍 {indent}   ✅ Marked child {child_id} as decorative")
                        
            except Exception as e:
                if debug:
                    logger.debug(f"🔍 {indent}   ❌ Error marking child as decorative: {e}")
                continue
    
    def _create_title_from_alt_text(self, alt_text: str) -> str:
        """
        Create a short title from the full ALT text for PowerPoint Reading Order.
        Same logic as in pptx_alt_injector.py
        
        Args:
            alt_text: Full ALT text
            
        Returns:
            Shortened title (60-80 characters)
        """
        if not alt_text or not alt_text.strip():
            return ""
        
        clean_text = alt_text.strip()
        
        # Remove common prefixes to make title more concise
        prefixes_to_remove = [
            "This is a PowerPoint shape. It is ",
            "This is a PowerPoint ",
            "This is a ",
            "A PowerPoint ",
            "PowerPoint ",
            "Group containing ",
            "Group of "
        ]
        
        for prefix in prefixes_to_remove:
            if clean_text.startswith(prefix):
                clean_text = clean_text[len(prefix):]
                break
        
        # Capitalize first letter
        if clean_text:
            clean_text = clean_text[0].upper() + clean_text[1:]
        
        # Truncate to reasonable title length (60-80 chars)
        max_title_length = 70
        if len(clean_text) <= max_title_length:
            return clean_text
        
        # Find a good breaking point
        truncated = clean_text[:max_title_length]
        
        # Try to break at sentence boundaries
        for break_char in ['. ', '! ', '? ']:
            last_break = truncated.rfind(break_char)
            if last_break > max_title_length * 0.7:
                return truncated[:last_break + 1].strip()
        
        # Fall back to word boundary
        last_space = truncated.rfind(' ')
        if last_space > max_title_length * 0.7:
            return truncated[:last_space].strip()
        
        # Hard truncate with ellipsis
        return truncated.strip() + "..."
    
    def _generate_alt_text_for_visual_element(self, visual_element: PPTXVisualElement, debug: bool = False) -> Tuple[Optional[str], Optional[str]]:
        """
        Generate ALT text for any type of visual element.
        
        Args:
            visual_element: Visual element information
            debug: Enable debug logging
            
        Returns:
            Tuple of (alt_text, failure_reason)
        """
        try:
            # For images, use the existing image generation logic
            if visual_element.element_type == "image" and visual_element.image_data:
                # Create a temporary PPTXImageInfo for compatibility
                temp_image_info = type('PPTXImageInfo', (), {
                    'image_data': visual_element.image_data,
                    'filename': visual_element.filename,
                    'slide_text': visual_element.slide_text,
                    'width_px': visual_element.width_px,
                    'height_px': visual_element.height_px,
                    'slide_idx': visual_element.slide_idx,
                    'shape_idx': visual_element.shape_idx,
                    'shape': visual_element.shape,
                    'image_key': visual_element.element_key,  # Use the stable element key
                    'image_hash': visual_element.element_hash  # Include the hash too
                })()
                
                return self._generate_alt_text_for_image_with_validation(temp_image_info, debug)
            
            # For non-image visual elements, create a screenshot and analyze it
            else:
                return self._generate_alt_text_for_shape_element(visual_element, debug)
                
        except Exception as e:
            error_msg = f"Error generating ALT text for {visual_element.element_type}: {str(e)}"
            if debug:
                logger.error(f"Exception in _generate_alt_text_for_visual_element: {e}", exc_info=True)
            return None, error_msg
    
    def _generate_alt_text_for_shape_element(self, visual_element: PPTXVisualElement, debug: bool = False) -> Tuple[Optional[str], Optional[str]]:
        """
        Generate ALT text for shape elements by creating descriptive text based on element type and properties.
        
        Args:
            visual_element: Visual element information
            debug: Enable debug logging
            
        Returns:
            Tuple of (alt_text, failure_reason)
        """
        try:
            # Create context-aware ALT text based on element type
            element_description = self._create_element_description(visual_element)
            
            if not element_description:
                return None, f"Could not create description for {visual_element.element_type}"
            
            # Enhance with context from slide
            if visual_element.slide_text and len(visual_element.slide_text.strip()) > 0:
                context_prompt = f"Slide context: {visual_element.slide_text[:200]}...\n\nShape: {element_description}\n\nCreate appropriate ALT text for this visual element considering the slide context. If it appears decorative, respond with 'decorative [element type]':"
            else:
                context_prompt = f"Shape: {element_description}\n\nCreate appropriate ALT text for this visual element. If it appears decorative, respond with 'decorative [element type]':"
            
            # Use the text generator to create ALT text
            alt_text = self.alt_generator.generate_text_response(context_prompt)
            
            if alt_text and alt_text.strip():
                # PHASE 1: Apply universal normalization before returning
                normalized_alt_text = self._normalize_alt(alt_text.strip())
                return normalized_alt_text, None
            else:
                return None, "Empty response from text generator"
                
        except Exception as e:
            error_msg = f"Error generating ALT text for shape: {str(e)}"
            if debug:
                logger.error(f"Exception in _generate_alt_text_for_shape_element: {e}", exc_info=True)
            return None, error_msg
    
    def _create_element_description(self, visual_element: PPTXVisualElement) -> str:
        """
        Create a descriptive text for a visual element based on its properties.
        
        Args:
            visual_element: Visual element to describe
            
        Returns:
            Descriptive text string
        """
        # Special handling for connectors and lines - bypass raster checks
        if visual_element.element_type in ['connector', 'line']:
            return self._create_connector_line_description(visual_element)
        
        description_parts = []
        
        # Add element type
        description_parts.append(f"A {visual_element.element_type}")
        
        # Add dimensions if available
        if visual_element.width_px > 0 and visual_element.height_px > 0:
            description_parts.append(f"sized {visual_element.width_px}x{visual_element.height_px} pixels")
        
        # Add text content if available
        if visual_element.has_text and visual_element.text_content:
            text_preview = visual_element.text_content[:100] + "..." if len(visual_element.text_content) > 100 else visual_element.text_content
            description_parts.append(f"containing text: '{text_preview}'")
        
        # Add shape name if available and meaningful
        if visual_element.shape_name and visual_element.shape_name.lower() not in ['unnamed', 'shape', 'autoshape']:
            description_parts.append(f"named '{visual_element.shape_name}'")
        
        # Add position context
        if visual_element.top_px < 200:  # Likely in title area
            description_parts.append("located in the upper area of the slide")
        elif visual_element.top_px > 400:  # Likely in lower area
            description_parts.append("located in the lower area of the slide")
        
        return " ".join(description_parts)
    
    def _create_connector_line_description(self, visual_element: PPTXVisualElement) -> str:
        """
        Create direct descriptive text for connectors and lines, bypassing raster checks.
        
        Args:
            visual_element: Connector or line visual element
            
        Returns:
            Direct descriptive ALT text
        """
        try:
            # Get shape properties for enhanced description
            shape = visual_element.shape
            width_px = visual_element.width_px
            height_px = visual_element.height_px
            
            # Determine line/connector type and orientation
            if hasattr(shape, 'connector_type'):
                element_name = "connector"
            else:
                element_name = "line"
            
            # Determine orientation
            orientation = self._get_line_orientation(width_px, height_px)
            
            # Get stroke/line properties if available
            color_info = ""
            try:
                if hasattr(shape, 'line') and shape.line:
                    # Try to get line color
                    if hasattr(shape.line, 'color') and shape.line.color:
                        # Extract color information if available
                        color_info = " colored"
                    # Could add more detailed stroke analysis here
            except:
                pass
            
            # Create minimal semantic description (boilerplate removed)
            if orientation:
                return f"{orientation}{color_info} {element_name}"
            else:
                return f"{color_info} {element_name}".strip()
                
        except Exception as e:
            logger.debug(f"Error creating connector/line description: {e}")
            # Fallback description (no boilerplate)
            return visual_element.element_type
    
    def _create_descriptive_shape_alt_text(self, shape: BaseShape, width_px: int, height_px: int) -> str:
        """
        Create descriptive ALT text for PowerPoint shapes when LLaVa processing fails.
        Provides meaningful information for screen readers.
        
        Args:
            shape: PowerPoint shape object
            width_px: Shape width in pixels
            height_px: Shape height in pixels
            
        Returns:
            Descriptive ALT text string
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
            
            # Get shape type and auto-shape type
            shape_type = getattr(shape, 'shape_type', None)
            auto_shape_type = getattr(shape, 'auto_shape_type', None)
            
            # Determine the specific shape name and appropriate article
            shape_name, article = self._get_shape_name_and_article(shape_type, auto_shape_type)
            
            # Format dimensions
            dimension_text = f"({width_px}x{height_px}px)" if width_px > 0 and height_px > 0 else ""
            
            # Create minimal semantic text (no boilerplate - injector handles that)
            alt_text = f"{article} {shape_name} {dimension_text}".strip()
            
            # Add orientation information for lines and rectangles
            if shape_type == MSO_SHAPE_TYPE.LINE or "line" in shape_name.lower():
                orientation = self._get_line_orientation(width_px, height_px)
                if orientation:
                    alt_text = alt_text.replace(f"{shape_name}", f"{orientation} {shape_name}")
            
            return alt_text
            
        except Exception as e:
            logger.debug(f"Error creating descriptive shape ALT text: {e}")
            # Ultimate fallback (minimal semantic)
            if width_px > 0 and height_px > 0:
                return f"shape ({width_px}x{height_px}px)"
            else:
                return "shape"
    
    def _get_shape_name_and_article(self, shape_type, auto_shape_type) -> Tuple[str, str]:
        """
        Get the shape name and appropriate article (a/an) for a PowerPoint shape.
        
        Args:
            shape_type: MSO_SHAPE_TYPE enumeration value
            auto_shape_type: MSO_AUTO_SHAPE_TYPE enumeration value
            
        Returns:
            Tuple of (shape_name, article)
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
            
            # Handle auto shapes with specific types
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and auto_shape_type:
                # Create shape mapping using only existing enum values
                shape_mapping = {}
                
                # Basic shapes
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'RECTANGLE'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.RECTANGLE] = ("rectangle", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'ROUNDED_RECTANGLE'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE] = ("rounded rectangle", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'OVAL'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.OVAL] = ("oval", "an")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'DIAMOND'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.DIAMOND] = ("diamond", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'RIGHT_TRIANGLE'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE] = ("right triangle", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'PARALLELOGRAM'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM] = ("parallelogram", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'TRAPEZOID'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.TRAPEZOID] = ("trapezoid", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'HEXAGON'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.HEXAGON] = ("hexagon", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'OCTAGON'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.OCTAGON] = ("octagon", "an")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'DECAGON'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.DECAGON] = ("decagon", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'DODECAGON'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.DODECAGON] = ("dodecagon", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'PENTAGON'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.PENTAGON] = ("pentagon", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'REGULAR_PENTAGON'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.REGULAR_PENTAGON] = ("regular pentagon", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'HEPTAGON'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.HEPTAGON] = ("heptagon", "a")
                
                # Stars (using actual enum names)
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'STAR_4_POINT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.STAR_4_POINT] = ("4-point star", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'STAR_5_POINT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.STAR_5_POINT] = ("5-point star", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'STAR_6_POINT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.STAR_6_POINT] = ("6-point star", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'STAR_7_POINT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.STAR_7_POINT] = ("7-point star", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'STAR_8_POINT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.STAR_8_POINT] = ("8-point star", "an")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'STAR_10_POINT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.STAR_10_POINT] = ("10-point star", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'STAR_12_POINT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.STAR_12_POINT] = ("12-point star", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'STAR_16_POINT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.STAR_16_POINT] = ("16-point star", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'STAR_24_POINT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.STAR_24_POINT] = ("24-point star", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'STAR_32_POINT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.STAR_32_POINT] = ("32-point star", "a")
                    
                # Other common shapes
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'CROSS'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.CROSS] = ("cross", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'HEART'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.HEART] = ("heart", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'LIGHTNING_BOLT'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.LIGHTNING_BOLT] = ("lightning bolt", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'MOON'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.MOON] = ("moon", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'SUN'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.SUN] = ("sun", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'CLOUD'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.CLOUD] = ("cloud", "a")
                
                # Arrows
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'LEFT_ARROW'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.LEFT_ARROW] = ("left arrow", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'RIGHT_ARROW'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW] = ("right arrow", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'UP_ARROW'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.UP_ARROW] = ("up arrow", "an")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'DOWN_ARROW'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.DOWN_ARROW] = ("down arrow", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'LEFT_RIGHT_ARROW'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_ARROW] = ("left-right arrow", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'UP_DOWN_ARROW'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.UP_DOWN_ARROW] = ("up-down arrow", "an")
                
                # Flowchart shapes (common ones)
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'FLOWCHART_PROCESS'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS] = ("flowchart process box", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'FLOWCHART_DECISION'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.FLOWCHART_DECISION] = ("flowchart decision diamond", "a")
                if hasattr(MSO_AUTO_SHAPE_TYPE, 'FLOWCHART_TERMINATOR'):
                    shape_mapping[MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR] = ("flowchart terminator", "a")
                
                if auto_shape_type in shape_mapping:
                    return shape_mapping[auto_shape_type]
            
            # Handle general shape types
            if shape_type == MSO_SHAPE_TYPE.LINE:
                return ("line", "a")
            elif shape_type == MSO_SHAPE_TYPE.FREEFORM:
                return ("freeform shape", "a")
            elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return ("auto shape", "an")
            elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                return ("text box", "a")
            elif shape_type == MSO_SHAPE_TYPE.CHART:
                return ("chart", "a")
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                return ("table", "a")
            elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                return ("picture", "a")
            elif shape_type == MSO_SHAPE_TYPE.MEDIA:
                return ("media object", "a")
            elif shape_type == MSO_SHAPE_TYPE.OLE_OBJECT:
                return ("embedded object", "an")
            elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                return ("placeholder", "a")
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                return ("group of shapes", "a")
            else:
                return ("shape", "a")
                
        except Exception as e:
            logger.debug(f"Error mapping shape type: {e}")
            return ("shape", "a")
    
    def _get_line_orientation(self, width_px: int, height_px: int) -> str:
        """
        Determine line orientation based on dimensions.
        
        Args:
            width_px: Width in pixels
            height_px: Height in pixels
            
        Returns:
            Orientation description or empty string
        """
        if width_px <= 0 or height_px <= 0:
            return ""
        
        # Calculate aspect ratio
        aspect_ratio = width_px / height_px
        
        if aspect_ratio > 3:  # Much wider than tall
            return "horizontal"
        elif aspect_ratio < 0.33:  # Much taller than wide
            return "vertical"
        elif 0.8 <= aspect_ratio <= 1.2:  # Roughly square
            return "diagonal"
        else:
            return ""  # Don't specify if unclear
    
    def _extract_visual_elements_from_shapes(self, shapes, slide_idx: int, slide_text: str, debug: bool = False) -> List[PPTXVisualElement]:
        """
        Recursively extract visual elements from shapes collection using stable shape IDs.
        
        Args:
            shapes: Collection of shapes to process
            slide_idx: Slide index
            slide_text: Slide text context
            debug: Enable debug logging
            
        Returns:
            List of PPTXVisualElement objects
        """
        visual_elements = []
        
        for shape_idx, shape in enumerate(shapes):
            try:
                # Handle grouped shapes recursively - ALWAYS PROCESS BOTH CHILDREN AND PARENT
                if hasattr(shape, 'shapes') and shape.shapes:
                    # First, process shapes within groups using recursive DFS
                    if debug:
                        logger.debug(f"Processing group shape with {len(shape.shapes)} children")
                    group_elements = self._extract_visual_elements_from_shapes(
                        shape.shapes, slide_idx, slide_text, debug
                    )
                    visual_elements.extend(group_elements)
                    
                    # DESCRIBE EVERYTHING POLICY: Also create element for group parent
                    # This ensures PowerPoint Reading Order has ALT text at group level
                    if debug:
                        logger.debug(f"Creating group parent element for {len(shape.shapes)} children")
                    
                    # Don't continue - process the group parent as well
                
                # Determine element type
                element_type = self._classify_visual_element(shape)
                
                # Skip text-only placeholders unless they have visual significance
                if element_type == "text_placeholder" and not self._has_visual_significance(shape):
                    if debug:
                        logger.debug(f"Skipping text-only placeholder: {getattr(shape, 'name', 'unnamed')}")
                    continue
                
                # Use stable shape ID when available, with enhanced XML extraction for grouped children
                shape_identifier = self._extract_robust_shape_id(shape, shape_idx, debug)
                
                if shape_identifier == shape_idx:
                    if debug:
                        logger.debug(f"Using fallback index {shape_idx} for shape without ID")
                else:
                    if debug:
                        logger.debug(f"Using stable shape ID {shape_identifier}")
                
                # Create visual element with stable identifier
                visual_element = PPTXVisualElement(shape, slide_idx, shape_identifier, slide_text, element_type)
                visual_elements.append(visual_element)
                
                if debug:
                    logger.debug(f"Added {element_type}: {visual_element.element_key} ({visual_element.width_px}x{visual_element.height_px}px)")
                
            except Exception as e:
                logger.warning(f"Error processing shape {shape_idx} on slide {slide_idx}: {e}")
                continue
        
        return visual_elements
    
    def _classify_visual_element(self, shape) -> str:
        """
        Classify the type of visual element.
        
        Args:
            shape: Shape object
            
        Returns:
            String classification of the element type
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            shape_type = getattr(shape, 'shape_type', None)
            
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                return "image"
            elif shape_type == MSO_SHAPE_TYPE.CHART:
                return "chart" 
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                return "table"
            elif shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]:
                return "shape"
            elif shape_type == MSO_SHAPE_TYPE.LINE:
                return "line"
            elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                return "text_box"
            elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                return "text_placeholder"
            elif shape_type == MSO_SHAPE_TYPE.MEDIA:
                return "media"
            elif shape_type == MSO_SHAPE_TYPE.OLE_OBJECT:
                return "embedded_object"
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                return "group"
            # Add connector support
            elif hasattr(MSO_SHAPE_TYPE, 'CONNECTOR') and shape_type == MSO_SHAPE_TYPE.CONNECTOR:
                return "connector"
            # Check for connector-like behavior in shapes without explicit CONNECTOR type
            elif hasattr(shape, 'connector_type') or 'Connector' in type(shape).__name__:
                return "connector"
            # Check if this is a group shape (fallback)
            elif hasattr(shape, 'shapes') and shape.shapes:
                return "group"
            else:
                return "unknown"
                
        except Exception:
            return "unknown"
    
    def _has_visual_significance(self, shape) -> bool:
        """
        Determine if a shape has visual significance beyond just text.
        
        Args:
            shape: Shape object
            
        Returns:
            Boolean indicating if shape is visually significant
        """
        try:
            # Check for fills, borders, or visual styling
            if hasattr(shape, 'fill') and shape.fill:
                return True
            if hasattr(shape, 'line') and shape.line:
                return True
            
            # Check dimensions - very large elements might be visually significant
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                try:
                    width_px = int(shape.width.emu / 914400 * 96) if shape.width else 0
                    height_px = int(shape.height.emu / 914400 * 96) if shape.height else 0
                    if width_px > 200 or height_px > 200:
                        return True
                except:
                    pass
            
            return False
        except:
            return False
    
    def detect_decorative_shapes(self, presentation: Presentation, debug: bool = False) -> List[PPTXShapeInfo]:
        """
        Detect decorative shapes (basic geometric shapes without meaningful content) in PPTX.
        
        Args:
            presentation: PowerPoint presentation
            debug: Enable debug logging
            
        Returns:
            List of PPTXShapeInfo objects representing decorative shapes
        """
        decorative_shapes = []
        
        # Define comprehensive decorative shape types with version compatibility
        decorative_shape_types = {
            MSO_SHAPE_TYPE.AUTO_SHAPE,    # Most geometric shapes (rectangles, ovals, etc.)
            MSO_SHAPE_TYPE.LINE,          # Lines and arrows
            MSO_SHAPE_TYPE.FREEFORM,      # Freeform drawings and custom shapes
            MSO_SHAPE_TYPE.CALLOUT,       # Callout shapes
            MSO_SHAPE_TYPE.TEXT_EFFECT,   # WordArt/text effects (often decorative)
        }
        
        # Add additional shape types that might be considered decorative
        # Check for existence first as different python-pptx versions have different attributes
        try:
            if hasattr(MSO_SHAPE_TYPE, 'CONNECTOR'):
                decorative_shape_types.add(MSO_SHAPE_TYPE.CONNECTOR)  # Connector lines between shapes
            if hasattr(MSO_SHAPE_TYPE, 'CANVAS'):
                decorative_shape_types.add(MSO_SHAPE_TYPE.CANVAS)
            if hasattr(MSO_SHAPE_TYPE, 'DIAGRAM'):
                decorative_shape_types.add(MSO_SHAPE_TYPE.DIAGRAM)  # SmartArt diagrams
            if hasattr(MSO_SHAPE_TYPE, 'INK'):
                decorative_shape_types.add(MSO_SHAPE_TYPE.INK)      # Ink annotations
            if hasattr(MSO_SHAPE_TYPE, 'MEDIA'):
                decorative_shape_types.add(MSO_SHAPE_TYPE.MEDIA)    # Media objects (might be decorative)
            if hasattr(MSO_SHAPE_TYPE, 'OLE_OBJECT'):
                decorative_shape_types.add(MSO_SHAPE_TYPE.OLE_OBJECT)  # OLE objects (might be decorative)
        except AttributeError:
            pass
        
        if debug:
            logger.debug(f"🎨 Using decorative shape types: {[t for t in decorative_shape_types]}")
        
        for slide_idx, slide in enumerate(presentation.slides):
            if debug:
                logger.info(f"🔍 DEBUG: Scanning slide {slide_idx + 1} for decorative shapes")
            
            # Extract slide text for context
            slide_text = self._extract_slide_text(slide) if self.include_slide_text else ""
            
            # Process all shapes recursively with enhanced detection
            decorative_on_slide = self._detect_decorative_shapes_recursive(
                slide.shapes, slide_idx, slide_text, decorative_shape_types, debug
            )
            
            decorative_shapes.extend(decorative_on_slide)
            
            if debug and decorative_on_slide:
                logger.info(f"🔍 DEBUG: Found {len(decorative_on_slide)} decorative shapes on slide {slide_idx + 1}")
        
        logger.info(f"Detected {len(decorative_shapes)} potentially decorative shapes")
        return decorative_shapes
    
    def _detect_decorative_shapes_recursive(self, shapes, slide_idx: int, slide_text: str, 
                                          decorative_types: set, debug: bool = False, 
                                          parent_group_idx: str = None, depth: int = 0) -> List[PPTXShapeInfo]:
        """
        Recursively detect decorative shapes with enhanced traversal and fallback detection.
        
        Args:
            shapes: Collection of shapes to process
            slide_idx: Slide index
            slide_text: Slide text context
            decorative_types: Set of shape types considered potentially decorative
            debug: Enable debug logging
            parent_group_idx: Parent group identifier for nested shapes
            depth: Current recursion depth for debugging
            
        Returns:
            List of decorative PPTXShapeInfo objects
        """
        decorative_shapes = []
        indent = "  " + "  " * depth
        
        if debug and depth == 0:
            logger.debug(f"🎨 Starting decorative shape detection on slide {slide_idx + 1}")
        
        total_shapes = len(shapes) if hasattr(shapes, '__len__') else 0
        if debug:
            context = f"group {parent_group_idx}" if parent_group_idx else "slide"
            logger.debug(f"{indent}🔍 Checking {total_shapes} shapes in {context} (depth {depth})")
        
        # Enhanced debug: Count shapes by type to understand what we're working with
        shape_type_counts = {}
        if debug and total_shapes > 0:
            for shape in shapes:
                try:
                    shape_type = getattr(shape, 'shape_type', None)
                    shape_class = shape.__class__.__name__
                    type_key = f"{shape_type} ({shape_class})" if shape_type else f"None ({shape_class})"
                    shape_type_counts[type_key] = shape_type_counts.get(type_key, 0) + 1
                except:
                    shape_type_counts['Unknown'] = shape_type_counts.get('Unknown', 0) + 1
            
            logger.debug(f"{indent}📊 Shape type distribution: {dict(sorted(shape_type_counts.items()))}")
        
        for shape_idx, shape in enumerate(shapes):
            try:
                # Create hierarchical shape identifier
                if parent_group_idx is not None:
                    shape_id = f"{parent_group_idx}_{shape_idx}"
                else:
                    shape_id = shape_idx
                
                shape_class = type(shape).__name__
                shape_type = getattr(shape, 'shape_type', None)
                
                if debug:
                    logger.debug(f"{indent}  🔹 Shape {shape_id}: {shape_class} (type={shape_type})")
                
                # Skip images (handled separately)
                if hasattr(shape, 'image') and shape.image:
                    if debug:
                        logger.debug(f"{indent}    ⏭️ Skipping image shape (handled separately)")
                    continue
                
                # Enhanced group detection and traversal
                if hasattr(shape, 'shapes') and shape.shapes:
                    group_size = len(shape.shapes)
                    if debug:
                        logger.debug(f"{indent}    📁 Group with {group_size} children - traversing recursively")
                    
                    group_decorative = self._detect_decorative_shapes_recursive(
                        shape.shapes, slide_idx, slide_text, decorative_types, debug, shape_id, depth + 1
                    )
                    
                    # Enhanced group analysis
                    if self._is_group_decorative(shape, group_decorative, debug, indent):
                        shape_info = PPTXShapeInfo(shape, slide_idx, shape_id, slide_text)
                        decorative_shapes.append(shape_info)
                        if debug:
                            logger.debug(f"{indent}    ✅ Marked entire group {shape_id} as decorative")
                    else:
                        # Add individual decorative shapes from within the group
                        decorative_shapes.extend(group_decorative)
                        if debug and group_decorative:
                            logger.debug(f"{indent}    📝 Added {len(group_decorative)} decorative shapes from group")
                    
                    continue
                
                # Enhanced shape type checking with fallback detection
                is_potentially_decorative = False
                fallback_reason = None
                
                try:
                    if shape_type is not None:
                        # Standard type checking
                        if shape_type in decorative_types:
                            is_potentially_decorative = True
                            if debug:
                                logger.debug(f"{indent}    ✅ Shape type {shape_type} is in decorative types")
                        else:
                            if debug:
                                logger.debug(f"{indent}    ❌ Shape type {shape_type} not in decorative types")
                    else:
                        # Fallback detection for shapes without standard type
                        fallback_result = self._detect_decorative_fallback(shape, shape_class, debug, indent)
                        is_potentially_decorative = fallback_result['is_decorative']
                        fallback_reason = fallback_result['reason']
                        
                        if debug:
                            if is_potentially_decorative:
                                logger.debug(f"{indent}    ⚠️ FALLBACK: Detected as decorative - {fallback_reason}")
                            else:
                                logger.debug(f"{indent}    ⚠️ FALLBACK: Not decorative - {fallback_reason}")
                except Exception as e:
                    if debug:
                        logger.debug(f"{indent}    ❌ Error in type checking: {e}")
                    # Try fallback detection even on errors
                    fallback_result = self._detect_decorative_fallback(shape, shape_class, debug, indent)
                    is_potentially_decorative = fallback_result['is_decorative']
                    fallback_reason = f"Error in standard detection ({e}), used fallback: {fallback_result['reason']}"
                
                if not is_potentially_decorative:
                    if debug:
                        rejection_reason = fallback_reason or f"Shape type {shape_type} not in decorative types"
                        logger.debug(f"{indent}    ❌ REJECTED: {shape_id} - {rejection_reason}")
                    continue
                
                # Create shape info for detailed analysis with comprehensive error handling
                try:
                    shape_info = PPTXShapeInfo(shape, slide_idx, shape_id, slide_text)
                    
                    # Validate shape info was created successfully
                    if shape_info is None:
                        if debug:
                            logger.debug(f"{indent}    ❌ Failed to create shape info for {shape_id}")
                        continue
                    
                    if debug:
                        try:
                            size_info = f"({shape_info.width_px}x{shape_info.height_px}px)"
                            type_info = f"{shape_info.shape_type_name}"
                            if fallback_reason:
                                type_info += f" [FALLBACK: {fallback_reason}]"
                            logger.debug(f"{indent}    📋 Analyzing: {type_info} {size_info}")
                        except Exception as debug_error:
                            logger.debug(f"{indent}    📋 Analyzing: {shape_id} (error formatting debug info: {debug_error})")
                    
                    # Apply comprehensive decorative detection heuristics with error handling
                    try:
                        is_decorative = self._is_shape_decorative(shape_info, debug, indent)
                        if is_decorative:
                            decorative_shapes.append(shape_info)
                            if debug:
                                logger.debug(f"{indent}    ✅ DECORATIVE: {shape_id} marked as decorative")
                        elif debug:
                            logger.debug(f"{indent}    ❌ MEANINGFUL: {shape_id} has meaningful content")
                    except Exception as detection_error:
                        if debug:
                            logger.debug(f"{indent}    ⚠️ Error in decorative detection for {shape_id}: {detection_error}")
                        # Assume not decorative if detection fails
                        if debug:
                            logger.debug(f"{indent}    ❌ FALLBACK: {shape_id} assumed meaningful due to detection error")
                
                except Exception as e:
                    logger.warning(f"Error creating or analyzing shape info for {shape_id}: {e}")
                    if debug:
                        logger.debug(f"{indent}    ❌ Skipping {shape_id} due to critical error: {e}")
                    continue
                
            except Exception as e:
                logger.warning(f"Error analyzing shape {shape_idx} on slide {slide_idx}: {e}")
                if debug:
                    logger.debug(f"{indent}  ❌ Error in shape {shape_idx}: {e}")
                continue
        
        if debug:
            context = f"group {parent_group_idx}" if parent_group_idx else "slide"
            logger.debug(f"{indent}📊 Found {len(decorative_shapes)} decorative shapes in {context} (depth {depth})")
        
        return decorative_shapes
    
    def _is_shape_decorative(self, shape_info: PPTXShapeInfo, debug: bool = False, indent: str = "") -> bool:
        """
        Enhanced decorative detection with educational content awareness and comprehensive heuristics.
        
        Args:
            shape_info: Shape information
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            bool: True if shape appears to be decorative
        """
        # PRIORITY RULE: Educational content shapes are never decorative
        if self._is_educational_shape(shape_info, debug, indent):
            if debug:
                logger.debug(f"{indent}      🎓 Educational Rule: Shape contains educational content - NOT decorative")
            return False
        
        # PRIORITY RULE: Text boxes with meaningful content are NEVER decorative
        if self._is_text_box_with_content(shape_info, debug, indent):
            if debug:
                logger.debug(f"{indent}      📝 Text Box Rule: Has meaningful text content - NOT decorative")
            return False
        
        # PRIORITY RULE: Title areas and slide headers are NEVER decorative
        if self._is_title_or_header_area(shape_info, debug, indent):
            if debug:
                logger.debug(f"{indent}      📜 Title Rule: Title or header area - NOT decorative")
            return False
        
        # Rule 1: Shapes with meaningful text content are not decorative (fallback)
        try:
            if shape_info.has_text and shape_info.text_content and len(shape_info.text_content.strip()) > 2:
                text_preview = shape_info.text_content[:30] + ('...' if len(shape_info.text_content) > 30 else '')
                if debug:
                    logger.debug(f"{indent}      ✅ Rule 1: Has meaningful text ({len(shape_info.text_content)} chars): '{text_preview}'")
                return False
        except Exception as e:
            if debug:
                logger.debug(f"{indent}      ⚠️ Error checking text content: {e}")
        
        # Get dimensions for multiple rules with safe fallbacks
        try:
            min_dimension = min(shape_info.width_px, shape_info.height_px)
            max_dimension = max(shape_info.width_px, shape_info.height_px)
            
            # Sanity check dimensions
            if min_dimension < 0 or max_dimension < 0:
                if debug:
                    logger.debug(f"{indent}      ⚠️ Invalid dimensions detected, using fallback values")
                min_dimension = max(0, min_dimension)
                max_dimension = max(0, max_dimension)
        except Exception as e:
            if debug:
                logger.debug(f"{indent}      ⚠️ Error calculating dimensions: {e}, using defaults")
            min_dimension, max_dimension = 0, 0
        
        # Rule 2: Very small shapes are likely decorative, but check context first
        if min_dimension < self.decorative_size_threshold and min_dimension > 0:
            # Exception: Small shapes in educational context might be important (labels, annotations)
            try:
                has_educational_context = self._has_educational_context(shape_info)
                if has_educational_context:
                    if debug:
                        logger.debug(f"{indent}      ❌ Rule 2a: Small shape but in educational context - NOT decorative")
                    return False
                else:
                    if debug:
                        logger.debug(f"{indent}      ✅ Rule 2b: Very small shape ({min_dimension}px < {self.decorative_size_threshold}px)")
                    return True
            except Exception as e:
                if debug:
                    logger.debug(f"{indent}      ⚠️ Error checking educational context: {e}, assuming decorative")
                if debug:
                    logger.debug(f"{indent}      ✅ Rule 2b (fallback): Very small shape ({min_dimension}px < {self.decorative_size_threshold}px)")
                return True
        
        # Rule 3: Lines are typically decorative unless they have text
        if shape_info.shape_type == MSO_SHAPE_TYPE.LINE:
            if debug:
                logger.debug(f"{indent}      ✅ Rule 3: Line shape")
            return True
        
        # Rule 4: Connectors are typically decorative (if supported)
        if hasattr(MSO_SHAPE_TYPE, 'CONNECTOR') and shape_info.shape_type == MSO_SHAPE_TYPE.CONNECTOR:
            if debug:
                logger.debug(f"{indent}      ✅ Rule 4: Connector shape")
            return True
        
        # Rule 5: Freeform shapes are often decorative drawings
        if shape_info.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            if debug:
                logger.debug(f"{indent}      ✅ Rule 5: Freeform drawing shape")
            return True
        
        # Rule 6: Text effects are often decorative WordArt
        if shape_info.shape_type == MSO_SHAPE_TYPE.TEXT_EFFECT:
            # Exception: if it has substantial text, it might be functional
            if shape_info.has_text and len(shape_info.text_content) > 10:
                if debug:
                    logger.debug(f"{indent}      ❌ Rule 6a: Text effect with substantial content (not decorative)")
                return False
            else:
                if debug:
                    logger.debug(f"{indent}      ✅ Rule 6b: Text effect without substantial content (decorative)")
                return True
        
        # Rule 7: Callouts without text are often decorative
        if shape_info.shape_type == MSO_SHAPE_TYPE.CALLOUT and not shape_info.has_text:
            if debug:
                logger.debug(f"{indent}      ✅ Rule 7: Empty callout shape")
            return True
        
        # Rule 8: Auto shapes - enhanced analysis with geometric shape detection
        if shape_info.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and not shape_info.has_text:
            # Identify simple geometric shapes by their auto shape type
            try:
                auto_shape_details = self._get_auto_shape_details(shape_info.shape, debug, indent)
                
                # Simple geometric shapes without text are often decorative
                simple_geometric_shapes = [
                    'OVAL', 'CIRCLE', 'RECTANGLE', 'SQUARE', 'TRIANGLE', 
                    'HEXAGON', 'OCTAGON', 'DECAGON', 'PENTAGON', 'DIAMOND',
                    'PARALLELOGRAM', 'TRAPEZOID', 'STAR', 'HEART', 'CROSS'
                ]
                
                if auto_shape_details:  # Check if we got valid details
                    shape_type_upper = auto_shape_details.upper()
                    is_simple_geometric = any(geo_shape in shape_type_upper for geo_shape in simple_geometric_shapes)
                    
                    if is_simple_geometric:
                        if debug:
                            logger.debug(f"{indent}      ✅ Rule 8a: Simple geometric shape detected ({auto_shape_details}) - decorative")
                        return True
                elif debug:
                    logger.debug(f"{indent}      ⚠️ Could not get auto shape details for geometric analysis")
            except Exception as e:
                if debug:
                    logger.debug(f"{indent}      ⚠️ Error in geometric shape detection: {e}")
                # Continue with other rules if geometric detection fails
            
            # Very large background elements (likely decorative backgrounds)
            if max_dimension > 500:  # Increased threshold for educational content
                # Exception: Large shapes in educational presentations might be important backgrounds
                if self._has_educational_context(shape_info):
                    if debug:
                        logger.debug(f"{indent}      ❌ Rule 8b: Large auto shape in educational context - might be content")
                    return False
                else:
                    if debug:
                        logger.debug(f"{indent}      ✅ Rule 8c: Very large auto shape background ({max_dimension}px)")
                    return True
            
            # Medium-large shapes - check for elongation (decorative dividers/lines)
            elif max_dimension > 200:
                # Check aspect ratio - very wide or tall shapes might be decorative
                if min_dimension > 0:
                    aspect_ratio = max_dimension / min_dimension
                    if aspect_ratio > 4:  # Increased threshold for educational content
                        if debug:
                            logger.debug(f"{indent}      ✅ Rule 8d: Very elongated auto shape (ratio {aspect_ratio:.1f})")
                        return True
                if debug:
                    logger.debug(f"{indent}      ❌ Rule 8e: Medium auto shape - checking if decorative ({auto_shape_details})")
            
            # Small geometric shapes - still likely decorative
            elif min_dimension < 80 and min_dimension > 0:  # Slightly more generous
                if debug:
                    logger.debug(f"{indent}      ✅ Rule 8f: Small auto shape decoration ({min_dimension}px)")
                return True
        
        # Rule 9: Very thin shapes (likely decorative lines/dividers)
        try:
            if min_dimension > 0 and max_dimension > 0:
                aspect_ratio = max_dimension / min_dimension
                if aspect_ratio > 15:  # Very thin shapes
                    if debug:
                        logger.debug(f"{indent}      ✅ Rule 9: Very thin shape (aspect ratio {aspect_ratio:.1f})")
                    return True
        except (ZeroDivisionError, TypeError, ValueError) as e:
            if debug:
                logger.debug(f"{indent}      ⚠️ Error calculating aspect ratio: {e}")
        
        # Rule 10: Tiny shapes (likely bullets or decorative elements)
        if max_dimension > 0 and max_dimension < 15:
            if debug:
                logger.debug(f"{indent}      ✅ Rule 10: Tiny shape ({max_dimension}px max dimension)")
            return True
        
        # Rule 11: Shapes with only very short text (1-2 characters) might be decorative
        try:
            if shape_info.has_text and shape_info.text_content:
                text_length = len(shape_info.text_content.strip())
                if 0 < text_length <= 2:
                    # Common decorative text patterns
                    decorative_text = shape_info.text_content.strip().lower()
                    decorative_symbols = ['•', '●', '○', '■', '□', '★', '☆', '*', '-', '_', '|', '/', '\\']
                    if decorative_text in decorative_symbols:
                        if debug:
                            logger.debug(f"{indent}      ✅ Rule 11: Decorative symbol text: '{shape_info.text_content}'")
                        return True
        except (AttributeError, TypeError) as e:
            if debug:
                logger.debug(f"{indent}      ⚠️ Error checking decorative text: {e}")
        
        # Default: not decorative
        if debug:
            type_name = shape_info.shape_type_name if hasattr(shape_info, 'shape_type_name') else 'unknown'
            logger.debug(f"{indent}      ❌ No decorative rules matched for {type_name} ({shape_info.width_px}x{shape_info.height_px}px) - has meaningful content")
        return False
    
    def _detect_decorative_fallback(self, shape, shape_class: str, debug: bool = False, indent: str = "") -> Dict[str, Any]:
        """
        Fallback detection for shapes that don't match standard MSO_SHAPE_TYPE categories.
        
        Args:
            shape: Shape object to analyze
            shape_class: Class name of the shape
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            Dictionary with 'is_decorative' bool and 'reason' string
        """
        result = {'is_decorative': False, 'reason': 'No fallback criteria matched'}
        
        if debug:
            logger.debug(f"{indent}    🔍 FALLBACK: Analyzing {shape_class} with fallback detection")
        
        try:
            # Check XML element for additional clues
            xml_tag = None
            if hasattr(shape, '_element') and shape._element is not None:
                try:
                    xml_tag = shape._element.tag
                    if debug:
                        logger.debug(f"{indent}      📄 XML tag: {xml_tag}")
                except:
                    if debug:
                        logger.debug(f"{indent}      ⚠️ Could not extract XML tag")
                    pass
            else:
                if debug:
                    logger.debug(f"{indent}      ❌ No XML element available")
            
            # Fallback criteria based on class name and properties
            
            # 1. Check for geometric shape classes
            geometric_classes = ['Shape', 'AutoShape', 'Rectangle', 'Oval', 'Triangle', 'Line']
            geometric_match = [cls for cls in geometric_classes if cls in shape_class]
            if geometric_match:
                result['is_decorative'] = True
                result['reason'] = f"Geometric shape class: {shape_class} (matched: {geometric_match})"
                if debug:
                    logger.debug(f"{indent}      ✅ Geometric match: {geometric_match}")
                return result
            elif debug:
                logger.debug(f"{indent}      ❌ No geometric class match in {shape_class}")
            
            # 2. Check for connector-related classes
            connector_classes = ['Connector', 'Connection']
            connector_match = [cls for cls in connector_classes if cls in shape_class]
            if connector_match:
                result['is_decorative'] = True
                result['reason'] = f"Connector class: {shape_class} (matched: {connector_match})"
                if debug:
                    logger.debug(f"{indent}      ✅ Connector match: {connector_match}")
                return result
            elif debug:
                logger.debug(f"{indent}      ❌ No connector class match in {shape_class}")
            
            # 3. Check for drawing/freeform classes
            drawing_classes = ['Freeform', 'Drawing', 'Path', 'Curve']
            drawing_match = [cls for cls in drawing_classes if cls in shape_class]
            if drawing_match:
                result['is_decorative'] = True
                result['reason'] = f"Drawing class: {shape_class} (matched: {drawing_match})"
                if debug:
                    logger.debug(f"{indent}      ✅ Drawing match: {drawing_match}")
                return result
            elif debug:
                logger.debug(f"{indent}      ❌ No drawing class match in {shape_class}")
            
            # 4. XML-based detection
            if xml_tag:
                decorative_xml_patterns = ['line', 'rect', 'ellipse', 'path', 'polygon', 'connector']
                xml_lower = xml_tag.lower()
                xml_matches = [pattern for pattern in decorative_xml_patterns if pattern in xml_lower]
                if xml_matches:
                    result['is_decorative'] = True
                    result['reason'] = f"XML tag contains decorative patterns {xml_matches}: {xml_tag}"
                    if debug:
                        logger.debug(f"{indent}      ✅ XML pattern match: {xml_matches} in {xml_tag}")
                    return result
                elif debug:
                    logger.debug(f"{indent}      ❌ No XML pattern match in {xml_tag} (checked: {decorative_xml_patterns})")
            elif debug:
                logger.debug(f"{indent}      ❌ No XML tag available for pattern matching")
            
            # 5. Check for size-based heuristics (very small shapes are likely decorative)
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                try:
                    width_px = int(shape.width.emu / 914400 * 96) if shape.width else 0
                    height_px = int(shape.height.emu / 914400 * 96) if shape.height else 0
                    
                    # Very small shapes (< 20px in any dimension) are likely decorative
                    if (width_px > 0 and width_px < 20) or (height_px > 0 and height_px < 20):
                        result['is_decorative'] = True
                        result['reason'] = f"Very small shape: {width_px}x{height_px}px"
                        return result
                    
                    # Very thin shapes (aspect ratio > 10:1) might be decorative lines
                    if width_px > 0 and height_px > 0:
                        aspect_ratio = max(width_px, height_px) / min(width_px, height_px)
                        if aspect_ratio > 10:
                            result['is_decorative'] = True
                            result['reason'] = f"Very thin shape (aspect ratio {aspect_ratio:.1f}): {width_px}x{height_px}px"
                            return result
                except:
                    pass
            
            # 6. Check for absence of text content in text-capable shapes
            if hasattr(shape, 'text_frame') or hasattr(shape, 'text'):
                has_text = False
                try:
                    if hasattr(shape, 'text') and shape.text:
                        text_content = shape.text.strip()
                        has_text = len(text_content) > 0
                    elif hasattr(shape, 'text_frame') and shape.text_frame:
                        # Check if text frame has content
                        if hasattr(shape.text_frame, 'text') and shape.text_frame.text:
                            text_content = shape.text_frame.text.strip()
                            has_text = len(text_content) > 0
                except:
                    pass
                
                if not has_text:
                    result['is_decorative'] = True
                    result['reason'] = f"Text-capable shape without content: {shape_class}"
                    return result
            
            # 7. Check for fill patterns that suggest decoration
            if hasattr(shape, 'fill'):
                try:
                    from pptx.dml.fill import MSO_FILL_TYPE
                    fill_type = getattr(shape.fill, 'type', None)
                    if fill_type == MSO_FILL_TYPE.GRADIENT:
                        result['is_decorative'] = True
                        result['reason'] = "Shape with gradient fill (likely decorative)"
                        return result
                    elif fill_type == MSO_FILL_TYPE.TEXTURED:
                        result['is_decorative'] = True
                        result['reason'] = "Shape with texture fill (likely decorative)"
                        return result
                except:
                    pass
            
            # Default: not decorative
            result['reason'] = f"No fallback criteria matched for {shape_class}"
            if debug:
                logger.debug(f"{indent}      ❌ FALLBACK FAILED: All criteria failed for {shape_class} - not decorative")
            return result
            
        except Exception as e:
            result['reason'] = f"Error in fallback detection: {e}"
            if debug:
                logger.debug(f"{indent}      Fallback error: {e}")
            return result
    
    def _is_group_decorative(self, group_shape, group_decorative_shapes: List[PPTXShapeInfo], debug: bool = False, indent: str = "") -> bool:
        """
        Enhanced group decorative analysis with better heuristics.
        
        Args:
            group_shape: The group shape object
            group_decorative_shapes: List of decorative shapes found within the group
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            bool: True if the entire group should be marked as decorative
        """
        try:
            total_shapes_in_group = len(group_shape.shapes) if hasattr(group_shape, 'shapes') else 0
            decorative_shapes_in_group = len(group_decorative_shapes)
            
            if debug:
                logger.debug(f"{indent}      📋 Group analysis: {decorative_shapes_in_group}/{total_shapes_in_group} shapes are decorative")
            
            # No shapes in group - not decorative
            if total_shapes_in_group == 0:
                return False
            
            # If all shapes in the group are decorative, mark the whole group as decorative
            if decorative_shapes_in_group == total_shapes_in_group:
                if debug:
                    logger.debug(f"{indent}      ✅ All shapes in group are decorative")
                return True
            
            # Enhanced heuristics for partial decorative groups
            decorative_ratio = decorative_shapes_in_group / total_shapes_in_group
            
            # Small groups with high decorative ratio
            if total_shapes_in_group <= 3 and decorative_ratio >= 0.67:  # 2/3 or more
                if debug:
                    logger.debug(f"{indent}      ✅ Small group ({total_shapes_in_group}) is {decorative_ratio:.1%} decorative")
                return True
            
            # Medium groups with very high decorative ratio
            if total_shapes_in_group <= 5 and decorative_ratio >= 0.8:  # 4/5 or more
                if debug:
                    logger.debug(f"{indent}      ✅ Medium group ({total_shapes_in_group}) is {decorative_ratio:.1%} decorative")
                return True
            
            # Check if remaining shapes are just text boxes or have minimal content
            non_decorative_count = total_shapes_in_group - decorative_shapes_in_group
            meaningful_content_count = 0
            
            try:
                for shape in group_shape.shapes:
                    # Skip already identified decorative shapes
                    is_decorative_shape = any(
                        ds.shape == shape for ds in group_decorative_shapes
                    )
                    if is_decorative_shape:
                        continue
                    
                    # Check if this shape has meaningful content
                    has_meaningful_content = False
                    
                    # Check for substantial text
                    if hasattr(shape, 'text') and shape.text:
                        text_length = len(shape.text.strip())
                        if text_length > 10:  # More than just a few characters
                            has_meaningful_content = True
                    
                    # Check for images
                    if hasattr(shape, 'image') and shape.image:
                        has_meaningful_content = True
                    
                    # Check for charts
                    if hasattr(shape, 'chart'):
                        has_meaningful_content = True
                    
                    if has_meaningful_content:
                        meaningful_content_count += 1
                
                if debug:
                    logger.debug(f"{indent}      📋 Non-decorative shapes with meaningful content: {meaningful_content_count}/{non_decorative_count}")
                
                # If no non-decorative shapes have meaningful content, group might be decorative
                if meaningful_content_count == 0 and decorative_ratio >= 0.5:
                    if debug:
                        logger.debug(f"{indent}      ✅ No meaningful content in remaining shapes, group is {decorative_ratio:.1%} decorative")
                    return True
                
            except Exception as content_error:
                if debug:
                    logger.debug(f"{indent}      ⚠️ Error analyzing group content: {content_error}")
            
            if debug:
                logger.debug(f"{indent}      ❌ Group not decorative: {decorative_ratio:.1%} decorative ratio with meaningful content")
            return False
            
        except Exception as e:
            if debug:
                logger.debug(f"{indent}      ❌ Error analyzing group decorativeness: {e}")
            return False
    
    def set_decorative_flag(self, decorative_shapes: List[PPTXShapeInfo], debug: bool = False) -> int:
        """
        Mark shapes as decorative in the PPTX XML structure using Office 2019+ decorative attribute.
        
        Args:
            decorative_shapes: List of shapes to mark as decorative
            debug: Enable debug logging
            
        Returns:
            int: Number of shapes successfully marked as decorative
        """
        marked_count = 0
        
        # Register decorative namespace if not already done
        try:
            _nsmap["adec"] = "http://schemas.microsoft.com/office/drawing/2017/decorative"
        except:
            pass  # May already be registered
        
        for shape_info in decorative_shapes:
            try:
                if debug:
                    logger.debug(f"🔍 DEBUG: Setting decorative flag for {shape_info.shape_key}")
                
                success = self._set_shape_decorative_xml(shape_info.shape, debug)
                
                if success:
                    marked_count += 1
                    if debug:
                        logger.debug(f"    ✅ Successfully marked {shape_info.shape_key} as decorative")
                else:
                    logger.warning(f"Failed to mark shape as decorative: {shape_info.shape_key} "
                                 f"(type: {shape_info.shape_type_name}, size: {shape_info.width_px}x{shape_info.height_px}px)")
                    if debug:
                        logger.debug(f"    ❌ XML marking failed for {shape_info.shape_key}")
                        # Try to provide more diagnostic info
                        if hasattr(shape_info.shape, '_element'):
                            logger.debug(f"    Shape has _element: {shape_info.shape._element}")
                            logger.debug(f"    Element tag: {getattr(shape_info.shape._element, 'tag', 'no_tag')}")
                        else:
                            logger.debug(f"    Shape has no _element attribute")
                
            except Exception as e:
                logger.warning(f"Error setting decorative flag for {shape_info.shape_key}: {e}")
                continue
        
        logger.info(f"Successfully marked {marked_count}/{len(decorative_shapes)} shapes as decorative")
        return marked_count
    
    def _set_shape_decorative_xml(self, shape: BaseShape, debug: bool = False) -> bool:
        """
        Set the decorative attribute in the shape's XML structure.
        
        Args:
            shape: Shape to mark as decorative
            debug: Enable debug logging
            
        Returns:
            bool: True if successfully set
        """
        success_methods = []
        
        try:
            # Method 1: Try to set decorative attribute on cNvPr element
            if hasattr(shape, '_element') and shape._element is not None:
                element = shape._element
                
                # Find the cNvPr (non-visual properties) element with multiple namespace tries
                cnvpr_elements = []
                try:
                    cnvpr_elements = _safe_xpath(element, './/p:cNvPr | .//pic:cNvPr | .//a:cNvPr',
                                                 namespaces={
                                                     'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                                                     'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
                                                     'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
                                                 })
                except Exception as xpath_error:
                    if debug:
                        logger.debug(f"      XPath query failed: {xpath_error}")
                
                if cnvpr_elements:
                    for i, cnvpr in enumerate(cnvpr_elements):
                        try:
                            # Set the decorative attribute (Office 2019+ feature)
                            decorative_attr = '{http://schemas.microsoft.com/office/drawing/2017/decorative}decorative'
                            cnvpr.set(decorative_attr, '1')
                            success_methods.append(f"cNvPr_element_{i}")
                            
                            if debug:
                                logger.debug(f"      Set decorative='1' on cNvPr element {i}")
                        except Exception as cnvpr_error:
                            if debug:
                                logger.debug(f"      Failed to set decorative on cNvPr {i}: {cnvpr_error}")
                
                # Method 2: Try alternative XML approaches
                try:
                    # Try setting decorative directly on the shape element
                    decorative_attr = '{http://schemas.microsoft.com/office/drawing/2017/decorative}decorative'
                    element.set(decorative_attr, '1')
                    success_methods.append("shape_element")
                    
                    if debug:
                        logger.debug(f"      Set decorative='1' on shape element")
                except Exception as element_error:
                    if debug:
                        logger.debug(f"      Failed to set decorative on shape element: {element_error}")
                
                # Method 3: Fallback - set a custom attribute for tracking
                try:
                    # Set a custom attribute that can be used for identification
                    custom_attr = '{http://schemas.anthropic.com/accessibility/2024}decorative'
                    element.set(custom_attr, '1')
                    success_methods.append("custom_fallback")
                    
                    if debug:
                        logger.debug(f"      Set custom decorative attribute as fallback")
                except Exception as custom_error:
                    if debug:
                        logger.debug(f"      Failed to set custom decorative attribute: {custom_error}")
                
                # Method 4: Try setting ALT text to empty (accessibility best practice for decorative)
                if cnvpr_elements:
                    for i, cnvpr in enumerate(cnvpr_elements):
                        try:
                            cnvpr.set('descr', '')  # Empty ALT text for decorative images
                            success_methods.append(f"empty_alt_text_{i}")
                            
                            if debug:
                                logger.debug(f"      Set empty ALT text on cNvPr element {i} as fallback")
                        except Exception as alt_error:
                            if debug:
                                logger.debug(f"      Failed to set empty ALT text on cNvPr {i}: {alt_error}")
            
            # If we had any success, return True
            if success_methods:
                if debug:
                    logger.debug(f"      Decorative marking succeeded via: {', '.join(success_methods)}")
                return True
            
            if debug:
                logger.debug(f"      No suitable XML element found for decorative marking")
            return False
            
        except Exception as e:
            if debug:
                logger.debug(f"      Error in decorative XML marking: {e}")
            return len(success_methods) > 0  # Return True if we had any success before the error
    
    def _count_all_shapes(self, presentation: Presentation) -> Tuple[int, int]:
        """
        Count total shapes and shapes with text content for coverage reporting.
        
        Args:
            presentation: PowerPoint presentation
            
        Returns:
            Tuple of (total_shapes, shapes_with_content)
        """
        total_shapes = 0
        shapes_with_content = 0
        
        for slide in presentation.slides:
            shapes_on_slide, content_shapes_on_slide = self._count_shapes_recursive(slide.shapes)
            total_shapes += shapes_on_slide
            shapes_with_content += content_shapes_on_slide
        
        return total_shapes, shapes_with_content
    
    def _count_shapes_recursive(self, shapes) -> Tuple[int, int]:
        """
        Recursively count shapes and those with meaningful content.
        
        Args:
            shapes: Collection of shapes to count
            
        Returns:
            Tuple of (total_shapes, shapes_with_content)
        """
        total_shapes = 0
        shapes_with_content = 0
        
        for shape in shapes:
            # Skip images (counted separately)
            if hasattr(shape, 'image') and shape.image:
                continue
            
            total_shapes += 1
            
            # Check if shape has meaningful content
            has_text = hasattr(shape, 'text') and shape.text and shape.text.strip()
            if has_text and len(shape.text.strip()) > 2:
                shapes_with_content += 1
            
            # Recursively count grouped shapes
            if hasattr(shape, 'shapes'):
                group_total, group_content = self._count_shapes_recursive(shape.shapes)
                total_shapes += group_total
                shapes_with_content += group_content
        
        return total_shapes, shapes_with_content
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text content from a slide."""
        text_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                text_parts.append(shape.text.strip())
        
        return " ".join(text_parts)
    
    def _extract_slide_notes(self, slide) -> str:
        """Extract notes from a slide."""
        try:
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text
                return notes_text.strip()
        except Exception as e:
            logger.debug(f"Failed to extract slide notes: {e}")
        
        return ""
    
    def _should_generate_alt_text(self, image_info: PPTXImageInfo, 
                                 image_tracker: defaultdict) -> bool:
        """
        Enhanced determination of whether ALT text should be generated, with strong bias towards 
        educational and technical content.
        
        Args:
            image_info: Image information
            image_tracker: Dictionary tracking image occurrences
            
        Returns:
            bool: True if ALT text should be generated
        """
        # Check configuration-based decorative rules first
        if is_force_decorative_by_filename(image_info.filename, self.config_manager.config):
            logger.debug(f"Image marked as decorative by config rules: {image_info.filename}")
            return False
        
        # Check if decorative detection is disabled
        if not self.skip_decorative:
            return True
        
        # ENHANCED: Priority check for educational/technical content
        # Images with educational content should NEVER be decorative
        is_educational = self._is_educational_content(image_info)
        if is_educational:
            logger.debug(f"Image identified as educational content, generating ALT text: {image_info.filename}")
            return True
        
        # ENHANCED: Refined size-based content detection  
        # Medium to large images are more likely to be meaningful content
        dimensions = (image_info.width_px, image_info.height_px)
        if self._is_content_by_size_and_context(image_info, dimensions):
            logger.debug(f"Image identified as content by size/context analysis: {image_info.filename}")
            return True
        
        # Use the existing heuristic-based decorative detection only for remaining cases
        position = (image_info.left_px, image_info.top_px)
        slide_shapes = []  # Not used by current heuristics
        
        is_decorative, notes = is_decorative_image(
            image_bytes=image_info.image_data,
            image_name=image_info.filename,
            position=position,
            dimensions=dimensions,
            slide_shapes=slide_shapes,
            image_hash=image_info.image_hash,
            image_tracker=image_tracker
        )
        
        if is_decorative:
            logger.debug(f"Image marked as decorative by heuristics: {image_info.filename} - {', '.join(notes)}")
            return False
        
        return True
    
    def _is_educational_content(self, image_info: PPTXImageInfo) -> bool:
        """
        Determine if an image contains educational/technical content that should never be decorative.
        Uses advanced heuristics for scientific, medical, and technical content detection.
        
        Args:
            image_info: Image information including context and properties
            
        Returns:
            bool: True if image contains educational content
        """
        # Combine all available context
        combined_context = f"{image_info.filename} {image_info.slide_text}".lower()
        
        # 1. SCIENTIFIC/MEDICAL CONTENT PATTERNS
        scientific_patterns = [
            # Anatomical and biological
            r'\b(anatomy|anatomical|organ|cell|tissue|muscle|bone|nerve|neuron|brain|heart|lung|kidney|liver)\b',
            r'\b(molecule|molecular|protein|dna|rna|gene|chromosome|mitochondria|nucleus)\b',
            r'\b(bacteria|virus|pathogen|infection|immune|antibody|enzyme|hormone)\b',
            
            # Medical diagnostics and procedures
            r'\b(x-?ray|ct scan|mri|ultrasound|ekg|ecg|radiograph|tomography|imaging|scan)\b',
            r'\b(diagnosis|diagnostic|clinical|medical|surgical|procedure|operation|treatment)\b',
            r'\b(patient|case study|syndrome|disease|disorder|condition|symptom)\b',
            
            # Scientific processes and concepts
            r'\b(mechanism|pathway|process|cycle|reaction|synthesis|metabolism|photosynthesis)\b',
            r'\b(experiment|study|research|analysis|method|procedure|protocol|technique)\b',
            r'\b(hypothesis|theory|model|simulation|calculation|measurement)\b'
        ]
        
        for pattern in scientific_patterns:
            if re.search(pattern, combined_context):
                logger.debug(f"Educational content detected - scientific pattern: {pattern}")
                return True
        
        # 2. TECHNICAL DIAGRAMS AND CHARTS
        technical_patterns = [
            # Diagrams and technical illustrations
            r'\b(diagram|schematic|flowchart|blueprint|circuit|wiring|layout)\b',
            r'\b(graph|chart|plot|data|statistics|results|analysis|trend)\b',
            r'\b(figure|illustration|drawing|sketch|model|3d|cross-section)\b',
            
            # Engineering and technical
            r'\b(engineering|mechanical|electrical|structural|design|system|component)\b',
            r'\b(algorithm|flow|pipeline|architecture|framework|structure)\b',
            r'\b(specification|requirements|standards|protocol|interface)\b'
        ]
        
        for pattern in technical_patterns:
            if re.search(pattern, combined_context):
                logger.debug(f"Educational content detected - technical pattern: {pattern}")
                return True
        
        # 3. EDUCATIONAL CONTEXT INDICATORS
        educational_patterns = [
            r'\b(learn|learning|teach|education|course|lesson|lecture|tutorial)\b',
            r'\b(example|demonstration|illustration|explanation|concept|principle)\b',
            r'\b(objective|goal|overview|summary|conclusion|key point)\b',
            r'\b(step|phase|stage|part [abc]|section|chapter)\b'
        ]
        
        for pattern in educational_patterns:
            if re.search(pattern, combined_context):
                logger.debug(f"Educational content detected - educational pattern: {pattern}")
                return True
        
        # 4. MATHEMATICAL AND QUANTITATIVE CONTENT
        math_patterns = [
            r'\b(equation|formula|calculation|mathematics|statistics|probability)\b',
            r'\b(graph|plot|curve|distribution|correlation|regression|analysis)\b',
            r'\b(data|dataset|measurement|metric|value|parameter|variable)\b'
        ]
        
        for pattern in math_patterns:
            if re.search(pattern, combined_context):
                logger.debug(f"Educational content detected - mathematical pattern: {pattern}")
                return True
        
        # 5. FILENAME-BASED EDUCATIONAL INDICATORS
        filename_lower = image_info.filename.lower()
        educational_filename_patterns = [
            # Common educational image naming
            r'fig(ure)?[_-]?\d+',  # figure1, fig_2, etc.
            r'(table|chart|graph|plot)[_-]?\d*',
            r'(diagram|schematic|flow)[_-]?(chart)?',
            r'(slide|page)[_-]?\d+',
            
            # Scientific naming conventions
            r'(anatomy|medical|clinical|scientific)',
            r'(experiment|study|research|analysis)',
            r'(process|mechanism|pathway|cycle)',
            r'(structure|system|model|simulation)'
        ]
        
        for pattern in educational_filename_patterns:
            if re.search(pattern, filename_lower):
                logger.debug(f"Educational content detected - filename pattern: {pattern}")
                return True
        
        # 6. SLIDE POSITION AND CONTEXT ANALYSIS
        # Images in educational slides are often centrally positioned and substantial
        if image_info.slide_text:
            slide_context = image_info.slide_text.lower()
            
            # Check for surrounding educational context
            if len(slide_context) > 50:  # Substantial educational content on slide
                # Common educational slide patterns
                educational_slide_indicators = [
                    r'\b(definition|overview|introduction|explanation)\b',
                    r'\b(example|case|instance|application)\b',
                    r'\b(compare|contrast|difference|similarity)\b',
                    r'\b(function|role|purpose|importance)\b',
                    r'\b(feature|characteristic|property|attribute)\b'
                ]
                
                for pattern in educational_slide_indicators:
                    if re.search(pattern, slide_context):
                        logger.debug(f"Educational content detected - slide context pattern: {pattern}")
                        return True
        
        return False
    
    def _is_content_by_size_and_context(self, image_info: PPTXImageInfo, dimensions: Tuple[int, int]) -> bool:
        """
        Determine if an image is likely content based on enhanced size analysis and context.
        
        Args:
            image_info: Image information
            dimensions: (width, height) in pixels
            
        Returns:
            bool: True if image is likely meaningful content
        """
        width, height = dimensions
        
        # Ignore zero or invalid dimensions
        if width <= 0 or height <= 0:
            return False
        
        min_dimension = min(width, height)
        max_dimension = max(width, height)
        
        # 1. ENHANCED SIZE THRESHOLDS for educational content
        # Medium-to-large images are more likely to be content in educational presentations
        
        # Large images (>300px) are almost certainly content
        if min_dimension > 300:
            logger.debug(f"Content by size - large image: {width}x{height}px")
            return True
        
        # Medium images (>150px) are likely content, especially in educational context
        if min_dimension > 150:
            # Additional checks for medium-sized images
            
            # Check if image is reasonably proportioned (not extremely thin)
            aspect_ratio = max_dimension / min_dimension if min_dimension > 0 else float('inf')
            if aspect_ratio < 5:  # Not extremely thin
                logger.debug(f"Content by size - medium proportioned image: {width}x{height}px, ratio: {aspect_ratio:.1f}")
                return True
            
            # Medium images with educational context are likely content
            if image_info.slide_text and len(image_info.slide_text) > 30:
                educational_context_indicators = [
                    'figure', 'diagram', 'chart', 'graph', 'illustration',
                    'example', 'model', 'structure', 'process', 'system'
                ]
                
                slide_text_lower = image_info.slide_text.lower()
                for indicator in educational_context_indicators:
                    if indicator in slide_text_lower:
                        logger.debug(f"Content by size - medium image with educational context: {indicator}")
                        return True
        
        # 2. CONTEXTUAL SIZE ANALYSIS
        # Images that are prominent on slide (good size relative to slide)
        # Standard slide dimensions are approximately 960x720px
        slide_area_estimate = 960 * 720
        image_area = width * height
        area_ratio = image_area / slide_area_estimate
        
        # Images taking up significant slide space (>5%) are likely content
        if area_ratio > 0.05:  # More than 5% of slide area
            logger.debug(f"Content by size - significant slide coverage: {area_ratio:.1%} of slide")
            return True
        
        # Images with reasonable aspect ratio and decent size are likely content
        if min_dimension > 100:  # Reasonable minimum size
            aspect_ratio = max_dimension / min_dimension if min_dimension > 0 else float('inf')
            if aspect_ratio < 3:  # Not extremely elongated
                # Check for central positioning (content images are often centered)
                if image_info.left_px > 50 and image_info.top_px > 50:  # Not in corner
                    logger.debug(f"Content by size - well-positioned medium image: {width}x{height}px")
                    return True
        
        return False
    
    def _is_educational_shape(self, shape_info: PPTXShapeInfo, debug: bool = False, indent: str = "") -> bool:
        """
        Determine if a shape contains educational content and should never be considered decorative.
        
        Args:
            shape_info: Shape information
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            bool: True if shape contains educational content
        """
        # Check shape name for educational indicators
        shape_name = shape_info.shape_name.lower() if shape_info.shape_name else ""
        
        educational_shape_names = [
            'diagram', 'chart', 'graph', 'figure', 'illustration',
            'flowchart', 'schematic', 'model', 'structure', 'process',
            'table', 'data', 'analysis', 'result', 'example'
        ]
        
        for name_pattern in educational_shape_names:
            if name_pattern in shape_name:
                if debug:
                    logger.debug(f"{indent}        🎓 Educational shape name: '{shape_name}' contains '{name_pattern}'")
                return True
        
        # Check text content for educational patterns
        if shape_info.has_text and shape_info.text_content:
            text_lower = shape_info.text_content.lower()
            
            # Educational text patterns
            educational_text_patterns = [
                r'\b(figure|fig|diagram|chart|graph|table)\s*\d+',  # Figure 1, Chart 2, etc.
                r'\b(step|phase|stage)\s*\d+',  # Step 1, Phase 2, etc.
                r'\b(example|case|instance)\s*\d*',
                r'\b(definition|explain|describe|illustrate)\b',
                r'\b(process|mechanism|pathway|cycle)\b',
                r'\b(structure|component|element|part)\b',
                r'\b(analysis|result|conclusion|finding)\b'
            ]
            
            for pattern in educational_text_patterns:
                if re.search(pattern, text_lower):
                    if debug:
                        logger.debug(f"{indent}        🎓 Educational text pattern: '{pattern}' in '{text_lower[:50]}'")
                    return True
        
        # Check slide context for educational content
        if self._has_educational_context(shape_info):
            # If shape is in educational context and has reasonable size, it might be educational
            min_dim = min(shape_info.width_px, shape_info.height_px)
            if min_dim > 20:  # Not tiny
                if debug:
                    logger.debug(f"{indent}        🎓 Educational context with reasonable size: {min_dim}px")
                return True
        
        return False
    
    def _has_educational_context(self, shape_info: PPTXShapeInfo) -> bool:
        """
        Check if a shape exists in an educational context based on slide content.
        
        Args:
            shape_info: Shape information including slide context
            
        Returns:
            bool: True if shape is in educational context
        """
        slide_context = shape_info.slide_text.lower() if shape_info.slide_text else ""
        
        # Quick educational context indicators
        educational_indicators = [
            'learn', 'teach', 'education', 'course', 'lesson', 'lecture',
            'objective', 'concept', 'principle', 'theory', 'method',
            'example', 'demonstration', 'illustration', 'explanation',
            'diagram', 'figure', 'chart', 'graph', 'table', 'data',
            'analysis', 'study', 'research', 'experiment', 'result',
            'process', 'mechanism', 'structure', 'function', 'system'
        ]
        
        # Check for multiple educational indicators (stronger signal)
        indicator_count = sum(1 for indicator in educational_indicators if indicator in slide_context)
        
        return indicator_count >= 2  # At least 2 educational indicators on slide
    
    def _is_text_box_with_content(self, shape_info: PPTXShapeInfo, debug: bool = False, indent: str = "") -> bool:
        """
        Determine if a shape is a text box with meaningful content that should never be decorative.
        
        Args:
            shape_info: Shape information
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            bool: True if shape is a text box with meaningful content
        """
        # Check if this is a TEXT_BOX shape type
        if hasattr(MSO_SHAPE_TYPE, 'TEXT_BOX') and shape_info.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # Any text box with content should not be decorative
            if shape_info.has_text and shape_info.text_content:
                text_length = len(shape_info.text_content.strip())
                
                # Even short text in text boxes can be meaningful (labels, captions, etc.)
                if text_length > 0:
                    if debug:
                        logger.debug(f"{indent}        ✅ TEXT BOX DETECTED: {text_length} chars: '{shape_info.text_content[:40]}...'")
                    return True
                elif debug:
                    logger.debug(f"{indent}        ❌ Text box empty or whitespace only")
            elif debug:
                logger.debug(f"{indent}        ❌ Text box has no text content (has_text={shape_info.has_text})")
        elif debug:
            textbox_available = hasattr(MSO_SHAPE_TYPE, 'TEXT_BOX')
            is_textbox = shape_info.shape_type == MSO_SHAPE_TYPE.TEXT_BOX if textbox_available else False
            logger.debug(f"{indent}        ❌ Not a text box (TEXT_BOX available: {textbox_available}, is_textbox: {is_textbox}, shape_type: {shape_info.shape_type})")
        
        # Check for shapes that function as text boxes (have text_frame)
        if hasattr(shape_info.shape, 'text_frame') and shape_info.has_text:
            text_content = shape_info.text_content.strip()
            if debug:
                logger.debug(f"{indent}        🔍 Checking text frame content: '{text_content[:40]}...'")
            
            # Functional text elements should not be decorative
            functional_text_patterns = [
                r'\b(title|header|caption|label|note|description)\b',
                r'\b(step|phase|instruction|guideline|tip)\b', 
                r'\b(name|date|author|source|reference)\b',
                r'\b(copyright|\u00a9|trademark|\u2122)\b',
                r'\d+',  # Numbers often indicate functional content
                r'[a-zA-Z]{3,}',  # Words of 3+ characters are likely meaningful
            ]
            
            text_lower = text_content.lower()
            for pattern in functional_text_patterns:
                if re.search(pattern, text_lower):
                    if debug:
                        logger.debug(f"{indent}        ✅ FUNCTIONAL TEXT DETECTED: Pattern '{pattern}' in: '{text_content[:40]}...'")
                    return True
                    
            # Any text content longer than 1 character in a text-capable shape is likely functional
            if len(text_content) > 1:
                if debug:
                    logger.debug(f"{indent}        ✅ MULTI-CHAR TEXT DETECTED: '{text_content[:40]}...'")
                return True
            elif debug:
                logger.debug(f"{indent}        ❌ Text content too short or pattern not matched: '{text_content}'")
        elif debug:
            has_text_frame = hasattr(shape_info.shape, 'text_frame')
            logger.debug(f"{indent}        ❌ No qualifying text frame (has_text_frame: {has_text_frame}, has_text: {shape_info.has_text})")
        
        if debug:
            logger.debug(f"{indent}        ❌ NOT A TEXT BOX: No qualifying text content found")
        return False
    
    def _is_title_or_header_area(self, shape_info: PPTXShapeInfo, debug: bool = False, indent: str = "") -> bool:
        """
        Determine if a shape is in a title or header area and should never be decorative.
        
        Args:
            shape_info: Shape information
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            bool: True if shape is in title/header area
        """
        # Check shape name for title/header indicators
        shape_name = shape_info.shape_name.lower() if shape_info.shape_name else ""
        
        title_name_patterns = [
            'title', 'header', 'heading', 'banner', 'top', 'head',
            'slide title', 'presentation title', 'main title', 'subtitle'
        ]
        
        for pattern in title_name_patterns:
            if pattern in shape_name:
                if debug:
                    logger.debug(f"{indent}        ✅ TITLE NAME DETECTED: Pattern '{pattern}' in: '{shape_name}'")
                return True
        
        if debug and shape_name:
            logger.debug(f"{indent}        ❌ No title name patterns in: '{shape_name}'")
        
        # Check position-based title detection
        # Title areas are typically in the upper portion of slides
        slide_height_estimate = 720  # Standard slide height
        title_area_threshold = slide_height_estimate * 0.25  # Top 25% of slide
        
        if shape_info.top_px <= title_area_threshold:
            if debug:
                logger.debug(f"{indent}        🔍 Shape in title area (y={shape_info.top_px} <= {title_area_threshold})")
            
            # Large shapes in title area are likely titles/headers
            if shape_info.width_px > 200 and shape_info.height_px > 30:
                if debug:
                    logger.debug(f"{indent}        ✅ TITLE AREA DETECTED: Large shape {shape_info.width_px}x{shape_info.height_px}px at y={shape_info.top_px}")
                return True
            elif debug:
                logger.debug(f"{indent}        ❌ Shape too small for title area: {shape_info.width_px}x{shape_info.height_px}px (need >200x30)")
            
            # Any text content in the title area is likely a title/header
            if shape_info.has_text and len(shape_info.text_content.strip()) > 2:
                # Check if text looks like a title
                text_content = shape_info.text_content.strip()
                if debug:
                    logger.debug(f"{indent}        🔍 Checking title text patterns in: '{text_content[:40]}...'")
                
                title_text_indicators = [
                    # Typical title patterns
                    r'^[A-Z][^.!?]*$',  # Starts with capital, no sentence punctuation
                    r'^(Chapter|Section|Part|Slide|Lesson|Topic)\s+\d+',
                    r'^(Introduction|Overview|Summary|Conclusion|Objectives?)',
                    r':\s*$',  # Ends with colon (common in titles)
                    # Short phrases that are likely titles
                    lambda t: len(t.split()) <= 8 and len(t) <= 60,  # Short title-like text
                ]
                
                text_lower = text_content.lower()
                for indicator in title_text_indicators:
                    try:
                        if callable(indicator):
                            if indicator(text_content):
                                if debug:
                                    logger.debug(f"{indent}        ✅ TITLE TEXT DETECTED: Short title pattern in: '{text_content[:40]}...'")
                                return True
                        elif re.search(indicator, text_content):
                            if debug:
                                logger.debug(f"{indent}        ✅ TITLE TEXT DETECTED: Pattern '{indicator}' in: '{text_content[:40]}...'")
                            return True
                    except:
                        continue
                        
                if debug:
                    logger.debug(f"{indent}        ❌ No title text patterns matched in: '{text_content[:40]}...'")
            elif debug:
                has_text_info = f"has_text={shape_info.has_text}"
                text_len = len(shape_info.text_content.strip()) if shape_info.has_text else 0
                logger.debug(f"{indent}        ❌ No qualifying text in title area ({has_text_info}, len={text_len})")
        elif debug:
            logger.debug(f"{indent}        ❌ Shape not in title area (y={shape_info.top_px} > {title_area_threshold})")
        
        if debug:
            logger.debug(f"{indent}        ❌ NOT A TITLE/HEADER: No qualifying criteria met")
        return False
    
    def _get_auto_shape_details(self, shape, debug: bool = False, indent: str = "") -> str:
        """
        Get detailed information about AUTO_SHAPE subtypes.
        
        Args:
            shape: Shape object to analyze
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            str: Description of auto shape details
        """
        details = []
        
        try:
            # Try to get auto_shape_type
            if hasattr(shape, 'auto_shape_type'):
                auto_shape_type = shape.auto_shape_type
                if auto_shape_type is not None:
                    # Try to get the name of the auto shape type
                    try:
                        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                        for attr_name in dir(MSO_AUTO_SHAPE_TYPE):
                            if not attr_name.startswith('_') and getattr(MSO_AUTO_SHAPE_TYPE, attr_name) == auto_shape_type:
                                details.append(attr_name)
                                break
                        else:
                            details.append(f"type_{auto_shape_type}")
                    except ImportError:
                        details.append(f"type_{auto_shape_type}")
            
            # Check for common geometric properties
            if hasattr(shape, 'adjustments'):
                try:
                    adj_count = len(shape.adjustments) if shape.adjustments else 0
                    if adj_count > 0:
                        details.append(f"{adj_count}_adjustments")
                except:
                    pass
                    
            # Check shape geometry hints from name
            shape_name = getattr(shape, 'name', '').lower()
            geometric_hints = {
                'rectangle': 'RECTANGLE', 'rect': 'RECTANGLE', 'square': 'RECTANGLE',
                'circle': 'OVAL', 'oval': 'OVAL', 'ellipse': 'OVAL',
                'triangle': 'TRIANGLE', 'arrow': 'ARROW', 'star': 'STAR',
                'diamond': 'DIAMOND', 'pentagon': 'PENTAGON', 'hexagon': 'HEXAGON'
            }
            
            for hint, shape_type in geometric_hints.items():
                if hint in shape_name:
                    details.append(f"named_{shape_type}")
                    break
            
            if debug and details:
                logger.debug(f"{indent}        AUTO_SHAPE details: {', '.join(details)}")
                
        except Exception as e:
            if debug:
                logger.debug(f"{indent}        Error getting AUTO_SHAPE details: {e}")
        
        return ', '.join(details) if details else "basic"
    
    def _get_connector_details(self, shape, debug: bool = False, indent: str = "") -> str:
        """
        Get detailed information about CONNECTOR shapes.
        
        Args:
            shape: Shape object to analyze
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            str: Description of connector details
        """
        details = []
        
        try:
            # Try to get connector_type
            if hasattr(shape, 'connector_type'):
                connector_type = shape.connector_type
                if connector_type is not None:
                    try:
                        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
                        for attr_name in dir(MSO_CONNECTOR_TYPE):
                            if not attr_name.startswith('_') and getattr(MSO_CONNECTOR_TYPE, attr_name) == connector_type:
                                details.append(attr_name)
                                break
                        else:
                            details.append(f"type_{connector_type}")
                    except ImportError:
                        details.append(f"type_{connector_type}")
            
            # Check connection points
            if hasattr(shape, 'begin_connect'):
                begin_connected = shape.begin_connect is not None
                details.append(f"begin_{'connected' if begin_connected else 'free'}")
            
            if hasattr(shape, 'end_connect'):
                end_connected = shape.end_connect is not None
                details.append(f"end_{'connected' if end_connected else 'free'}")
            
            if debug and details:
                logger.debug(f"{indent}        CONNECTOR details: {', '.join(details)}")
                
        except Exception as e:
            if debug:
                logger.debug(f"{indent}        Error getting CONNECTOR details: {e}")
        
        return ', '.join(details) if details else "basic"
    
    def _get_line_details(self, shape, debug: bool = False, indent: str = "") -> str:
        """
        Get detailed information about LINE shapes.
        
        Args:
            shape: Shape object to analyze
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            str: Description of line details
        """
        details = []
        
        try:
            # Check if it behaves like a connector
            if hasattr(shape, 'begin_connect') and hasattr(shape, 'end_connect'):
                details.append("connector_style")
            
            # Check line properties
            if hasattr(shape, 'line'):
                try:
                    line = shape.line
                    if hasattr(line, 'color'):
                        details.append("styled")
                    if hasattr(line, 'width'):
                        width = getattr(line, 'width', None)
                        if width:
                            width_pt = width.pt if hasattr(width, 'pt') else width
                            details.append(f"width_{width_pt}pt")
                except:
                    pass
            
            # Check for arrow properties
            shape_name = getattr(shape, 'name', '').lower()
            if 'arrow' in shape_name:
                details.append("arrow")
            
            if debug and details:
                logger.debug(f"{indent}        LINE details: {', '.join(details)}")
                
        except Exception as e:
            if debug:
                logger.debug(f"{indent}        Error getting LINE details: {e}")
        
        return ', '.join(details) if details else "basic"
    
    def _detect_shape_by_properties(self, shape, debug: bool = False, indent: str = "") -> str:
        """
        Attempt to detect shape type by analyzing properties when standard type detection fails.
        
        Args:
            shape: Shape object to analyze
            debug: Enable debug logging
            indent: Logging indentation
            
        Returns:
            str: Detected shape type description
        """
        try:
            shape_class = type(shape).__name__
            
            # Check for specific shape classes
            if 'Picture' in shape_class:
                return "PICTURE"
            elif 'Shape' in shape_class and hasattr(shape, 'text_frame'):
                return "TEXT_BOX"
            elif 'GroupShape' in shape_class:
                return "GROUP"
            elif 'Connector' in shape_class:
                return "CONNECTOR"
            elif 'FreeformBuilder' in shape_class or 'Freeform' in shape_class:
                return "FREEFORM"
            
            # Property-based detection
            properties = []
            
            if hasattr(shape, 'image') and shape.image:
                properties.append("has_image")
            
            if hasattr(shape, 'chart'):
                properties.append("has_chart")
                
            if hasattr(shape, 'text_frame') and hasattr(shape, 'text'):
                properties.append("has_text")
                
            if hasattr(shape, 'shapes'):
                properties.append("has_shapes")
            
            if hasattr(shape, 'adjustments'):
                properties.append("adjustable")
                
            if hasattr(shape, 'line'):
                properties.append("has_line")
                
            if hasattr(shape, 'fill'):
                properties.append("has_fill")
            
            # XML-based detection as fallback
            xml_detection = ""
            if hasattr(shape, '_element') and shape._element is not None:
                try:
                    xml_tag = shape._element.tag
                    if 'rect' in xml_tag.lower():
                        xml_detection = "XML_RECTANGLE"
                    elif 'ellipse' in xml_tag.lower():
                        xml_detection = "XML_ELLIPSE"
                    elif 'line' in xml_tag.lower():
                        xml_detection = "XML_LINE"
                    elif 'path' in xml_tag.lower():
                        xml_detection = "XML_PATH"
                except:
                    pass
            
            detection_parts = []
            if xml_detection:
                detection_parts.append(xml_detection)
            if properties:
                detection_parts.append("_".join(properties[:3]))  # Limit to avoid long names
            
            if debug:
                logger.debug(f"{indent}        Property detection: class={shape_class}, props={properties}, xml={xml_detection}")
            
            return "_".join(detection_parts) if detection_parts else "UNKNOWN"
                
        except Exception as e:
            if debug:
                logger.debug(f"{indent}        Error in property detection: {e}")
            return "DETECTION_ERROR"
    
    def _generate_alt_text_for_image(self, image_info: PPTXImageInfo) -> Optional[str]:
        """
        Generate ALT text for a single image using the existing ALT text generator.
        
        Args:
            image_info: Image information
            
        Returns:
            Generated ALT text or None if generation failed
        """
        try:
            # Normalize image format before processing
            try:
                normalized_image_data = self._normalize_image_format(image_info.image_data, image_info.filename)
                
                # Save normalized image to temporary file for ALT text generation
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                    temp_file.write(normalized_image_data)
                    temp_image_path = temp_file.name
                    
            except Exception as norm_error:
                # Check if this is a vector format conversion failure OR any WMF/EMF processing failure
                if (("Vector format conversion failed" in str(norm_error) or 
                     "format normalization failed" in str(norm_error) or
                     "Cannot process" in str(norm_error)) and 
                    image_info.filename.lower().endswith(('.wmf', '.emf'))):
                    # Generate contextual fallback ALT text
                    format_name = "WMF" if image_info.filename.lower().endswith('.wmf') else "EMF"
                    logger.info(f"Generating contextual fallback ALT for {image_info.filename}")
                    return self._generate_vector_fallback_alt(image_info, format_name)
                else:
                    # For other normalization failures, re-raise
                    raise norm_error
            
            try:
                # Build context for better ALT text generation
                context = self._build_generation_context(image_info)
                
                # Determine appropriate prompt type based on content
                prompt_type = self._determine_prompt_type(image_info)
                
                # Generate ALT text using the configured generator
                alt_text = self.alt_generator.generate_alt_text(
                    image_path=temp_image_path,
                    prompt_type=prompt_type,
                    context=context
                )
                
                return alt_text
                
            finally:
                # Clean up temporary file
                try:
                    os.unlink(temp_image_path)
                except OSError:
                    pass  # File cleanup failure is not critical
        
        except Exception as e:
            logger.error(f"Failed to generate ALT text for {image_info.image_key}: {e}")
            return None
    
    def _normalize_image_format(self, image_data: bytes, filename: str, debug: bool = False) -> bytes:
        """
        Normalize image format to prevent TIFF/WMF/EMF crashes with LLaVA.
        Converts problematic formats to PNG and optionally resizes large images.
        
        Args:
            image_data: Original image data
            filename: Original filename for format detection
            debug: Enable debug logging
            
        Returns:
            Normalized image data (PNG format)
        """
        try:
            if not PIL_AVAILABLE:
                logger.warning("PIL not available - cannot normalize image format")
                return image_data
            
            # Detect problematic formats that crash LLaVA
            filename_lower = filename.lower()
            is_problematic_format = (
                filename_lower.endswith(('.tiff', '.tif', '.wmf', '.emf')) or
                b'TIFF' in image_data[:100] or
                b'WMF' in image_data[:100] or 
                b'EMF' in image_data[:100]
            )
            
            if debug:
                logger.debug(f"Image format check: {filename} -> problematic: {is_problematic_format}")
            
            # Try to open the image with PIL
            try:
                with io.BytesIO(image_data) as img_buffer:
                    img = Image.open(img_buffer)
                    original_format = img.format
                    original_size = img.size
                    
                    if debug:
                        logger.debug(f"Original image: format={original_format}, size={original_size}")
                    
                    # Convert to RGB if needed (handles RGBA, CMYK, etc.)
                    if img.mode not in ('RGB', 'L'):  # L for grayscale
                        if debug:
                            logger.debug(f"Converting from {img.mode} to RGB")
                        img = img.convert('RGB')
                    
                    # Check if image is very large and resize if configured
                    max_dimension = self.processing_config.get('max_image_dimension', 1600)
                    if max(original_size) > max_dimension:
                        # Calculate new size maintaining aspect ratio
                        width, height = original_size
                        if width > height:
                            new_width = max_dimension
                            new_height = int(height * (max_dimension / width))
                        else:
                            new_height = max_dimension
                            new_width = int(width * (max_dimension / height))
                        
                        img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                        if debug:
                            logger.debug(f"Resized image from {original_size} to {img.size}")
                    
                    # Save as PNG
                    output_buffer = io.BytesIO()
                    img.save(output_buffer, format='PNG', optimize=True)
                    normalized_data = output_buffer.getvalue()
                    
                    if debug:
                        logger.debug(f"Normalized: {len(image_data)} -> {len(normalized_data)} bytes")
                    
                    return normalized_data
                    
            except Exception as pil_error:
                logger.warning(f"PIL failed to process image {filename}: {pil_error}")
                
                # For WMF/EMF formats, try external converters
                if filename_lower.endswith(('.wmf', '.emf')):
                    logger.info(f"Attempting external conversion for {filename}")
                    try:
                        converted_data = self._convert_vector_image_external(image_data, filename, debug)
                        if converted_data:
                            logger.info(f"External conversion successful for {filename}: {len(converted_data)} bytes")
                            return converted_data
                        else:
                            logger.warning(f"External conversion returned no data for {filename}")
                    except Exception as ext_error:
                        logger.warning(f"External conversion failed for {filename}: {ext_error}")
                        # Continue to contextual fallback instead of failing
                
                # If PIL fails but it's a problematic format, we can't process it
                if is_problematic_format:
                    logger.error(f"Cannot process problematic format {filename} - all conversion methods failed")
                    
                    # For WMF/EMF, provide contextual fallback instead of failing completely
                    if filename_lower.endswith(('.wmf', '.emf')):
                        logger.info(f"Using contextual fallback for unsupported vector format {filename}")
                        # Signal that contextual fallback should be used
                        raise Exception(f"Vector format conversion failed: {filename}")
                    
                    raise Exception(f"Unsupported image format: {filename}")
                
                # For other formats that PIL can't handle, return original data
                return image_data
                
        except Exception as e:
            logger.error(f"Image normalization failed for {filename}: {e}")
            # If normalization fails and it's a problematic format, we should fail
            if filename.lower().endswith(('.tiff', '.tif', '.wmf', '.emf')):
                raise Exception(f"Cannot process {filename}: format normalization failed")
            # Otherwise return original data and hope for the best
            return image_data
    
    def _convert_vector_image_external(self, image_data: bytes, filename: str, debug: bool = False) -> bytes:
        """
        Convert WMF/EMF images using external tools when PIL fails.
        Tries multiple conversion strategies in order of preference.
        
        Args:
            image_data: Original WMF/EMF image data
            filename: Original filename for logging
            debug: Enable debug logging
            
        Returns:
            Converted PNG image data
            
        Raises:
            Exception: If all conversion methods fail
        """
        import subprocess
        import tempfile
        import shutil
        
        # Create temporary files for input and output
        input_suffix = '.wmf' if filename.lower().endswith('.wmf') else '.emf'
        
        with tempfile.NamedTemporaryFile(suffix=input_suffix, delete=False) as input_file:
            input_file.write(image_data)
            input_path = input_file.name
            
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as output_file:
            output_path = output_file.name
        
        try:
            # Strategy 1: Inkscape (best quality for vector formats)
            if shutil.which('inkscape'):
                try:
                    logger.info(f"Trying Inkscape conversion for {filename}")
                    
                    cmd = [
                        'inkscape',
                        '--export-type=png',
                        '--export-dpi=300',
                        '--export-background=white',
                        '--export-background-opacity=1.0',
                        '--export-filename', output_path,
                        input_path
                    ]
                    
                    logger.info(f"Running command: {' '.join(cmd)}")
                    
                    result = subprocess.run(
                        cmd, 
                        capture_output=True, 
                        text=True, 
                        timeout=30,
                        check=False
                    )
                    
                    logger.info(f"Inkscape result: returncode={result.returncode}")
                    if result.stdout:
                        logger.info(f"Inkscape stdout: {result.stdout}")
                    if result.stderr:
                        logger.info(f"Inkscape stderr: {result.stderr}")
                    
                    if result.returncode == 0 and os.path.exists(output_path):
                        with open(output_path, 'rb') as f:
                            converted_data = f.read()
                        logger.info(f"Inkscape output file size: {len(converted_data)} bytes")
                        if len(converted_data) > 100:  # Sanity check for valid PNG
                            logger.info(f"Inkscape conversion successful: {len(converted_data)} bytes")
                            return converted_data
                    else:
                        logger.warning(f"Inkscape failed or no output file: returncode={result.returncode}, exists={os.path.exists(output_path)}")
                        
                except subprocess.TimeoutExpired:
                    logger.warning(f"Inkscape conversion timed out for {filename}")
                except Exception as e:
                    logger.warning(f"Inkscape conversion error: {e}")
            else:
                logger.info("Inkscape not available")
            
            # Strategy 2: ImageMagick/GraphicsMagick
            magick_commands = ['magick', 'convert']  # Try both names
            for magick_cmd in magick_commands:
                if shutil.which(magick_cmd):
                    try:
                        if debug:
                            logger.debug(f"Trying {magick_cmd} conversion for {filename}")
                        
                        cmd = [
                            magick_cmd,
                            '-density', '300',
                            '-background', 'white',
                            '-alpha', 'remove',
                            input_path,
                            output_path
                        ]
                        
                        result = subprocess.run(
                            cmd,
                            capture_output=True,
                            text=True,
                            timeout=30,
                            check=False
                        )
                        
                        if result.returncode == 0 and os.path.exists(output_path):
                            with open(output_path, 'rb') as f:
                                converted_data = f.read()
                            if len(converted_data) > 100:
                                if debug:
                                    logger.debug(f"{magick_cmd} conversion successful: {len(converted_data)} bytes")
                                return converted_data
                        elif debug:
                            logger.debug(f"{magick_cmd} failed: {result.stderr}")
                            
                    except subprocess.TimeoutExpired:
                        logger.warning(f"{magick_cmd} conversion timed out for {filename}")
                    except Exception as e:
                        if debug:
                            logger.debug(f"{magick_cmd} conversion error: {e}")
                    break  # Don't try other magick commands if one is found
            
            # Strategy 3: LibreOffice headless (last resort)
            if shutil.which('libreoffice'):
                try:
                    if debug:
                        logger.debug(f"Trying LibreOffice conversion for {filename}")
                    
                    # LibreOffice needs a directory to work in
                    with tempfile.TemporaryDirectory() as temp_dir:
                        temp_input = os.path.join(temp_dir, f"input{input_suffix}")
                        shutil.copy2(input_path, temp_input)
                        
                        cmd = [
                            'libreoffice',
                            '--headless',
                            '--convert-to', 'png',
                            '--outdir', temp_dir,
                            temp_input
                        ]
                        
                        result = subprocess.run(
                            cmd,
                            capture_output=True,
                            text=True,
                            timeout=60,
                            check=False
                        )
                        
                        # LibreOffice creates input.png
                        lo_output = os.path.join(temp_dir, "input.png")
                        if result.returncode == 0 and os.path.exists(lo_output):
                            with open(lo_output, 'rb') as f:
                                converted_data = f.read()
                            if len(converted_data) > 100:
                                if debug:
                                    logger.debug(f"LibreOffice conversion successful: {len(converted_data)} bytes")
                                return converted_data
                        elif debug:
                            logger.debug(f"LibreOffice failed: {result.stderr}")
                            
                except subprocess.TimeoutExpired:
                    logger.warning(f"LibreOffice conversion timed out for {filename}")
                except Exception as e:
                    if debug:
                        logger.debug(f"LibreOffice conversion error: {e}")
            
            # All external converters failed
            raise Exception(f"All external converters failed for {filename}")
            
        finally:
            # Clean up temporary files
            try:
                os.unlink(input_path)
                os.unlink(output_path)
            except OSError:
                pass  # Ignore cleanup failures
    
    def _generate_vector_fallback_alt(self, image_info: PPTXImageInfo, format_name: str, debug: bool = False) -> str:
        """
        Generate contextual fallback ALT text for vector images (WMF/EMF) that can't be converted.
        
        Args:
            image_info: Image information
            format_name: Format name (WMF or EMF)
            debug: Enable debug logging
            
        Returns:
            Contextual ALT text describing the vector image
        """
        # DETERMINISTIC WMF FALLBACK GUARD: Use metadata to prevent re-application
        if not hasattr(image_info, 'meta'):
            image_info.meta = {}
        
        fallback_key = f"wmf_fallback_applied_{format_name.lower()}"
        if image_info.meta.get(fallback_key):
            cached_text = image_info.meta.get(f"wmf_fallback_text_{format_name.lower()}")
            if debug:
                logger.debug(f"🔒 WMF fallback already computed for {image_info.filename}, returning cached")
            return cached_text or f"Vector diagram ({format_name} format)"
        
        context_parts = []
        
        # Get slide context
        slide_text = image_info.slide_text.strip() if image_info.slide_text else ""
        
        # Determine image type based on context and dimensions
        width_px = image_info.width_px or 0
        height_px = image_info.height_px or 0
        
        # Analyze context for diagram type hints
        context_lower = slide_text.lower()
        
        # Scientific/technical diagram indicators
        if any(word in context_lower for word in [
            'diagram', 'chart', 'graph', 'plot', 'circuit', 'schematic', 
            'flowchart', 'equation', 'formula', 'model', 'structure',
            'membrane', 'potential', 'channel', 'protein', 'cell'
        ]):
            if any(word in context_lower for word in ['membrane', 'potential', 'channel', 'cell', 'protein']):
                diagram_type = "scientific diagram"
            elif any(word in context_lower for word in ['circuit', 'electrical', 'voltage', 'current']):
                diagram_type = "electrical circuit diagram"  
            elif any(word in context_lower for word in ['flow', 'process', 'step']):
                diagram_type = "process flow diagram"
            else:
                diagram_type = "technical diagram"
        else:
            # Generic based on dimensions
            aspect_ratio = width_px / height_px if height_px > 0 else 1
            if aspect_ratio > 1.5:
                diagram_type = "horizontal diagram"
            elif aspect_ratio < 0.67:
                diagram_type = "vertical diagram"
            else:
                diagram_type = "diagram"
        
        # Build contextual description
        context_parts.append(f"Vector {diagram_type} ({format_name} format)")
        
        # Add size information
        if width_px and height_px:
            context_parts.append(f"({width_px}×{height_px} pixels)")
        
        # Add slide context if meaningful
        if slide_text:
            # Extract key terms from slide text (first meaningful sentence/phrase)
            slide_words = slide_text.split()[:15]  # First 15 words
            clean_text = ' '.join(slide_words)
            if len(clean_text) > 100:
                clean_text = clean_text[:97] + "..."
            context_parts.append(f"related to: {clean_text}")
        
        # Add hint about format limitation (avoid duplication)
        alt_text = ' '.join(context_parts)
        format_note = f". Note: Original {format_name} vector image could not be processed for detailed analysis."
        
        # DEDUPLICATION: Check if this note was already added to avoid "WMF format... WMF format" duplication
        if not alt_text.endswith(format_note.strip()) and format_note.strip() not in alt_text:
            alt_text += format_note
        
        # Additional deduplication: Remove any double format references
        alt_text = self._deduplicate_format_references(alt_text, format_name)
        
        # CACHE THE RESULT to prevent re-computation
        image_info.meta[fallback_key] = True
        image_info.meta[f"wmf_fallback_text_{format_name.lower()}"] = alt_text
        
        if debug:
            logger.debug(f"Generated and cached vector fallback ALT: {alt_text}")
        
        return alt_text
    
    def _check_element_bypass(self, visual_element: PPTXVisualElement) -> Optional[str]:
        """
        Check if a visual element should be bypassed and return the reason.
        
        Args:
            visual_element: Visual element to check
            
        Returns:
            Bypass reason string if element should be bypassed, None otherwise
        """
        # Connectors and lines get special handling but are not truly bypassed
        # They get descriptive text instead of AI-generated descriptions
        if visual_element.element_type in ['connector', 'line']:
            return f"Using descriptive text for {visual_element.element_type} (no AI analysis needed)"
        
        # Very small elements that might be decorative artifacts
        if visual_element.width_px < 5 or visual_element.height_px < 5:
            return f"Element too small ({visual_element.width_px}x{visual_element.height_px}px)"
        
        # Elements with no visual content
        if not hasattr(visual_element, 'shape') or visual_element.shape is None:
            return "No shape data available"
        
        return None
    
    def _extract_robust_shape_id(self, shape, fallback_idx: int, debug: bool = False) -> int:
        """
        Extract robust shape ID with enhanced XML-based extraction for grouped children.
        
        Args:
            shape: Shape object
            fallback_idx: Fallback index if no ID found
            debug: Enable debug logging
            
        Returns:
            Shape ID (integer) or fallback index
        """
        try:
            # Method 1: Standard python-pptx API
            shape_id = getattr(shape, 'shape_id', None)
            if shape_id is not None:
                if debug:
                    logger.debug(f"    Extracted ID via API: {shape_id}")
                return shape_id
            
            # Method 2: XML-based extraction for grouped children
            if hasattr(shape, '_element') and shape._element is not None:
                element = shape._element
                
                # Try different XML paths for shape IDs
                id_extraction_paths = [
                    # Standard shape ID paths
                    ('.//p:cNvPr/@id', 'p:cNvPr id attribute'),
                    ('.//pic:cNvPr/@id', 'pic:cNvPr id attribute'),
                    
                    # Group-specific paths
                    ('.//p:nvSpPr/p:cNvPr/@id', 'nvSpPr/cNvPr id'),
                    ('.//p:nvPicPr/p:cNvPr/@id', 'nvPicPr/cNvPr id'),
                    ('.//p:nvCxnSpPr/p:cNvPr/@id', 'nvCxnSpPr/cNvPr id'),
                    
                    # Fallback paths
                    ('.//@id', 'any id attribute')
                ]
                
                namespaces = {
                    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                    'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture'
                }
                
                for xpath, description in id_extraction_paths:
                    try:
                        id_results = _safe_xpath(element, xpath, namespaces=namespaces)
                        if id_results:
                            extracted_id = int(id_results[0])
                            if debug:
                                logger.debug(f"    Extracted ID via XML ({description}): {extracted_id}")
                            return extracted_id
                    except (ValueError, TypeError, Exception) as e:
                        if debug:
                            logger.debug(f"    XML extraction failed for {description}: {e}")
                        continue
            
            # Method 3: Shape name-based heuristic ID extraction
            shape_name = getattr(shape, 'name', '')
            if shape_name:
                # Look for ID patterns in shape names like "Shape 123" or "Picture 456"
                import re
                id_match = re.search(r'\d+$', shape_name)
                if id_match:
                    try:
                        name_based_id = int(id_match.group())
                        if debug:
                            logger.debug(f"    Extracted ID from shape name '{shape_name}': {name_based_id}")
                        return name_based_id
                    except ValueError:
                        pass
            
            if debug:
                logger.debug(f"    No stable ID found, using fallback index: {fallback_idx}")
            
        except Exception as e:
            if debug:
                logger.debug(f"    Error during robust ID extraction: {e}")
        
        return fallback_idx
    
    def _normalize_alt(self, txt: str) -> str:
        """
        PHASE 1: Enhanced normalization to eliminate all duplication patterns.
        Applied before EVERY injection path to ensure clean ALT text.
        
        Args:
            txt: ALT text that may contain duplications
            
        Returns:
            Cleaned and normalized ALT text with ALL duplications removed
        """
        if not txt or not txt.strip():
            return ""
        
        import re
        
        # Normalize whitespace first
        t = " ".join(txt.split())
        
        # PHASE 1 FIX 1: Remove duplicate size preambles
        # Target: "A shape (99x99px) This is a PowerPoint shape. It is a shape (99x99px)"
        t = re.sub(r"^(A (?:shape|line|connector) \(\d+x\d+px\))\s+(This is a PowerPoint shape\..*?)\s+\1", r"\2", t, flags=re.IGNORECASE)
        
        # PHASE 1 FIX 2: Remove repeated PowerPoint shape descriptions
        # Target: "This is a PowerPoint shape. It is a shape (99x99px). This is a PowerPoint shape. It is a shape (99x99px)"
        t = re.sub(r"(This is a PowerPoint shape\.[^.]*\.)\s*\1", r"\1", t, flags=re.IGNORECASE)
        
        # PHASE 1 FIX 3: Remove duplicate size information
        # Target: "This is a PowerPoint shape. It is a shape (99x99px) (99x99px)"
        t = re.sub(r"(\(\d+x\d+px\))\s*\1", r"\1", t)
        
        # PHASE 1 FIX 4: Handle repeated vector format descriptions
        # Target: "Vector scientific diagram... Vector scientific diagram..."
        vector_patterns = [
            r"(Vector [^.]+diagram[^.]*\.)\s*\1",
            r"(Vector [^.]+circuit[^.]*\.)\s*\1", 
            r"(Vector [^.]+process[^.]*\.)\s*\1"
        ]
        for pattern in vector_patterns:
            t = re.sub(pattern, r"\1", t, flags=re.IGNORECASE)
        
        # PHASE 1 FIX 5: Remove redundant "This is a PowerPoint shape" occurrences
        parts = t.split("This is a PowerPoint shape")
        if len(parts) > 2:  # More than one occurrence
            # Keep the first occurrence and the most detailed part
            detailed_part = max(parts[1:], key=len) if len(parts) > 1 else parts[1] if len(parts) > 1 else ""
            t = f"This is a PowerPoint shape{detailed_part}"
        
        # PHASE 1 FIX 6: Remove duplicate format references (WMF, EMF, etc.)
        format_types = ['WMF', 'EMF', 'TIFF', 'PNG', 'JPEG']
        for fmt in format_types:
            # Remove duplicate format names within close proximity
            pattern = rf'\b{fmt}\b(.{{0,50}}?)\b{fmt}\b'
            t = re.sub(pattern, f'{fmt}\\1', t, flags=re.IGNORECASE)
        
        # PHASE 1 FIX 7: Clean up whitespace and punctuation
        t = re.sub(r'\s+', ' ', t)  # Multiple spaces -> single space
        t = re.sub(r'\.+', '.', t)  # Multiple periods -> single period
        t = re.sub(r'\s+\.', '.', t)  # "space." -> "."
        t = re.sub(r'\s+,', ',', t)   # "space," -> ","
        t = self._ensure_terminal_punctuation(t)
        return t

    @staticmethod
    def _ensure_terminal_punctuation(text: str) -> str:
        """Ensure text ends with a terminal punctuation mark."""
        text = text.strip()
        if not text:
            return ""
        if text[-1] in ".!?":
            return text
        return f"{text}."
    
    def _compose_alt(self, parts: list[str]) -> str:
        """
        HOTPATCH FIX 2: Compose ALT text from multiple sources at join point.
        This addresses the root cause (concatenation) instead of only cleaning late.
        
        Args:
            parts: List of ALT text parts to combine
            
        Returns:
            Single clean composed ALT text string
        """
        if not parts:
            return ""
        
        # Filter out empty parts
        parts = [p.strip() for p in parts if p and p.strip()]
        if not parts:
            return ""
        
        # If only one part, just normalize and return
        if len(parts) == 1:
            return self._normalize_alt(parts[0])
        
        # Drop duplicates by stem (strip sizes, casefold)
        import re
        
        def get_stem(text: str) -> str:
            """Get the core content without sizes and formatting"""
            stem = text.strip().lower()
            # Remove size information
            stem = re.sub(r'\(\d+x\d+px\)', '', stem)
            stem = re.sub(r'\d+×\d+\s*pixels?', '', stem)
            # Remove common formatting
            stem = re.sub(r'\s+', ' ', stem).strip()
            return stem
        
        # Remove duplicate parts based on stem content
        seen_stems = set()
        unique_parts = []
        
        for part in parts:
            stem = get_stem(part)
            if stem and stem not in seen_stems:
                seen_stems.add(stem)
                unique_parts.append(part)
        
        if not unique_parts:
            return ""
        
        # Score parts to prefer semantic over boilerplate
        def score_semantic_value(text: str) -> int:
            """Score semantic value - higher is better"""
            text_lower = text.lower()
            score = 0
            
            # Penalty for boilerplate
            boilerplate_patterns = [
                'this is a powerpoint shape',
                'powerpoint element',  
                'visual element',
                'bypass:',
                r'^a (?:shape|line|connector) \(\d+x\d+px\)'
            ]
            
            for pattern in boilerplate_patterns:
                if re.search(pattern, text_lower):
                    score -= 2
            
            # Bonus for semantic content
            semantic_indicators = [
                'icon', 'arrow', 'circle', 'diagram', 'chart', 'graph', 'image',
                'photo', 'screenshot', 'illustration', 'logo', 'symbol', 'button'
            ]
            
            semantic_count = sum(1 for indicator in semantic_indicators if indicator in text_lower)
            score += semantic_count * 2
            
            # Bonus for good length (not too short, not too verbose)
            if 20 <= len(text.strip()) <= 150:
                score += 1
                
            return score
        
        # Sort by semantic value (highest first) and take the best one
        unique_parts.sort(key=score_semantic_value, reverse=True)
        best_part = unique_parts[0]
        
        # Apply final normalization
        composed = self._normalize_alt(best_part)
        
        return composed
    
    def _is_llava_error(self, description: str) -> bool:
        """
        Check if LLaVA returned an error response.
        
        Args:
            description: Generated description to check
            
        Returns:
            bool: True if this is an error response
        """
        if not description or not description.strip():
            return True
        
        error_patterns = [
            'error', 'failed', 'cannot', 'unable', 'sorry',
            'i cannot', 'i am unable', 'no description',
            'not available', 'description not available',
            'image could not be processed', 'cannot describe',
            'unable to describe', 'failed to process',
            'api error', 'request failed', 'timeout'
        ]
        
        description_lower = description.lower().strip()
        return any(pattern in description_lower for pattern in error_patterns)
    
    def _handle_llava_error_with_fallback(self, element, debug: bool = False) -> str:
        """
        Handle LLaVA errors with retry and fallback chain.
        
        Args:
            element: Visual element that failed
            debug: Enable debug logging
            
        Returns:
            str: Fallback description or empty string
        """
        if debug:
            logger.debug(f"    🔄 Handling LLaVA error for {element.element_key}")
        
        # Prefer vector-aware fallback for image elements with WMF/EMF
        try:
            if getattr(element, 'element_type', '') == 'image' and getattr(element, 'filename', ''):
                from pathlib import Path
                ext = str(Path(element.filename).suffix).lower()
                if ext in ('.wmf', '.emf'):
                    # Build a minimal PPTXImageInfo to reuse vector fallback
                    dummy = PPTXImageInfo(
                        shape=element.shape,
                        slide_idx=element.slide_idx,
                        shape_idx=element.shape_idx,
                        image_data=element.image_data or b'',
                        filename=element.filename,
                        slide_text=element.slide_text
                    )
                    fmt = 'WMF' if ext == '.wmf' else 'EMF'
                    return self._generate_vector_fallback_alt(dummy, fmt, debug=debug)
        except Exception as e:
            if debug:
                logger.debug(f"Vector-aware fallback path failed: {e}")

        # default: generic shape description
        return self._create_enhanced_fallback_description(element)
    
    def _create_enhanced_fallback_description(self, element) -> str:
        """
        Create enhanced fallback description with better freeform handling.
        
        Args:
            element: Visual element to describe
            
        Returns:
            str: Enhanced fallback description
        """
        try:
            # Handle freeform shapes explicitly
            if hasattr(element.shape, 'shape_type'):
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                if element.shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                    # Enhanced freeform description
                    width_px = getattr(element, 'width_px', 0)
                    height_px = getattr(element, 'height_px', 0)
                    size_info = f"({width_px}x{height_px}px)" if width_px > 0 and height_px > 0 else ""
                    
                    # Try to detect semantic meaning from context
                    semantic_hint = self._detect_freeform_semantic_meaning(element)
                    if semantic_hint:
                        return f"Freeform shape representing {semantic_hint} {size_info}".strip()
                    else:
                        return f"PowerPoint freeform shape {size_info}".strip()
            
            # Use existing describe_shape_with_details for other shapes
            return describe_shape_with_details(element.shape)
            
        except Exception as e:
            logger.debug(f"Error creating enhanced fallback: {e}")
            return f"PowerPoint visual element on slide {element.slide_idx + 1}"
    
    def _detect_freeform_semantic_meaning(self, element) -> str:
        """
        Detect semantic meaning of freeform shapes based on context.
        
        Args:
            element: Visual element with freeform shape
            
        Returns:
            str: Semantic description or empty string
        """
        try:
            # Check if it's part of a group that might indicate function
            slide_text = getattr(element, 'slide_text', '').lower()
            shape_name = getattr(element.shape, 'name', '').lower()
            
            # Common icon patterns
            icon_patterns = {
                'arrow': ['arrow', 'direction', 'next', 'previous'],
                'star': ['star', 'rating', 'favorite'],
                'heart': ['heart', 'love', 'like'],
                'check': ['check', 'tick', 'correct', 'done'],
                'cross': ['cross', 'x', 'close', 'delete'],
                'warning': ['warning', 'alert', 'caution'],
                'info': ['info', 'information', 'help']
            }
            
            for icon_type, keywords in icon_patterns.items():
                if any(keyword in shape_name for keyword in keywords):
                    return f"an {icon_type} icon"
                if any(keyword in slide_text for keyword in keywords):
                    return f"an {icon_type} icon"
            
            # Fallback based on size
            width_px = getattr(element, 'width_px', 0)
            height_px = getattr(element, 'height_px', 0)
            
            if width_px < 50 and height_px < 50:
                return "a small icon or symbol"
            elif width_px > 200 or height_px > 200:
                return "a large decorative shape"
            else:
                return "an icon or symbol"
                
        except Exception:
            return ""
    
    def _deduplicate_format_references(self, alt_text: str, format_name: str) -> str:
        """
        Remove duplicate format references from ALT text to prevent "WMF format... WMF format" patterns.
        
        Args:
            alt_text: ALT text that may contain duplicate format references
            format_name: Format name (WMF, EMF, etc.)
            
        Returns:
            Cleaned ALT text with deduplicated format references
        """
        try:
            # Simple approach: check for and remove obvious duplications
            import re
            
            # Check if format name appears multiple times in problematic patterns
            format_count = alt_text.upper().count(format_name.upper())
            
            if format_count > 1:
                # Remove duplicate format references in common patterns
                # Pattern 1: "WMF format... WMF format" -> "WMF format..."
                pattern1 = rf'({format_name}\s+format)(.+?)\1'
                alt_text = re.sub(pattern1, r'\1\2', alt_text, flags=re.IGNORECASE)
                
                # Pattern 2: Remove duplicate format names that are close together
                pattern2 = rf'({format_name})(\s+[^.]{0,30}?\s+){format_name}'
                alt_text = re.sub(pattern2, r'\1\2', alt_text, flags=re.IGNORECASE)
                
                # Pattern 3: Remove duplicate "Original X vector image" phrases
                pattern3 = rf'(Original {format_name} vector image[^.]*\.)(.+?)(Original {format_name} vector image[^.]*\.)'
                alt_text = re.sub(pattern3, r'\1\2', alt_text, flags=re.IGNORECASE)
            
            # Clean up any resulting double spaces or formatting issues
            alt_text = re.sub(r'\s+', ' ', alt_text)  # Multiple spaces -> single space
            alt_text = re.sub(r'\s+\.', '.', alt_text)  # "space." -> "."
            alt_text = re.sub(r'\s+,', ',', alt_text)   # "space," -> ","
            
            return alt_text.strip()
            
        except Exception as e:
            # If deduplication fails, return original text
            logger.debug(f"Error deduplicating format references: {e}")
            return alt_text
    
    def _generate_alt_text_for_image_with_validation(self, image_info: PPTXImageInfo, debug: bool = False) -> Tuple[Optional[str], Optional[str]]:
        """
        Generate ALT text with comprehensive validation and detailed failure tracking.
        
        Args:
            image_info: Image information
            debug: Whether to enable debug logging
            
        Returns:
            Tuple of (alt_text, failure_reason). If alt_text is None/empty, failure_reason explains why.
        """
        failure_reason = None
        
        try:
            # Normalize image format before processing
            try:
                normalized_image_data = self._normalize_image_format(image_info.image_data, image_info.filename, debug)
                
                # Save normalized image to temporary file for ALT text generation
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                    temp_file.write(normalized_image_data)
                    temp_image_path = temp_file.name
                    
            except Exception as norm_error:
                # Check if this is a vector format conversion failure OR any WMF/EMF processing failure
                if (("Vector format conversion failed" in str(norm_error) or 
                     "format normalization failed" in str(norm_error) or
                     "Cannot process" in str(norm_error)) and 
                    image_info.filename.lower().endswith(('.wmf', '.emf'))):
                    # Generate contextual fallback ALT text
                    format_name = "WMF" if image_info.filename.lower().endswith('.wmf') else "EMF"
                    logger.info(f"Generating contextual fallback ALT for {image_info.filename}")
                    contextual_alt = self._generate_vector_fallback_alt(image_info, format_name, debug)
                    return contextual_alt, None  # No failure reason since we provided fallback
                else:
                    # For other normalization failures, re-raise
                    raise norm_error
            
            try:
                # Build context for better ALT text generation
                context = self._build_generation_context(image_info)
                
                # Determine appropriate prompt type based on content
                prompt_type = self._determine_prompt_type(image_info)
                
                if debug:
                    logger.info(f"🔍 DEBUG: Using prompt type '{prompt_type}' with context: {context[:100]}...")
                
                # Generate ALT text using the configured generator
                alt_text = self.alt_generator.generate_alt_text(
                    image_path=temp_image_path,
                    prompt_type=prompt_type,
                    context=context
                )
                
                # Comprehensive validation of the generated ALT text
                if alt_text is None:
                    failure_reason = "Generator returned None"
                    return None, failure_reason
                
                if not isinstance(alt_text, str):
                    failure_reason = f"Generator returned non-string type: {type(alt_text)}"
                    return None, failure_reason
                
                alt_text_stripped = alt_text.strip()
                if not alt_text_stripped:
                    failure_reason = "Generator returned empty or whitespace-only string"
                    return None, failure_reason
                
                if len(alt_text_stripped) < 3:
                    failure_reason = f"Generator returned very short ALT text: '{alt_text_stripped}'"
                    return None, failure_reason
                
                # Check for common failure patterns and route to contextual fallback
                failure_patterns = [
                    'error', 'failed', 'cannot', 'unable', 'sorry', 
                    'i cannot', 'i am unable', 'no description',
                    'not available', 'description not available'
                ]
                
                alt_text_lower = alt_text_stripped.lower()
                for pattern in failure_patterns:
                    if pattern in alt_text_lower:
                        if debug:
                            logger.info(f"🔧 DEBUG: Error string detected ('{pattern}'), routing to contextual fallback")
                        # Route to contextual fallback instead of returning raw error
                        try:
                            fallback_alt = self._generate_vector_fallback_alt(image_info, "Unknown", debug)
                            if debug:
                                logger.info(f"✅ DEBUG: Used contextual fallback for error: '{fallback_alt[:50]}...'")
                            return fallback_alt, None  # Return fallback instead of failing
                        except Exception as fallback_error:
                            if debug:
                                logger.error(f"❌ DEBUG: Contextual fallback also failed: {fallback_error}")
                            failure_reason = f"Generator returned failure message containing '{pattern}' and fallback failed: {fallback_error}"
                            return None, failure_reason
                
                # ALT text passed all validation checks
                # PHASE 1: Apply universal normalization before returning
                normalized_alt_text = self._normalize_alt(alt_text_stripped)
                
                # HOTPATCH FIX 3: IMPOSSIBLE-TO-MISS GUARD for any remaining error strings
                # Final safety check after normalization to catch any errors that slipped through
                if not normalized_alt_text or 'error' in normalized_alt_text.lower():
                    if debug:
                        logger.warning(f"🚨 IMPOSSIBLE-TO-MISS GUARD: Error detected in final ALT text: '{normalized_alt_text}'")
                    # Force contextual fallback as last resort
                    try:
                        fallback_alt = self._generate_vector_fallback_alt(image_info, "Unknown", debug)
                        if debug:
                            logger.info(f"🔧 IMPOSSIBLE-TO-MISS: Used emergency fallback: '{fallback_alt}'")
                        return fallback_alt, None
                    except Exception as emergency_error:
                        if debug:
                            logger.error(f"❌ IMPOSSIBLE-TO-MISS: Emergency fallback failed: {emergency_error}")
                        return None, f"Emergency fallback failed: {emergency_error}"
                
                if debug:
                    logger.info(f"✅ DEBUG: Generated and normalized ALT text: '{normalized_alt_text[:50]}...'")
                
                return normalized_alt_text, None
                
            finally:
                # Clean up temporary file
                try:
                    os.unlink(temp_image_path)
                except OSError:
                    pass  # File cleanup failure is not critical
        
        except Exception as e:
            failure_reason = f"Exception during generation: {str(e)}"
            if debug:
                logger.error(f"💥 DEBUG: Exception in ALT text generation: {e}", exc_info=True)
            
            # Try to provide a descriptive fallback for rendered shapes
            if hasattr(image_info, 'is_rendered') and image_info.is_rendered:
                try:
                    fallback_alt_text = self._create_descriptive_shape_alt_text(
                        image_info.shape, image_info.width_px, image_info.height_px
                    )
                    # PHASE 1: Apply universal normalization to fallback ALT text too
                    normalized_fallback = self._normalize_alt(fallback_alt_text)
                    if debug:
                        logger.info(f"🔧 DEBUG: Using normalized descriptive shape fallback: '{normalized_fallback}'")
                    return normalized_fallback, None
                except Exception as fallback_error:
                    if debug:
                        logger.error(f"Fallback generation also failed: {fallback_error}")
                    
            return None, failure_reason
    
    def _build_generation_context(self, image_info: PPTXImageInfo) -> Optional[str]:
        """
        Build context string for ALT text generation.
        
        Args:
            image_info: Image information
            
        Returns:
            Context string or None
        """
        context_parts = []
        
        # Add slide text context
        if image_info.slide_text:
            context_parts.append(f"Slide content: {image_info.slide_text}")
        
        # Add slide number
        context_parts.append(f"Slide {image_info.slide_idx + 1}")
        
        # Add image filename if it provides context
        if image_info.filename and not image_info.filename.startswith('slide_'):
            context_parts.append(f"File: {image_info.filename}")
        
        return ". ".join(context_parts) if context_parts else None
    
    def _determine_prompt_type(self, image_info: PPTXImageInfo) -> str:
        """
        Determine the appropriate prompt type based on image and context.
        
        Args:
            image_info: Image information
            
        Returns:
            Prompt type string
        """
        # Special handling for rendered shapes
        if hasattr(image_info, 'is_rendered') and image_info.is_rendered:
            # For rendered shapes, use a shape-specific prompt or diagram prompt
            return 'diagram'
        
        # Check filename and context for medical content indicators
        text_to_check = (image_info.filename + " " + image_info.slide_text).lower()
        
        # Medical-specific prompt detection
        anatomical_keywords = ['anatomy', 'organ', 'body', 'muscle', 'bone', 'tissue']
        diagnostic_keywords = ['xray', 'x-ray', 'ct', 'mri', 'ultrasound', 'scan', 'radiograph']
        clinical_keywords = ['patient', 'clinical', 'medical', 'surgery', 'procedure', 'treatment']
        chart_keywords = ['chart', 'graph', 'data', 'results', 'statistics', 'plot']
        diagram_keywords = ['diagram', 'flowchart', 'process', 'workflow', 'schematic']
        
        if any(keyword in text_to_check for keyword in anatomical_keywords):
            return 'anatomical'
        elif any(keyword in text_to_check for keyword in diagnostic_keywords):
            return 'diagnostic'
        elif any(keyword in text_to_check for keyword in clinical_keywords):
            return 'clinical_photo'
        elif any(keyword in text_to_check for keyword in chart_keywords):
            return 'chart'
        elif any(keyword in text_to_check for keyword in diagram_keywords):
            return 'diagram'
        else:
            # Check if this appears to be medical content in general
            medical_keywords = ['medical', 'health', 'doctor', 'hospital', 'clinic']
            if any(keyword in text_to_check for keyword in medical_keywords):
                return 'unified_medical'
        
        return 'default'
    
    def _validate_visual_element_coverage(self, visual_elements: List[PPTXVisualElement], alt_text_mapping: Dict[str, Any], debug: bool = False) -> Dict[str, Any]:
        """
        Validate ALT text coverage for visual elements.
        
        Args:
            visual_elements: List of all visual elements found in the PPTX
            alt_text_mapping: Current ALT text mappings
            debug: Enable debug logging
            
        Returns:
            Dictionary with validation results
        """
        total_elements = len(visual_elements)
        covered_elements = len(alt_text_mapping)
        missing_elements = []
        
        # Check each visual element has ALT text
        for visual_element in visual_elements:
            if visual_element.element_key not in alt_text_mapping:
                missing_elements.append(visual_element.element_key)
        
        missing_count = len(missing_elements)
        coverage_percent = (covered_elements / total_elements * 100) if total_elements > 0 else 0
        complete_coverage = missing_count == 0
        
        # Count descriptive vs decorative
        descriptive_count = 0
        decorative_count = 0
        for element_key, element_data in alt_text_mapping.items():
            alt_text = element_data.get('alt_text', '').lower()
            if 'decorative' in alt_text or alt_text.startswith('[decorative'):
                decorative_count += 1
            else:
                descriptive_count += 1
        
        validation_result = {
            'complete_coverage': complete_coverage,
            'total_elements': total_elements,
            'covered_elements': covered_elements,
            'missing_count': missing_count,
            'missing_elements': missing_elements,
            'coverage_percent': coverage_percent,
            'descriptive_count': descriptive_count,
            'decorative_count': decorative_count,
            'total_coverage_percent': coverage_percent
        }
        
        if debug:
            logger.info(f"🔍 DEBUG: Visual element coverage validation results:")
            logger.info(f"   Total elements: {total_elements}")
            logger.info(f"   Elements with ALT text: {covered_elements}")
            logger.info(f"   Descriptive ALT text: {descriptive_count}")
            logger.info(f"   Decorative ALT text: {decorative_count}")
            logger.info(f"   Coverage: {coverage_percent:.1f}%")
            
            if missing_count > 0:
                logger.info(f"   Missing ALT text: {missing_count}")
                if debug and missing_count <= 5:  # Show first few missing
                    for missing_key in missing_elements[:5]:
                        logger.info(f"     - {missing_key}")
                    if missing_count > 5:
                        logger.info(f"     ... and {missing_count - 5} more")
        
        # Log validation summary
        if complete_coverage:
            logger.info(f"✅ Complete visual element ALT text coverage achieved: {covered_elements}/{total_elements} elements")
            logger.info(f"   Descriptive: {descriptive_count}, Decorative: {decorative_count}")
        else:
            logger.warning(f"⚠️ Partial visual element ALT text coverage: {covered_elements}/{total_elements} elements ({coverage_percent:.1f}%)")
            logger.warning(f"   Missing ALT text for {missing_count} elements")
        
        return validation_result
    
    def _validate_complete_coverage(self, image_infos: List[PPTXImageInfo], alt_text_mapping: Dict[str, Any], 
                                  force_decorative: bool = False, debug: bool = False) -> Dict[str, Any]:
        """
        Validate that every image has ALT text (either real or decorative fallback).
        
        Args:
            image_infos: List of all images found in the PPTX
            alt_text_mapping: Current ALT text mappings
            force_decorative: Whether decorative fallback was enabled
            debug: Whether to enable debug logging
            
        Returns:
            Dictionary with validation results
        """
        total_images = len(image_infos)
        covered_images = len(alt_text_mapping)
        missing_elements = []
        
        # Check each image has ALT text
        for image_info in image_infos:
            if image_info.image_key not in alt_text_mapping:
                missing_elements.append(image_info.image_key)
        
        missing_count = len(missing_elements)
        coverage_percent = (covered_images / total_images * 100) if total_images > 0 else 0
        complete_coverage = missing_count == 0
        
        # Count ALT text types
        descriptive_count = 0
        decorative_count = 0
        
        for key, info in alt_text_mapping.items():
            alt_text = info['alt_text']
            if alt_text == '[Decorative image]':
                decorative_count += 1
            else:
                descriptive_count += 1
        
        validation_result = {
            'complete_coverage': complete_coverage,
            'total_images': total_images,
            'covered_images': covered_images,
            'missing_count': missing_count,
            'missing_elements': missing_elements,
            'descriptive_count': descriptive_count,
            'decorative_count': decorative_count,
            'total_coverage_percent': coverage_percent
        }
        
        if debug:
            logger.info(f"🔍 DEBUG: Coverage validation results:")
            logger.info(f"   Total images: {total_images}")
            logger.info(f"   Images with ALT text: {covered_images}")
            logger.info(f"   Descriptive ALT text: {descriptive_count}")
            logger.info(f"   Decorative ALT text: {decorative_count}")
            logger.info(f"   Missing ALT text: {missing_count}")
            logger.info(f"   Coverage: {coverage_percent:.1f}%")
            
            if missing_count > 0:
                logger.warning(f"❌ DEBUG: {missing_count} images missing ALT text:")
                for missing_key in missing_elements[:5]:  # Show first 5
                    logger.warning(f"   - {missing_key}")
                if len(missing_elements) > 5:
                    logger.warning(f"   ... and {len(missing_elements) - 5} more")
        
        # Log validation summary
        if complete_coverage:
            logger.info(f"✅ Complete ALT text coverage achieved: {covered_images}/{total_images} images")
            logger.info(f"   Descriptive: {descriptive_count}, Decorative: {decorative_count}")
        else:
            logger.error(f"❌ Incomplete ALT text coverage: {covered_images}/{total_images} images ({coverage_percent:.1f}%)")
            logger.error(f"   Missing ALT text for {missing_count} images")
            if not force_decorative:
                logger.error("   💡 Consider using --force-decorative to ensure 100% coverage")
        
        return validation_result
    
    def _inject_alt_text_to_pptx(self, presentation: Presentation, 
                               alt_text_mapping: Dict[str, Any], output_path: str, debug: bool = False) -> tuple[bool, dict]:
        """
        Inject ALT text into PPTX presentation using the dedicated injector.
        
        Args:
            presentation: Presentation object
            alt_text_mapping: Mapping of image keys to ALT text and shape info
            output_path: Path to save modified PPTX
            
        Returns:
            bool: True if injection succeeded
        """
        try:
            # Import the dedicated ALT text injector
            from pptx_alt_injector import PPTXAltTextInjector
            
            # Create injector instance
            injector = PPTXAltTextInjector(self.config_manager)
            
            # Convert mapping format to match injector expectations
            simple_mapping = {}
            logger.debug("Converting ALT text mapping from processor format:")
            for image_key, info in alt_text_mapping.items():
                simple_mapping[image_key] = info['alt_text']
                logger.debug(f"  Processor key: {image_key} -> ALT: '{info['alt_text'][:50]}...'")
            
            logger.debug(f"Created simple mapping with {len(simple_mapping)} entries for injector")
            
            # Save presentation to temp file for injector processing
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_path = temp_file.name
            
            presentation.save(temp_path)
            
            # Use injector to perform robust ALT text injection
            result = injector.inject_alt_text_from_mapping(temp_path, simple_mapping, output_path)
            
            # Clean up temp file
            try:
                os.unlink(temp_path)
            except OSError:
                pass
            
            # Log injector statistics
            stats = result['statistics']
            logger.info(f"ALT text injection via dedicated injector:")
            logger.info(f"  Successfully injected: {stats['injected_successfully']}")
            logger.info(f"  Skipped (existing): {stats['skipped_existing']}")
            logger.info(f"  Failed: {stats['failed_injection']}")
            
            return result['success'], simple_mapping
            
        except Exception as e:
            logger.error(f"Failed to inject ALT text via dedicated injector: {e}")
            # Fallback to original simple method
            fallback_success = self._inject_alt_text_simple(presentation, alt_text_mapping, output_path)
            # Create simple mapping for fallback compatibility
            simple_mapping = {image_key: info['alt_text'] for image_key, info in alt_text_mapping.items()}
            return fallback_success, simple_mapping
    
    def _inject_alt_text_simple(self, presentation: Presentation, 
                              alt_text_mapping: Dict[str, Any], output_path: str) -> bool:
        """
        Fallback simple ALT text injection method.
        
        Args:
            presentation: Presentation object
            alt_text_mapping: Mapping of image keys to ALT text and shape info
            output_path: Path to save modified PPTX
            
        Returns:
            bool: True if injection succeeded
        """
        try:
            logger.info(f"Using fallback injection method for {len(alt_text_mapping)} images")
            
            for image_key, info in alt_text_mapping.items():
                try:
                    shape = info['shape']
                    alt_text = info['alt_text']
                    
                    # Set ALT text on the shape
                    if hasattr(shape, '_element'):
                        # Access the underlying XML element and set the description
                        shape._element.set('descr', alt_text)
                        logger.debug(f"Set ALT text for {image_key}: {alt_text[:50]}...")
                    else:
                        logger.warning(f"Cannot set ALT text for {image_key}: shape has no _element")
                        
                except Exception as e:
                    logger.error(f"Failed to set ALT text for {image_key}: {e}")
                    continue
            
            # Save the modified presentation
            presentation.save(output_path)
            logger.info(f"Saved PPTX with ALT text to: {output_path}")
            
            return True
            
        except Exception as e:
            logger.error(f"Fallback ALT text injection failed: {e}")
            return False
    
    def _log_processing_summary(self, result: Dict[str, Any]):
        """Log a summary of the processing results."""
        logger.info("\n" + "="*60)
        logger.info("📊 PPTX PROCESSING SUMMARY")
        logger.info("="*60)
        
        # File information
        logger.info(f"📁 Input file: {result['input_file']}")
        logger.info(f"💾 Output file: {result['output_file']}")
        logger.info(f"📄 Total slides: {result['total_slides']}")
        
        # Visual element processing summary
        total_elements = result.get('total_visual_elements', 0)
        processed_elements = result.get('processed_visual_elements', 0)
        failed_elements = result.get('failed_visual_elements', 0)
        
        logger.info(f"\n🎯 Visual Element Processing:")
        logger.info(f"   📊 Total elements found: {total_elements}")
        logger.info(f"   ✅ Elements processed: {processed_elements}")
        logger.info(f"   ❌ Failed elements: {failed_elements}")
        
        if total_elements > 0:
            success_rate = (processed_elements / total_elements) * 100
            logger.info(f"   🎯 Success rate: {success_rate:.1f}%")
            logger.info(f"   📈 Elements per slide: {total_elements / result['total_slides']:.1f}")
        
        # Timing information
        logger.info(f"  Generation time: {result['generation_time']:.2f}s")
        logger.info(f"  Injection time: {result['injection_time']:.2f}s")
        if 'decorative_marking_time' in result:
            logger.info(f"  Decorative marking time: {result['decorative_marking_time']:.2f}s")
        logger.info(f"  Total processing time: {result['total_time']:.2f}s")
        logger.info(f"  Success: {result['success']}")
        
        # Calculate and log coverage
        if result.get('total_visual_elements', 0) > 0:
            coverage_percent = (result.get('processed_visual_elements', 0) / result['total_visual_elements'] * 100)
            logger.info(f"  Visual Element Coverage: {result.get('processed_visual_elements', 0)}/{result['total_visual_elements']} ({coverage_percent:.1f}%)")
        
        if result['errors']:
            logger.warning(f"Errors encountered: {len(result['errors'])}")
            for error in result['errors']:
                logger.warning(f"  - {error}")


def debug_image_extraction(pptx_path: str):
    """
    Debug function to test comprehensive image extraction with detailed logging.
    
    Args:
        pptx_path: Path to PPTX file to analyze
    """
    # Set up debug logging
    logging.basicConfig(
        level=logging.DEBUG,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )
    
    print(f"🔍 DEBUG MODE: Analyzing image extraction in {pptx_path}")
    print("=" * 80)
    
    try:
        # Initialize processor with debug enabled
        config_manager = ConfigManager()
        processor = PPTXAccessibilityProcessor(config_manager, debug=True)
        
        # Extract images with comprehensive logging
        presentation, image_infos = processor._extract_images_from_pptx(pptx_path)
        
        # Summary
        print("\n" + "=" * 80)
        print("🔍 DEBUG SUMMARY:")
        print(f"   Total slides: {len(presentation.slides)}")
        print(f"   Total images found: {len(image_infos)}")
        
        if image_infos:
            print("\n📋 Image Details:")
            for i, img_info in enumerate(image_infos, 1):
                print(f"   {i:2d}. {img_info.filename} ({img_info.width_px}x{img_info.height_px}px)")
                print(f"       Key: {img_info.image_key}")
                print(f"       Slide: {img_info.slide_idx + 1}, Shape: {img_info.shape_idx}")
                print(f"       Hash: {img_info.image_hash[:8]}...")
                if img_info.slide_text:
                    print(f"       Context: {img_info.slide_text[:60]}...")
                print()
        else:
            print("   ❌ No images were detected!")
            print("   Check the debug logs above to see what shapes were found.")
        
    except Exception as e:
        print(f"❌ Error during debug extraction: {e}")
        raise

def main():
    """Test the PPTX accessibility processor."""
    import sys
    
    # Check for debug flag
    debug_mode = '--debug' in sys.argv
    if debug_mode:
        sys.argv.remove('--debug')
    
    # Set up logging
    log_level = logging.DEBUG if debug_mode else logging.INFO
    logging.basicConfig(
        level=log_level,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )
    
    if len(sys.argv) not in [2, 3]:
        print("Usage: python pptx_processor.py [--debug] <pptx_file> [output_file]")
        print("\nOptions:")
        print("  --debug    Enable debug mode with comprehensive image extraction logging")
        print("\nThis will process a PPTX to add ALT text to images.")
        print("If output_file is not specified, the original file will be overwritten.")
        print("\nFor debugging image extraction issues, use --debug flag.")
        return
    
    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) == 3 else None
    
    # If debug mode is enabled, just run image extraction debugging
    if debug_mode:
        debug_image_extraction(pptx_path)
        return 0
    
    try:
        print("PPTX Accessibility Processor Test")
        print("=" * 50)
        print(f"Processing: {pptx_path}")
        if output_path:
            print(f"Output: {output_path}")
        else:
            print("Output: Overwriting original file")
        print()
        
        # Initialize processor
        config_manager = ConfigManager()
        processor = PPTXAccessibilityProcessor(config_manager)
        
        # Process PPTX
        result = processor.process_pptx(pptx_path, output_path)
        
        # Display results
        print("Processing Results:")
        print(f"  Success: {result['success']}")
        print(f"  Total slides: {result['total_slides']}")
        print(f"  Total visual elements: {result.get('total_visual_elements', 0)}")
        print(f"  Processed: {result.get('processed_visual_elements', 0)}")
        print(f"  Failed: {result.get('failed_visual_elements', 0)}")
        
        print(f"  Total time: {result['total_time']:.2f}s")
        
        if result['errors']:
            print(f"  Errors: {len(result['errors'])}")
            for error in result['errors']:
                print(f"    - {error}")
        
        print()
        if result['success']:
            print("✅ PPTX processing completed successfully!")
            print(f"Modified PPTX saved to: {result['output_file']}")
        else:
            print("❌ PPTX processing failed!")
            return 1
        
    except Exception as e:
        logger.error(f"Processing failed: {e}")
        print(f"Error: {e}")
        return 1
    
    return 0


if __name__ == "__main__":
    exit(main())